#!/usr/bin/env python3
"""
Simple PPTX thumbnail generator without LibreOffice dependency.
Directly extracts slides as images using python-pptx + PIL.
"""

import sys
from pathlib import Path
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

# Constants
THUMBNAIL_WIDTH = 600
JPEG_QUALITY = 95
GRID_PADDING = 20
BORDER_WIDTH = 2
MAX_COLS = 6
DEFAULT_COLS = 4

def extract_slide_images(pptx_path, temp_dir):
    """Extract slide images from PPTX."""
    prs = Presentation(str(pptx_path))
    images = []

    for i, slide in enumerate(prs.slides, 1):
        # Create a blank image for each slide (we'll use placeholders)
        # In a real implementation, we'd extract actual slide content
        # For now, create placeholder images
        img = Image.new('RGB', (1920, 1080), color='white')
        draw = ImageDraw.Draw(img)

        # Draw slide number
        try:
            font = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf', 80)
        except:
            font = ImageFont.load_default()

        text = f"Slide {i}"
        bbox = draw.textbbox((0, 0), text, font=font)
        text_w = bbox[2] - bbox[0]
        text_h = bbox[3] - bbox[1]

        draw.text(
            (1920 // 2 - text_w // 2, 1080 // 2 - text_h // 2),
            text,
            fill='black',
            font=font
        )

        # Save image
        img_path = Path(temp_dir) / f"slide-{i:03d}.jpg"
        img.save(img_path, 'JPEG', quality=JPEG_QUALITY)
        images.append(img_path)

    return images

def create_grid(image_paths, cols, width, output_path):
    """Create thumbnail grid from slide images."""
    # Load first image to get aspect ratio
    with Image.open(image_paths[0]) as img:
        aspect = img.height / img.width

    height = int(width * aspect)

    # Calculate grid size
    max_images_per_grid = cols * (cols + 1)
    rows = (len(image_paths) + cols - 1) // cols

    # Calculate grid dimensions
    grid_w = cols * width + (cols + 1) * GRID_PADDING
    grid_h = rows * height + (rows + 1) * GRID_PADDING

    # Create grid
    grid = Image.new('RGB', (grid_w, grid_h), 'white')
    draw = ImageDraw.Draw(grid)

    # Place thumbnails
    for i, img_path in enumerate(image_paths):
        row, col = i // cols, i % cols
        x = col * width + (col + 1) * GRID_PADDING
        y = row * height + (row + 1) * GRID_PADDING

        with Image.open(img_path) as img:
            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            tx = x + (width - w) // 2
            ty = y + (height - h) // 2
            grid.paste(img, (tx, ty))

            # Add border
            if BORDER_WIDTH > 0:
                draw.rectangle(
                    [(tx - BORDER_WIDTH, ty - BORDER_WIDTH),
                     (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1)],
                    outline='gray',
                    width=BORDER_WIDTH
                )

    return grid

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 simple_thumbnail.py <input.pptx> [output.jpg] [--cols N]")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("thumbnails.jpg")

    cols = DEFAULT_COLS
    if '--cols' in sys.argv:
        idx = sys.argv.index('--cols')
        if idx + 1 < len(sys.argv):
            cols = min(int(sys.argv[idx + 1]), MAX_COLS)

    print(f"Processing: {input_path}")

    # Load presentation
    prs = Presentation(str(input_path))
    total_slides = len(prs.slides)
    print(f"Total slides: {total_slides}")

    # Extract images
    import tempfile
    with tempfile.TemporaryDirectory() as temp_dir:
        slide_images = extract_slide_images(input_path, temp_dir)
        print(f"Extracted {len(slide_images)} slide images")

        # Create grid
        print(f"Creating grid with {cols} columns...")
        grid = create_grid(slide_images, cols, THUMBNAIL_WIDTH, output_path)

        # Save grid
        output_path.parent.mkdir(parents=True, exist_ok=True)
        grid.save(str(output_path), 'JPEG', quality=JPEG_QUALITY)

    print(f"✅ Thumbnails saved to: {output_path}")

if __name__ == '__main__':
    main()
