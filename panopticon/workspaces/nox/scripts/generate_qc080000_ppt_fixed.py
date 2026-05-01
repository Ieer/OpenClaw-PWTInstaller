# generate_qc080000_ppt.py
# 需求: python3, pip install python-pptx pillow
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# 幻灯片内容：标题 / 要点 / 讲师台词（Notes）
slide_contents = [
    {
        "title": "QC080000 有害物质过程管理（HSPM）培训",
        "bullets": [
            "适用对象：一线生产员工（装配/焊接/仓库/检验）",
            "厂名：______    日期：______    讲师：______"
        ],
        "notes": "大家好，欢迎参加今天的 QC080000 有害物质过程管理培训。我是今天的讲师______。本次培训面向装配线、焊接、仓库和检验等一线同事，目的是让大家在日常工作中识别并控制有害物质，保证产品合规并保护自己与环境的安全。我们先看培训安排。"
    },
    {
        "title": "培训目标与安排",
        "bullets": [
            "培训目标：了解目的、识别受限物质、掌握关键控制点、正确处理异常、完成测验",
            "安排：简介 → 受限物质 → 来料与工序控制 → 异常处理 → 实操演示 → 测验"
        ],
        "notes": "今天我们要达成四个目标：一，知道 QC080000 为什么重要；二，能认出常见受限物质；三，懂得材料与工序的关键控制动作；四，能按流程处理可疑/不合格物料。课程结束有一个简短测验，请认真听并参与实操。"
    },
    {
        "title": "什么是 QC080000？",
        "bullets": [
            "企业级有害物质过程管理体系（IECQ HSPM）",
            "目标：通过流程控制，确保产品受限物质符合客户与法规要求",
            "与我们的工作关系：材料、工具、工序、记录、隔离"
        ],
        "notes": "简单说，QC080000 要求我们在\'从进料到出货'的过程中，把受限物质控制好。它不只是质量部的事，而是每个人的责任：仓库要核材料，线员要用对材料，检验要做记录，遇到问题要隔离上报。"
    },
    {
        "title": "为什么要控制有害物质？",
        "bullets": [
            "保护你我健康与环境",
            "满足客户与法规，避免退货与罚款",
            "维护公司声誉与订单稳定"
        ],
        "notes": "若产品含有受限物质被客户发现，会退货、罚款，甚至丢订单。更重要的是，有些物质对人体有害，长期接触会影响健康。做好控制，就是保护自己、保护家人，也保护公司工作机会。"
    },
    {
        "title": "常见受限/有害物质 — 简单列举",
        "bullets": [
            "铅 (Pb) —— 来自旧式焊料、某些零件",
            "汞 (Hg)、镉 (Cd)、六价铬 (Cr6+)",
            "多溴联苯（PBB）/ 多溴二苯醚（PBDE）",
            "邻苯二甲酸酯（塑料增塑剂）"
        ],
        "notes": "这些是我们最常碰到需要注意的物质。你不必记住化学式，但要知道这些物质可能来自焊料、塑件、涂层、电池或添加剂。工作时要只使用\'批准'的材料和工具，别随意替换。"
    },
    {
        "title": "这些物质如何进入产品",
        "bullets": [
            "来料（芯片、外壳、线束、焊料）",
            "临时替代材料或来路不明材料",
            "工具/容器污染与交叉混料",
            "仓库与线边管理不善"
        ],
        "notes": "常见情形包括：仓库放行错误材料、线边有人临时用外面买的清洗剂或焊料、或者工具上有残留交叉污染。发生这些事情，受限物质就可能进入产品。"
    },
    {
        "title": "一线员工的五项关键职责",
        "bullets": [
            "只使用批准的物料、工具与耗材",
            "检查物料标签、批号与 CoC（合格证）",
            "正确存放与标识（按颜色/位置）",
            "发现可疑/不合格物料立即隔离并报告",
            "按规定填写和保存记录，保持可追溯性"
        ],
        "notes": "这五条是我们每天必须做到的。特别强调：见到不完整标签或CoC缺失，不要继续使用；把它隔离并通知质检。不要自行做决定或继续使用\'试一下'。"
    },
    {
        "title": "来料控制要点 — 仓库与收料",
        "bullets": [
            "供应商与物料需在\'批准清单'中",
            "入库核对：型号/数量/批号/CoC",
            "打标签、扫描条码、放到指定位置",
            "缺少文件或外观异常要隔离并报告"
        ],
        "notes": "仓库收到料要和采购单逐项核对，确认 CoC 与标签是否齐全。若发现标签缺批号或外包装破损，要贴\'隔离'标签并通知质检和相关负责人。"
    },
    {
        "title": "来料检验与线边投料",
        "bullets": [
            "检验内容：型号、批次、外观、标签、CoC 检查",
            "抽样检测按 SOP 执行（检验员负责）",
            "线边只放批准放行的料（先不要混放）"
        ],
        "notes": "检验员按标准抽样并记录结果。线长与线员不要把来路不明的料放在线边，也不要用口头确认代替书面 CoC。若不确定就隔离。"
    },
    {
        "title": "生产过程控制 — 关键点",
        "bullets": [
            "使用批准工艺、焊料与助焊膏",
            "避免混用（如铅/无铅分区、颜色标识工具）",
            "工具与容器清洁、避免交叉污染",
            "设备维护与校准按计划执行"
        ],
        "notes": "在无铅区不能用含铅焊料，工具要颜色或标签区分。操作工具或助焊膏后要及时清洁，避免把一种物质带到另一条线或另一批料上。"
    },
    {
        "title": "标识与隔离 — 现场规则",
        "bullets": [
            "红色标签 = 不合格 / 禁用",
            "黄色标签 = 隔离 / 待检",
            "绿色标签 = 合格 / 可用",
            "隔离区要整洁并有明显标识"
        ],
        "notes": "工厂用颜色区分状态。看到红色标签不要动、不拿、不用；看到黄色标签说明需要检验或等待处理。把隔离物料放在规定位置，不要随意移动。"
    },
    {
        "title": "不合格/可疑物料处理流程",
        "bullets": [
            "识别 → 隔离并贴标签 → 填写不合格单 → 通知质检/线长 → 等待处理决定（退货/返工/报废）"
        ],
        "notes": "一旦发现可疑物料，第一步是隔离并贴标签；第二步填写不合格单，记录是谁、何时、什么问题；第三步通知质检并等待处理。不要擅自处理或继续使用。"
    },
    {
        "title": "记录与可追溯性",
        "bullets": [
            "需要记录的关键项：批次号、供应商、生产日期、操作员、检测结果",
            "目的：追溯问题来源、支持客户调查、减少范围影响"
        ],
        "notes": "记录看似\'多此一举'，但如果客户发现问题，我们需要通过记录迅速定位问题批次并采取补救。务必填写批号、时间和你的工号，签名或扫码留痕。"
    },
    {
        "title": "常见违规行为与后果",
        "bullets": [
            "私自使用替代材料或来路清洗剂",
            "未核对标签放行来料",
            "混用含铅/无铅焊料",
            "不填写不合格单或擅自丢弃隔离物料"
        ],
        "notes": "这些违规行为会导致客户退货、罚款或停单。有的还会危害个人健康。请不要为了图省事而做违规操作，遵守流程对每个人都有好处。"
    },
    {
        "title": "个人防护与安全注意",
        "bullets": [
            "根据岗位佩戴 PPE（手套、口罩、护目、工作服）",
            "正确处置溶剂与化学品，遵守废弃物处理规则",
            "若接触或溅到皮肤/眼睛，立刻用清水冲洗并报告"
        ],
        "notes": "除了保护产品合规，我们也要保护自己。做清洗或接触化学品时请戴手套与护目镜。如发生溅到皮肤或眼睛，要立即按应急流程处理并报告。"
    },
    {
        "title": "线边日常自检清单",
        "bullets": [
            "物料标签完整（型号/批号/供应商） □",
            "使用批准焊料/助焊膏 □",
            "线边无红色/黄色隔离物料（或有则说明位置） □",
            "关键记录已填写（批号/时间/操作员） □"
        ],
        "notes": "上班后、换班或停线前做一次自检，只需几分钟。把清单贴在工位，线长在班前交接时确认签字。这样可以及早发现问题。"
    },
    {
        "title": "现场演示：来料标签检查",
        "bullets": [
            "找到标签 → 核对型号与批号 → 检查 CoC → 放入合格区或隔离"
        ],
        "notes": "现在我演示一次：这是一卷元件，先看包装上的标签，找型号、批号和供应商信息；再看随货的 CoC；若信息齐全，放到绿色合格区并扫码；若缺失或与订单不符，就贴黄色/红色隔离标签并填写不合格单。接下来每组请上来练习一次。"
    },
    {
        "title": "现场演示：贴隔离标签与填写不合格单",
        "bullets": [
            "隔离标签字段：物料编码 / 数量 / 发现人 / 日期 / 原因",
            "不合格单核心字段：单号 / 发现时间 / 发现人 / 物料编码 / 描述 / 建议处置"
        ],
        "notes": "我现在示范如何贴隔离标签并填写不合格单。标签要贴在明显位置；不合格单要完整写明原因与数量。练习时请按模板填写——每组演练 5 分钟，我会巡视指导。"
    },
    {
        "title": "小组练习说明",
        "bullets": [
            "每组给一套来料样本（合格/缺失 CoC/外观异常）",
            "判断是否放行、隔离或退回，并填写隔离单/不合格单",
            "每组演示并说明理由（2–3 分钟）"
        ],
        "notes": "现在分小组练习。我发每组一套样本，请按流程判断并填写表单。每组演示后我会点评正确或错误之处。时间不多，请抓紧完成。"
    },
    {
        "title": "测验说明与答题",
        "bullets": [
            "测验形式：10 题（选择/判断）——答对 8 题及格（80%）",
            "时间：5–8 分钟",
            "交卷与签名：通过者记录入培训档案"
        ],
        "notes": "下面是短测验，内容都是今天讲过的要点。答题时间 5–8 分钟，答对 8/10 算通过。答卷收回后我们会录入培训档案。"
    },
    {
        "title": "测验题示例（供打印）",
        "bullets": [
            "（示例题略 — 实际文件中会打印完整 10 题）"
        ],
        "notes": "参考测验题请参照培训文本（10 题）。运行脚本后我也会在 Notes 中包含完整题目与答案。"
    },
    {
        "title": "测验答案与讲解",
        "bullets": [
            "答案：1:C 2:错 3:B 4:C 5:对 6:A 7:B 8:错 9:B 10:对"
        ],
        "notes": "答案如屏幕所示。逐题讲解一下：第1题看到标签缺批号要隔离并通知质检……若有人不通过，我们会安排补训并复测。"
    },
    {
        "title": "关键记忆点 — 三条必须做",
        "bullets": [
            "只用批准材料与工具（查标签、看 CoC）",
            "发现异常立即隔离并通知质检（贴标签、填单）",
            "认真填写记录，保持可追溯性（批号/操作员/时间）"
        ],
        "notes": "把这三条记在心里：只用批准材料、异常隔离上报、认真记录。这就是我们日常能做到、也能防大问题的三件事。"
    },
    {
        "title": "常用表单与标签模板文本（供打印）",
        "bullets": [
            "隔离标签字段：物料编码/描述/数量/发现日期/发现人/原因/隔离位置",
            "不合格单字段：单号/发现时间/发现人/物料编码/描述/数量/建议处置/处理结果",
            "线边自检表字段：日期/班次/物料编码/是否使用批准焊料/隔离物料/检查人签名"
        ],
        "notes": "这些是我们要用的标准字段。培训结束后我会把电子版模板发给你们，大家可以直接打印到不干胶或表单纸上。"
    },
    {
        "title": "培训总结与行动要求",
        "bullets": [
            "关键回顾：三条必须做",
            "现场要求：线长每天签字自检、仓库严格核验、不合格 24 小时内反馈",
            "联系人：质检 ______ 工程 ______ 仓库 ______"
        ],
        "notes": "总结一下：掌握三条必须做、使用表单并做好记录。从明天开始，线长要每天做自检并签字。若发现不合格，24 小时内要有初步处理反馈。培训负责人会把联系方式贴到公告栏。"
    },
    {
        "title": "答疑与签到",
        "bullets": [
            "提问时间（请就地提问）",
            "培训签到（请在签到表签名并填写工号）"
        ],
        "notes": "现在是答疑时间，大家有什么问题可以提。答疑后请在培训记录表上签名并写上你的工号，测验通过的会记入你培训档案。谢谢大家配合。"
    }
]

def add_logo(slide, logo_path, left=Inches(11.0), top=Inches(0.2), height=Inches(0.6)):
    try:
        slide.shapes.add_picture(logo_path, left, top, height=height)
    except Exception as e:
        # 若插入失败则添加文本占位
        txBox = slide.shapes.add_textbox(left, top, Inches(2.0), height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "LOGO"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(90, 90, 90)

def add_slide(prs, title, bullets, notes, logo_path=None):
    # 使用 Title and Content 布局
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    # 标题
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    # 内容
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    # Notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes
    # logo 或占位
    if logo_path:
        add_logo(slide, logo_path)
    else:
        left = Inches(11.0)
        top = Inches(0.2)
        txBox = slide.shapes.add_textbox(left, top, Inches(2.0), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "LOGO 占位"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(120, 120, 120)

def main():
    parser = argparse.ArgumentParser(description="生成 QC080000 培训 PPT (16:9)，含 Presenter Notes")
    parser.add_argument("--logo", help="logo 图片路径（可选 PNG/JPG）", default=None)
    parser.add_argument("--company", help="公司名称（可选，会写入封面）", default=None)
    parser.add_argument("--date", help="培训日期（可选，会写入封面）", default=None)
    parser.add_argument("--presenter", help="讲师姓名（可选，会写入封面）", default=None)
    parser.add_argument("--output", help="输出 PPT 文件名（可选）", default="QC080000_HSPM_Training.pptx")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # 如果提供封面信息，替换第一张幻灯片 bullets 中占位
    if args.company or args.date or args.presenter:
        first = slide_contents[0]
        company_text = args.company if args.company else "______"
        date_text = args.date if args.date else "______"
        presenter_text = args.presenter if args.presenter else "______"
        first["bullets"][1] = f"厂名：{company_text}    日期：{date_text}    讲师：{presenter_text}"

    # 添加所有幻灯片
    for s in slide_contents:
        add_slide(prs, s["title"], s["bullets"], s["notes"], logo_path=args.logo)

    prs.save(args.output)
    print(f"已生成 PPT 文件: {args.output}")

if __name__ == "__main__":
    main()
