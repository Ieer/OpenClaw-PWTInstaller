#!/usr/bin/env python3
"""生成下周工作汇报 PPTX — 纯 python-pptx 可编辑版本"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import shutil, os

prs = Presentation()
prs.slide_width = Emu(12192000)  # 16:9
prs.slide_height = Emu(6858000)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Colors ──
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)
WHITE        = RGBColor(0xF1, 0xF5, 0xF9)
GRAY         = RGBColor(0x94, 0xA3, 0xB8)
DIM_GRAY     = RGBColor(0x64, 0x74, 0x8B)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LIGHT   = RGBColor(0x60, 0xA5, 0xFA)
GREEN        = RGBColor(0x22, 0xC5, 0x5E)
GREEN_LIGHT  = RGBColor(0x4A, 0xDE, 0x80)
RED          = RGBColor(0xEF, 0x44, 0x44)
RED_LIGHT    = RGBColor(0xFC, 0xA5, 0xA5)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
BORDER       = RGBColor(0x33, 0x43, 0x55)

# Background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = BG_DARK

def add_rect(left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(left, top, width, height, text, font_size=10, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_multiline_textbox(left, top, width, height, lines, default_size=10, default_color=WHITE):
    """lines: list of (text, font_size, bold, color)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line[0]
        size = line[1] if len(line) > 1 else default_size
        bold = line[2] if len(line) > 2 else False
        color = line[3] if len(line) > 3 else default_color
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(2)
    return txBox

# ═══════════════════ Top Decoration Line ═══════════════════
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(457200), 0, Emu(11289600), Emu(4572))
deco.fill.solid()
deco.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6)
deco.line.fill.background()

# ═══════════════════ Title Area ═══════════════════
# Main title
add_textbox(Emu(457200), Emu(274320), Emu(7000000), Emu(365760),
            "📋  WCQ 数位工具 下周工作汇报", 24, True, WHITE)
# Subtitle
add_textbox(Emu(457200), Emu(640080), Emu(9000000), Emu(201168),
            "汇报人: Ieer Qin  |  团队: Ieer + Bella Wu  |  周期: 2026.05.18 ~ 05.22", 
            11, False, GRAY)

# Badges
badge_data = [
    (Emu(8000000), Emu(310000), "🔴 P0 冲刺 2 项", RGBColor(0x1E, 0x3A, 0x5A), BLUE_LIGHT),
    (Emu(9000000), Emu(310000), "🟢 Bella 交接启动", RGBColor(0x1A, 0x3A, 0x2A), GREEN_LIGHT),
    (Emu(10000000), Emu(310000), "📖 Week3 培训", RGBColor(0x2A, 0x1A, 0x3A), PURPLE_LIGHT),
]
for x, y, text, bg_c, fg_c in badge_data:
    shape = add_rect(x, y, Emu(850000), Emu(280000), bg_c, radius=None)
    add_textbox(x + Emu(50000), y + Emu(20000), Emu(800000), Emu(240000),
                text, 8, True, fg_c, PP_ALIGN.CENTER)

# ═══════════════════ Summary Cards ═══════════════════
cards = [
    ("⚡", "2", "P0 任务", "专业训上架 + 报告Approve", RED),
    ("🛠", "40H", "PCBA钢板开发", "整合测试目标 50%", AMBER),
    ("👤", "3项", "Bella 交接", "Delay/SAP/评审跟催", GREEN),
    ("📖", "Week3", "培训进度", "项目管理 + 案例说明", BLUE),
]

card_w = Emu(2646000)
card_h = Emu(686000)
card_gap = Emu(120000)
start_x = Emu(457200)
card_y = Emu(975000)

for i, (icon, num, label, sub, accent) in enumerate(cards):
    cx = start_x + i * (card_w + card_gap)
    
    # Card background
    card_shape = add_rect(cx, card_y, card_w, card_h, CARD_BG, BORDER)
    
    # Icon circle
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Emu(100000), card_y + Emu(80000), Emu(280000), Emu(280000))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = RGBColor(0x2A, 0x1A, 0x3A)
    icon_shape.line.fill.background()
    add_textbox(cx + Emu(100000), card_y + Emu(100000), Emu(280000), Emu(240000),
                icon, 14, False, accent, PP_ALIGN.CENTER)
    
    # Number
    add_textbox(cx + Emu(430000), card_y + Emu(60000), Emu(2000000), Emu(280000),
                num, 22, True, WHITE)
    # Label
    add_textbox(cx + Emu(430000), card_y + Emu(320000), Emu(2000000), Emu(160000),
                label, 11, False, GRAY)
    # Sub
    add_textbox(cx + Emu(100000), card_y + Emu(480000), Emu(2400000), Emu(160000),
                sub, 9, False, DIM_GRAY)

# ═══════════════════ Two Column Content ═══════════════════
col_top = Emu(1820000)
col_h = Emu(4600000)

# Left Column: 本周重点任务
left_x = Emu(457200)
col_w = Emu(5320000)

# Column header
add_rect(left_x, col_top, col_w, Emu(32000), BLUE, BLUE)
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, col_top + Emu(40000), col_w, Emu(12000))
divider.fill.solid()
divider.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
divider.line.fill.background()

add_textbox(left_x + Emu(10000), col_top - Emu(60000), Emu(4000000), Emu(280000),
            "🎯  本周重点任务", 14, True, BLUE_LIGHT)

# Task cards - left column
tasks_left = [
    ("🔴", RED, "各部门专业训自动收集整理", 
     "最后 20% 校对 → 上架部署，本周闭环", "P0 · 开发中", RGBColor(0x3A, 0x1A, 0x1A)),
    ("🔴", RED, "数位交响乐-制造类 Columbus 报告",
     "5/20 前完成内部 Approve · 5/29 发表", "P0 · 报告迭代中", RGBColor(0x3A, 0x1A, 0x1A)),
    ("🟡", AMBER, "PCBA 钢板自动绑定",
     "各模块代码完成 · 整合测试目标过半", "P1 · 整合测试", RGBColor(0x3A, 0x2A, 0x1A)),
    ("🟢", GREEN, "Bella 三件事交接",
     "Delay追踪 / SAP进度 / 评审跟催 — 周五独立", "交接中", RGBColor(0x1A, 0x3A, 0x1A)),
    ("⚪", DIM_GRAY, "NPI RFQ Agent 确认",
     "确认资料到位时间，到则启动开发", "待确认", RGBColor(0x2A, 0x2A, 0x2A)),
]

task_start_y = col_top + Emu(60000)
task_h = Emu(760000)
task_gap_y = Emu(80000)

for i, (dot, dot_color, title, desc, tag, tag_bg) in enumerate(tasks_left):
    ty = task_start_y + i * (task_h + task_gap_y)
    
    # Task card
    tc = add_rect(left_x, ty, col_w, task_h, RGBColor(0x14, 0x1E, 0x33), BORDER)
    
    # Dot
    dot_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_x + Emu(100000), ty + Emu(80000), Emu(50000), Emu(50000))
    dot_shape.fill.solid()
    dot_shape.fill.fore_color.rgb = dot_color
    dot_shape.line.fill.background()
    
    # Title
    add_textbox(left_x + Emu(200000), ty + Emu(50000), Emu(4800000), Emu(220000),
                title, 12, True, WHITE)
    
    # Description
    add_textbox(left_x + Emu(200000), ty + Emu(280000), Emu(4800000), Emu(200000),
                desc, 9, False, GRAY)
    
    # Tag
    tag_shape = add_rect(left_x + Emu(200000), ty + Emu(500000), Emu(1200000), Emu(180000), tag_bg)
    add_textbox(left_x + Emu(220000), ty + Emu(510000), Emu(1160000), Emu(160000),
                tag, 8, True, dot_color, PP_ALIGN.CENTER)

# Right Column: 下周迭代 + 风险
right_x = Emu(6240000)
rcol_w = Emu(5480000)

add_textbox(right_x, col_top - Emu(60000), Emu(4000000), Emu(280000),
            "🔄  下周迭代计划", 14, True, GREEN_LIGHT)
add_rect(right_x, col_top, rcol_w, Emu(32000), GREEN, GREEN)

# Task cards - right column (迭代)
tasks_right = [
    ("🟢", GREEN, "SAP 升级 RPA Phase 1", 
     "目标 15 支 · Bella 接手进度追踪", "等各部门排程", RGBColor(0x1A, 0x3A, 0x1A)),
    ("🟢", GREEN, "AI Agent 课程通知准备",
     "6/1 发送 MFG. 课程通知 · 本周备料", "准备中", RGBColor(0x1A, 0x3A, 0x1A)),
    ("🟡", AMBER, "月KPI报表 + PMO周报合并设计",
     "一次开发两次复用，评估可行性", "P2 · 待启动", RGBColor(0x3A, 0x2A, 0x1A)),
]

for i, (dot, dot_color, title, desc, tag, tag_bg) in enumerate(tasks_right):
    ty = task_start_y + i * (task_h + task_gap_y)
    
    tc = add_rect(right_x, ty, rcol_w, task_h, RGBColor(0x14, 0x1E, 0x33), BORDER)
    
    dot_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Emu(100000), ty + Emu(80000), Emu(50000), Emu(50000))
    dot_shape.fill.solid()
    dot_shape.fill.fore_color.rgb = dot_color
    dot_shape.line.fill.background()
    
    add_textbox(right_x + Emu(200000), ty + Emu(50000), Emu(5000000), Emu(220000),
                title, 12, True, WHITE)
    add_textbox(right_x + Emu(200000), ty + Emu(280000), Emu(5000000), Emu(200000),
                desc, 9, False, GRAY)
    
    tag_shape = add_rect(right_x + Emu(200000), ty + Emu(500000), Emu(1200000), Emu(180000), tag_bg)
    add_textbox(right_x + Emu(220000), ty + Emu(510000), Emu(1160000), Emu(160000),
                tag, 8, True, dot_color, PP_ALIGN.CENTER)

# Risk section (右栏下方)
risk_y = col_top + Emu(800000) + 3 * (task_h + task_gap_y)
add_textbox(right_x, risk_y - Emu(60000), Emu(4000000), Emu(280000),
            "⚠️  风险与关注", 14, True, RED_LIGHT)

risks = [
    ("🔴", RED, "吸嘴比对暂停 → 黑客松交付缺口", "原4人黑客松团队解散，PMO独力开发周期长"),
    ("🟡", AMBER, "Bella 培训进度", "Week3 Open，减压预期依赖培训完成"),
    ("🟡", AMBER, "SAP 70件 + WiMES 17件 潜在撞车", "若同期到达将打满 Ieer 工时"),
]

for i, (dot, dot_color, title, desc) in enumerate(risks):
    ty = risk_y + Emu(60000) + i * Emu(260000)
    tc = add_rect(right_x, ty, rcol_w, Emu(230000), RGBColor(0x18, 0x18, 0x18), RGBColor(0x3A, 0x1A, 0x1A))
    
    dot_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Emu(100000), ty + Emu(30000), Emu(40000), Emu(40000))
    dot_shape.fill.solid()
    if "🔴" in dot:
        dot_shape.fill.fore_color.rgb = RED
    else:
        dot_shape.fill.fore_color.rgb = AMBER
    dot_shape.line.fill.background()
    
    add_textbox(right_x + Emu(180000), ty + Emu(15000), Emu(5000000), Emu(160000),
                title, 11, True, WHITE)
    add_textbox(right_x + Emu(180000), ty + Emu(140000), Emu(5000000), Emu(80000),
                desc, 8, False, GRAY)

# ═══════════════════ Footer ═══════════════════
footer_y = Emu(6500000)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(457200), footer_y, Emu(11289600), Emu(6000))
line.fill.solid()
line.fill.fore_color.rgb = BORDER
line.line.fill.background()

add_textbox(Emu(457200), footer_y + Emu(30000), Emu(6000000), Emu(200000),
            "下周 P0: 2/2 ✅   |   PCBA 整合: 50% 🟡   |   Bella: 3/3 ✅", 
            9, False, GRAY)
add_textbox(Emu(8500000), footer_y + Emu(30000), Emu(3400000), Emu(200000),
            "生成: 2026-05-16  |  WCQ 数位转型专案室",
            9, False, DIM_GRAY, PP_ALIGN.RIGHT)

# ── Save ──
out_path = "/mnt/usb/WCQ_下周工作汇报_2026-05-18.pptx"
prs.save(out_path)
print(f"✅ PPTX saved: {out_path}")

shutil.copy(out_path, "/home/node/.openclaw/workspace/exports/WCQ_下周工作汇报_2026-05-18.pptx")
print("✅ Backup to exports/ OK")
