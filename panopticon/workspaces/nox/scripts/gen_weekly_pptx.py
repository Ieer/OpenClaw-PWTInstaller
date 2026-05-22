#!/usr/bin/env python3
"""Take screenshot of HTML report and create PPTX"""
import subprocess, sys

html_path = "/home/node/.openclaw/workspace/artifacts/下周工作计划-20260518-0522/weekly-report.html"
png_dir = "/home/node/.openclaw/workspace/artifacts/下周工作计划-20260518-0522"
png_path = f"{png_dir}/weekly-report.png"

# Step 1: Screenshot using browser tool (headless approach)
print("Step 1: Taking screenshot...")

# Try Playwright first
try:
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(f"file://{html_path}")
        page.screenshot(path=png_path, full_page=False)
        browser.close()
    print("Playwright screenshot OK")
except ImportError:
    print("Playwright not available, trying alternative...")
    # Use wkhtmltoimage if available
    result = subprocess.run(
        ["wkhtmltoimage", "--width", "1920", "--height", "1080", html_path, png_path],
        capture_output=True, text=True, timeout=30
    )
    if result.returncode == 0:
        print("wkhtmltoimage OK")
    else:
        print(f"wkhtmltoimage failed: {result.stderr}")
        sys.exit(1)

# Verify
import os
if os.path.exists(png_path):
    size = os.path.getsize(png_path)
    print(f"PNG created: {size} bytes")
else:
    print("PNG not created!")
    sys.exit(1)

# Step 2: Create PPTX
print("\nStep 2: Creating PPTX...")
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Emu(12192000)  # 16:9
prs.slide_height = Emu(6858000)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# Background color (dark)
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

# Add screenshot image
img = slide.shapes.add_picture(png_path, 0, 0, prs.slide_width, prs.slide_height)

out_path = "/mnt/usb/WCQ_下周工作汇报_2026-05-18.pptx"
prs.save(out_path)
print(f"PPTX saved: {out_path}")

# Also copy to exports
import shutil
shutil.copy(out_path, "/home/node/.openclaw/workspace/exports/WCQ_下周工作汇报_2026-05-18.pptx")
print("Backup to exports/ OK")
