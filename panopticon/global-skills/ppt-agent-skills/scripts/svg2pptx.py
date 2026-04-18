#!/usr/bin/env python3
"""SVG to PPTX -- 将 SVG 元素解析为原生 OOXML 形状

支持: rect, text+tspan, circle, ellipse, line, path, image(data URI + file)
      linearGradient, radialGradient, transform(translate/scale/matrix)
      group opacity 传递, 首屏 rect 自动设为幻灯片背景

用法:
  python3 scripts/svg2pptx.py <svg_dir_or_file> -o output.pptx
"""

import argparse
import base64
import collections.abc
import io
import json
import math
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree, html as lhtml
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from speech_script import load_speech_page_entries, write_slide_speaker_notes

# -------------------------------------------------------------------
# 常量
# -------------------------------------------------------------------
SVG_NS = 'http://www.w3.org/2000/svg'
XLINK_NS = 'http://www.w3.org/1999/xlink'
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}
EMU_PX = 9525
SLIDE_W = 12192000
SLIDE_H = 6858000
SOURCE_REPORT_FILE_NAME = 'svg-export-report.json'
SEMANTIC_FILE_SUFFIX = '.semantic.json'
CARD_ROOT_XPATH = '//*[@data-card-id]'
CHART_GROUP_NAME_PREFIX = 'ChartGroup:'
NATIVE_CHART_NAME_PREFIX = 'NativeChart:'
MANAGED_BLOCK_NAME_PREFIX = 'Block:'
BLOCK_SLOT_NAME_PREFIXES = ('BlockSlot:', 'ManagedBlock:', 'ReplaceBlock:')
BLOCK_NAME_SUFFIX_MARKER = ':block='
HTML_SEMANTIC_XPATH = (
    '//h1[not(ancestor::li)]|//h2[not(ancestor::li)]|//h3[not(ancestor::li)]|'
    '//h4[not(ancestor::li)]|//h5[not(ancestor::li)]|//h6[not(ancestor::li)]|'
    '//p[not(ancestor::li or ancestor::td or ancestor::th)]|//li|'
    '//caption|//figcaption|//th|//td'
)

# CSS 完整命名颜色表（常用子集）
CSS_COLORS = {
    'aliceblue': 'f0f8ff', 'antiquewhite': 'faebd7', 'aqua': '00ffff',
    'aquamarine': '7fffd4', 'azure': 'f0ffff', 'beige': 'f5f5dc',
    'bisque': 'ffe4c4', 'black': '000000', 'blanchedalmond': 'ffebcd',
    'blue': '0000ff', 'blueviolet': '8a2be2', 'brown': 'a52a2a',
    'burlywood': 'deb887', 'cadetblue': '5f9ea0', 'chartreuse': '7fff00',
    'chocolate': 'd2691e', 'coral': 'ff7f50', 'cornflowerblue': '6495ed',
    'cornsilk': 'fff8dc', 'crimson': 'dc143c', 'cyan': '00ffff',
    'darkblue': '00008b', 'darkcyan': '008b8b', 'darkgoldenrod': 'b8860b',
    'darkgray': 'a9a9a9', 'darkgreen': '006400', 'darkgrey': 'a9a9a9',
    'darkkhaki': 'bdb76b', 'darkmagenta': '8b008b', 'darkolivegreen': '556b2f',
    'darkorange': 'ff8c00', 'darkorchid': '9932cc', 'darkred': '8b0000',
    'darksalmon': 'e9967a', 'darkseagreen': '8fbc8f', 'darkslateblue': '483d8b',
    'darkslategray': '2f4f4f', 'darkturquoise': '00ced1', 'darkviolet': '9400d3',
    'deeppink': 'ff1493', 'deepskyblue': '00bfff', 'dimgray': '696969',
    'dodgerblue': '1e90ff', 'firebrick': 'b22222', 'floralwhite': 'fffaf0',
    'forestgreen': '228b22', 'fuchsia': 'ff00ff', 'gainsboro': 'dcdcdc',
    'ghostwhite': 'f8f8ff', 'gold': 'ffd700', 'goldenrod': 'daa520',
    'gray': '808080', 'green': '008000', 'greenyellow': 'adff2f',
    'grey': '808080', 'honeydew': 'f0fff0', 'hotpink': 'ff69b4',
    'indianred': 'cd5c5c', 'indigo': '4b0082', 'ivory': 'fffff0',
    'khaki': 'f0e68c', 'lavender': 'e6e6fa', 'lawngreen': '7cfc00',
    'lemonchiffon': 'fffacd', 'lightblue': 'add8e6', 'lightcoral': 'f08080',
    'lightcyan': 'e0ffff', 'lightgoldenrodyellow': 'fafad2', 'lightgray': 'd3d3d3',
    'lightgreen': '90ee90', 'lightpink': 'ffb6c1', 'lightsalmon': 'ffa07a',
    'lightseagreen': '20b2aa', 'lightskyblue': '87cefa', 'lightslategray': '778899',
    'lightsteelblue': 'b0c4de', 'lightyellow': 'ffffe0', 'lime': '00ff00',
    'limegreen': '32cd32', 'linen': 'faf0e6', 'magenta': 'ff00ff',
    'maroon': '800000', 'mediumaquamarine': '66cdaa', 'mediumblue': '0000cd',
    'mediumorchid': 'ba55d3', 'mediumpurple': '9370db', 'mediumseagreen': '3cb371',
    'mediumslateblue': '7b68ee', 'mediumspringgreen': '00fa9a',
    'mediumturquoise': '48d1cc', 'mediumvioletred': 'c71585', 'midnightblue': '191970',
    'mintcream': 'f5fffa', 'mistyrose': 'ffe4e1', 'moccasin': 'ffe4b5',
    'navajowhite': 'ffdead', 'navy': '000080', 'oldlace': 'fdf5e6',
    'olive': '808000', 'olivedrab': '6b8e23', 'orange': 'ffa500',
    'orangered': 'ff4500', 'orchid': 'da70d6', 'palegoldenrod': 'eee8aa',
    'palegreen': '98fb98', 'paleturquoise': 'afeeee', 'palevioletred': 'db7093',
    'papayawhip': 'ffefd5', 'peachpuff': 'ffdab9', 'peru': 'cd853f',
    'pink': 'ffc0cb', 'plum': 'dda0dd', 'powderblue': 'b0e0e6',
    'purple': '800080', 'rebeccapurple': '663399', 'red': 'ff0000',
    'rosybrown': 'bc8f8f', 'royalblue': '4169e1', 'saddlebrown': '8b4513',
    'salmon': 'fa8072', 'sandybrown': 'f4a460', 'seagreen': '2e8b57',
    'seashell': 'fff5ee', 'sienna': 'a0522d', 'silver': 'c0c0c0',
    'skyblue': '87ceeb', 'slateblue': '6a5acd', 'slategray': '708090',
    'snow': 'fffafa', 'springgreen': '00ff7f', 'steelblue': '4682b4',
    'tan': 'd2b48c', 'teal': '008080', 'thistle': 'd8bfd8',
    'tomato': 'ff6347', 'turquoise': '40e0d0', 'violet': 'ee82ee',
    'wheat': 'f5deb3', 'white': 'ffffff', 'whitesmoke': 'f5f5f5',
    'yellow': 'ffff00', 'yellowgreen': '9acd32',
}

# 字体回退链
FONT_FALLBACK = {
    'PingFang SC': 'Microsoft YaHei',
    'SF Pro Display': 'Arial',
    'Helvetica Neue': 'Arial',
    'Helvetica': 'Arial',
    'system-ui': 'Microsoft YaHei',
    'sans-serif': 'Microsoft YaHei',
}


def iso_now():
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace('+00:00', 'Z')


def display_path(path, base_dir):
    resolved = Path(path).resolve()
    try:
        return str(resolved.relative_to(Path(base_dir).resolve()))
    except ValueError:
        return str(resolved)


def resolve_source_report_path(svg_input, output_path, explicit_path=None):
    candidates = []
    if explicit_path:
        candidates.append(Path(explicit_path).resolve())

    svg_input = Path(svg_input)
    output_path = Path(output_path)
    if svg_input.is_dir():
        candidates.extend([
            (svg_input.parent / SOURCE_REPORT_FILE_NAME).resolve(),
            (svg_input / SOURCE_REPORT_FILE_NAME).resolve(),
        ])
    else:
        candidates.extend([
            (svg_input.parent.parent / SOURCE_REPORT_FILE_NAME).resolve(),
            (svg_input.parent / SOURCE_REPORT_FILE_NAME).resolve(),
        ])
    candidates.append((output_path.parent / SOURCE_REPORT_FILE_NAME).resolve())

    seen = set()
    for candidate in candidates:
        key = str(candidate)
        if key in seen:
            continue
        seen.add(key)
        if candidate.exists():
            return candidate

    if explicit_path:
        return Path(explicit_path).resolve()
    return None


def load_source_report(report_path):
    if report_path is None or not report_path.exists():
        return None, [f'source_report_missing:{report_path}' if report_path else 'source_report_missing']
    try:
        payload = json.loads(report_path.read_text(encoding='utf-8'))
    except (OSError, json.JSONDecodeError) as exc:
        return None, [f'source_report_invalid:{exc}']
    if not isinstance(payload, dict):
        return None, ['source_report_invalid:root_not_object']
    return payload, []


def build_source_page_lookup(source_report):
    if not isinstance(source_report, dict):
        return {}
    pages = source_report.get('pages')
    if not isinstance(pages, list):
        return {}
    lookup = {}
    for page in pages:
        if not isinstance(page, dict):
            continue
        svg_name = page.get('source_svg')
        if isinstance(svg_name, str) and svg_name.strip():
            lookup[Path(svg_name).name] = page
            continue
        page_name = page.get('page_name')
        if isinstance(page_name, str) and page_name.strip():
            lookup[f'{page_name}.svg'] = page
    return lookup


def normalize_text(value):
    if value is None:
        return ''
    return re.sub(r'\s+', ' ', str(value).replace('\xa0', ' ')).strip()


def semantic_match_key(value):
    text = normalize_text(value)
    if not text:
        return ''
    text = re.sub(r'^[\u2022\u2023\u25E6\u25CF\u2219\-\u2013\u2014]+\s*', '', text)
    text = re.sub(r'^\d+[.)]\s*', '', text)
    return text.casefold()


def estimate_text_width(text, font_size):
    font_size = float(font_size)
    return sum((0.95 if ord(ch) > 0x2E7F else 0.6) * font_size for ch in text)


def resolve_html_source_path(svg_path, html_dir):
    if not html_dir:
        return None
    html_dir = Path(html_dir)
    if html_dir.is_file():
        return html_dir.resolve()
    candidates = [
        html_dir / f'{Path(svg_path).stem}.html',
        html_dir / f'{Path(svg_path).stem}.htm',
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    return None


def infer_block_role(card_id, node=None):
    if node is not None:
        explicit_role = node.get('data-role')
        if isinstance(explicit_role, str) and explicit_role.strip():
            return explicit_role.strip()
    if not isinstance(card_id, str) or not card_id.strip():
        return None
    parts = [part.strip() for part in card_id.split('-') if part.strip()]
    if len(parts) >= 2 and parts[1] in {'anchor', 'support', 'context'}:
        return parts[1]
    return None


def resolve_block_id_for_node(node):
    current = node
    while current is not None:
        if hasattr(current, 'get'):
            block_id = current.get('data-card-id')
            if isinstance(block_id, str) and block_id.strip():
                return block_id.strip()
        current = current.getparent()
    return None


def parse_int_attr(node, name, default=1):
    try:
        value = int(node.get(name, default))
        return value if value > 0 else default
    except (TypeError, ValueError):
        return default


def parse_html_table(table_node, table_index):
    row_nodes = [row for row in table_node.xpath('.//tr') if isinstance(getattr(row, 'tag', None), str)]
    parsed_rows = []
    occupied = set()
    max_columns = 0
    for row_index, row_node in enumerate(row_nodes):
        logical_column = 0
        parsed_cells = []
        cell_nodes = [cell for cell in row_node.xpath('./th|./td') if isinstance(getattr(cell, 'tag', None), str)]
        for cell_index, cell_node in enumerate(cell_nodes):
            while (row_index, logical_column) in occupied:
                logical_column += 1
            colspan = parse_int_attr(cell_node, 'colspan', default=1)
            rowspan = parse_int_attr(cell_node, 'rowspan', default=1)
            cell_info = {
                'row_index': row_index,
                'cell_index': cell_index,
                'column_index': logical_column,
                'colspan': colspan,
                'rowspan': rowspan,
                'tag': cell_node.tag.lower(),
                'text': normalize_text(cell_node.text_content()),
                'is_header': cell_node.tag.lower() == 'th' or bool(cell_node.xpath('ancestor::thead')),
                'match_key': semantic_match_key(cell_node.text_content()),
            }
            parsed_cells.append(cell_info)
            for row_offset in range(rowspan):
                for col_offset in range(colspan):
                    occupied.add((row_index + row_offset, logical_column + col_offset))
            logical_column += colspan
        max_columns = max(max_columns, logical_column)
        parsed_rows.append({'row_index': row_index, 'cells': parsed_cells})

    caption_nodes = table_node.xpath('./caption')
    caption = normalize_text(caption_nodes[0].text_content()) if caption_nodes else ''
    return {
        'table_id': f'table-{table_index}',
        'block_id': resolve_block_id_for_node(table_node),
        'caption': caption,
        'row_count': len(parsed_rows),
        'column_count': max_columns,
        'rows': parsed_rows,
    }


def load_html_semantics(html_path):
    if html_path is None:
        return None, ['html_semantics_missing']
    try:
        document = lhtml.fromstring(Path(html_path).read_bytes())
    except (OSError, etree.ParserError, etree.XMLSyntaxError, ValueError) as exc:
        return None, [f'html_semantics_invalid:{exc}']

    block_nodes = [node for node in document.xpath(CARD_ROOT_XPATH) if isinstance(getattr(node, 'tag', None), str)]
    blocks = []
    block_lookup = {}
    for index, block_node in enumerate(block_nodes, start=1):
        block_id = normalize_text(block_node.get('data-card-id')) or f'page-root-{index}'
        block_info = {
            'block_id': block_id,
            'card_id': normalize_text(block_node.get('data-card-id')) or None,
            'role_hint': infer_block_role(block_id, node=block_node),
            'tag': block_node.tag.lower(),
            'class_name': normalize_text(block_node.get('class')),
            'entries': [],
            'tables': [],
        }
        blocks.append(block_info)
        block_lookup[block_id] = block_info

    entries = []
    nodes = document.xpath(HTML_SEMANTIC_XPATH)
    for node in nodes:
        if not isinstance(getattr(node, 'tag', None), str):
            continue
        tag = node.tag.lower()
        text = normalize_text(node.text_content())
        if not text:
            continue
        role = {
            'p': 'paragraph',
            'li': 'list_item',
            'caption': 'caption',
            'figcaption': 'caption',
            'th': 'table_header',
            'td': 'table_cell',
        }.get(tag, 'heading' if tag.startswith('h') else 'text')
        entry = {
            'entry_id': len(entries) + 1,
            'order': len(entries) + 1,
            'tag': tag,
            'role': role,
            'text': text,
            'match_key': semantic_match_key(text),
            'source_line': getattr(node, 'sourceline', None),
            'block_id': resolve_block_id_for_node(node),
        }
        if role == 'list_item':
            list_container = None
            for ancestor in node.iterancestors():
                if isinstance(getattr(ancestor, 'tag', None), str) and ancestor.tag.lower() in ('ul', 'ol'):
                    list_container = ancestor
                    break
            list_type = list_container.tag.lower() if list_container is not None else 'ul'
            depth = max(
                0,
                sum(
                    1
                    for ancestor in node.iterancestors()
                    if isinstance(getattr(ancestor, 'tag', None), str)
                    and ancestor.tag.lower() in ('ul', 'ol')
                ) - 1,
            )
            start_at = 1
            if list_type == 'ol' and list_container is not None:
                try:
                    start_at = int(list_container.get('start', '1'))
                except (TypeError, ValueError):
                    start_at = 1
            preceding_items = sum(
                1
                for sibling in node.itersiblings(preceding=True)
                if isinstance(getattr(sibling, 'tag', None), str) and sibling.tag.lower() == 'li'
            )
            entry['list_type'] = list_type
            entry['list_depth'] = depth
            entry['list_index'] = start_at + preceding_items if list_type == 'ol' else None
        entries.append(entry)
        block_id = entry.get('block_id')
        if block_id and block_id in block_lookup:
            block_lookup[block_id]['entries'].append(entry)

    list_nodes = document.xpath('//ul|//ol')
    table_nodes = document.xpath('//table')
    tables = []
    for index, table_node in enumerate(table_nodes, start=1):
        table_info = parse_html_table(table_node, index)
        tables.append(table_info)
        block_id = table_info.get('block_id')
        if block_id and block_id in block_lookup:
            block_lookup[block_id]['tables'].append(table_info['table_id'])
    list_item_count = sum(1 for entry in entries if entry.get('role') == 'list_item')
    table_cell_count = sum(1 for entry in entries if entry.get('role') in ('table_header', 'table_cell'))
    return {
        'path': str(Path(html_path).resolve()),
        'summary': {
            'text_blocks': len(entries),
            'lists': len(list_nodes),
            'list_items': list_item_count,
            'tables': len(table_nodes),
            'table_cells': table_cell_count,
            'blocks': len(blocks),
        },
        'entries': entries,
        'blocks': blocks,
        'tables': tables,
    }, []


def load_rendered_semantics(semantic_path):
    if semantic_path is None:
        return None, ['rendered_semantics_missing']
    if not semantic_path.exists():
        return None, [f'rendered_semantics_missing:{semantic_path}']
    try:
        payload = json.loads(semantic_path.read_text(encoding='utf-8'))
    except (OSError, json.JSONDecodeError) as exc:
        return None, [f'rendered_semantics_invalid:{exc}']
    if not isinstance(payload, dict):
        return None, ['rendered_semantics_invalid:root_not_object']
    return payload, []


def resolve_semantic_sidecar_path(svg_file, source_page, source_report_path=None):
    candidates = []
    raw_path = source_page.get('semantic_path') if isinstance(source_page, dict) else None
    if isinstance(raw_path, str) and raw_path.strip():
        path_obj = Path(raw_path)
        if path_obj.is_absolute():
            candidates.append(path_obj.resolve())
        else:
            candidates.extend([
                (svg_file.parent / path_obj.name).resolve(),
                (svg_file.parent / path_obj).resolve(),
                (svg_file.parent.parent / path_obj).resolve(),
            ])
            if source_report_path is not None:
                candidates.append((source_report_path.parent / path_obj).resolve())
    candidates.append(svg_file.with_suffix(SEMANTIC_FILE_SUFFIX).resolve())

    seen = set()
    for candidate in candidates:
        key = str(candidate)
        if key in seen:
            continue
        seen.add(key)
        if candidate.exists():
            return candidate
    return None


def normalize_bbox(raw_bbox):
    if not isinstance(raw_bbox, dict):
        return None
    try:
        x = float(raw_bbox.get('x', 0))
        y = float(raw_bbox.get('y', 0))
        width = float(raw_bbox.get('width', 0))
        height = float(raw_bbox.get('height', 0))
    except (TypeError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return {
        'x': x,
        'y': y,
        'width': width,
        'height': height,
    }


def bbox_center(bbox):
    if not bbox:
        return None
    return (
        float(bbox['x']) + float(bbox['width']) / 2,
        float(bbox['y']) + float(bbox['height']) / 2,
    )


def merge_bboxes(bboxes):
    valid = [normalize_bbox(bbox) for bbox in bboxes]
    valid = [bbox for bbox in valid if bbox]
    if not valid:
        return None
    left = min(bbox['x'] for bbox in valid)
    top = min(bbox['y'] for bbox in valid)
    right = max(bbox['x'] + bbox['width'] for bbox in valid)
    bottom = max(bbox['y'] + bbox['height'] for bbox in valid)
    return normalize_bbox({
        'x': left,
        'y': top,
        'width': right - left,
        'height': bottom - top,
    })


def bbox_area(bbox):
    if not bbox:
        return 0.0
    return max(float(bbox.get('width', 0)), 0.0) * max(float(bbox.get('height', 0)), 0.0)


def parse_svg_point_pairs(raw_value):
    if not raw_value:
        return []
    coords = re.findall(r'[+-]?\d*\.?\d+', str(raw_value))
    if len(coords) < 4:
        return []
    if len(coords) % 2 == 1:
        coords = coords[:-1]
    points = []
    for index in range(0, len(coords), 2):
        try:
            points.append((float(coords[index]), float(coords[index + 1])))
        except (TypeError, ValueError):
            return []
    return points


def parse_path_point_pairs(path_data):
    if not path_data:
        return []
    return parse_svg_point_pairs(path_data)


def transform_point_pairs(points, ox=0.0, oy=0.0, scale=1.0):
    return [
        ((float(x) * scale) + ox, (float(y) * scale) + oy)
        for x, y in points
    ]


def bbox_from_point_pairs(points):
    if len(points) < 2:
        return None
    xs = [float(point[0]) for point in points]
    ys = [float(point[1]) for point in points]
    return normalize_bbox({
        'x': min(xs),
        'y': min(ys),
        'width': max(xs) - min(xs) or 1.0,
        'height': max(ys) - min(ys) or 1.0,
    })


def dedupe_adjacent_points(points, tolerance=0.25):
    deduped = []
    for point in points or []:
        if not deduped:
            deduped.append((float(point[0]), float(point[1])))
            continue
        prev_x, prev_y = deduped[-1]
        cur_x, cur_y = float(point[0]), float(point[1])
        if abs(prev_x - cur_x) <= tolerance and abs(prev_y - cur_y) <= tolerance:
            continue
        deduped.append((cur_x, cur_y))
    return deduped


def bbox_intersection_area(left, right):
    if not left or not right:
        return 0.0
    x1 = max(float(left['x']), float(right['x']))
    y1 = max(float(left['y']), float(right['y']))
    x2 = min(float(left['x']) + float(left['width']), float(right['x']) + float(right['width']))
    y2 = min(float(left['y']) + float(left['height']), float(right['y']) + float(right['height']))
    if x2 <= x1 or y2 <= y1:
        return 0.0
    return (x2 - x1) * (y2 - y1)


def clamp(value, minimum, maximum):
    return max(minimum, min(maximum, value))


def extract_numeric_value(text):
    normalized = normalize_text(text).replace('−', '-')
    if not normalized:
        return None
    match = re.search(r'-?\d+(?:,\d{3})*(?:\.\d+)?', normalized)
    if not match:
        return None
    try:
        return float(match.group(0).replace(',', ''))
    except ValueError:
        return None


def is_numeric_like_text(text):
    normalized = normalize_text(text)
    value = extract_numeric_value(normalized)
    if value is None:
        return False
    residual = re.sub(r'-?\d+(?:,\d{3})*(?:\.\d+)?', '', normalized)
    residual = residual.replace('%', '').replace('$', '').replace('¥', '').replace('€', '')
    residual = residual.replace('+', '').replace('-', '').replace('/', '').replace('x', '')
    residual = residual.strip()
    return len(residual) <= 3


def rgb_to_hex(value):
    if isinstance(value, RGBColor):
        return bytes(value).hex().upper()
    if isinstance(value, str) and value:
        parsed = parse_color(value)
        if parsed and parsed[0] != 'grad':
            return parsed[0].upper()
    return None


def merge_semantic_tables(static_tables, rendered_tables):
    merged = []
    rendered_tables = rendered_tables if isinstance(rendered_tables, list) else []
    for index, table in enumerate(static_tables or []):
        merged_table = {k: v for k, v in table.items() if k != 'rows'}
        rendered_table = rendered_tables[index] if index < len(rendered_tables) and isinstance(rendered_tables[index], dict) else None
        merged_table['bbox'] = normalize_bbox(rendered_table.get('bbox')) if rendered_table else None
        rendered_lookup = {}
        if rendered_table:
            for rendered_row in rendered_table.get('rows', []):
                if not isinstance(rendered_row, list):
                    continue
                for rendered_cell in rendered_row:
                    if not isinstance(rendered_cell, dict):
                        continue
                    key = (int(rendered_cell.get('row_index', 0)), int(rendered_cell.get('cell_index', 0)))
                    rendered_lookup[key] = rendered_cell
        merged_rows = []
        for row in table.get('rows', []):
            merged_cells = []
            for cell in row.get('cells', []):
                merged_cell = dict(cell)
                rendered_cell = rendered_lookup.get((cell['row_index'], cell['cell_index']))
                merged_cell['bbox'] = normalize_bbox(rendered_cell.get('bbox')) if rendered_cell else None
                merged_cell['styles'] = rendered_cell.get('styles') if rendered_cell else None
                merged_cells.append(merged_cell)
            merged_row = dict(row)
            merged_row['cells'] = merged_cells
            merged_rows.append(merged_row)
        merged_table['rows'] = merged_rows
        merged.append(merged_table)
    return merged


class HtmlSemanticMatcher:
    def __init__(self, entries, consumed_entry_ids=None):
        self.entries = list(entries or [])
        self.cursor = 0
        self.consumed_entry_ids = consumed_entry_ids if consumed_entry_ids is not None else set()

    def _entry_id(self, entry):
        if isinstance(entry, dict) and isinstance(entry.get('entry_id'), int):
            return entry['entry_id']
        return id(entry)

    def _is_consumed(self, entry):
        return self._entry_id(entry) in self.consumed_entry_ids

    def _find_match_index(self, key):
        if not key:
            return None
        upper = min(len(self.entries), self.cursor + 16)
        best_idx = None
        for idx in range(self.cursor, upper):
            entry = self.entries[idx]
            if self._is_consumed(entry):
                continue
            entry_key = entry.get('match_key', '')
            if not entry_key:
                continue
            if entry_key == key:
                return idx
            if len(key) >= 6 and (key in entry_key or entry_key in key):
                best_idx = idx
        if best_idx is not None:
            return best_idx
        for idx, entry in enumerate(self.entries):
            if self._is_consumed(entry):
                continue
            if entry.get('match_key') == key:
                return idx
        return None

    def find_match(self, text, consume=True):
        key = semantic_match_key(text)
        idx = self._find_match_index(key)
        if idx is None:
            return None
        entry = self.entries[idx]
        if consume:
            self.consumed_entry_ids.add(self._entry_id(entry))
            self.cursor = max(self.cursor, idx + 1)
        return entry

    def summary(self):
        total = len(self.entries)
        matched = sum(1 for entry in self.entries if self._is_consumed(entry))
        return {
            'entries': total,
            'matched': matched,
            'unmatched': max(total - matched, 0),
        }


class PageSemanticModel:
    def __init__(self, html_semantics=None, rendered_semantics=None):
        self.html_semantics = html_semantics if isinstance(html_semantics, dict) else {}
        self.rendered_semantics = rendered_semantics if isinstance(rendered_semantics, dict) else {}
        self.consumed_entry_ids = set()
        self.global_matcher = HtmlSemanticMatcher(
            self.html_semantics.get('entries', []),
            consumed_entry_ids=self.consumed_entry_ids,
        )
        rendered_blocks = {}
        for block in self.rendered_semantics.get('blocks', []) if isinstance(self.rendered_semantics.get('blocks'), list) else []:
            if not isinstance(block, dict):
                continue
            block_id = block.get('block_id')
            if isinstance(block_id, str) and block_id.strip():
                rendered_blocks[block_id] = block
        self.blocks = []
        for block in self.html_semantics.get('blocks', []) if isinstance(self.html_semantics.get('blocks'), list) else []:
            if not isinstance(block, dict):
                continue
            merged_block = dict(block)
            rendered_block = rendered_blocks.get(merged_block.get('block_id'), {})
            merged_block['bbox'] = normalize_bbox(rendered_block.get('bbox'))
            merged_block['contains_table'] = bool(rendered_block.get('contains_table'))
            merged_block['contains_chart_like'] = bool(rendered_block.get('contains_chart_like'))
            merged_block['matcher'] = HtmlSemanticMatcher(
                merged_block.get('entries', []),
                consumed_entry_ids=self.consumed_entry_ids,
            )
            self.blocks.append(merged_block)
        self.tables = merge_semantic_tables(
            self.html_semantics.get('tables', []),
            self.rendered_semantics.get('tables', []),
        )
        self.charts = []
        for chart in self.rendered_semantics.get('charts', []) if isinstance(self.rendered_semantics.get('charts'), list) else []:
            if not isinstance(chart, dict):
                continue
            merged_chart = dict(chart)
            merged_chart['bbox'] = normalize_bbox(chart.get('bbox'))
            if not merged_chart.get('bbox'):
                continue
            block_id = merged_chart.get('block_id')
            merged_chart['block'] = next((block for block in self.blocks if block.get('block_id') == block_id), None)
            self.charts.append(merged_chart)
        self.block_matches = 0

    def find_block_for_bbox(self, bbox):
        bbox = normalize_bbox(bbox)
        if bbox is None:
            return None
        best_block = None
        best_score = 0.0
        bbox_total_area = max(bbox_area(bbox), 1.0)
        center_x = bbox['x'] + bbox['width'] / 2
        center_y = bbox['y'] + bbox['height'] / 2
        for block in self.blocks:
            block_bbox = block.get('bbox')
            if not block_bbox:
                continue
            overlap_area = bbox_intersection_area(bbox, block_bbox)
            if overlap_area <= 0:
                continue
            score = overlap_area / bbox_total_area
            if (
                block_bbox['x'] <= center_x <= block_bbox['x'] + block_bbox['width']
                and block_bbox['y'] <= center_y <= block_bbox['y'] + block_bbox['height']
            ):
                score += 1.0
            if score > best_score:
                best_score = score
                best_block = block
        return best_block

    def match_text(self, text, bbox=None, consume=True):
        block = self.find_block_for_bbox(bbox)
        if block is not None:
            entry = block['matcher'].find_match(text, consume=consume)
            if entry is not None:
                if consume:
                    self.block_matches += 1
                return entry
        return self.global_matcher.find_match(text, consume=consume)

    def summary(self):
        summary = self.global_matcher.summary()
        summary['block_matches'] = self.block_matches
        summary['block_count'] = len([block for block in self.blocks if block.get('bbox')])
        summary['tables'] = len(self.tables)
        summary['charts'] = len(self.charts)
        return summary


def sanitize_name_token(value):
    if value is None:
        return None
    token = re.sub(r'[^0-9A-Za-z_.-]+', '_', str(value).strip())
    token = re.sub(r'_+', '_', token).strip('._')
    return token or None


def extract_block_name_tokens(shape_name):
    if not isinstance(shape_name, str) or not shape_name:
        return set()
    tokens = set()
    for prefix in (MANAGED_BLOCK_NAME_PREFIX,) + BLOCK_SLOT_NAME_PREFIXES:
        if not shape_name.startswith(prefix):
            continue
        remainder = shape_name[len(prefix):]
        token = sanitize_name_token(remainder.split(':', 1)[0])
        if token:
            tokens.add(token)
    for match in re.finditer(r':block=([0-9A-Za-z_.-]+)', shape_name):
        token = sanitize_name_token(match.group(1))
        if token:
            tokens.add(token)
    return tokens


def collect_update_blocks(html_semantics=None, rendered_semantics=None):
    blocks = []
    seen = set()
    page_semantics = PageSemanticModel(html_semantics, rendered_semantics)
    source_blocks = list(page_semantics.blocks)
    if not source_blocks and isinstance(rendered_semantics, dict):
        source_blocks = [
            block for block in rendered_semantics.get('blocks', [])
            if isinstance(block, dict)
        ]
    for block in source_blocks:
        block_id = normalize_text(block.get('block_id'))
        token = sanitize_name_token(block_id)
        if not token or token in seen:
            continue
        seen.add(token)
        blocks.append({
            'block_id': block_id,
            'name_token': token,
            'bbox': normalize_bbox(block.get('bbox')),
        })
    return blocks


def build_pptx_report(*, output_path, svg_input, html_dir, source_report_path, source_report, top_warnings, page_reports,
                      template_pptx_path=None, target_slide_numbers=None, update_mode='new_presentation',
                      preserve_template_background=False):
    editable_count = sum(1 for page in page_reports if page.get('source_method') == 'dom_to_svg_editable')
    raster_count = sum(1 for page in page_reports if page.get('source_method') == 'png_wrapper_raster')
    pathified_count = sum(1 for page in page_reports if page.get('source_method') == 'pdf2svg_pathified')
    failed_count = sum(1 for page in page_reports if page.get('source_method') == 'failed')
    unknown_count = sum(1 for page in page_reports if page.get('source_method') == 'unknown')
    pptx_shapes_total = sum(int(page.get('pptx_shapes', 0)) for page in page_reports)
    pptx_skipped_total = sum(int(page.get('pptx_skipped', 0)) for page in page_reports)
    pptx_errors_total = sum(int(page.get('pptx_errors', 0)) for page in page_reports)
    html_semantic_entries_total = sum(int(page.get('html_semantic_entries') or 0) for page in page_reports)
    html_semantic_matches_total = sum(int(page.get('html_semantic_matches') or 0) for page in page_reports)
    html_semantic_unmatched_total = sum(int(page.get('html_semantic_unmatched') or 0) for page in page_reports)
    html_block_matches_total = sum(int(page.get('html_block_matches') or 0) for page in page_reports)
    native_tables_total = sum(int(page.get('native_tables') or 0) for page in page_reports)
    semantic_blocks_total = sum(int(page.get('html_blocks') or 0) for page in page_reports)
    rendered_chart_regions_total = sum(int(page.get('rendered_charts') or 0) for page in page_reports)
    native_charts_total = sum(int(page.get('native_charts') or 0) for page in page_reports)
    structured_chart_groups_total = sum(int(page.get('structured_chart_groups') or 0) for page in page_reports)
    updated_blocks_total = sum(len(page.get('updated_block_ids') or []) for page in page_reports)
    template_removed_shapes_total = sum(int(page.get('template_removed_shapes') or 0) for page in page_reports)
    template_removed_slot_shapes_total = sum(int(page.get('template_removed_slot_shapes') or 0) for page in page_reports)
    template_removed_managed_shapes_total = sum(int(page.get('template_removed_managed_shapes') or 0) for page in page_reports)
    structured_chart_hits_total = native_charts_total + structured_chart_groups_total
    html_semantic_slides = sum(1 for page in page_reports if page.get('html_source_path'))
    source_summary = source_report.get('summary') if isinstance(source_report, dict) else None
    return {
        'generated_at': iso_now(),
        'presentation_path': str(Path(output_path).resolve()),
        'source_svg_input': str(Path(svg_input).resolve()),
        'source_html_dir': str(Path(html_dir).resolve()) if html_dir else None,
        'source_report_path': str(source_report_path.resolve()) if source_report_path else None,
        'update_mode': update_mode,
        'template_update_scope': 'block_update' if update_mode == 'template_update' else None,
        'template_pptx_path': str(Path(template_pptx_path).resolve()) if template_pptx_path else None,
        'target_slide_numbers': list(target_slide_numbers or []),
        'preserve_template_background': bool(preserve_template_background),
        'summary': {
            'total_slides': len(page_reports),
            'source_editable_slides': editable_count,
            'source_raster_slides': raster_count,
            'source_pathified_slides': pathified_count,
            'source_failed_slides': failed_count,
            'source_unknown_slides': unknown_count,
            'pptx_shapes_total': pptx_shapes_total,
            'pptx_skipped_total': pptx_skipped_total,
            'pptx_errors_total': pptx_errors_total,
            'html_semantic_slides': html_semantic_slides,
            'html_semantic_entries_total': html_semantic_entries_total,
            'html_semantic_matches_total': html_semantic_matches_total,
            'html_semantic_unmatched_total': html_semantic_unmatched_total,
            'html_block_matches_total': html_block_matches_total,
            'html_blocks_total': semantic_blocks_total,
            'native_tables_total': native_tables_total,
            'rendered_chart_regions_total': rendered_chart_regions_total,
            'native_charts_total': native_charts_total,
            'structured_chart_groups_total': structured_chart_groups_total,
            'structured_chart_hits_total': structured_chart_hits_total,
            'updated_blocks_total': updated_blocks_total,
            'template_removed_shapes_total': template_removed_shapes_total,
            'template_removed_slot_shapes_total': template_removed_slot_shapes_total,
            'template_removed_managed_shapes_total': template_removed_managed_shapes_total,
            'structured_chart_hit_rate': round(
                structured_chart_hits_total / rendered_chart_regions_total, 4
            ) if rendered_chart_regions_total else None,
            'html_semantic_match_rate': round(
                html_semantic_matches_total / html_semantic_entries_total, 4
            ) if html_semantic_entries_total else None,
            'warnings': len(top_warnings),
            'source_total_pages': source_summary.get('total_pages') if isinstance(source_summary, dict) else None,
        },
        'warnings': top_warnings,
        'pages': page_reports,
    }


def write_pptx_report(report_path, payload):
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding='utf-8')


def _natural_key(path):
    return [int(x) if x.isdigit() else x.lower() for x in re.split(r'(\d+)', path.stem)]


def resolve_svg_files(svg_input):
    svg_input = Path(svg_input)
    if svg_input.is_file():
        svg_files = [svg_input]
    elif svg_input.is_dir():
        svg_files = sorted(svg_input.glob('*.svg'), key=_natural_key)
    else:
        print(f"Error: {svg_input} not found", file=sys.stderr)
        sys.exit(1)
    if not svg_files:
        print("Error: No SVG files found", file=sys.stderr)
        sys.exit(1)
    return svg_files


def parse_target_slide_numbers(raw_value, expected_count):
    if not raw_value:
        return None
    values = []
    seen = set()
    for token in str(raw_value).split(','):
        cleaned = token.strip()
        if not cleaned:
            continue
        try:
            slide_number = int(cleaned)
        except ValueError as exc:
            raise ValueError(f'invalid target slide number: {cleaned}') from exc
        if slide_number <= 0:
            raise ValueError(f'target slide number must be >= 1: {cleaned}')
        if slide_number in seen:
            raise ValueError(f'duplicate target slide number: {slide_number}')
        seen.add(slide_number)
        values.append(slide_number)
    if len(values) != expected_count:
        raise ValueError(f'target slide count mismatch: expected {expected_count}, got {len(values)}')
    return values


def clear_slide_shapes(slide):
    sp_tree = slide.shapes._spTree
    for element in list(sp_tree)[2:]:
        sp_tree.remove(element)


def clear_target_block_shapes(slide, block_ids):
    target_tokens = {
        token for token in (sanitize_name_token(block_id) for block_id in block_ids or [])
        if token
    }
    removal_stats = {
        'removed_total': 0,
        'removed_slot_shapes': 0,
        'removed_managed_shapes': 0,
        'target_blocks': sorted(target_tokens),
    }
    if not target_tokens:
        return removal_stats

    for shape in reversed(list(slide.shapes)):
        try:
            shape_name = shape.name or ''
        except Exception:
            shape_name = ''
        shape_tokens = extract_block_name_tokens(shape_name)
        if not shape_tokens or shape_tokens.isdisjoint(target_tokens):
            continue
        element = getattr(shape, '_element', None)
        parent = element.getparent() if element is not None else None
        if parent is None:
            continue
        parent.remove(element)
        removal_stats['removed_total'] += 1
        if any(shape_name.startswith(prefix) for prefix in BLOCK_SLOT_NAME_PREFIXES):
            removal_stats['removed_slot_shapes'] += 1
        else:
            removal_stats['removed_managed_shapes'] += 1
    return removal_stats


def px(v):
    return int(float(v) * EMU_PX)

def font_sz(svg_px):
    return max(100, int(float(svg_px) * 75))

def strip_unit(v):
    return re.sub(r'[a-z%]+', '', str(v))

def resolve_font(ff_str):
    """解析 font-family 字符串，返回 PPT 可用字体。"""
    ff_str = ff_str.replace('&quot;', '').replace('"', '').replace("'", '')
    fonts = [f.strip() for f in ff_str.split(',') if f.strip()]
    for f in fonts:
        if f in FONT_FALLBACK:
            return FONT_FALLBACK[f]
        if f and f not in ('sans-serif', 'serif', 'monospace', 'system-ui'):
            return f
    return 'Microsoft YaHei'


# -------------------------------------------------------------------
# 颜色解析（完整 CSS 命名颜色）
# -------------------------------------------------------------------
def parse_color(s):
    if not s or s.strip() == 'none':
        return None
    s = s.strip()
    if s.startswith('url('):
        m = re.search(r'#([\w-]+)', s)
        return ('grad', m.group(1)) if m else None
    m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*(?:,\s*([\d.]+))?\s*\)', s)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        return (f'{r:02x}{g:02x}{b:02x}', int(a * 100000))
    if s.startswith('#'):
        h = s[1:]
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        return (h.lower().ljust(6, '0')[:6], 100000)
    c = CSS_COLORS.get(s.lower())
    return (c, 100000) if c else None


# -------------------------------------------------------------------
# OOXML 元素构造
# -------------------------------------------------------------------
def _el(tag, attrib=None, text=None, children=None):
    pre, local = tag.split(':') if ':' in tag else ('a', tag)
    el = etree.Element(f'{{{NS[pre]}}}{local}')
    if attrib:
        for k, v in attrib.items():
            el.set(k, str(v))
    if text is not None:
        el.text = str(text)
    for c in (children or []):
        if c is not None:
            el.append(c)
    return el

def _srgb(hex6, alpha=100000):
    el = _el('a:srgbClr', {'val': hex6})
    if alpha < 100000:
        el.append(_el('a:alpha', {'val': str(alpha)}))
    return el

def make_fill(fill_str, grads, opacity=1.0):
    c = parse_color(fill_str)
    if c is None:
        return _el('a:noFill')
    if c[0] == 'grad':
        gdef = grads.get(c[1])
        return _make_grad(gdef) if gdef else _el('a:noFill')
    hex6, alpha = c
    alpha = int(alpha * opacity)
    return _el('a:solidFill', children=[_srgb(hex6, alpha)])

def _make_grad(gdef):
    gs_lst = _el('a:gsLst')
    for stop in gdef['stops']:
        pos = int(stop['offset'] * 1000)
        sc = parse_color(stop['color_str'])
        if not sc or sc[0] == 'grad':
            continue
        hex6, alpha = sc
        alpha = int(alpha * stop.get('opacity', 1.0))
        gs_lst.append(_el('a:gs', {'pos': str(pos)}, children=[_srgb(hex6, alpha)]))

    if gdef.get('type') == 'radial':
        # 径向渐变
        path = _el('a:path', {'path': 'circle'}, children=[
            _el('a:fillToRect', {'l': '50000', 't': '50000', 'r': '50000', 'b': '50000'})
        ])
        return _el('a:gradFill', {'rotWithShape': '1'}, children=[gs_lst, path])
    else:
        # 线性渐变
        dx = gdef.get('x2', 1) - gdef.get('x1', 0)
        dy = gdef.get('y2', 1) - gdef.get('y1', 0)
        ang = int(math.degrees(math.atan2(dy, dx)) * 60000)
        if ang < 0:
            ang += 21600000
        lin = _el('a:lin', {'ang': str(ang), 'scaled': '0'})
        return _el('a:gradFill', children=[gs_lst, lin])

def make_line(stroke_str, stroke_w=1):
    c = parse_color(stroke_str)
    if not c or c[0] == 'grad':
        return None
    hex6, alpha = c
    w = max(1, int(float(strip_unit(stroke_w)) * 12700))
    return _el('a:ln', {'w': str(w)},
               children=[_el('a:solidFill', children=[_srgb(hex6, alpha)])])

def make_shape(sid, name, x, y, cx, cy, preset='rect',
               fill_el=None, line_el=None, rx=0, geom_el=None):
    sp = _el('p:sp')
    sp.append(_el('p:nvSpPr', children=[
        _el('p:cNvPr', {'id': str(sid), 'name': name}),
        _el('p:cNvSpPr'), _el('p:nvPr'),
    ]))
    sp_pr = _el('p:spPr')
    sp_pr.append(_el('a:xfrm', children=[
        _el('a:off', {'x': str(max(0, int(x))), 'y': str(max(0, int(y)))}),
        _el('a:ext', {'cx': str(max(0, int(cx))), 'cy': str(max(0, int(cy)))}),
    ]))
    if geom_el is not None:
        sp_pr.append(geom_el)
    else:
        geom = _el('a:prstGeom', {'prst': preset})
        av = _el('a:avLst')
        if preset == 'roundRect' and rx > 0:
            shorter = max(min(cx, cy), 1)
            adj = min(50000, int(rx / (shorter / 2) * 50000))
            av.append(_el('a:gd', {'name': 'adj', 'fmla': f'val {adj}'}))
        geom.append(av)
        sp_pr.append(geom)
    sp_pr.append(fill_el if fill_el is not None else _el('a:noFill'))
    if line_el is not None:
        sp_pr.append(line_el)
    sp.append(sp_pr)
    return sp

def make_textbox(sid, name, x, y, cx, cy, paragraphs, anchor='t'):
    """paragraphs = [[{text,sz,bold,hex,alpha,font}, ...], ...] or
    [{runs, align, bullet_char, number, indent_level, continuation}, ...]
    anchor: 't'=top, 'ctr'=center, 'b'=bottom
    """
    sp = _el('p:sp')
    sp.append(_el('p:nvSpPr', children=[
        _el('p:cNvPr', {'id': str(sid), 'name': name}),
        _el('p:cNvSpPr', {'txBox': '1'}), _el('p:nvPr'),
    ]))
    sp.append(_el('p:spPr', children=[
        _el('a:xfrm', children=[
            _el('a:off', {'x': str(max(0, int(x))), 'y': str(max(0, int(y)))}),
            _el('a:ext', {'cx': str(max(0, int(cx))), 'cy': str(max(0, int(cy)))}),
        ]),
        _el('a:prstGeom', {'prst': 'rect'}, children=[_el('a:avLst')]),
        _el('a:noFill'), _el('a:ln', children=[_el('a:noFill')]),
    ]))
    tx = _el('p:txBody', children=[
        _el('a:bodyPr', {'wrap': 'none', 'lIns': '0', 'tIns': '0',
                         'rIns': '0', 'bIns': '0', 'anchor': anchor}),
        _el('a:lstStyle'),
    ])
    for paragraph in paragraphs:
        if isinstance(paragraph, dict):
            runs = paragraph.get('runs', [])
            align = paragraph.get('align')
            bullet_char = paragraph.get('bullet_char')
            number = paragraph.get('number')
            indent_level = max(0, int(paragraph.get('indent_level', 0)))
            continuation = bool(paragraph.get('continuation'))
        else:
            runs = paragraph
            align = None
            bullet_char = None
            number = None
            indent_level = 0
            continuation = False
        p_el = _el('a:p')
        # 段落属性: 行距=90%, 段前距=0, 段后距=0
        p_pr_attrib = {}
        if align:
            p_pr_attrib['algn'] = align
        if bullet_char or number is not None or continuation:
            margin_left = 228600 * max(1, indent_level + 1)
            p_pr_attrib['marL'] = str(margin_left)
            p_pr_attrib['indent'] = '0' if continuation else str(-228600)
        p_pr = _el('a:pPr', p_pr_attrib)
        if continuation or (not bullet_char and number is None):
            p_pr.append(_el('a:buNone'))
        elif number is not None:
            p_pr.append(_el('a:buAutoNum', {
                'type': 'arabicPeriod',
                'startAt': str(max(1, int(number))),
            }))
        else:
            p_pr.append(_el('a:buChar', {'char': bullet_char}))
        p_pr.append(_el('a:lnSpc', children=[_el('a:spcPct', {'val': '90000'})]))
        p_pr.append(_el('a:spcBef', children=[_el('a:spcPts', {'val': '0'})]))
        p_pr.append(_el('a:spcAft', children=[_el('a:spcPts', {'val': '0'})]))
        p_el.append(p_pr)
        for run in runs:
            rpr_a = {'lang': 'zh-CN', 'dirty': '0'}
            if run.get('sz'):
                rpr_a['sz'] = str(run['sz'])
            if run.get('bold'):
                rpr_a['b'] = '1'
            rpr = _el('a:rPr', rpr_a)
            rpr.append(_el('a:solidFill', children=[
                _srgb(run.get('hex', '000000'), run.get('alpha', 100000))
            ]))
            font = run.get('font', 'Microsoft YaHei')
            rpr.append(_el('a:latin', {'typeface': font}))
            rpr.append(_el('a:ea', {'typeface': font}))
            p_el.append(_el('a:r', children=[rpr, _el('a:t', text=run.get('text', ''))]))
        tx.append(p_el)
    sp.append(tx)
    return sp


def parse_css_px(value, default=0.0):
    if value is None:
        return float(default)
    match = re.search(r'-?[\d.]+', str(value))
    if not match:
        return float(default)
    try:
        return float(match.group(0))
    except ValueError:
        return float(default)


def parse_css_font_weight(value, default=False):
    if value is None:
        return default
    text = str(value).strip().lower()
    if text in {'bold', 'bolder'}:
        return True
    if text in {'normal', 'lighter'}:
        return False
    try:
        return int(text) >= 600
    except ValueError:
        return default


def to_pp_align(value):
    text = str(value or '').strip().lower()
    if text in {'center', 'middle'}:
        return PP_ALIGN.CENTER
    if text in {'right', 'end'}:
        return PP_ALIGN.RIGHT
    if text in {'justify'}:
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def set_table_cell_borders(cell, color_hex='D1D5DB', width=12700):
    tc_pr = cell._tc.get_or_add_tcPr()
    for side in ('lnL', 'lnR', 'lnT', 'lnB'):
        existing = tc_pr.find(f'{{{NS["a"]}}}{side}')
        if existing is not None:
            tc_pr.remove(existing)
        tc_pr.append(_el(f'a:{side}', {'w': str(width)}, children=[
            _el('a:solidFill', children=[_srgb(color_hex)]),
            _el('a:prstDash', {'val': 'solid'}),
            _el('a:round'),
            _el('a:headEnd', {'type': 'none', 'w': 'med', 'len': 'med'}),
            _el('a:tailEnd', {'type': 'none', 'w': 'med', 'len': 'med'}),
        ]))


def apply_table_cell_text(cell, cell_semantics):
    styles = cell_semantics.get('styles') or {}
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = px(8)
    text_frame.margin_right = px(8)
    text_frame.margin_top = px(5)
    text_frame.margin_bottom = px(5)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = to_pp_align(styles.get('text_align'))
    run = paragraph.add_run()
    run.text = cell_semantics.get('text', '')
    font = run.font
    font.bold = cell_semantics.get('is_header', False) or parse_css_font_weight(styles.get('font_weight'))
    font.size = Pt(max(9.0, parse_css_px(styles.get('font_size'), default=14.0) * 0.75))
    font.name = 'Microsoft YaHei'
    text_color = parse_color(styles.get('color'))
    if text_color and text_color[0] != 'grad':
        font.color.rgb = RGBColor.from_string(text_color[0].upper())

    fill_color = parse_color(styles.get('background_color'))
    if fill_color and fill_color[0] != 'grad':
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill_color[0].upper())
    elif cell_semantics.get('is_header'):
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string('F3F4F6')

    border_color = None
    for key in ('border_top_color', 'border_right_color', 'border_bottom_color', 'border_left_color'):
        parsed = parse_color(styles.get(key))
        if parsed and parsed[0] != 'grad':
            border_color = parsed[0].upper()
            break
    set_table_cell_borders(cell, color_hex=border_color or 'D1D5DB')


# -------------------------------------------------------------------
# SVG Path 解析器 -> OOXML custGeom
# -------------------------------------------------------------------
_PATH_RE = re.compile(r'([mMzZlLhHvVcCsSqQtTaA])|([+-]?(?:\d+\.?\d*|\.\d+)(?:[eE][+-]?\d+)?)')

def parse_path_to_custgeom(d_str, bbox):
    """SVG path d -> OOXML a:custGeom 元素。bbox=(x,y,w,h) 用于坐标偏移。"""
    bx, by, bw, bh = bbox
    scale = 100000  # OOXML 路径坐标空间

    def coord(v, is_x=True):
        base = bw if is_x else bh
        offset = bx if is_x else by
        if base <= 0:
            return 0
        return int((float(v) - offset) / base * scale)

    tokens = _PATH_RE.findall(d_str)
    items = []
    for cmd_match, num_match in tokens:
        if cmd_match:
            items.append(cmd_match)
        elif num_match:
            items.append(float(num_match))

    path_el = _el('a:path', {'w': str(scale), 'h': str(scale)})
    i = 0
    cx_p, cy_p = 0, 0  # current point (absolute)
    last_cx2, last_cy2 = 0, 0  # 上一个 C/S 的第二控制点（S 命令反射用）
    cmd = None
    rel = False

    while i < len(items):
        if isinstance(items[i], str):
            cmd = items[i].lower()
            rel = items[i].islower()
            i += 1
            if cmd == 'z':
                path_el.append(_el('a:close'))
                continue

        if cmd is None:
            i += 1
            continue

        try:
            if cmd == 'm':
                x, y = float(items[i]), float(items[i+1])
                if rel:
                    x += cx_p; y += cy_p
                cx_p, cy_p = x, y
                path_el.append(_el('a:moveTo', children=[
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))})
                ]))
                i += 2
                cmd = 'l'  # implicit lineTo after moveTo

            elif cmd == 'l':
                x, y = float(items[i]), float(items[i+1])
                if rel:
                    x += cx_p; y += cy_p
                cx_p, cy_p = x, y
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))})
                ]))
                i += 2

            elif cmd == 'h':
                x = float(items[i])
                if rel:
                    x += cx_p
                cx_p = x
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(cx_p, True)), 'y': str(coord(cy_p, False))})
                ]))
                i += 1

            elif cmd == 'v':
                y = float(items[i])
                if rel:
                    y += cy_p
                cy_p = y
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(cx_p, True)), 'y': str(coord(cy_p, False))})
                ]))
                i += 1

            elif cmd == 'c':
                x1, y1 = float(items[i]), float(items[i+1])
                x2, y2 = float(items[i+2]), float(items[i+3])
                x, y = float(items[i+4]), float(items[i+5])
                if rel:
                    x1 += cx_p; y1 += cy_p
                    x2 += cx_p; y2 += cy_p
                    x += cx_p; y += cy_p
                last_cx2, last_cy2 = x2, y2
                cx_p, cy_p = x, y
                path_el.append(_el('a:cubicBezTo', children=[
                    _el('a:pt', {'x': str(coord(x1, True)), 'y': str(coord(y1, False))}),
                    _el('a:pt', {'x': str(coord(x2, True)), 'y': str(coord(y2, False))}),
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))}),
                ]))
                i += 6

            elif cmd == 's':
                # 平滑三次贝塞尔：反射上一个 C/S 的第二控制点
                x2, y2 = float(items[i]), float(items[i+1])
                x, y = float(items[i+2]), float(items[i+3])
                if rel:
                    x2 += cx_p; y2 += cy_p
                    x += cx_p; y += cy_p
                x1 = 2 * cx_p - last_cx2
                y1 = 2 * cy_p - last_cy2
                last_cx2, last_cy2 = x2, y2
                cx_p, cy_p = x, y
                path_el.append(_el('a:cubicBezTo', children=[
                    _el('a:pt', {'x': str(coord(x1, True)), 'y': str(coord(y1, False))}),
                    _el('a:pt', {'x': str(coord(x2, True)), 'y': str(coord(y2, False))}),
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))}),
                ]))
                i += 4

            elif cmd == 'q':
                # 二次贝塞尔 -> 三次贝塞尔近似
                qx, qy = float(items[i]), float(items[i+1])
                x, y = float(items[i+2]), float(items[i+3])
                if rel:
                    qx += cx_p; qy += cy_p
                    x += cx_p; y += cy_p
                x1 = cx_p + 2/3 * (qx - cx_p)
                y1 = cy_p + 2/3 * (qy - cy_p)
                x2 = x + 2/3 * (qx - x)
                y2 = y + 2/3 * (qy - y)
                last_cx2, last_cy2 = qx, qy
                cx_p, cy_p = x, y
                path_el.append(_el('a:cubicBezTo', children=[
                    _el('a:pt', {'x': str(coord(x1, True)), 'y': str(coord(y1, False))}),
                    _el('a:pt', {'x': str(coord(x2, True)), 'y': str(coord(y2, False))}),
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))}),
                ]))
                i += 4

            elif cmd == 't':
                # 平滑二次贝塞尔：反射上一个 Q 控制点
                x, y = float(items[i]), float(items[i+1])
                if rel:
                    x += cx_p; y += cy_p
                qx = 2 * cx_p - last_cx2
                qy = 2 * cy_p - last_cy2
                x1 = cx_p + 2/3 * (qx - cx_p)
                y1 = cy_p + 2/3 * (qy - cy_p)
                x2 = x + 2/3 * (qx - x)
                y2 = y + 2/3 * (qy - y)
                last_cx2, last_cy2 = qx, qy
                cx_p, cy_p = x, y
                path_el.append(_el('a:cubicBezTo', children=[
                    _el('a:pt', {'x': str(coord(x1, True)), 'y': str(coord(y1, False))}),
                    _el('a:pt', {'x': str(coord(x2, True)), 'y': str(coord(y2, False))}),
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))}),
                ]))
                i += 2

            elif cmd == 'a':
                # 弧线：完整转换复杂，降级为直线连终点保持路径连续
                # 参数: rx ry x-rot large-arc sweep x y
                _rx = float(items[i])
                _ry = float(items[i+1])
                # items[i+2] = x-rotation, items[i+3] = large-arc, items[i+4] = sweep
                x, y = float(items[i+5]), float(items[i+6])
                if rel:
                    x += cx_p; y += cy_p
                last_cx2, last_cy2 = x, y
                cx_p, cy_p = x, y
                path_el.append(_el('a:lnTo', children=[
                    _el('a:pt', {'x': str(coord(x, True)), 'y': str(coord(y, False))}),
                ]))
                i += 7
            else:
                i += 1
        except (IndexError, ValueError):
            i += 1

    cust_geom = _el('a:custGeom', children=[
        _el('a:avLst'), _el('a:gdLst'), _el('a:ahLst'), _el('a:cxnLst'),
        _el('a:rect', {'l': 'l', 't': 't', 'r': 'r', 'b': 'b'}),
        _el('a:pathLst', children=[path_el]),
    ])
    return cust_geom


# -------------------------------------------------------------------
# SVG -> PPTX 转换器
# -------------------------------------------------------------------
class SvgConverter:
    def __init__(self, on_progress=None):
        self.sid = 100
        self.grads = {}
        self.bg_set = False  # 是否已设置幻灯片背景
        self.on_progress = on_progress  # 进度回调 (i, total, filename)
        self.page_semantics = None
        self.semantic_tables = []
        self.semantic_charts = []
        self.chart_groups = {}
        self.native_chart_groups = {}
        self.native_chart_ids = set()
        self.current_slide = None
        self.svg_primitives = []
        self.preserve_background = False
        self.stats = {'shapes': 0, 'skipped': 0, 'errors': 0}

    def _id(self):
        self.sid += 1
        return self.sid

    def _resolve_block_id(self, bbox=None, block_id=None, semantic=None):
        resolved = normalize_text(block_id)
        if resolved:
            return resolved
        if isinstance(semantic, dict):
            resolved = normalize_text(semantic.get('block_id'))
            if resolved:
                return resolved
        if bbox is None or self.page_semantics is None:
            return None
        block = self.page_semantics.find_block_for_bbox(bbox)
        if block is None:
            return None
        return normalize_text(block.get('block_id'))

    def _managed_shape_name(self, base_name, bbox=None, block_id=None, semantic=None):
        token = sanitize_name_token(self._resolve_block_id(bbox=bbox, block_id=block_id, semantic=semantic))
        if not token:
            return base_name
        return f'{MANAGED_BLOCK_NAME_PREFIX}{token}:{base_name}'

    def _managed_chart_name(self, prefix, identifier, chart_type=None, bbox=None, block_id=None):
        name = f'{prefix}{identifier}'
        if chart_type:
            name += f':{chart_type}'
        token = sanitize_name_token(self._resolve_block_id(bbox=bbox, block_id=block_id))
        if token:
            name += f'{BLOCK_NAME_SUFFIX_MARKER}{token}'
        return name

    def _set_element_name(self, shape_el, name):
        c_nv_pr = shape_el.find(f'.//{{{NS["p"]}}}cNvPr')
        if c_nv_pr is not None:
            c_nv_pr.set('name', name)

    def _append_shape_element(self, sp, shape_el, base_name, bbox=None, block_id=None, semantic=None):
        self._set_element_name(
            shape_el,
            self._managed_shape_name(base_name, bbox=bbox, block_id=block_id, semantic=semantic),
        )
        sp.append(shape_el)
        self.stats['shapes'] += 1

    def _assign_shape_name(self, shape_obj, base_name, bbox=None, block_id=None, semantic=None):
        try:
            shape_obj.name = self._managed_shape_name(base_name, bbox=bbox, block_id=block_id, semantic=semantic)
        except Exception:
            pass

    def _assign_chart_name(self, shape_obj, identifier, chart_type, bbox=None, block_id=None, prefix=NATIVE_CHART_NAME_PREFIX):
        try:
            shape_obj.name = self._managed_chart_name(
                prefix,
                identifier,
                chart_type=chart_type,
                bbox=bbox,
                block_id=block_id,
            )
        except Exception:
            pass

    def convert(self, svg_path, slide, html_semantics=None, rendered_semantics=None, preserve_existing_background=False):
        self.bg_set = False
        self.current_slide = slide
        self.chart_groups = {}
        self.native_chart_groups = {}
        self.native_chart_ids = set()
        self.preserve_background = bool(preserve_existing_background)
        self.page_semantics = PageSemanticModel(html_semantics, rendered_semantics)
        self.semantic_tables = [table for table in self.page_semantics.tables if table.get('bbox')]
        self.semantic_charts = [chart for chart in self.page_semantics.charts if chart.get('bbox')]
        semantic_entries = self.page_semantics.html_semantics.get('entries', [])
        self.stats = {
            'shapes': 0,
            'skipped': 0,
            'errors': 0,
            'semantic_entries': len(semantic_entries),
            'semantic_matches': 0,
            'semantic_unmatched': len(semantic_entries),
            'semantic_blocks': len(self.page_semantics.blocks),
            'block_matches': 0,
            'native_tables': 0,
            'native_charts': 0,
            'structured_chart_groups': 0,
            'rendered_charts': len(self.semantic_charts),
        }
        tree = etree.parse(str(svg_path))
        root = tree.getroot()
        self.root = root  # 供 _use 方法查找引用元素
        self._parse_grads(root)
        self.svg_primitives = self._collect_svg_primitives(root)
        sp_tree = None
        for d in slide._element.iter():
            if d.tag.endswith('}spTree'):
                sp_tree = d
                break
        if sp_tree is None:
            return
        self._render_semantic_charts(slide)
        self._walk(root, sp_tree, 0, 0, 1.0, 1.0, slide)
        self._render_semantic_tables(slide)
        self._finalize_chart_groups()
        semantic_summary = self.page_semantics.summary()
        self.stats['semantic_entries'] = semantic_summary['entries']
        self.stats['semantic_matches'] = semantic_summary['matched']
        self.stats['semantic_unmatched'] = semantic_summary['unmatched']
        self.stats['semantic_blocks'] = semantic_summary['block_count']
        self.stats['block_matches'] = semantic_summary['block_matches']

    def _parse_grads(self, root):
        self.grads = {}
        pct = lambda v: float(v.rstrip('%')) / 100 if '%' in str(v) else float(v)
        for g in root.iter(f'{{{SVG_NS}}}linearGradient'):
            gid = g.get('id')
            if not gid:
                continue
            stops = []
            for s in g.findall(f'{{{SVG_NS}}}stop'):
                off = s.get('offset', '0%')
                off = float(off.rstrip('%')) if '%' in off else float(off) * 100
                stops.append({'offset': off, 'color_str': s.get('stop-color', '#000'),
                              'opacity': float(s.get('stop-opacity', '1'))})
            self.grads[gid] = {
                'type': 'linear', 'stops': stops,
                'x1': pct(g.get('x1', '0%')), 'y1': pct(g.get('y1', '0%')),
                'x2': pct(g.get('x2', '100%')), 'y2': pct(g.get('y2', '100%')),
            }
        for g in root.iter(f'{{{SVG_NS}}}radialGradient'):
            gid = g.get('id')
            if not gid:
                continue
            stops = []
            for s in g.findall(f'{{{SVG_NS}}}stop'):
                off = s.get('offset', '0%')
                off = float(off.rstrip('%')) if '%' in off else float(off) * 100
                stops.append({'offset': off, 'color_str': s.get('stop-color', '#000'),
                              'opacity': float(s.get('stop-opacity', '1'))})
            self.grads[gid] = {'type': 'radial', 'stops': stops}

    def _match_semantic_text(self, text, bbox=None, consume=True):
        return self.page_semantics.match_text(text, bbox=bbox, consume=consume) if self.page_semantics else None

    def _table_region_for_bbox(self, bbox):
        bbox = normalize_bbox(bbox)
        if bbox is None:
            return None
        area = max(bbox_area(bbox), 1.0)
        center_x = bbox['x'] + bbox['width'] / 2
        center_y = bbox['y'] + bbox['height'] / 2
        for table in self.semantic_tables:
            table_bbox = table.get('bbox')
            if not table_bbox:
                continue
            overlap = bbox_intersection_area(bbox, table_bbox)
            if overlap <= 0:
                continue
            overlap_ratio = overlap / area
            contains_center = (
                table_bbox['x'] <= center_x <= table_bbox['x'] + table_bbox['width']
                and table_bbox['y'] <= center_y <= table_bbox['y'] + table_bbox['height']
            )
            if overlap_ratio >= 0.6 or (contains_center and overlap_ratio >= 0.3):
                return table
        return None

    def _chart_region_for_bbox(self, bbox):
        bbox = normalize_bbox(bbox)
        if bbox is None:
            return None
        area = max(bbox_area(bbox), 1.0)
        center_x = bbox['x'] + bbox['width'] / 2
        center_y = bbox['y'] + bbox['height'] / 2
        best_chart = None
        best_score = 0.0
        best_area = None
        for chart in self.semantic_charts:
            chart_bbox = chart.get('bbox')
            if not chart_bbox:
                continue
            overlap = bbox_intersection_area(bbox, chart_bbox)
            if overlap <= 0:
                continue
            score = overlap / area
            if (
                chart_bbox['x'] <= center_x <= chart_bbox['x'] + chart_bbox['width']
                and chart_bbox['y'] <= center_y <= chart_bbox['y'] + chart_bbox['height']
            ):
                score += 1.0
            chart_area = bbox_area(chart_bbox)
            if score > best_score or (score == best_score and best_area is not None and chart_area < best_area):
                best_chart = chart
                best_score = score
                best_area = chart_area
        return best_chart

    def _native_chart_region_for_bbox(self, bbox):
        chart_region = self._chart_region_for_bbox(bbox)
        if chart_region is None:
            return None
        chart_id = chart_region.get('chart_id')
        if not isinstance(chart_id, str) or chart_id not in self.native_chart_ids:
            return None
        return chart_region

    def _ensure_native_chart_group(self, chart_region, chart_type):
        chart_id = chart_region.get('chart_id') or f'native-chart-{len(self.native_chart_groups) + 1}'
        group = self.native_chart_groups.get(chart_id)
        if group is not None:
            return group
        group = self.current_slide.shapes.add_group_shape()
        try:
            group.name = self._managed_chart_name(
                NATIVE_CHART_NAME_PREFIX,
                chart_id,
                chart_type=chart_type,
                bbox=chart_region.get('bbox'),
                block_id=chart_region.get('block_id'),
            )
        except Exception:
            pass
        self.native_chart_groups[chart_id] = group
        return group

    def _append_group_shape_element(self, group, shape_el):
        group.shapes._spTree.append(shape_el)
        self.stats['shapes'] += 1

    def _solid_fill_from_rgb(self, color_value, opacity=1.0):
        hex6 = rgb_to_hex(color_value)
        if not hex6:
            return _el('a:noFill')
        alpha = max(0, min(100000, int(opacity * 100000)))
        return _el('a:solidFill', children=[_srgb(hex6, alpha)])

    def _line_from_rgb(self, color_value, stroke_width_px, opacity=1.0):
        hex6 = rgb_to_hex(color_value)
        if not hex6:
            return None
        alpha = max(0, min(100000, int(opacity * 100000)))
        return _el('a:ln', {'w': str(max(1, int(float(stroke_width_px) * 12700)))}, children=[
            _el('a:solidFill', children=[_srgb(hex6, alpha)]),
            _el('a:round'),
        ])

    def _recalculate_structured_groups(self):
        for group in self.chart_groups.values():
            group.shapes._recalculate_extents()
        for group in self.native_chart_groups.values():
            group.shapes._recalculate_extents()

    def _paint_to_rgb(self, paint_value, default=None):
        color = parse_color(paint_value)
        if not color:
            return default
        if color[0] == 'grad':
            gradient = self.grads.get(color[1])
            if not gradient:
                return default
            for stop in gradient.get('stops', []):
                stop_color = parse_color(stop.get('color_str'))
                if stop_color and stop_color[0] != 'grad':
                    return RGBColor.from_string(stop_color[0].upper())
            return default
        return RGBColor.from_string(color[0].upper())

    def _normalized_fill_present(self, paint_value):
        normalized = normalize_text(paint_value).lower()
        return normalized not in {'', 'none', 'transparent'}

    def _collect_svg_primitives(self, el, ox=0.0, oy=0.0, scale=1.0, sink=None):
        sink = sink if sink is not None else []
        tag = self._tag(el)

        if tag == 'rect':
            x = (float(el.get('x', 0)) * scale) + ox
            y = (float(el.get('y', 0)) * scale) + oy
            w = float(el.get('width', 0)) * scale
            h = float(el.get('height', 0)) * scale
            bbox = normalize_bbox({'x': x, 'y': y, 'width': w, 'height': h})
            if bbox is not None:
                sink.append({
                    'kind': 'rect',
                    'bbox': bbox,
                    'fill': el.get('fill', ''),
                    'stroke': el.get('stroke', ''),
                    'stroke_width': strip_unit(el.get('stroke-width', '1')),
                    'rx': float(el.get('rx', 0) or 0),
                    'ry': float(el.get('ry', 0) or 0),
                })
            return sink

        if tag == 'circle':
            cx = float(el.get('cx', 0)) * scale + ox
            cy = float(el.get('cy', 0)) * scale + oy
            r = float(el.get('r', 0)) * scale
            bbox = normalize_bbox({'x': cx - r, 'y': cy - r, 'width': 2 * r, 'height': 2 * r})
            if bbox is not None:
                sink.append({
                    'kind': 'circle',
                    'bbox': bbox,
                    'cx': cx,
                    'cy': cy,
                    'r': r,
                    'fill': el.get('fill', ''),
                    'stroke': el.get('stroke', ''),
                    'stroke_width': strip_unit(el.get('stroke-width', '1')),
                    'dasharray': el.get('stroke-dasharray', ''),
                    'dashoffset': el.get('stroke-dashoffset', ''),
                    'transform': el.get('transform', ''),
                    'opacity': float(el.get('opacity', '1') or 1.0),
                })
            return sink

        if tag == 'line':
            x1 = float(el.get('x1', 0)) * scale + ox
            y1 = float(el.get('y1', 0)) * scale + oy
            x2 = float(el.get('x2', 0)) * scale + ox
            y2 = float(el.get('y2', 0)) * scale + oy
            bbox = normalize_bbox({
                'x': min(x1, x2),
                'y': min(y1, y2),
                'width': abs(x2 - x1) or 1.0,
                'height': abs(y2 - y1) or 1.0,
            })
            if bbox is not None:
                sink.append({
                    'kind': 'line',
                    'bbox': bbox,
                    'points': [(x1, y1), (x2, y2)],
                    'stroke': el.get('stroke', ''),
                    'stroke_width': strip_unit(el.get('stroke-width', '1')),
                    'opacity': float(el.get('opacity', '1') or 1.0),
                })
            return sink

        if tag == 'polyline':
            raw_points = parse_svg_point_pairs(el.get('points', ''))
            points = transform_point_pairs(raw_points, ox=ox, oy=oy, scale=scale)
            bbox = bbox_from_point_pairs(points)
            if bbox is not None:
                sink.append({
                    'kind': 'polyline',
                    'bbox': bbox,
                    'points': points,
                    'fill': el.get('fill', ''),
                    'stroke': el.get('stroke', ''),
                    'stroke_width': strip_unit(el.get('stroke-width', '1')),
                    'opacity': float(el.get('opacity', '1') or 1.0),
                })
            return sink

        if tag == 'polygon':
            raw_points = parse_svg_point_pairs(el.get('points', ''))
            points = transform_point_pairs(raw_points, ox=ox, oy=oy, scale=scale)
            bbox = bbox_from_point_pairs(points)
            if bbox is not None:
                sink.append({
                    'kind': 'polygon',
                    'bbox': bbox,
                    'points': points,
                    'fill': el.get('fill', ''),
                    'stroke': el.get('stroke', ''),
                    'stroke_width': strip_unit(el.get('stroke-width', '1')),
                    'opacity': float(el.get('opacity', '1') or 1.0),
                    'closed': True,
                })
            return sink

        if tag == 'path':
            raw_points = parse_path_point_pairs(el.get('d', ''))
            points = transform_point_pairs(raw_points, ox=ox, oy=oy, scale=scale)
            bbox = bbox_from_point_pairs(points)
            if bbox is not None:
                sink.append({
                    'kind': 'path',
                    'bbox': bbox,
                    'points': points,
                    'fill': el.get('fill', ''),
                    'stroke': el.get('stroke', ''),
                    'stroke_width': strip_unit(el.get('stroke-width', '1')),
                    'opacity': float(el.get('opacity', '1') or 1.0),
                    'closed': 'z' in str(el.get('d', '')).lower(),
                })
            return sink

        if tag == 'text':
            fill_s = el.get('fill', el.get('color', ''))
            fsz = el.get('font-size', '14px').replace('px', '')
            fw = el.get('font-weight', '')
            ff = el.get('font-family', '')
            baseline = el.get('dominant-baseline', '')
            anchor = el.get('text-anchor', 'start')
            tspans = list(el.findall(f'{{{SVG_NS}}}tspan'))
            if tspans:
                for ts in tspans:
                    txt = normalize_text(ts.text)
                    if not txt:
                        continue
                    x = float(ts.get('x', el.get('x', 0))) * scale + ox
                    y = float(ts.get('y', el.get('y', 0))) * scale + oy
                    tlen = float(ts.get('textLength', 0))
                    fragment = self._build_text_fragment(
                        text=txt,
                        x=x,
                        y=y,
                        text_length=tlen,
                        font_size=ts.get('font-size', fsz).replace('px', ''),
                        font_weight=ts.get('font-weight', fw),
                        fill_s=ts.get('fill', fill_s),
                        font_family=ts.get('font-family', ff),
                        baseline=baseline,
                        anchor=anchor,
                        opacity=1.0,
                    )
                    if fragment is not None:
                        sink.append({
                            'kind': 'text',
                            'bbox': fragment.get('bbox'),
                            'text': fragment['text'],
                            'fill': ts.get('fill', fill_s),
                            'run': fragment['run'],
                            'align': fragment['align'],
                            'anchor': fragment['anchor'],
                            'font_size_px': fragment['font_size_px'],
                            'numeric_value': extract_numeric_value(fragment['text']),
                            'numeric_like': is_numeric_like_text(fragment['text']),
                        })
            elif el.text and normalize_text(el.text):
                x = float(el.get('x', 0)) * scale + ox
                y = float(el.get('y', 0)) * scale + oy
                fragment = self._build_text_fragment(
                    text=el.text,
                    x=x,
                    y=y,
                    text_length=0,
                    font_size=fsz,
                    font_weight=fw,
                    fill_s=fill_s,
                    font_family=ff,
                    baseline=baseline,
                    anchor=anchor,
                    opacity=1.0,
                )
                if fragment is not None:
                    sink.append({
                        'kind': 'text',
                        'bbox': fragment.get('bbox'),
                        'text': fragment['text'],
                        'fill': fill_s,
                        'run': fragment['run'],
                        'align': fragment['align'],
                        'anchor': fragment['anchor'],
                        'font_size_px': fragment['font_size_px'],
                        'numeric_value': extract_numeric_value(fragment['text']),
                        'numeric_like': is_numeric_like_text(fragment['text']),
                    })
            return sink

        if tag == 'g':
            dx, dy, sx, _sy = self._parse_transform(el)
            child_scale = scale * sx
            new_ox = ox + dx * scale
            new_oy = oy + dy * scale
            for child in el:
                self._collect_svg_primitives(child, new_ox, new_oy, child_scale, sink)
            return sink

        if tag == 'use':
            href = el.get(f'{{{XLINK_NS}}}href') or el.get('href', '')
            if href.startswith('#'):
                ref_id = href[1:]
                ref_el = None
                for candidate in self.root.iter():
                    if candidate.get('id') == ref_id:
                        ref_el = candidate
                        break
                if ref_el is not None:
                    use_x = float(el.get('x', 0)) * scale
                    use_y = float(el.get('y', 0)) * scale
                    self._collect_svg_primitives(ref_el, ox + use_x, oy + use_y, scale, sink)
            return sink

        for child in el:
            self._collect_svg_primitives(child, ox, oy, scale, sink)
        return sink

    def _primitive_in_bbox(self, primitive, region_bbox, min_overlap=0.35):
        primitive_bbox = primitive.get('bbox')
        if not primitive_bbox:
            return False
        overlap = bbox_intersection_area(primitive_bbox, region_bbox)
        if overlap <= 0:
            return False
        primitive_area = max(bbox_area(primitive_bbox), 1.0)
        center = bbox_center(primitive_bbox)
        contains_center = False
        if center is not None:
            contains_center = (
                region_bbox['x'] <= center[0] <= region_bbox['x'] + region_bbox['width']
                and region_bbox['y'] <= center[1] <= region_bbox['y'] + region_bbox['height']
            )
        return contains_center or (overlap / primitive_area) >= min_overlap

    def _primitives_for_chart_region(self, chart_region, kind=None):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return []
        result = []
        for primitive in self.svg_primitives:
            if kind and primitive.get('kind') != kind:
                continue
            if self._primitive_in_bbox(primitive, region_bbox):
                result.append(primitive)
        return result

    def _nearest_text_for_rect(self, rect_bbox, text_items, prefer='below'):
        rect_center = bbox_center(rect_bbox)
        if rect_center is None:
            return None
        best_item = None
        best_score = None
        rect_top = rect_bbox['y']
        rect_bottom = rect_bbox['y'] + rect_bbox['height']
        for item in text_items:
            text_bbox = item.get('bbox')
            text_center = bbox_center(text_bbox)
            if not text_bbox or text_center is None:
                continue
            dx = abs(text_center[0] - rect_center[0])
            max_dx = max(rect_bbox['width'] * 1.8, text_bbox['width'] * 2.0, 48.0)
            if dx > max_dx:
                continue
            if prefer == 'above':
                penalty = rect_bbox['height'] if text_center[1] > rect_top + rect_bbox['height'] * 0.75 else 0.0
                dy = abs((text_bbox['y'] + text_bbox['height']) - rect_top)
            else:
                penalty = rect_bbox['height'] if text_center[1] < rect_bottom - text_bbox['height'] else 0.0
                dy = abs(text_bbox['y'] - rect_bottom)
            score = dx + dy + penalty
            if best_score is None or score < best_score:
                best_item = item
                best_score = score
        if best_score is None or best_score > max(rect_bbox['height'] * 2.5, rect_bbox['width'] * 2.0, 96.0):
            return None
        return best_item

    def _nearest_text_same_row(self, rect_bbox, text_items, side='left'):
        rect_center = bbox_center(rect_bbox)
        if rect_center is None:
            return None
        best_item = None
        best_score = None
        for item in text_items:
            text_bbox = item.get('bbox')
            text_center = bbox_center(text_bbox)
            if not text_bbox or text_center is None:
                continue
            dy = abs(text_center[1] - rect_center[1])
            max_dy = max(rect_bbox['height'] * 2.5, text_bbox['height'] * 2.5, 40.0)
            if dy > max_dy:
                continue
            if side == 'right':
                if text_bbox['x'] + text_bbox['width'] < rect_bbox['x'] + rect_bbox['width'] * 0.5:
                    continue
                dx = abs(text_bbox['x'] - (rect_bbox['x'] + rect_bbox['width']))
            else:
                if text_bbox['x'] > rect_bbox['x'] + rect_bbox['width'] * 0.35:
                    continue
                dx = abs(rect_bbox['x'] - (text_bbox['x'] + text_bbox['width']))
            score = dy + dx
            if best_score is None or score < best_score:
                best_item = item
                best_score = score
        if best_score is None or best_score > max(rect_bbox['width'] * 1.5, 120.0):
            return None
        return best_item

    def _extract_comparison_bar_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        rects = self._primitives_for_chart_region(chart_region, kind='rect')
        texts = self._primitives_for_chart_region(chart_region, kind='text')
        chart_w = region_bbox['width']
        chart_h = region_bbox['height']
        candidate_rects = []
        for rect in rects:
            bbox = rect.get('bbox')
            if not bbox:
                continue
            if bbox['width'] < max(10.0, chart_w * 0.08) or bbox['height'] < max(14.0, chart_h * 0.18):
                continue
            if bbox['width'] > chart_w * 0.62 or bbox['height'] > chart_h * 0.98:
                continue
            if bbox['height'] < max(28.0, bbox['width'] * 0.75):
                continue
            candidate_rects.append(rect)
        if len(candidate_rects) < 2:
            return None

        max_bottom = max(rect['bbox']['y'] + rect['bbox']['height'] for rect in candidate_rects)
        bottom_tolerance = max(10.0, chart_h * 0.12)
        bottom_aligned = [
            rect for rect in candidate_rects
            if abs((rect['bbox']['y'] + rect['bbox']['height']) - max_bottom) <= bottom_tolerance
        ]
        if len(bottom_aligned) >= 2:
            candidate_rects = bottom_aligned

        candidate_rects.sort(key=lambda rect: (rect['bbox']['height'], rect['bbox']['width']), reverse=True)
        selected = []
        min_center_gap = max(18.0, chart_w * 0.12)
        for rect in candidate_rects:
            center = bbox_center(rect['bbox'])
            if center is None:
                continue
            if all(abs(center[0] - bbox_center(existing['bbox'])[0]) >= min_center_gap for existing in selected):
                selected.append(rect)
            if len(selected) == 2:
                break
        if len(selected) < 2:
            return None

        selected.sort(key=lambda rect: bbox_center(rect['bbox'])[0])
        numeric_texts = [item for item in texts if item.get('numeric_like')]
        label_texts = [item for item in texts if not item.get('numeric_like')]
        numeric_values_found = 0
        categories = []
        values = []
        point_colors = []
        for index, rect in enumerate(selected, start=1):
            rect_bbox = rect['bbox']
            value_item = self._nearest_text_for_rect(rect_bbox, numeric_texts, prefer='above')
            if value_item and value_item.get('numeric_value') is not None:
                numeric_values_found += 1
                value = max(float(value_item['numeric_value']), 0.0)
            else:
                value = max(float(rect_bbox['height']), 1.0)
            values.append(value)

            label_item = self._nearest_text_for_rect(rect_bbox, label_texts, prefer='below')
            categories.append(label_item.get('text') if label_item else f'Item {index}')
            default_color = RGBColor(203, 213, 225) if index == 1 else RGBColor(37, 99, 235)
            point_colors.append(self._paint_to_rgb(rect.get('fill'), default=default_color))

        if numeric_values_found == 0:
            max_value = max(values) if values else 0.0
            if max_value <= 0:
                return None
            values = [round((value / max_value) * 100, 2) for value in values]

        if len(set(categories)) != len(categories):
            categories = [f'Item {index}' for index in range(1, len(values) + 1)]

        plot_bbox = merge_bboxes([rect['bbox'] for rect in selected])
        if plot_bbox is None:
            return None
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-comparison',
            'chart_type': 'comparison_bar',
            'plot_bbox': plot_bbox,
            'values': values,
            'categories': categories,
            'point_colors': point_colors,
        }

    def _extract_progress_bar_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        rects = self._primitives_for_chart_region(chart_region, kind='rect')
        texts = self._primitives_for_chart_region(chart_region, kind='text')
        chart_w = region_bbox['width']
        chart_h = region_bbox['height']
        candidate_rects = []
        for rect in rects:
            bbox = rect.get('bbox')
            if not bbox:
                continue
            if bbox['width'] < max(32.0, chart_w * 0.25) or bbox['height'] < max(6.0, chart_h * 0.05):
                continue
            if bbox['width'] <= bbox['height'] * 1.8:
                continue
            if bbox['height'] > chart_h * 0.55:
                continue
            candidate_rects.append(rect)
        if not candidate_rects:
            return None

        candidate_rects.sort(key=lambda rect: (rect['bbox']['width'], rect['bbox']['height']), reverse=True)
        track = candidate_rects[0]
        track_bbox = track['bbox']
        fill_candidates = []
        for rect in candidate_rects[1:]:
            bbox = rect['bbox']
            vertical_overlap = bbox_intersection_area(
                {'x': track_bbox['x'], 'y': track_bbox['y'], 'width': track_bbox['width'], 'height': track_bbox['height']},
                {'x': bbox['x'], 'y': bbox['y'], 'width': bbox['width'], 'height': bbox['height']},
            )
            if vertical_overlap <= 0:
                continue
            if abs(bbox['y'] - track_bbox['y']) > max(8.0, track_bbox['height'] * 0.9):
                continue
            if abs(bbox['height'] - track_bbox['height']) > max(8.0, track_bbox['height'] * 0.9):
                continue
            if bbox['width'] > track_bbox['width'] * 1.02:
                continue
            if bbox['x'] > track_bbox['x'] + max(10.0, track_bbox['width'] * 0.15):
                continue
            fill_candidates.append(rect)
        fill = max(fill_candidates, key=lambda rect: rect['bbox']['width']) if fill_candidates else None

        percent_texts = [item for item in texts if '%' in str(item.get('text', '')) and item.get('numeric_value') is not None]
        numeric_texts = percent_texts or [item for item in texts if item.get('numeric_like') and item.get('numeric_value') is not None]
        progress_value = None
        if numeric_texts:
            region_center = bbox_center(region_bbox)
            progress_item = min(
                numeric_texts,
                key=lambda item: (
                    abs(bbox_center(item['bbox'])[1] - region_center[1]) if bbox_center(item['bbox']) else 9999,
                    abs(bbox_center(item['bbox'])[0] - region_center[0]) if bbox_center(item['bbox']) else 9999,
                ),
            )
            progress_value = progress_item.get('numeric_value')
        if progress_value is None and fill is not None:
            progress_value = round((fill['bbox']['width'] / max(track_bbox['width'], 1.0)) * 100, 2)
        if progress_value is None:
            return None
        progress_value = clamp(float(progress_value), 0.0, 100.0)

        fill_color = self._paint_to_rgb(fill.get('fill') if fill else None, default=RGBColor(37, 99, 235))
        track_color = self._paint_to_rgb(track.get('fill'), default=RGBColor(226, 232, 240))
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-progress',
            'chart_type': 'progress_bar',
            'plot_bbox': track_bbox,
            'progress_value': progress_value,
            'fill_color': fill_color,
            'track_color': track_color,
        }

    def _extract_stacked_bar_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        rects = self._primitives_for_chart_region(chart_region, kind='rect')
        texts = self._primitives_for_chart_region(chart_region, kind='text')
        chart_w = region_bbox['width']
        chart_h = region_bbox['height']
        candidate_rects = []
        for rect in rects:
            bbox = rect.get('bbox')
            if not bbox:
                continue
            if bbox['width'] < max(10.0, chart_w * 0.04):
                continue
            if bbox['height'] < max(8.0, chart_h * 0.03) or bbox['height'] > max(40.0, chart_h * 0.22):
                continue
            if bbox['width'] <= bbox['height'] * 1.2:
                continue
            candidate_rects.append(rect)
        if len(candidate_rects) < 2:
            return None

        candidate_rects.sort(key=lambda rect: bbox_center(rect['bbox'])[1])
        heights = sorted(rect['bbox']['height'] for rect in candidate_rects)
        median_height = heights[len(heights) // 2] if heights else 12.0
        row_tolerance = max(10.0, median_height * 1.4)
        rows = []
        for rect in candidate_rects:
            center_y = bbox_center(rect['bbox'])[1]
            if not rows or abs(center_y - rows[-1]['center_y']) > row_tolerance:
                rows.append({'center_y': center_y, 'items': [rect]})
            else:
                rows[-1]['items'].append(rect)
                item_count = len(rows[-1]['items'])
                rows[-1]['center_y'] = ((rows[-1]['center_y'] * (item_count - 1)) + center_y) / item_count

        best_segments = None
        best_span = 0.0
        for row in rows:
            items = sorted(row['items'], key=lambda rect: rect['bbox']['x'])
            filtered = []
            for rect in items:
                rect_bbox = rect['bbox']
                contains_others = 0
                rect_area = max(bbox_area(rect_bbox), 1.0)
                for other in items:
                    if other is rect:
                        continue
                    other_bbox = other['bbox']
                    overlap = bbox_intersection_area(rect_bbox, other_bbox)
                    other_area = max(bbox_area(other_bbox), 1.0)
                    if overlap / other_area >= 0.9 and rect_area > other_area * 1.25:
                        contains_others += 1
                if contains_others >= 2:
                    continue
                filtered.append(rect)
            items = filtered if len(filtered) >= 2 else items
            if len(items) < 2:
                continue

            clusters = []
            cluster = [items[0]]
            for rect in items[1:]:
                prev_bbox = cluster[-1]['bbox']
                bbox = rect['bbox']
                gap = bbox['x'] - (prev_bbox['x'] + prev_bbox['width'])
                height_delta = abs(bbox['height'] - prev_bbox['height'])
                if gap <= max(16.0, median_height * 1.4) and height_delta <= max(8.0, median_height * 0.8):
                    cluster.append(rect)
                else:
                    clusters.append(cluster)
                    cluster = [rect]
            if cluster:
                clusters.append(cluster)

            for segment_cluster in clusters:
                if len(segment_cluster) < 2:
                    continue
                span_start = segment_cluster[0]['bbox']['x']
                span_end = segment_cluster[-1]['bbox']['x'] + segment_cluster[-1]['bbox']['width']
                span = span_end - span_start
                total_width = sum(item['bbox']['width'] for item in segment_cluster)
                coverage = span / max(chart_w, 1.0)
                compactness = total_width / max(span, 1.0)
                if coverage < 0.25 or compactness < 0.72:
                    continue
                if span > best_span:
                    best_span = span
                    best_segments = segment_cluster

        if not best_segments or len(best_segments) < 2:
            return None

        def color_key(color_value):
            return rgb_to_hex(color_value) or ''

        plot_bbox = merge_bboxes([segment['bbox'] for segment in best_segments])
        if plot_bbox is None:
            return None

        non_numeric_texts = [item for item in texts if not item.get('numeric_like')]
        legend_labels = {}
        for rect in rects:
            if rect in best_segments:
                continue
            bbox = rect.get('bbox')
            if not bbox:
                continue
            if bbox['width'] > 28 or bbox['height'] > 28:
                continue
            if abs(bbox['width'] - bbox['height']) > max(6.0, bbox['height'] * 0.8):
                continue
            swatch_color = self._paint_to_rgb(rect.get('fill'))
            if swatch_color is None:
                continue
            label_item = self._nearest_text_same_row(bbox, non_numeric_texts, side='right')
            if label_item:
                legend_labels[color_key(swatch_color)] = label_item.get('text')

        total_width = sum(segment['bbox']['width'] for segment in best_segments)
        if total_width <= 0:
            return None
        series = []
        fallback_colors = [
            RGBColor(37, 99, 235),
            RGBColor(14, 165, 233),
            RGBColor(34, 197, 94),
            RGBColor(245, 158, 11),
            RGBColor(239, 68, 68),
        ]
        for index, segment in enumerate(best_segments, start=1):
            fill_color = self._paint_to_rgb(segment.get('fill'), default=fallback_colors[(index - 1) % len(fallback_colors)])
            series.append({
                'name': legend_labels.get(color_key(fill_color)) or f'Segment {index}',
                'value': round((segment['bbox']['width'] / total_width) * 100, 2),
                'color': fill_color,
            })
        total_value = sum(item['value'] for item in series)
        if total_value <= 0:
            return None
        scale = 100.0 / total_value
        for item in series:
            item['value'] = round(item['value'] * scale, 2)
        remainder = round(100.0 - sum(item['value'] for item in series), 2)
        series[-1]['value'] = max(0.0, round(series[-1]['value'] + remainder, 2))
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-stacked-bar',
            'chart_type': 'stacked_bar',
            'plot_bbox': plot_bbox,
            'series': series,
        }

    def _extract_kpi_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        text_items = self._primitives_for_chart_region(chart_region, kind='text')
        if not text_items:
            return None
        numeric_items = [item for item in text_items if item.get('numeric_like')]
        main_value = max(
            numeric_items or text_items,
            key=lambda item: (float(item.get('font_size_px') or 0), bbox_area(item.get('bbox') or {})),
        )
        delta_candidates = [
            item for item in text_items
            if item is not main_value and normalize_text(item.get('text'))
            and (
                '%' in str(item.get('text', ''))
                or str(item.get('text', '')).strip().startswith('+')
                or str(item.get('text', '')).strip().startswith('-')
            )
        ]
        main_bbox = main_value.get('bbox') or region_bbox
        delta_item = None
        if delta_candidates:
            main_center = bbox_center(main_bbox) or bbox_center(region_bbox)
            delta_item = min(
                delta_candidates,
                key=lambda item: (
                    abs((bbox_center(item.get('bbox')) or (0, 0))[1] - main_center[1]),
                    abs((bbox_center(item.get('bbox')) or (0, 0))[0] - main_center[0]),
                ),
            )
        trend = 'flat'
        trend_text = normalize_text(delta_item.get('text')) if delta_item else ''
        if trend_text.startswith('+') or '增长' in trend_text or '上升' in trend_text:
            trend = 'up'
        elif trend_text.startswith('-') or '下降' in trend_text or '下滑' in trend_text:
            trend = 'down'

        trend_color = self._paint_to_rgb(delta_item.get('fill') if delta_item else None)
        if trend_color is None:
            if trend == 'up':
                trend_color = RGBColor(34, 197, 94)
            elif trend == 'down':
                trend_color = RGBColor(239, 68, 68)
            else:
                trend_color = RGBColor(100, 116, 139)

        if delta_item and delta_item.get('bbox'):
            delta_bbox = delta_item['bbox']
            size = clamp(min(delta_bbox['height'] * 0.85, region_bbox['height'] * 0.18), 10.0, 24.0)
            arrow_bbox = normalize_bbox({
                'x': max(region_bbox['x'], delta_bbox['x'] - size - 6.0),
                'y': delta_bbox['y'] + max((delta_bbox['height'] - size) / 2, 0.0),
                'width': size,
                'height': size,
            })
        else:
            size = clamp(min(main_bbox['height'] * 0.35, region_bbox['height'] * 0.18), 10.0, 24.0)
            arrow_bbox = normalize_bbox({
                'x': min(region_bbox['x'] + region_bbox['width'] - size, main_bbox['x'] + main_bbox['width'] + 8.0),
                'y': main_bbox['y'] + max(main_bbox['height'] * 0.25, 0.0),
                'width': size,
                'height': size,
            })
        if arrow_bbox is None:
            return None
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-kpi',
            'chart_type': 'kpi',
            'plot_bbox': arrow_bbox,
            'trend': trend,
            'trend_color': trend_color,
        }

    def _extract_ring_chart(self, chart_region):
        circles = self._primitives_for_chart_region(chart_region, kind='circle')
        if len(circles) < 2:
            return None
        progress_candidates = [circle for circle in circles if normalize_text(circle.get('dasharray'))]
        if not progress_candidates:
            return None
        progress_circle = max(progress_candidates, key=lambda circle: circle.get('r') or 0)
        progress_center = (progress_circle.get('cx'), progress_circle.get('cy'))
        track_candidates = [
            circle for circle in circles
            if circle is not progress_circle and not normalize_text(circle.get('dasharray')) and circle.get('stroke')
        ]
        track_circle = None
        if track_candidates:
            track_circle = min(
                track_candidates,
                key=lambda circle: (
                    abs((circle.get('cx') or 0) - (progress_center[0] or 0))
                    + abs((circle.get('cy') or 0) - (progress_center[1] or 0))
                    + abs((circle.get('r') or 0) - (progress_circle.get('r') or 0))
                ),
            )
        plot_bbox = merge_bboxes([
            progress_circle.get('bbox'),
            track_circle.get('bbox') if track_circle else None,
        ])
        if plot_bbox is None:
            return None
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-ring',
            'chart_type': 'ring',
            'plot_bbox': plot_bbox,
            'progress_circle': progress_circle,
            'track_circle': track_circle,
        }

    def _extract_metric_row_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        rects = self._primitives_for_chart_region(chart_region, kind='rect')
        texts = self._primitives_for_chart_region(chart_region, kind='text')
        chart_w = region_bbox['width']
        chart_h = region_bbox['height']
        candidate_rects = []
        for rect in rects:
            bbox = rect.get('bbox')
            if not bbox:
                continue
            if bbox['width'] < max(40.0, chart_w * 0.16):
                continue
            if bbox['height'] < max(5.0, chart_h * 0.03) or bbox['height'] > max(26.0, chart_h * 0.2):
                continue
            if bbox['width'] <= bbox['height'] * 2.2:
                continue
            candidate_rects.append(rect)
        if len(candidate_rects) < 4:
            return None

        candidate_rects.sort(key=lambda rect: bbox_center(rect['bbox'])[1])
        heights = sorted(rect['bbox']['height'] for rect in candidate_rects)
        median_height = heights[len(heights) // 2] if heights else 8.0
        row_tolerance = max(10.0, median_height * 1.5)
        rows = []
        for rect in candidate_rects:
            center_y = bbox_center(rect['bbox'])[1]
            if not rows or abs(center_y - rows[-1]['center_y']) > row_tolerance:
                rows.append({'center_y': center_y, 'items': [rect]})
            else:
                rows[-1]['items'].append(rect)
                item_count = len(rows[-1]['items'])
                rows[-1]['center_y'] = ((rows[-1]['center_y'] * (item_count - 1)) + center_y) / item_count

        row_specs = []
        for row in rows:
            items = sorted(row['items'], key=lambda rect: rect['bbox']['width'], reverse=True)
            if not items:
                continue
            track = items[0]
            track_bbox = track['bbox']
            fill_candidates = []
            for rect in items[1:]:
                bbox = rect['bbox']
                if abs(bbox['y'] - track_bbox['y']) > max(8.0, track_bbox['height'] * 0.9):
                    continue
                if abs(bbox['height'] - track_bbox['height']) > max(8.0, track_bbox['height'] * 0.9):
                    continue
                if bbox['width'] > track_bbox['width'] * 1.02:
                    continue
                if bbox['x'] > track_bbox['x'] + max(10.0, track_bbox['width'] * 0.12):
                    continue
                fill_candidates.append(rect)
            fill = max(fill_candidates, key=lambda rect: rect['bbox']['width']) if fill_candidates else None
            radius_px = max(
                float(track.get('rx') or 0),
                float(track.get('ry') or 0),
                min(track_bbox['height'] / 2, 8.0),
            )
            progress_value = None
            if fill is not None and track_bbox['width'] > 0:
                progress_value = round((fill['bbox']['width'] / track_bbox['width']) * 100, 2)
            label_item = self._nearest_text_same_row(
                track_bbox,
                [item for item in texts if not item.get('numeric_like')],
                side='left',
            )
            value_item = self._nearest_text_same_row(
                track_bbox,
                [item for item in texts if item.get('numeric_like')],
                side='right',
            )
            row_specs.append({
                'track_bbox': track_bbox,
                'fill_bbox': fill.get('bbox') if fill else None,
                'track_color': self._paint_to_rgb(track.get('fill'), default=RGBColor(226, 232, 240)),
                'fill_color': self._paint_to_rgb(fill.get('fill') if fill else None, default=RGBColor(37, 99, 235)),
                'radius_px': radius_px,
                'label': label_item.get('text') if label_item else None,
                'value_text': value_item.get('text') if value_item else None,
                'progress_value': progress_value,
            })
        if len(row_specs) < 2:
            return None
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-metric-row',
            'chart_type': 'metric_row',
            'rows': row_specs,
            'plot_bbox': merge_bboxes([row['track_bbox'] for row in row_specs]),
        }

    def _ring_arc_geometry(self, circle):
        dasharray = normalize_text(circle.get('dasharray'))
        if not dasharray:
            return None
        dash_parts = [float(strip_unit(part.strip())) for part in dasharray.replace(',', ' ').split() if part.strip()]
        radius = float(circle.get('r') or 0)
        if len(dash_parts) < 2 or radius <= 0:
            return None
        circumference = 2 * math.pi * radius
        arc_len = dash_parts[0]
        sweep_angle = clamp((arc_len / circumference) * 360, 0.0, 360.0)
        start_angle = 0.0
        transform = circle.get('transform', '') or ''
        rot_m = re.search(r'rotate\(\s*([\d.\-]+)', transform)
        if rot_m:
            start_angle = float(rot_m.group(1))
        dashoffset = circle.get('dashoffset')
        if dashoffset:
            try:
                start_angle += -(float(strip_unit(dashoffset)) / circumference) * 360
            except (ValueError, ZeroDivisionError):
                pass
        ppt_start = (start_angle + 90) % 360
        ppt_end = (ppt_start + sweep_angle) % 360
        return ppt_start, ppt_end

    def _append_ring_outline_shape(self, group, circle, color_value):
        circle_bbox = circle.get('bbox')
        if not circle_bbox:
            return
        line_el = self._line_from_rgb(color_value, float(circle.get('stroke_width') or 1), opacity=float(circle.get('opacity') or 1.0))
        if line_el is None:
            return
        shape = make_shape(
            self._id(),
            f'NR{self.sid}',
            px(circle_bbox['x']),
            px(circle_bbox['y']),
            px(circle_bbox['width']),
            px(circle_bbox['height']),
            preset='ellipse',
            fill_el=_el('a:noFill'),
            line_el=line_el,
        )
        self._append_group_shape_element(group, shape)

    def _append_ring_arc_shape(self, group, circle, color_value):
        circle_bbox = circle.get('bbox')
        arc_geometry = self._ring_arc_geometry(circle)
        if not circle_bbox or arc_geometry is None:
            return
        ppt_start, ppt_end = arc_geometry
        geom = _el('a:prstGeom', {'prst': 'arc'})
        av = _el('a:avLst')
        av.append(_el('a:gd', {'name': 'adj1', 'fmla': f'val {int(ppt_start * 60000)}'}))
        av.append(_el('a:gd', {'name': 'adj2', 'fmla': f'val {int(ppt_end * 60000)}'}))
        geom.append(av)
        line_el = self._line_from_rgb(color_value, float(circle.get('stroke_width') or 1), opacity=float(circle.get('opacity') or 1.0))
        if line_el is None:
            return
        shape = _el('p:sp')
        shape.append(_el('p:nvSpPr', children=[
            _el('p:cNvPr', {'id': str(self._id()), 'name': f'NA{self.sid}'}),
            _el('p:cNvSpPr'), _el('p:nvPr'),
        ]))
        sp_pr = _el('p:spPr')
        sp_pr.append(_el('a:xfrm', children=[
            _el('a:off', {'x': str(px(circle_bbox['x'])), 'y': str(px(circle_bbox['y']))}),
            _el('a:ext', {'cx': str(px(circle_bbox['width'])), 'cy': str(px(circle_bbox['height']))}),
        ]))
        sp_pr.append(geom)
        sp_pr.append(_el('a:noFill'))
        sp_pr.append(line_el)
        shape.append(sp_pr)
        self._append_group_shape_element(group, shape)

    def _append_metric_row_bar_shape(self, group, bbox, color_value, radius_px):
        if not bbox:
            return
        preset = 'roundRect' if radius_px > 0 else 'rect'
        shape = make_shape(
            self._id(),
            f'MR{self.sid}',
            px(bbox['x']),
            px(bbox['y']),
            px(bbox['width']),
            px(bbox['height']),
            preset=preset,
            fill_el=self._solid_fill_from_rgb(color_value),
            line_el=None,
            rx=px(radius_px),
        )
        self._append_group_shape_element(group, shape)

    def _append_native_basic_shape(self, group, bbox, *, preset='ellipse', fill_color=None,
                                   fill_opacity=1.0, line_color=None, line_width_px=1.0,
                                   radius_px=0.0, name_prefix='NS'):
        if not bbox:
            return
        fill_el = self._solid_fill_from_rgb(fill_color, opacity=fill_opacity) if fill_color else _el('a:noFill')
        line_el = self._line_from_rgb(line_color, line_width_px) if line_color else None
        shape = make_shape(
            self._id(),
            f'{name_prefix}{self.sid}',
            px(bbox['x']),
            px(bbox['y']),
            px(bbox['width']),
            px(bbox['height']),
            preset=preset,
            fill_el=fill_el,
            line_el=line_el,
            rx=px(radius_px),
        )
        self._append_group_shape_element(group, shape)

    def _append_native_path_shape(self, group, points, *, close=False, fill_color=None,
                                  fill_opacity=1.0, line_color=None, line_width_px=1.0,
                                  line_opacity=1.0, name_prefix='NP'):
        points = dedupe_adjacent_points(points)
        if len(points) < 2:
            return None
        bbox = bbox_from_point_pairs(points)
        if bbox is None:
            return None
        d_parts = [f'M {points[0][0]} {points[0][1]}']
        for point in points[1:]:
            d_parts.append(f'L {point[0]} {point[1]}')
        if close:
            d_parts.append('Z')
        geom_el = parse_path_to_custgeom(' '.join(d_parts), (bbox['x'], bbox['y'], bbox['width'], bbox['height']))
        fill_el = self._solid_fill_from_rgb(fill_color, opacity=fill_opacity) if fill_color else _el('a:noFill')
        line_el = self._line_from_rgb(line_color, line_width_px, opacity=line_opacity) if line_color else None
        shape = make_shape(
            self._id(),
            f'{name_prefix}{self.sid}',
            px(bbox['x']),
            px(bbox['y']),
            px(bbox['width']),
            px(bbox['height']),
            fill_el=fill_el,
            line_el=line_el,
            geom_el=geom_el,
        )
        self._append_group_shape_element(group, shape)
        return bbox

    def _append_native_text_primitive(self, group, primitive):
        bbox = primitive.get('bbox')
        run = primitive.get('run')
        if not bbox or not run:
            return
        shape = make_textbox(
            self._id(),
            f'NT{self.sid}',
            px(bbox['x']),
            px(bbox['y']),
            max(1, px(bbox['width'])),
            max(1, px(bbox['height'])),
            [self._build_paragraph([run], primitive.get('align', 'l'))],
            anchor=primitive.get('anchor', 't'),
        )
        self._append_group_shape_element(group, shape)

    def _extract_native_group_chart(self, chart_region, chart_type, *, allowed_kinds, min_primitives,
                                    require_text=False, require_any_of=None, require_closed_paths=0):
        primitives = [
            primitive
            for primitive in self._primitives_for_chart_region(chart_region)
            if primitive.get('kind') in allowed_kinds
        ]
        if len(primitives) < min_primitives:
            return None
        kinds = {primitive.get('kind') for primitive in primitives}
        if require_text and 'text' not in kinds:
            return None
        if require_any_of and not any(any(kind in kinds for kind in option_group) for option_group in require_any_of):
            return None
        if require_closed_paths:
            closed_count = sum(
                1
                for primitive in primitives
                if primitive.get('kind') in {'polygon', 'path'} and primitive.get('closed')
            )
            if closed_count < require_closed_paths:
                return None
        plot_bbox = merge_bboxes([primitive.get('bbox') for primitive in primitives if primitive.get('bbox')])
        if plot_bbox is None:
            return None
        return {
            'chart_id': chart_region.get('chart_id') or f'chart-{chart_type}',
            'chart_type': chart_type,
            'plot_bbox': plot_bbox,
            'primitives': primitives,
        }

    def _extract_timeline_chart(self, chart_region):
        return self._extract_native_group_chart(
            chart_region,
            'timeline',
            allowed_kinds={'line', 'polyline', 'path', 'circle', 'rect', 'text'},
            min_primitives=5,
            require_text=True,
            require_any_of=[{'line', 'polyline', 'path'}, {'circle', 'rect'}],
        )

    def _extract_funnel_chart(self, chart_region):
        return self._extract_native_group_chart(
            chart_region,
            'funnel',
            allowed_kinds={'rect', 'polygon', 'path', 'text'},
            min_primitives=5,
            require_text=True,
            require_any_of=[{'rect', 'polygon', 'path'}],
            require_closed_paths=1,
        )

    def _extract_radar_chart(self, chart_region):
        return self._extract_native_group_chart(
            chart_region,
            'radar',
            allowed_kinds={'polygon', 'path', 'line', 'circle', 'text'},
            min_primitives=6,
            require_any_of=[{'polygon', 'path'}, {'line'}],
            require_closed_paths=2,
        )

    def _extract_waffle_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        region_area = max(bbox_area(region_bbox), 1.0)
        candidates = []
        for primitive in self._primitives_for_chart_region(chart_region):
            if primitive.get('kind') not in {'rect', 'circle'}:
                continue
            bbox = primitive.get('bbox')
            if not bbox:
                continue
            if bbox_area(bbox) > region_area * 0.03:
                continue
            if bbox['width'] < 4.0 or bbox['height'] < 4.0:
                continue
            if bbox['width'] > max(18.0, region_bbox['width'] * 0.12):
                continue
            if bbox['height'] > max(18.0, region_bbox['height'] * 0.12):
                continue
            ratio = bbox['width'] / max(bbox['height'], 1.0)
            if ratio < 0.55 or ratio > 1.8:
                continue
            candidates.append(primitive)
        if len(candidates) < 12 or len(candidates) > 140:
            return None
        text_items = self._primitives_for_chart_region(chart_region, kind='text')
        plot_bbox = merge_bboxes([primitive.get('bbox') for primitive in candidates])
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-waffle',
            'chart_type': 'waffle',
            'plot_bbox': plot_bbox,
            'primitives': candidates + text_items,
        }

    def _extract_treemap_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        region_area = max(bbox_area(region_bbox), 1.0)
        rects = self._primitives_for_chart_region(chart_region, kind='rect')
        candidate_rects = []
        for rect in rects:
            bbox = rect.get('bbox')
            if not bbox:
                continue
            if bbox_area(bbox) < region_area * 0.025:
                continue
            if bbox['width'] < max(24.0, region_bbox['width'] * 0.08):
                continue
            if bbox['height'] < max(18.0, region_bbox['height'] * 0.08):
                continue
            if bbox['width'] > region_bbox['width'] * 0.99 and bbox['height'] > region_bbox['height'] * 0.99:
                continue
            candidate_rects.append(rect)
        if len(candidate_rects) < 3:
            return None

        filtered_rects = []
        for rect in candidate_rects:
            rect_bbox = rect['bbox']
            rect_area = max(bbox_area(rect_bbox), 1.0)
            contained = 0
            for other in candidate_rects:
                if other is rect:
                    continue
                other_bbox = other['bbox']
                overlap = bbox_intersection_area(rect_bbox, other_bbox)
                other_area = max(bbox_area(other_bbox), 1.0)
                if overlap / other_area >= 0.88 and rect_area > other_area * 1.15:
                    contained += 1
            if contained >= 2:
                continue
            filtered_rects.append(rect)
        if len(filtered_rects) < 3:
            filtered_rects = candidate_rects

        text_items = []
        for primitive in self._primitives_for_chart_region(chart_region, kind='text'):
            text_bbox = primitive.get('bbox')
            if not text_bbox:
                continue
            if any(
                bbox_intersection_area(text_bbox, rect['bbox']) / max(bbox_area(text_bbox), 1.0) >= 0.5
                for rect in filtered_rects
            ):
                text_items.append(primitive)

        plot_bbox = merge_bboxes([primitive.get('bbox') for primitive in filtered_rects])
        if plot_bbox is None:
            return None
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-treemap',
            'chart_type': 'treemap',
            'plot_bbox': plot_bbox,
            'primitives': filtered_rects + text_items,
        }

    def _sparkline_area_candidate(self, chart_region, line_bbox):
        candidates = []
        for primitive in self._primitives_for_chart_region(chart_region, kind='path'):
            bbox = primitive.get('bbox')
            if not bbox or not primitive.get('closed'):
                continue
            if not self._normalized_fill_present(primitive.get('fill')):
                continue
            if bbox['width'] < line_bbox['width'] * 0.85:
                continue
            overlap = bbox_intersection_area(bbox, line_bbox)
            if overlap <= 0:
                continue
            candidates.append((overlap / max(bbox_area(line_bbox), 1.0), primitive))
        if not candidates:
            return None
        candidates.sort(key=lambda item: item[0], reverse=True)
        return candidates[0][1]

    def _extract_sparkline_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        chart_w = region_bbox['width']
        chart_h = region_bbox['height']
        path_candidates = [
            primitive
            for primitive in self._primitives_for_chart_region(chart_region, kind='path')
            if primitive.get('stroke') and not primitive.get('closed')
        ]
        line_candidates = self._primitives_for_chart_region(chart_region, kind='polyline') + path_candidates
        scored = []
        for primitive in line_candidates:
            bbox = primitive.get('bbox')
            points = dedupe_adjacent_points(primitive.get('points'))
            if not bbox or len(points) < 3:
                continue
            if bbox['width'] < max(32.0, chart_w * 0.35):
                continue
            if bbox['height'] > chart_h * 0.95:
                continue
            score = (bbox['width'] / max(chart_w, 1.0)) + min(len(points), 12) * 0.04
            if primitive.get('stroke'):
                score += 0.25
            if self._normalized_fill_present(primitive.get('fill')):
                score -= 0.1
            scored.append((score, primitive, points))
        if not scored:
            return None
        scored.sort(key=lambda item: item[0], reverse=True)
        _, line_primitive, line_points = scored[0]
        line_bbox = line_primitive.get('bbox')
        if line_bbox is None:
            return None

        line_color = self._paint_to_rgb(
            line_primitive.get('stroke') or line_primitive.get('fill'),
            default=RGBColor(37, 99, 235),
        )
        area_primitive = self._sparkline_area_candidate(chart_region, line_bbox)
        area_points = dedupe_adjacent_points(area_primitive.get('points')) if area_primitive else None
        area_color = self._paint_to_rgb(area_primitive.get('fill') if area_primitive else None, default=line_color)
        area_opacity = float(area_primitive.get('opacity') or 0.12) if area_primitive else 0.12

        last_point = line_points[-1]
        marker_candidates = [
            primitive for primitive in self._primitives_for_chart_region(chart_region, kind='circle')
            if float(primitive.get('r') or 0) <= max(8.0, chart_h * 0.1)
        ]
        marker = None
        if marker_candidates:
            marker = min(
                marker_candidates,
                key=lambda primitive: abs((primitive.get('cx') or 0) - last_point[0]) + abs((primitive.get('cy') or 0) - last_point[1]),
            )
            if abs((marker.get('cx') or 0) - last_point[0]) > max(12.0, chart_w * 0.08):
                marker = None
        marker_bbox = marker.get('bbox') if marker else normalize_bbox({
            'x': last_point[0] - clamp(float(line_primitive.get('stroke_width') or 2.0) * 1.4, 3.0, 6.0),
            'y': last_point[1] - clamp(float(line_primitive.get('stroke_width') or 2.0) * 1.4, 3.0, 6.0),
            'width': clamp(float(line_primitive.get('stroke_width') or 2.0) * 2.8, 6.0, 12.0),
            'height': clamp(float(line_primitive.get('stroke_width') or 2.0) * 2.8, 6.0, 12.0),
        })
        return {
            'chart_id': chart_region.get('chart_id') or 'chart-sparkline',
            'chart_type': 'sparkline',
            'plot_bbox': merge_bboxes([line_bbox, marker_bbox]),
            'line_points': line_points,
            'line_color': line_color,
            'line_width_px': max(float(line_primitive.get('stroke_width') or 2.0), 1.5),
            'area_points': area_points,
            'area_color': area_color,
            'area_opacity': clamp(area_opacity, 0.05, 0.35),
            'marker_bbox': marker_bbox,
        }

    def _extract_rating_chart(self, chart_region):
        region_bbox = chart_region.get('bbox')
        if not region_bbox:
            return None
        chart_w = region_bbox['width']
        chart_h = region_bbox['height']

        def select_row(primitives, primitive_kind):
            filtered = []
            for primitive in primitives:
                bbox = primitive.get('bbox')
                if not bbox:
                    continue
                if primitive_kind == 'circle':
                    radius = float(primitive.get('r') or 0)
                    if radius < 3.0 or radius > max(12.0, chart_h * 0.16):
                        continue
                else:
                    if bbox['width'] < 6.0 or bbox['height'] < 6.0:
                        continue
                    if bbox['width'] > max(24.0, chart_w * 0.12) or bbox['height'] > max(24.0, chart_h * 0.2):
                        continue
                filtered.append(primitive)
            if len(filtered) < 3:
                return None

            filtered.sort(key=lambda primitive: bbox_center(primitive['bbox'])[1])
            heights = sorted(primitive['bbox']['height'] for primitive in filtered)
            median_height = heights[len(heights) // 2] if heights else 12.0
            row_tolerance = max(8.0, median_height * 1.2)
            rows = []
            for primitive in filtered:
                center_y = bbox_center(primitive['bbox'])[1]
                if not rows or abs(center_y - rows[-1]['center_y']) > row_tolerance:
                    rows.append({'center_y': center_y, 'items': [primitive]})
                else:
                    rows[-1]['items'].append(primitive)
                    count = len(rows[-1]['items'])
                    rows[-1]['center_y'] = ((rows[-1]['center_y'] * (count - 1)) + center_y) / count

            best = None
            best_score = None
            for row in rows:
                items = sorted(row['items'], key=lambda primitive: bbox_center(primitive['bbox'])[0])
                if len(items) < 3 or len(items) > 7:
                    continue
                widths = [item['bbox']['width'] for item in items]
                heights = [item['bbox']['height'] for item in items]
                if min(widths) <= 0 or min(heights) <= 0:
                    continue
                if (max(widths) / min(widths)) > 2.2 or (max(heights) / min(heights)) > 2.2:
                    continue
                span = items[-1]['bbox']['x'] + items[-1]['bbox']['width'] - items[0]['bbox']['x']
                if span < max(40.0, chart_w * 0.18) or span > chart_w * 0.96:
                    continue
                score = (len(items) * 10.0) + span - abs(len(items) - 5) * 3.0
                if best_score is None or score > best_score:
                    best = items
                    best_score = score
            return best

        circle_row = select_row(self._primitives_for_chart_region(chart_region, kind='circle'), 'circle')
        rect_row = select_row(self._primitives_for_chart_region(chart_region, kind='rect'), 'rect')
        if circle_row and (not rect_row or len(circle_row) >= len(rect_row)):
            items = circle_row
            preset = 'ellipse'
        elif rect_row:
            items = rect_row
            preset = 'roundRect'
        else:
            return None

        item_specs = []
        for primitive in items:
            bbox = primitive.get('bbox')
            fill_color = self._paint_to_rgb(primitive.get('fill')) if self._normalized_fill_present(primitive.get('fill')) else None
            stroke_color = self._paint_to_rgb(primitive.get('stroke'), default=fill_color or RGBColor(37, 99, 235))
            item_specs.append({
                'bbox': bbox,
                'preset': preset,
                'filled': fill_color is not None,
                'fill_color': fill_color,
                'line_color': stroke_color,
                'line_width_px': max(float(primitive.get('stroke_width') or 1.6), 1.2),
                'radius_px': min(bbox['width'], bbox['height']) / 2 if preset == 'roundRect' else 0.0,
            })

        return {
            'chart_id': chart_region.get('chart_id') or 'chart-rating',
            'chart_type': 'rating',
            'plot_bbox': merge_bboxes([item['bbox'] for item in item_specs]),
            'items': item_specs,
        }

    def _render_native_ring_group(self, chart_region, chart_spec):
        group = self._ensure_native_chart_group(chart_region, chart_spec['chart_type'])
        track_circle = chart_spec.get('track_circle')
        progress_circle = chart_spec.get('progress_circle')
        if track_circle is not None:
            track_color = self._paint_to_rgb(track_circle.get('stroke'), default=RGBColor(226, 232, 240))
            self._append_ring_outline_shape(group, track_circle, track_color)
        if progress_circle is None:
            return False
        progress_color = self._paint_to_rgb(progress_circle.get('stroke'), default=RGBColor(37, 99, 235))
        self._append_ring_arc_shape(group, progress_circle, progress_color)
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        group.shapes._recalculate_extents()
        return True

    def _render_native_kpi_shape(self, slide, chart_region, chart_spec):
        plot_bbox = chart_spec.get('plot_bbox')
        if not plot_bbox:
            return False
        fill_el = self._solid_fill_from_rgb(chart_spec.get('trend_color'), opacity=1.0)
        trend = chart_spec.get('trend') or 'flat'
        if trend == 'flat':
            bar_height = max(plot_bbox['height'] * 0.35, 4.0)
            bar_bbox = normalize_bbox({
                'x': plot_bbox['x'],
                'y': plot_bbox['y'] + (plot_bbox['height'] - bar_height) / 2,
                'width': plot_bbox['width'],
                'height': bar_height,
            })
            if bar_bbox is None:
                return False
            shape = make_shape(
                self._id(),
                f'KPI{self.sid}',
                px(bar_bbox['x']),
                px(bar_bbox['y']),
                px(bar_bbox['width']),
                px(bar_bbox['height']),
                preset='roundRect',
                fill_el=fill_el,
                line_el=None,
                rx=px(min(bar_bbox['height'] / 2, 8.0)),
            )
        else:
            shape = make_shape(
                self._id(),
                f'KPI{self.sid}',
                px(plot_bbox['x']),
                px(plot_bbox['y']),
                px(plot_bbox['width']),
                px(plot_bbox['height']),
                preset='triangle',
                fill_el=fill_el,
                line_el=None,
            )
            if trend == 'down':
                xfrm = shape.find(f'.//{{{NS["a"]}}}xfrm')
                if xfrm is not None:
                    xfrm.set('flipV', '1')
        sp_tree = next((node for node in slide._element.iter() if node.tag.endswith('}spTree')), None)
        if sp_tree is None:
            return False
        self._set_element_name(
            shape,
            self._managed_chart_name(
                NATIVE_CHART_NAME_PREFIX,
                chart_spec['chart_id'],
                chart_type=chart_spec['chart_type'],
                bbox=plot_bbox,
                block_id=chart_region.get('block_id'),
            ),
        )
        sp_tree.append(shape)
        self.stats['shapes'] += 1
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        return True

    def _render_native_metric_row_group(self, chart_region, chart_spec):
        rows = chart_spec.get('rows') or []
        if len(rows) < 2:
            return False
        group = self._ensure_native_chart_group(chart_region, chart_spec['chart_type'])
        for row in rows:
            self._append_metric_row_bar_shape(group, row.get('track_bbox'), row.get('track_color'), row.get('radius_px', 0.0))
            self._append_metric_row_bar_shape(group, row.get('fill_bbox'), row.get('fill_color'), row.get('radius_px', 0.0))
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        group.shapes._recalculate_extents()
        return True

    def _render_native_sparkline_group(self, chart_region, chart_spec):
        line_points = chart_spec.get('line_points') or []
        if len(line_points) < 3:
            return False
        group = self._ensure_native_chart_group(chart_region, chart_spec['chart_type'])
        area_points = chart_spec.get('area_points') or []
        if len(area_points) >= 3:
            self._append_native_path_shape(
                group,
                area_points,
                close=True,
                fill_color=chart_spec.get('area_color'),
                fill_opacity=chart_spec.get('area_opacity', 0.12),
                line_color=None,
                name_prefix='SPA',
            )
        self._append_native_path_shape(
            group,
            line_points,
            close=False,
            fill_color=None,
            line_color=chart_spec.get('line_color'),
            line_width_px=chart_spec.get('line_width_px', 2.0),
            name_prefix='SPL',
        )
        self._append_native_basic_shape(
            group,
            chart_spec.get('marker_bbox'),
            preset='ellipse',
            fill_color=chart_spec.get('line_color'),
            line_color=None,
            name_prefix='SPD',
        )
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        group.shapes._recalculate_extents()
        return True

    def _render_native_rating_group(self, chart_region, chart_spec):
        items = chart_spec.get('items') or []
        if len(items) < 3:
            return False
        group = self._ensure_native_chart_group(chart_region, chart_spec['chart_type'])
        for item in items:
            self._append_native_basic_shape(
                group,
                item.get('bbox'),
                preset=item.get('preset', 'ellipse'),
                fill_color=item.get('fill_color') if item.get('filled') else None,
                line_color=item.get('line_color'),
                line_width_px=item.get('line_width_px', 1.4),
                radius_px=item.get('radius_px', 0.0),
                name_prefix='RAT',
            )
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        group.shapes._recalculate_extents()
        return True

    def _render_native_primitive_group(self, chart_region, chart_spec):
        primitives = chart_spec.get('primitives') or []
        if len(primitives) < 3:
            return False
        group = self._ensure_native_chart_group(chart_region, chart_spec['chart_type'])
        for primitive in primitives:
            kind = primitive.get('kind')
            bbox = primitive.get('bbox')
            if kind == 'text':
                self._append_native_text_primitive(group, primitive)
            elif kind == 'rect':
                self._append_native_basic_shape(
                    group,
                    bbox,
                    preset='roundRect' if max(float(primitive.get('rx') or 0), float(primitive.get('ry') or 0)) > 0 else 'rect',
                    fill_color=self._paint_to_rgb(primitive.get('fill')),
                    fill_opacity=float(primitive.get('opacity') or 1.0),
                    line_color=self._paint_to_rgb(primitive.get('stroke')),
                    line_width_px=max(float(primitive.get('stroke_width') or 1.0), 1.0),
                    radius_px=max(float(primitive.get('rx') or 0), float(primitive.get('ry') or 0)),
                    name_prefix='NGR',
                )
            elif kind == 'circle':
                self._append_native_basic_shape(
                    group,
                    bbox,
                    preset='ellipse',
                    fill_color=self._paint_to_rgb(primitive.get('fill')),
                    fill_opacity=float(primitive.get('opacity') or 1.0),
                    line_color=self._paint_to_rgb(primitive.get('stroke')),
                    line_width_px=max(float(primitive.get('stroke_width') or 1.0), 1.0),
                    name_prefix='NGC',
                )
            elif kind in {'line', 'polyline'}:
                self._append_native_path_shape(
                    group,
                    primitive.get('points') or [],
                    close=False,
                    fill_color=None,
                    line_color=self._paint_to_rgb(primitive.get('stroke')),
                    line_width_px=max(float(primitive.get('stroke_width') or 1.0), 1.0),
                    line_opacity=float(primitive.get('opacity') or 1.0),
                    name_prefix='NGL',
                )
            elif kind in {'polygon', 'path'}:
                self._append_native_path_shape(
                    group,
                    primitive.get('points') or [],
                    close=bool(primitive.get('closed') or kind == 'polygon'),
                    fill_color=self._paint_to_rgb(primitive.get('fill')),
                    fill_opacity=float(primitive.get('opacity') or 1.0),
                    line_color=self._paint_to_rgb(primitive.get('stroke')),
                    line_width_px=max(float(primitive.get('stroke_width') or 1.0), 1.0),
                    line_opacity=float(primitive.get('opacity') or 1.0),
                    name_prefix='NGP',
                )
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        group.shapes._recalculate_extents()
        return True

    def _render_native_stacked_bar(self, slide, chart_region, chart_spec):
        plot_bbox = chart_spec.get('plot_bbox')
        series_specs = chart_spec.get('series') or []
        if not plot_bbox or len(series_specs) < 2:
            return False
        data = CategoryChartData()
        data.categories = ['Composition']
        for series_spec in series_specs:
            data.add_series(series_spec['name'], (float(series_spec['value']),))
        chart_type = getattr(XL_CHART_TYPE, 'BAR_STACKED_100', XL_CHART_TYPE.BAR_STACKED)
        chart_shape = slide.shapes.add_chart(
            chart_type,
            px(plot_bbox['x']),
            px(plot_bbox['y']),
            px(plot_bbox['width']),
            px(plot_bbox['height']),
            data,
        )
        self._assign_chart_name(
            chart_shape,
            chart_spec['chart_id'],
            chart_spec['chart_type'],
            bbox=plot_bbox,
            block_id=chart_region.get('block_id'),
        )
        chart = chart_shape.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.gap_width = 0
        plot.overlap = 100
        chart.category_axis.visible = False
        chart.value_axis.visible = False
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.maximum_scale = 100.0
        for series, series_spec in zip(chart.series, series_specs):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = series_spec['color']
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        return True

    def _render_native_progress_bar(self, slide, chart_region, chart_spec):
        plot_bbox = chart_spec.get('plot_bbox')
        if not plot_bbox:
            return False
        data = CategoryChartData()
        data.categories = ['Progress']
        value = float(chart_spec['progress_value'])
        data.add_series('Done', (value,))
        data.add_series('Remaining', (max(0.0, 100.0 - value),))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            px(plot_bbox['x']),
            px(plot_bbox['y']),
            px(plot_bbox['width']),
            px(plot_bbox['height']),
            data,
        )
        self._assign_chart_name(
            chart_shape,
            chart_spec['chart_id'],
            chart_spec['chart_type'],
            bbox=plot_bbox,
            block_id=chart_region.get('block_id'),
        )
        chart = chart_shape.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.gap_width = 0
        plot.overlap = 100
        chart.category_axis.visible = False
        chart.value_axis.visible = False
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.maximum_scale = 100.0
        for series, color in zip(chart.series, (chart_spec['fill_color'], chart_spec['track_color'])):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        return True

    def _render_native_comparison_bar(self, slide, chart_region, chart_spec):
        plot_bbox = chart_spec.get('plot_bbox')
        values = chart_spec.get('values') or []
        categories = chart_spec.get('categories') or []
        if not plot_bbox or len(values) < 2 or len(categories) < 2:
            return False
        data = CategoryChartData()
        data.categories = categories
        data.add_series('Value', tuple(values))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            px(plot_bbox['x']),
            px(plot_bbox['y']),
            px(plot_bbox['width']),
            px(plot_bbox['height']),
            data,
        )
        self._assign_chart_name(
            chart_shape,
            chart_spec['chart_id'],
            chart_spec['chart_type'],
            bbox=plot_bbox,
            block_id=chart_region.get('block_id'),
        )
        chart = chart_shape.chart
        chart.has_legend = False
        plot = chart.plots[0]
        plot.gap_width = 60
        chart.category_axis.visible = False
        chart.value_axis.visible = False
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.maximum_scale = max(max(values) * 1.15, 1.0)
        series = chart.series[0]
        for point, color in zip(series.points, chart_spec.get('point_colors') or []):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
        self.native_chart_ids.add(chart_region.get('chart_id') or chart_spec['chart_id'])
        return True

    def _chart_promotion_registry(self):
        return {
            'comparison_bar': {
                'extract': self._extract_comparison_bar_chart,
                'render': lambda slide, chart_region, chart_spec: self._render_native_comparison_bar(slide, chart_region, chart_spec),
            },
            'progress_bar': {
                'extract': self._extract_progress_bar_chart,
                'render': lambda slide, chart_region, chart_spec: self._render_native_progress_bar(slide, chart_region, chart_spec),
            },
            'stacked_bar': {
                'extract': self._extract_stacked_bar_chart,
                'render': lambda slide, chart_region, chart_spec: self._render_native_stacked_bar(slide, chart_region, chart_spec),
            },
            'sparkline': {
                'extract': self._extract_sparkline_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_sparkline_group(chart_region, chart_spec),
            },
            'rating': {
                'extract': self._extract_rating_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_rating_group(chart_region, chart_spec),
            },
            'kpi': {
                'extract': self._extract_kpi_chart,
                'render': lambda slide, chart_region, chart_spec: self._render_native_kpi_shape(slide, chart_region, chart_spec),
            },
            'ring': {
                'extract': self._extract_ring_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_ring_group(chart_region, chart_spec),
            },
            'metric_row': {
                'extract': self._extract_metric_row_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_metric_row_group(chart_region, chart_spec),
            },
            'timeline': {
                'extract': self._extract_timeline_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_primitive_group(chart_region, chart_spec),
            },
            'funnel': {
                'extract': self._extract_funnel_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_primitive_group(chart_region, chart_spec),
            },
            'radar': {
                'extract': self._extract_radar_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_primitive_group(chart_region, chart_spec),
            },
            'waffle': {
                'extract': self._extract_waffle_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_primitive_group(chart_region, chart_spec),
            },
            'treemap': {
                'extract': self._extract_treemap_chart,
                'render': lambda _slide, chart_region, chart_spec: self._render_native_primitive_group(chart_region, chart_spec),
            },
        }

    def _render_semantic_charts(self, slide):
        rendered = 0
        registry = self._chart_promotion_registry()
        for chart_region in self.semantic_charts:
            chart_type = normalize_text(chart_region.get('chart_type_hint')).lower()
            promoter = registry.get(chart_type)
            if promoter is None:
                continue
            chart_spec = promoter['extract'](chart_region)
            rendered_ok = chart_spec is not None and promoter['render'](slide, chart_region, chart_spec)
            if rendered_ok:
                rendered += 1
        self.stats['native_charts'] = rendered

    def _ensure_chart_group(self, chart_region):
        chart_id = chart_region.get('chart_id') or f'chart-{len(self.chart_groups) + 1}'
        group = self.chart_groups.get(chart_id)
        if group is not None:
            return group
        group = self.current_slide.shapes.add_group_shape()
        try:
            group._element.nvGrpSpPr.cNvPr.set(
                'name',
                self._managed_chart_name(
                    CHART_GROUP_NAME_PREFIX,
                    chart_id,
                    bbox=chart_region.get('bbox'),
                    block_id=chart_region.get('block_id'),
                ),
            )
        except Exception:
            pass
        self.chart_groups[chart_id] = group
        return group

    def _target_tree_for_bbox(self, default_sp, bbox):
        chart_region = self._chart_region_for_bbox(bbox)
        if chart_region is None:
            return default_sp
        chart_id = chart_region.get('chart_id')
        if isinstance(chart_id, str) and chart_id in self.native_chart_ids:
            return default_sp
        return self._ensure_chart_group(chart_region).shapes._spTree

    def _grid_from_table(self, table_semantics):
        widths = [0.0] * max(int(table_semantics.get('column_count') or 0), 0)
        heights = [0.0] * max(int(table_semantics.get('row_count') or 0), 0)
        for row in table_semantics.get('rows', []):
            for cell in row.get('cells', []):
                cell_bbox = cell.get('bbox')
                if not cell_bbox:
                    continue
                col_span = max(int(cell.get('colspan', 1)), 1)
                row_span = max(int(cell.get('rowspan', 1)), 1)
                width_share = cell_bbox['width'] / col_span
                height_share = cell_bbox['height'] / row_span
                for col_index in range(cell['column_index'], min(cell['column_index'] + col_span, len(widths))):
                    widths[col_index] = max(widths[col_index], width_share)
                for row_index in range(cell['row_index'], min(cell['row_index'] + row_span, len(heights))):
                    heights[row_index] = max(heights[row_index], height_share)

        table_bbox = table_semantics.get('bbox') or {}
        if widths and not any(widths):
            widths = [float(table_bbox.get('width', 0)) / len(widths)] * len(widths)
        elif widths and any(width <= 0 for width in widths):
            fallback = max(sum(widths) / max(sum(1 for width in widths if width > 0), 1), 80.0)
            widths = [width if width > 0 else fallback for width in widths]
        if heights and not any(heights):
            heights = [float(table_bbox.get('height', 0)) / len(heights)] * len(heights)
        elif heights and any(height <= 0 for height in heights):
            fallback = max(sum(heights) / max(sum(1 for height in heights if height > 0), 1), 28.0)
            heights = [height if height > 0 else fallback for height in heights]

        total_width = sum(widths)
        total_height = sum(heights)
        if widths and table_bbox.get('width') and total_width > 0:
            scale = table_bbox['width'] / total_width
            widths = [width * scale for width in widths]
        if heights and table_bbox.get('height') and total_height > 0:
            scale = table_bbox['height'] / total_height
            heights = [height * scale for height in heights]
        return widths, heights

    def _render_native_table(self, slide, table_semantics):
        table_bbox = table_semantics.get('bbox')
        row_count = int(table_semantics.get('row_count') or 0)
        col_count = int(table_semantics.get('column_count') or 0)
        if not table_bbox or row_count <= 0 or col_count <= 0:
            return False

        shape = slide.shapes.add_table(
            row_count,
            col_count,
            px(table_bbox['x']),
            px(table_bbox['y']),
            px(table_bbox['width']),
            px(table_bbox['height']),
        )
        self._assign_shape_name(
            shape,
            f'TBL{self.sid}',
            bbox=table_bbox,
            block_id=table_semantics.get('block_id'),
        )
        table = shape.table
        widths, heights = self._grid_from_table(table_semantics)
        for col_index, width_px in enumerate(widths):
            table.columns[col_index].width = max(1, px(max(width_px, 8.0)))
        for row_index, height_px in enumerate(heights):
            table.rows[row_index].height = max(1, px(max(height_px, 18.0)))

        for row in table_semantics.get('rows', []):
            for cell_semantics in row.get('cells', []):
                row_index = int(cell_semantics['row_index'])
                col_index = int(cell_semantics['column_index'])
                if row_index >= row_count or col_index >= col_count:
                    continue
                cell = table.cell(row_index, col_index)
                colspan = max(int(cell_semantics.get('colspan', 1)), 1)
                rowspan = max(int(cell_semantics.get('rowspan', 1)), 1)
                if colspan > 1 or rowspan > 1:
                    end_row = min(row_count - 1, row_index + rowspan - 1)
                    end_col = min(col_count - 1, col_index + colspan - 1)
                    if end_row != row_index or end_col != col_index:
                        cell = cell.merge(table.cell(end_row, end_col))
                apply_table_cell_text(cell, cell_semantics)
        return True

    def _render_semantic_tables(self, slide):
        rendered = 0
        for table in self.semantic_tables:
            if self._render_native_table(slide, table):
                rendered += 1
        self.stats['native_tables'] = rendered

    def _finalize_chart_groups(self):
        self._recalculate_structured_groups()
        self.stats['structured_chart_groups'] = len(self.chart_groups)

    def _text_align(self, anchor):
        if anchor == 'middle':
            return 'ctr'
        if anchor == 'end':
            return 'r'
        return 'l'

    def _ascent_ratio(self, font_height):
        return 0.80 if font_height >= 32 else (0.85 if font_height >= 16 else 0.88)

    def _resolve_vertical_anchor(self, baseline, y, font_height):
        if 'after-edge' in baseline:
            return 'b', y
        if baseline in ('central', 'middle'):
            return 'ctr', y
        return 't', y - font_height * self._ascent_ratio(font_height)

    def _build_text_fragment(self, *, text, x, y, text_length, font_size, font_weight,
                             fill_s, font_family, baseline, anchor, opacity):
        txt = normalize_text(text)
        if not txt:
            return None
        font_height = float(font_size)
        text_anchor_v, y = self._resolve_vertical_anchor(baseline, y, font_height)
        c = parse_color(fill_s)
        hex6 = c[0] if c and c[0] != 'grad' else '000000'
        alpha = c[1] if c and c[0] != 'grad' else 100000
        alpha = int(alpha * opacity)
        width_px = float(text_length) if float(text_length or 0) > 0 else estimate_text_width(txt, font_size)
        cx_v = px(width_px)
        cy_v = px(font_height * 1.5)
        if anchor == 'middle':
            x -= width_px / 2
        elif anchor == 'end':
            x -= width_px
        if text_anchor_v == 'b':
            tb_y = px(y) - cy_v
        elif text_anchor_v == 'ctr':
            tb_y = px(y) - cy_v // 2
        else:
            tb_y = px(y)
        return {
            'text': txt,
            'run': {
                'text': txt,
                'sz': font_sz(font_size),
                'bold': font_weight in ('bold', '700', '800', '900'),
                'hex': hex6,
                'alpha': alpha,
                'font': resolve_font(font_family),
            },
            'x_px': x,
            'y_emu': tb_y,
            'width_emu': cx_v,
            'height_emu': cy_v,
            'anchor': text_anchor_v,
            'horizontal_anchor': anchor,
            'align': self._text_align(anchor),
            'font_size_px': font_height,
            'bbox': normalize_bbox({
                'x': x,
                'y': tb_y / EMU_PX,
                'width': width_px,
                'height': cy_v / EMU_PX,
            }),
        }

    def _should_group_text_fragments(self, fragments, semantic):
        if not semantic or len(fragments) < 2:
            return False
        if semantic.get('role') not in {'heading', 'paragraph', 'list_item', 'caption'}:
            return False
        alignment_points = []
        for fragment in fragments:
            width_px = fragment['width_emu'] / EMU_PX
            if fragment['horizontal_anchor'] == 'middle':
                alignment_points.append(fragment['x_px'] + width_px / 2)
            elif fragment['horizontal_anchor'] == 'end':
                alignment_points.append(fragment['x_px'] + width_px)
            else:
                alignment_points.append(fragment['x_px'])
        tolerance = max(12.0, max(fragment['font_size_px'] for fragment in fragments) * 0.8)
        return (max(alignment_points) - min(alignment_points)) <= tolerance

    def _build_paragraph(self, runs, align, semantic=None, continuation=False):
        paragraph = {'runs': runs, 'align': align}
        if semantic and semantic.get('role') == 'list_item':
            paragraph['indent_level'] = semantic.get('list_depth', 0)
            if continuation:
                paragraph['continuation'] = True
            elif semantic.get('list_type') == 'ol':
                paragraph['number'] = semantic.get('list_index', 1)
            else:
                paragraph['bullet_char'] = '•'
        return paragraph

    def _append_text_shape(self, sp, fragment, semantic=None):
        shape_id = self._id()
        base_name = f'T{self.sid}'
        shape = make_textbox(
            shape_id,
            base_name,
            px(fragment['x_px']),
            fragment['y_emu'],
            fragment['width_emu'],
            fragment['height_emu'],
            [self._build_paragraph([fragment['run']], fragment['align'], semantic=semantic)],
            anchor=fragment['anchor'],
        )
        self._append_shape_element(sp, shape, base_name, bbox=fragment.get('bbox'), semantic=semantic)

    def _append_grouped_text_shape(self, sp, fragments, semantic):
        left = min(px(fragment['x_px']) for fragment in fragments)
        top = min(fragment['y_emu'] for fragment in fragments)
        right = max(px(fragment['x_px']) + fragment['width_emu'] for fragment in fragments)
        bottom = max(fragment['y_emu'] + fragment['height_emu'] for fragment in fragments)
        padding = px(max(2.0, max(fragment['font_size_px'] for fragment in fragments) * 0.25))
        left = max(0, left - padding)
        top = max(0, top - padding // 2)
        right += padding
        bottom += padding
        paragraphs = []
        for idx, fragment in enumerate(fragments):
            paragraphs.append(
                self._build_paragraph(
                    [fragment['run']],
                    fragment['align'],
                    semantic=semantic,
                    continuation=idx > 0 and semantic.get('role') == 'list_item',
                )
            )
        shape_id = self._id()
        base_name = f'T{self.sid}'
        shape = make_textbox(
            shape_id,
            base_name,
            left,
            top,
            max(1, right - left),
            max(1, bottom - top),
            paragraphs,
            anchor='t',
        )
        self._append_shape_element(
            sp,
            shape,
            base_name,
            bbox=self._fragments_bbox(fragments),
            semantic=semantic,
        )

    def _fragments_bbox(self, fragments):
        if not fragments:
            return None
        left = min(fragment['x_px'] for fragment in fragments)
        top = min(fragment['bbox']['y'] for fragment in fragments if fragment.get('bbox'))
        right = max(fragment['x_px'] + fragment['width_emu'] / EMU_PX for fragment in fragments)
        bottom = max(fragment['bbox']['y'] + fragment['bbox']['height'] for fragment in fragments if fragment.get('bbox'))
        return normalize_bbox({'x': left, 'y': top, 'width': right - left, 'height': bottom - top})

    def _tag(self, el):
        t = el.tag
        return t.split('}')[1] if isinstance(t, str) and '}' in t else (t if isinstance(t, str) else '')

    def _parse_transform(self, el):
        """解析 transform -> (dx, dy, sx, sy)。"""
        t = el.get('transform', '')
        dx, dy, sx, sy = 0.0, 0.0, 1.0, 1.0
        # translate
        m = re.search(r'translate\(\s*([\d.\-]+)[,\s]+([\d.\-]+)', t)
        if m:
            dx, dy = float(m.group(1)), float(m.group(2))
        # scale
        m = re.search(r'scale\(\s*([\d.\-]+)(?:[,\s]+([\d.\-]+))?\s*\)', t)
        if m:
            sx = float(m.group(1))
            sy = float(m.group(2)) if m.group(2) else sx
        # matrix(a,b,c,d,e,f) -> e=translateX, f=translateY
        m = re.search(r'matrix\(\s*([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)[,\s]+([\d.\-]+)', t)
        if m:
            dx = float(m.group(5))
            dy = float(m.group(6))
            sx = float(m.group(1))
            sy = float(m.group(4))
        return dx, dy, sx, sy

    def _walk(self, el, sp, ox, oy, group_opacity, scale, slide):
        tag = self._tag(el)
        try:
            if tag == 'rect':
                self._rect(el, sp, ox, oy, group_opacity, scale, slide)
            elif tag == 'text':
                self._text(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'circle':
                self._circle(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'ellipse':
                self._ellipse(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'line':
                self._line(el, sp, ox, oy, scale)
            elif tag == 'path':
                self._path(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'image':
                self._image(el, sp, ox, oy, group_opacity, scale, slide)
            elif tag == 'polygon':
                self._polygon(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'polyline':
                self._polyline(el, sp, ox, oy, group_opacity, scale)
            elif tag == 'use':
                self._use(el, sp, ox, oy, group_opacity, scale, slide)
            elif tag == 'g':
                dx, dy, sx, sy = self._parse_transform(el)
                el_opacity = float(el.get('opacity', '1'))
                child_opacity = group_opacity * el_opacity
                # scale 累积：父级scale * 当前g的scale
                child_scale = scale * sx  # 假设sx==sy（等比缩放）
                new_ox = ox + dx * scale
                new_oy = oy + dy * scale
                for c in el:
                    self._walk(c, sp, new_ox, new_oy,
                               child_opacity, child_scale, slide)
            elif tag in ('defs', 'style', 'linearGradient', 'radialGradient',
                         'stop', 'pattern', 'clipPath', 'filter', 'mask'):
                pass
            else:
                for c in el:
                    self._walk(c, sp, ox, oy, group_opacity, scale, slide)
        except Exception as e:
            self.stats['errors'] += 1
            print(f"    Warning: {tag} element failed: {e}", file=sys.stderr)

    def _rect(self, el, sp, ox, oy, opacity, scale, slide):
        x = (float(el.get('x', 0)) * scale) + ox
        y = (float(el.get('y', 0)) * scale) + oy
        w = float(el.get('width', 0)) * scale
        h = float(el.get('height', 0)) * scale
        if w <= 0 or h <= 0:
            return

        # 过滤面积 < 4px 的纯装饰元素
        if w < 4 and h < 4:
            self.stats['skipped'] += 1
            return

        if self._table_region_for_bbox({'x': x, 'y': y, 'width': w, 'height': h}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': x, 'y': y, 'width': w, 'height': h}):
            self.stats['skipped'] += 1
            return

        fill_s = el.get('fill', '')
        stroke_s = el.get('stroke', '')
        c = parse_color(fill_s)

        # 跳过全透明无边框矩形
        if c and c[0] != 'grad' and c[1] == 0 and not stroke_s:
            return

        el_opacity = float(el.get('opacity', '1')) * opacity

        # 首个全屏 rect -> 幻灯片背景
        if not self.bg_set and w >= 1270 and h >= 710:
            if self.preserve_background:
                self.bg_set = True
                self.stats['skipped'] += 1
                return
            self.bg_set = True
            bg = slide._element.find(f'.//{{{NS["p"]}}}bg')
            if bg is None:
                cSld = slide._element.find(f'{{{NS["p"]}}}cSld')
                if cSld is not None:
                    bg_el = _el('p:bg', children=[
                        _el('p:bgPr', children=[
                            make_fill(fill_s, self.grads, el_opacity),
                            _el('a:effectLst'),
                        ])
                    ])
                    cSld.insert(0, bg_el)
            return  # 不再作为形状添加

        r = max(float(el.get('rx', 0)), float(el.get('ry', 0)))
        preset = 'roundRect' if r > 0 else 'rect'
        fill_el = make_fill(fill_s, self.grads, el_opacity)
        line_el = make_line(stroke_s, el.get('stroke-width', '1')) if stroke_s else None
        bbox = {'x': x, 'y': y, 'width': w, 'height': h}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        shape_id = self._id()
        base_name = f'R{self.sid}'
        shape = make_shape(shape_id, base_name,
                           px(x), px(y), px(w), px(h),
                           preset=preset, fill_el=fill_el, line_el=line_el, rx=px(r))
        self._append_shape_element(target_sp, shape, base_name, bbox=bbox)

    def _text(self, el, sp, ox, oy, opacity, scale):
        """优先利用 HTML 语义合并多行文本，并为列表项生成原生段落语义。"""
        fill_s = el.get('fill', el.get('color', ''))
        fsz = el.get('font-size', '14px').replace('px', '')
        fw = el.get('font-weight', '')
        ff = el.get('font-family', '')
        baseline = el.get('dominant-baseline', '')
        anchor = el.get('text-anchor', 'start')

        tspans = list(el.findall(f'{{{SVG_NS}}}tspan'))
        fragments = []

        if tspans:
            for ts in tspans:
                txt = normalize_text(ts.text)
                if not txt:
                    continue
                x = float(ts.get('x', el.get('x', 0))) * scale + ox
                y = float(ts.get('y', el.get('y', 0))) * scale + oy
                tlen = float(ts.get('textLength', 0))
                ts_fsz = ts.get('font-size', fsz).replace('px', '')
                ts_fw = ts.get('font-weight', fw)
                ts_fill = ts.get('fill', fill_s)
                ts_ff = ts.get('font-family', ff)
                fragment = self._build_text_fragment(
                    text=txt,
                    x=x,
                    y=y,
                    text_length=tlen,
                    font_size=ts_fsz,
                    font_weight=ts_fw,
                    fill_s=ts_fill,
                    font_family=ts_ff,
                    baseline=baseline,
                    anchor=anchor,
                    opacity=opacity,
                )
                if fragment is not None:
                    fragments.append(fragment)

        elif el.text and normalize_text(el.text):
            x = float(el.get('x', 0)) * scale + ox
            y = float(el.get('y', 0)) * scale + oy
            fragment = self._build_text_fragment(
                text=el.text,
                x=x,
                y=y,
                text_length=0,
                font_size=fsz,
                font_weight=fw,
                fill_s=fill_s,
                font_family=ff,
                baseline=baseline,
                anchor=anchor,
                opacity=opacity,
            )
            if fragment is not None:
                fragments.append(fragment)

        if not fragments:
            return

        fragments_bbox = self._fragments_bbox(fragments)
        if self._table_region_for_bbox(fragments_bbox):
            self.stats['skipped'] += len(fragments)
            return

        if len(fragments) > 1:
            combined_text = ' '.join(fragment['text'] for fragment in fragments)
            semantic = self._match_semantic_text(combined_text, bbox=fragments_bbox, consume=False)
            if self._should_group_text_fragments(fragments, semantic):
                semantic = self._match_semantic_text(combined_text, bbox=fragments_bbox, consume=True)
                target_sp = self._target_tree_for_bbox(sp, fragments_bbox)
                self._append_grouped_text_shape(target_sp, fragments, semantic)
                return

        for fragment in fragments:
            semantic = self._match_semantic_text(fragment['text'], bbox=fragment.get('bbox'))
            target_sp = self._target_tree_for_bbox(sp, fragment.get('bbox'))
            self._append_text_shape(target_sp, fragment, semantic=semantic)

    def _circle(self, el, sp, ox, oy, opacity, scale):
        cx_v = float(el.get('cx', 0)) * scale + ox
        cy_v = float(el.get('cy', 0)) * scale + oy
        r = float(el.get('r', 0)) * scale
        if r <= 0 or r < 2:
            self.stats['skipped'] += 1
            return

        if self._table_region_for_bbox({'x': cx_v - r, 'y': cy_v - r, 'width': 2 * r, 'height': 2 * r}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': cx_v - r, 'y': cy_v - r, 'width': 2 * r, 'height': 2 * r}):
            self.stats['skipped'] += 1
            return

        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_s = el.get('fill', '')
        stroke_s = el.get('stroke', '')
        stroke_w_s = el.get('stroke-width', '1')
        dasharray = el.get('stroke-dasharray', '')
        dashoffset = el.get('stroke-dashoffset', '')

        # stroke-dashoffset 兼容：转为等效旋转角度
        extra_rotate = 0
        if dashoffset and dasharray:
            try:
                offset_val = float(strip_unit(dashoffset))
                circumference = 2 * math.pi * r
                extra_rotate = -(offset_val / circumference) * 360  # 负偏移 = 顺时针旋转
            except (ValueError, ZeroDivisionError):
                pass

        # 环形图特殊处理：fill=none + stroke + dasharray -> OOXML arc + 粗描边
        if (fill_s == 'none' or not fill_s) and stroke_s and dasharray:
            sw = float(strip_unit(stroke_w_s))
            # 解析 dasharray (格式: "188.1 188.5" 或 "113.097px, 150.796px")
            dash_parts = [float(strip_unit(p.strip())) for p in dasharray.replace(',', ' ').split() if p.strip()]
            if len(dash_parts) >= 2:
                circumference = 2 * math.pi * r
                arc_len = dash_parts[0]
                angle_pct = min(arc_len / circumference, 1.0)

                # 检查 rotate transform
                transform = el.get('transform', '')
                start_angle = 0
                rot_m = re.search(r'rotate\(\s*([\d.\-]+)', transform)
                if rot_m:
                    start_angle = float(rot_m.group(1))
                start_angle += extra_rotate  # 合并 dashoffset 等效旋转

                # SVG -> PowerPoint 角度转换
                # SVG rotate(-90) = 从 12 点钟方向开始
                # PowerPoint arc: adj1=startAngle, adj2=endAngle (从3点钟顺时针, 60000单位/度)
                ppt_start = (start_angle + 90) % 360
                sweep = angle_pct * 360
                ppt_end = (ppt_start + sweep) % 360

                adj1 = int(ppt_start * 60000)
                adj2 = int(ppt_end * 60000)

                # 用 arc 预设 (只画弧线轮廓) + 粗描边 = 环形弧
                geom = _el('a:prstGeom', {'prst': 'arc'})
                av = _el('a:avLst')
                av.append(_el('a:gd', {'name': 'adj1', 'fmla': f'val {adj1}'}))
                av.append(_el('a:gd', {'name': 'adj2', 'fmla': f'val {adj2}'}))
                geom.append(av)

                # 描边颜色 = SVG 的 stroke 颜色（支持渐变引用）
                stroke_color = parse_color(stroke_s)
                ln_children = []
                if stroke_color and stroke_color[0] == 'grad':
                    # stroke 引用渐变 -> 提取渐变的第一个 stop 颜色作为实色
                    gdef = self.grads.get(stroke_color[1])
                    if gdef and gdef.get('stops'):
                        first_stop = gdef['stops'][0]
                        sc = parse_color(first_stop['color_str'])
                        if sc and sc[0] != 'grad':
                            ln_children.append(_el('a:solidFill', children=[
                                _srgb(sc[0], int(sc[1] * el_opacity))
                            ]))
                    # 也尝试用渐变填充（OOXML线条支持渐变）
                    if not ln_children and gdef:
                        grad_fill = _make_grad(gdef)
                        if grad_fill is not None:
                            ln_children.append(grad_fill)
                elif stroke_color and stroke_color[0] != 'grad':
                    ln_children.append(_el('a:solidFill', children=[
                        _srgb(stroke_color[0], int(stroke_color[1] * el_opacity))
                    ]))
                ln_children.append(_el('a:round'))
                line_el = _el('a:ln', {'w': str(int(sw * 12700))}, children=ln_children)

                shape = _el('p:sp')
                shape.append(_el('p:nvSpPr', children=[
                    _el('p:cNvPr', {'id': str(self._id()), 'name': f'Arc{self.sid}'}),
                    _el('p:cNvSpPr'), _el('p:nvPr'),
                ]))
                sp_pr = _el('p:spPr')
                sp_pr.append(_el('a:xfrm', children=[
                    _el('a:off', {'x': str(max(0, px(cx_v - r))),
                                  'y': str(max(0, px(cy_v - r)))}),
                    _el('a:ext', {'cx': str(px(2 * r)),
                                  'cy': str(px(2 * r))}),
                ]))
                sp_pr.append(geom)
                sp_pr.append(_el('a:noFill'))
                sp_pr.append(line_el)
                shape.append(sp_pr)
                bbox = {'x': cx_v - r, 'y': cy_v - r, 'width': 2 * r, 'height': 2 * r}
                target_sp = self._target_tree_for_bbox(sp, bbox)
                self._append_shape_element(target_sp, shape, f'Arc{self.sid}', bbox=bbox)
                return

        # fill=none + stroke (无dasharray) -> 空心圆 + 粗描边
        if (fill_s == 'none' or not fill_s) and stroke_s and stroke_s != 'none':
            sw = float(strip_unit(stroke_w_s))
            stroke_color = parse_color(stroke_s)
            ln_children = []
            if stroke_color and stroke_color[0] != 'grad':
                ln_children.append(_el('a:solidFill', children=[
                    _srgb(stroke_color[0], int(stroke_color[1] * el_opacity))
                ]))
            ln_children.append(_el('a:round'))
            line_el = _el('a:ln', {'w': str(int(sw * 12700))}, children=ln_children)

            bbox = {'x': cx_v - r, 'y': cy_v - r, 'width': 2 * r, 'height': 2 * r}
            target_sp = self._target_tree_for_bbox(sp, bbox)
            shape_id = self._id()
            base_name = f'C{self.sid}'
            shape = make_shape(shape_id, base_name,
                               px(cx_v - r), px(cy_v - r), px(2*r), px(2*r),
                               preset='ellipse',
                               fill_el=_el('a:noFill'),
                               line_el=line_el)
            self._append_shape_element(target_sp, shape, base_name, bbox=bbox)
            return

        # 普通圆形
        fill_el = make_fill(fill_s, self.grads, el_opacity)
        line_el = make_line(stroke_s, stroke_w_s) if stroke_s and stroke_s != 'none' else None
        bbox = {'x': cx_v - r, 'y': cy_v - r, 'width': 2 * r, 'height': 2 * r}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        shape_id = self._id()
        base_name = f'C{self.sid}'
        shape = make_shape(shape_id, base_name,
                   px(cx_v - r), px(cy_v - r), px(2*r), px(2*r),
                   preset='ellipse', fill_el=fill_el, line_el=line_el)
        self._append_shape_element(target_sp, shape, base_name, bbox=bbox)

    def _ellipse(self, el, sp, ox, oy, opacity, scale):
        cx_v = float(el.get('cx', 0)) * scale + ox
        cy_v = float(el.get('cy', 0)) * scale + oy
        rx = float(el.get('rx', 0)) * scale
        ry = float(el.get('ry', 0)) * scale
        if rx <= 0 or ry <= 0:
            return
        if self._table_region_for_bbox({'x': cx_v - rx, 'y': cy_v - ry, 'width': 2 * rx, 'height': 2 * ry}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': cx_v - rx, 'y': cy_v - ry, 'width': 2 * rx, 'height': 2 * ry}):
            self.stats['skipped'] += 1
            return
        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_el = make_fill(el.get('fill', ''), self.grads, el_opacity)
        bbox = {'x': cx_v - rx, 'y': cy_v - ry, 'width': 2 * rx, 'height': 2 * ry}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        shape_id = self._id()
        base_name = f'E{self.sid}'
        shape = make_shape(shape_id, base_name,
                           px(cx_v - rx), px(cy_v - ry), px(2*rx), px(2*ry),
                           preset='ellipse', fill_el=fill_el)
        self._append_shape_element(target_sp, shape, base_name, bbox=bbox)

    def _line(self, el, sp, ox, oy, scale):
        x1 = float(el.get('x1', 0)) * scale + ox
        y1 = float(el.get('y1', 0)) * scale + oy
        x2 = float(el.get('x2', 0)) * scale + ox
        y2 = float(el.get('y2', 0)) * scale + oy
        if self._table_region_for_bbox({'x': min(x1, x2), 'y': min(y1, y2), 'width': abs(x2 - x1) or 1, 'height': abs(y2 - y1) or 1}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': min(x1, x2), 'y': min(y1, y2), 'width': abs(x2 - x1) or 1, 'height': abs(y2 - y1) or 1}):
            self.stats['skipped'] += 1
            return
        line_el = make_line(el.get('stroke', '#000'), el.get('stroke-width', '1'))
        if line_el is None:
            return
        mx, my = min(x1, x2), min(y1, y2)
        w, h = abs(x2 - x1) or 1, abs(y2 - y1) or 1
        shape = make_shape(self._id(), f'L{self.sid}',
                           px(mx), px(my), px(w), px(h),
                           preset='line', fill_el=_el('a:noFill'), line_el=line_el)
        xfrm = shape.find(f'.//{{{NS["a"]}}}xfrm')
        if x1 > x2:
            xfrm.set('flipH', '1')
        if y1 > y2:
            xfrm.set('flipV', '1')
        bbox = {'x': mx, 'y': my, 'width': w, 'height': h}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        self._append_shape_element(target_sp, shape, f'L{self.sid}', bbox=bbox)

    def _polygon(self, el, sp, ox, oy, opacity, scale):
        """SVG <polygon> -> OOXML custGeom。"""
        points_str = el.get('points', '')
        if not points_str:
            return
        coords = re.findall(r'[+-]?\d*\.?\d+', points_str)
        if len(coords) < 6:  # 至少 3 个点
            return
        xs = [float(coords[i]) for i in range(0, len(coords), 2)]
        ys = [float(coords[i]) for i in range(1, len(coords), 2)]
        d_parts = [f'M {xs[0]} {ys[0]}']
        for j in range(1, len(xs)):
            d_parts.append(f'L {xs[j]} {ys[j]}')
        d_parts.append('Z')
        d_str = ' '.join(d_parts)
        bx, by = min(xs), min(ys)
        bw = max(xs) - bx or 1
        bh = max(ys) - by or 1
        if bw < 4 and bh < 4:
            self.stats['skipped'] += 1
            return
        if self._table_region_for_bbox({'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}):
            self.stats['skipped'] += 1
            return
        geom_el = parse_path_to_custgeom(d_str, (bx, by, bw, bh))
        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_el = make_fill(el.get('fill', ''), self.grads, el_opacity)
        line_el = make_line(el.get('stroke', ''), el.get('stroke-width', '1')) if el.get('stroke') else None
        shape = make_shape(self._id(), f'PG{self.sid}',
                           px((bx + ox) * scale) if scale != 1.0 else px(bx + ox),
                           px((by + oy) * scale) if scale != 1.0 else px(by + oy),
                           px(bw * scale), px(bh * scale),
                           fill_el=fill_el, line_el=line_el, geom_el=geom_el)
        bbox = {'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        self._append_shape_element(target_sp, shape, f'PG{self.sid}', bbox=bbox)

    def _polyline(self, el, sp, ox, oy, opacity, scale):
        """SVG <polyline> -> OOXML custGeom（不闭合）。"""
        points_str = el.get('points', '')
        if not points_str:
            return
        coords = re.findall(r'[+-]?\d*\.?\d+', points_str)
        if len(coords) < 4:  # 至少 2 个点
            return
        xs = [float(coords[i]) for i in range(0, len(coords), 2)]
        ys = [float(coords[i]) for i in range(1, len(coords), 2)]
        d_parts = [f'M {xs[0]} {ys[0]}']
        for j in range(1, len(xs)):
            d_parts.append(f'L {xs[j]} {ys[j]}')
        d_str = ' '.join(d_parts)  # 不加 Z（不闭合）
        bx, by = min(xs), min(ys)
        bw = max(xs) - bx or 1
        bh = max(ys) - by or 1
        if bw < 4 and bh < 4:
            self.stats['skipped'] += 1
            return
        if self._table_region_for_bbox({'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}):
            self.stats['skipped'] += 1
            return
        geom_el = parse_path_to_custgeom(d_str, (bx, by, bw, bh))
        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_el = make_fill(el.get('fill', 'none'), self.grads, el_opacity)
        line_el = make_line(el.get('stroke', '#000'), el.get('stroke-width', '1'))
        shape = make_shape(self._id(), f'PL{self.sid}',
                           px((bx + ox) * scale) if scale != 1.0 else px(bx + ox),
                           px((by + oy) * scale) if scale != 1.0 else px(by + oy),
                           px(bw * scale), px(bh * scale),
                           fill_el=fill_el, line_el=line_el, geom_el=geom_el)
        bbox = {'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        self._append_shape_element(target_sp, shape, f'PL{self.sid}', bbox=bbox)

    def _use(self, el, sp, ox, oy, opacity, scale, slide):
        """SVG <use> -> 解析引用并递归处理。"""
        href = el.get(f'{{{XLINK_NS}}}href') or el.get('href', '')
        if not href or not href.startswith('#'):
            return
        ref_id = href[1:]
        # 在整个 SVG 树中查找被引用元素
        ref_el = None
        for candidate in self.root.iter():
            if candidate.get('id') == ref_id:
                ref_el = candidate
                break
        if ref_el is None:
            return
        use_x = float(el.get('x', 0)) * scale
        use_y = float(el.get('y', 0)) * scale
        self._walk(ref_el, sp, ox + use_x, oy + use_y, opacity, scale, slide)

    def _path(self, el, sp, ox, oy, opacity, scale):
        """SVG <path> -> OOXML custGeom 形状。"""
        d = el.get('d', '')
        if not d or 'nan' in d:
            return
        # 计算 bounding box（简化：从 path 数据提取所有数字坐标）
        nums = re.findall(r'[+-]?(?:\d+\.?\d*|\.\d+)', d)
        if len(nums) < 4:
            return
        coords = [float(n) for n in nums]
        xs = coords[0::2]
        ys = coords[1::2] if len(coords) > 1 else [0]
        bx, by = min(xs), min(ys)
        bw = max(xs) - bx or 1
        bh = max(ys) - by or 1

        # 过滤极小路径
        if bw < 4 and bh < 4:
            self.stats['skipped'] += 1
            return

        if self._table_region_for_bbox({'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}):
            self.stats['skipped'] += 1
            return

        geom_el = parse_path_to_custgeom(d, (bx, by, bw, bh))
        el_opacity = float(el.get('opacity', '1')) * opacity
        fill_el = make_fill(el.get('fill', ''), self.grads, el_opacity)
        line_el = make_line(el.get('stroke', ''), el.get('stroke-width', '1')) if el.get('stroke') else None

        shape = make_shape(self._id(), f'P{self.sid}',
                           px((bx + ox) * scale) if scale != 1.0 else px(bx + ox),
                           px((by + oy) * scale) if scale != 1.0 else px(by + oy),
                           px(bw * scale), px(bh * scale),
                           fill_el=fill_el, line_el=line_el, geom_el=geom_el)
        bbox = {'x': bx + ox, 'y': by + oy, 'width': bw * scale, 'height': bh * scale}
        target_sp = self._target_tree_for_bbox(sp, bbox)
        self._append_shape_element(target_sp, shape, f'P{self.sid}', bbox=bbox)

    def _image(self, el, sp, ox, oy, opacity, scale, slide):
        href = el.get(f'{{{XLINK_NS}}}href') or el.get('href', '')
        x = float(el.get('x', 0)) * scale + ox
        y = float(el.get('y', 0)) * scale + oy
        w = float(el.get('width', 0)) * scale
        h = float(el.get('height', 0)) * scale
        el_opacity = float(el.get('opacity', '1')) * opacity
        if not href or w <= 0 or h <= 0:
            return
        if self._table_region_for_bbox({'x': x, 'y': y, 'width': w, 'height': h}):
            self.stats['skipped'] += 1
            return
        if self._native_chart_region_for_bbox({'x': x, 'y': y, 'width': w, 'height': h}):
            self.stats['skipped'] += 1
            return

        img_source = None
        if href.startswith('data:'):
            m = re.match(r'data:image/\w+;base64,(.*)', href, re.DOTALL)
            if m:
                img_source = io.BytesIO(base64.b64decode(m.group(1)))
        elif href.startswith('file://'):
            p = Path(href.replace('file://', ''))
            if p.exists():
                img_source = str(p)
        elif not href.startswith('http'):
            p = Path(href)
            if p.exists():
                img_source = str(p)

        if img_source is None:
            return

        # 获取图片原始尺寸以计算宽高比
        try:
            from PIL import Image as PILImage
            if isinstance(img_source, io.BytesIO):
                img_source.seek(0)
                pil_img = PILImage.open(img_source)
                img_w, img_h = pil_img.size
                # 不 close -- PIL close 会关掉底层 BytesIO
                del pil_img
                img_source.seek(0)
            else:
                with PILImage.open(img_source) as pil_img:
                    img_w, img_h = pil_img.size
        except ImportError:
            # 没有 PIL，退回直接拉伸
            picture_shapes = slide.shapes
            chart_region = self._chart_region_for_bbox({'x': x, 'y': y, 'width': w, 'height': h})
            if chart_region is not None:
                picture_shapes = self._ensure_chart_group(chart_region).shapes
            pic = picture_shapes.add_picture(img_source,
                                             Emu(px(x)), Emu(px(y)),
                                             Emu(px(w)), Emu(px(h)))
            if chart_region is None:
                self._assign_shape_name(pic, f'IMG{self.sid}', bbox={'x': x, 'y': y, 'width': w, 'height': h})
            self.stats['shapes'] += 1
            return

        # object-fit: cover -- 容器尺寸放置 + srcRect 源裁剪
        container_w = px(w)
        container_h = px(h)
        img_ratio = img_w / img_h
        container_ratio = container_w / container_h

        # 计算源裁剪区域 (srcRect, 百分比 0-100000)
        if img_ratio > container_ratio:
            # 图片更宽 -> 按高度填满，裁剪左右
            visible_w_pct = container_ratio / img_ratio  # 0~1
            crop_lr = int((1 - visible_w_pct) / 2 * 100000)
            crop_tb = 0
        else:
            # 图片更高 -> 按宽度填满，裁剪上下
            visible_h_pct = img_ratio / container_ratio  # 0~1
            crop_lr = 0
            crop_tb = int((1 - visible_h_pct) / 2 * 100000)

        # 以容器尺寸放置图片（不放大）
        picture_shapes = slide.shapes
        chart_region = self._chart_region_for_bbox({'x': x, 'y': y, 'width': w, 'height': h})
        if chart_region is not None:
            picture_shapes = self._ensure_chart_group(chart_region).shapes
        pic = picture_shapes.add_picture(img_source,
                                         Emu(px(x)), Emu(px(y)),
                                         Emu(container_w), Emu(container_h))
        if chart_region is None:
            self._assign_shape_name(pic, f'IMG{self.sid}', bbox={'x': x, 'y': y, 'width': w, 'height': h})

        # 用 srcRect 在 blipFill 内定义源裁剪区域（等效 object-fit: cover）
        if crop_lr > 0 or crop_tb > 0:
            from pptx.oxml.ns import qn
            blip_fill = pic._element.find(qn('p:blipFill'))
            if blip_fill is not None:
                src_rect = _el('a:srcRect', {
                    'l': str(crop_lr), 't': str(crop_tb),
                    'r': str(crop_lr), 'b': str(crop_tb)
                })
                # stretch 前面插入 srcRect
                stretch = blip_fill.find(qn('a:stretch'))
                if stretch is not None:
                    blip_fill.insert(list(blip_fill).index(stretch), src_rect)
                else:
                    blip_fill.append(src_rect)

        # 应用透明度（通过 OOXML alphaModFix）
        if el_opacity < 0.99:
            from pptx.oxml.ns import qn
            sp_pr = pic._element.find(qn('p:spPr'))
            if sp_pr is None:
                sp_pr = pic._element.find(qn('pic:spPr'))
            # 在 blipFill 的 blip 上设置 alphaModFix
            blip = pic._element.find('.//' + qn('a:blip'))
            if blip is not None:
                alpha_val = int(el_opacity * 100000)
                alpha_el = _el('a:alphaModFix', {'amt': str(alpha_val)})
                blip.append(alpha_el)

        self.stats['shapes'] += 1


# -------------------------------------------------------------------
# 主流程
# -------------------------------------------------------------------
def convert(svg_input, output_path, on_progress=None, export_report_path=None, html_dir=None, report_path=None,
            template_pptx=None, target_slides=None, preserve_template_background=False, speech_script=None):
    svg_input = Path(svg_input)
    output_path = Path(output_path)
    svg_files = resolve_svg_files(svg_input)

    template_pptx_path = Path(template_pptx).resolve() if template_pptx else None
    if template_pptx_path and not template_pptx_path.exists():
        print(f"Error: template PPTX not found: {template_pptx_path}", file=sys.stderr)
        sys.exit(1)

    if template_pptx_path:
        prs = Presentation(str(template_pptx_path))
        if int(prs.slide_width) != SLIDE_W or int(prs.slide_height) != SLIDE_H:
            print(
                f"Error: template PPTX slide size must be {SLIDE_W}x{SLIDE_H} EMU; "
                f"got {int(prs.slide_width)}x{int(prs.slide_height)}",
                file=sys.stderr,
            )
            sys.exit(1)
        try:
            target_slide_numbers = parse_target_slide_numbers(target_slides, len(svg_files))
        except ValueError as exc:
            print(f"Error: {exc}", file=sys.stderr)
            sys.exit(1)
        if target_slide_numbers is None:
            print("Error: --target-slides is required when --template-pptx is used", file=sys.stderr)
            sys.exit(1)
        if any(number > len(prs.slides) for number in target_slide_numbers):
            print(
                f"Error: target slide out of range (deck has {len(prs.slides)} slides): {target_slide_numbers}",
                file=sys.stderr,
            )
            sys.exit(1)
        blank = None
        update_mode = 'template_update'
    else:
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_W)
        prs.slide_height = Emu(SLIDE_H)
        blank = prs.slide_layouts[6]
        target_slide_numbers = list(range(1, len(svg_files) + 1))
        update_mode = 'new_presentation'

    converter = SvgConverter(on_progress=on_progress)
    total = len(svg_files)
    speech_pages = None
    if speech_script:
        try:
            speech_pages = load_speech_page_entries(speech_script, expected_pages=total)
        except ValueError as exc:
            print(f"Error: {exc}", file=sys.stderr)
            sys.exit(1)
    source_report_path = resolve_source_report_path(svg_input, output_path, export_report_path)
    source_report, top_warnings = load_source_report(source_report_path)
    html_dir_path = Path(html_dir).resolve() if html_dir else None
    if html_dir_path and not html_dir_path.exists():
        top_warnings.append(f'html_dir_missing:{html_dir_path}')
        html_dir_path = None
    source_pages = build_source_page_lookup(source_report)
    html_semantic_cache = {}
    rendered_semantic_cache = {}
    if isinstance(source_report, dict):
        source_page_count = len(source_report.get('pages', [])) if isinstance(source_report.get('pages'), list) else None
        if source_page_count is not None and source_page_count != total:
            top_warnings.append(f'source_page_count_mismatch:{source_page_count}!={total}')
    page_reports = []

    for i, svg_file in enumerate(svg_files):
        target_slide_number = target_slide_numbers[i]
        slide = prs.slides[target_slide_number - 1] if template_pptx_path else prs.slides.add_slide(blank)
        source_page = source_pages.get(svg_file.name, {})
        source_method = source_page.get('method', 'unknown') if isinstance(source_page, dict) else 'unknown'
        source_warning = source_page.get('warning') if isinstance(source_page, dict) else None
        html_semantics = None
        html_warnings = []
        html_source_path = resolve_html_source_path(svg_file, html_dir_path)
        if html_source_path:
            cache_key = str(html_source_path)
            if cache_key not in html_semantic_cache:
                html_semantic_cache[cache_key] = load_html_semantics(html_source_path)
            html_semantics, html_warnings = html_semantic_cache[cache_key]
        elif html_dir_path:
            html_warnings = [f'html_semantics_missing:{svg_file.stem}.html']

        rendered_semantics = None
        rendered_warnings = []
        rendered_semantic_path = resolve_semantic_sidecar_path(svg_file, source_page, source_report_path=source_report_path)
        if rendered_semantic_path:
            cache_key = str(rendered_semantic_path)
            if cache_key not in rendered_semantic_cache:
                rendered_semantic_cache[cache_key] = load_rendered_semantics(rendered_semantic_path)
            rendered_semantics, rendered_warnings = rendered_semantic_cache[cache_key]
        elif source_method == 'dom_to_svg_editable':
            rendered_warnings = [f'rendered_semantics_missing:{svg_file.name}']

        update_blocks = collect_update_blocks(html_semantics, rendered_semantics) if template_pptx_path else []
        updated_block_ids = [block.get('block_id') for block in update_blocks if block.get('block_id')]
        template_update_stats = {
            'removed_total': 0,
            'removed_slot_shapes': 0,
            'removed_managed_shapes': 0,
            'target_blocks': [],
        }
        if template_pptx_path:
            template_update_stats = clear_target_block_shapes(slide, updated_block_ids)

        converter.convert(
            svg_file,
            slide,
            html_semantics=html_semantics,
            rendered_semantics=rendered_semantics,
            preserve_existing_background=bool(template_pptx_path and preserve_template_background),
        )
        if speech_pages:
            write_slide_speaker_notes(slide, speech_pages[i]["speaker_notes"])
        s = converter.stats
        semantic_note = ''
        if s.get('semantic_entries'):
            semantic_chunks = [f"{s['semantic_matches']}/{s['semantic_entries']} semantic matches"]
            if s.get('block_matches'):
                semantic_chunks.append(f"{s['block_matches']} block-bound")
            if s.get('native_tables'):
                semantic_chunks.append(f"{s['native_tables']} native tables")
            if s.get('native_charts'):
                semantic_chunks.append(f"{s['native_charts']} native charts")
            if s.get('structured_chart_groups'):
                semantic_chunks.append(f"{s['structured_chart_groups']} chart groups")
            semantic_note = ', ' + ', '.join(semantic_chunks)
        print(f"  [{i+1}/{total}] {svg_file.name} via {source_method} "
              f"({s['shapes']} shapes, {s['skipped']} skipped, {s['errors']} errors{semantic_note})")
        page_warning_list = []
        if isinstance(source_warning, str) and source_warning.strip():
            page_warning_list.append(source_warning)
        page_warning_list.extend(html_warnings)
        page_warning_list.extend(rendered_warnings)
        if template_pptx_path and updated_block_ids and template_update_stats['removed_total'] == 0:
            page_warning_list.append(f"template_block_replace_no_matches:{','.join(updated_block_ids)}")
        html_summary = html_semantics.get('summary') if isinstance(html_semantics, dict) else {}
        rendered_summary = rendered_semantics.get('summary') if isinstance(rendered_semantics, dict) else {}
        page_reports.append({
            'slide_number': target_slide_number,
            'svg_file': svg_file.name,
            'source_method': source_method,
            'source_editable': source_page.get('editable') if isinstance(source_page, dict) else None,
            'source_text_count': source_page.get('text_count') if isinstance(source_page, dict) else None,
            'html_source_path': html_semantics.get('path') if isinstance(html_semantics, dict) else None,
            'rendered_semantic_path': str(rendered_semantic_path.resolve()) if rendered_semantic_path else None,
            'html_semantic_entries': s.get('semantic_entries'),
            'html_semantic_matches': s.get('semantic_matches'),
            'html_semantic_unmatched': s.get('semantic_unmatched'),
            'html_block_matches': s.get('block_matches'),
            'html_list_items': html_summary.get('list_items'),
            'html_tables': html_summary.get('tables'),
            'html_blocks': s.get('semantic_blocks') or html_summary.get('blocks'),
            'rendered_blocks': rendered_summary.get('blocks'),
            'rendered_tables': rendered_summary.get('tables'),
            'rendered_charts': s.get('rendered_charts') or rendered_summary.get('charts'),
            'native_charts': s.get('native_charts'),
            'native_tables': s.get('native_tables'),
            'structured_chart_groups': s.get('structured_chart_groups'),
            'pptx_shapes': s['shapes'],
            'pptx_skipped': s['skipped'],
            'pptx_errors': s['errors'],
            'updated_block_ids': updated_block_ids if template_pptx_path else [],
            'template_update_scope': 'block_update' if template_pptx_path else None,
            'template_removed_shapes': template_update_stats['removed_total'] if template_pptx_path else 0,
            'template_removed_slot_shapes': template_update_stats['removed_slot_shapes'] if template_pptx_path else 0,
            'template_removed_managed_shapes': template_update_stats['removed_managed_shapes'] if template_pptx_path else 0,
            'speaker_notes_chars': len(speech_pages[i]['speaker_notes']) if speech_pages else 0,
            'warnings': page_warning_list,
        })
        if on_progress:
            on_progress(i + 1, total, svg_file.name)

    prs.save(str(output_path))
    report_path = Path(report_path).resolve() if report_path else output_path.with_suffix('.report.json')
    write_pptx_report(
        report_path,
        build_pptx_report(
            output_path=output_path,
            svg_input=svg_input,
            html_dir=html_dir,
            source_report_path=source_report_path,
            source_report=source_report,
            top_warnings=top_warnings,
            page_reports=page_reports,
            template_pptx_path=template_pptx_path,
            target_slide_numbers=target_slide_numbers,
            update_mode=update_mode,
            preserve_template_background=bool(template_pptx_path and preserve_template_background),
        ),
    )
    print(f"Saved: {output_path} ({len(prs.slides) if template_pptx_path else total} slides)")
    print(f"Report: {report_path}")


def main():
    parser = argparse.ArgumentParser(description="SVG to PPTX (native shapes)")
    parser.add_argument('svg', help='SVG file or directory')
    parser.add_argument('-o', '--output', default='presentation.pptx')
    parser.add_argument('--html-dir', default=None,
                        help='HTML source directory or file used for semantic text reconstruction')
    parser.add_argument('--export-report', default=None,
                        help='Path to html2svg report JSON (auto-detected if omitted)')
    parser.add_argument('--report-path', default=None,
                        help='Path to write PPTX export report JSON')
    parser.add_argument('--template-pptx', default=None,
                        help='Existing PPTX to update in-place while preserving theme/master/layout')
    parser.add_argument('--target-slides', default=None,
                        help='Comma-separated 1-based target slide numbers for template update mode')
    parser.add_argument('--speech-script', default=None,
                        help='Optional speech-script.json used to write speaker notes')
    parser.add_argument('--replace-template-background', action='store_true',
                        help='In template update mode, replace the target slide background using the SVG full-slide background rect')
    args = parser.parse_args()
    convert(
        args.svg,
        args.output,
        export_report_path=args.export_report,
        html_dir=args.html_dir,
        report_path=args.report_path,
        template_pptx=args.template_pptx,
        target_slides=args.target_slides,
        preserve_template_background=not args.replace_template_background,
        speech_script=args.speech_script,
    )


if __name__ == '__main__':
    main()
