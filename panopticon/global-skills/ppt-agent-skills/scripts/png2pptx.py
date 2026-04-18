#!/usr/bin/env python3
"""PNG to PPTX -- 将 PNG 截图嵌入 PowerPoint 幻灯片

每张 PNG 作为全屏图片覆盖整个幻灯片（16:9, 1280x720 逻辑像素）。
比 SVG 管线兼容性更好：
- 任何版本 PowerPoint / WPS / Keynote / Google Slides 均可打开
- 所有 CSS 效果完美保留（渐变、阴影、模糊等）
- 缺点：文字不可编辑（转为像素）

用法：
    python3 scripts/png2pptx.py <png_dir_or_file> -o output.pptx
"""

import argparse
import collections.abc
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from speech_script import load_speech_page_entries, write_slide_speaker_notes

# 16:9 标准尺寸（EMU）
SLIDE_W = 12192000
SLIDE_H = 6858000


def _natural_key(p):
    """自然排序 key：slide2 排在 slide10 前面。"""
    parts = re.split(r'(\d+)', p.stem)
    return [int(x) if x.isdigit() else x.lower() for x in parts]


def convert(png_input, output_path, on_progress=None, speech_script=None):
    """PNG 文件 -> PPTX 幻灯片。"""
    png_input = Path(png_input)
    if png_input.is_file():
        png_files = [png_input]
    elif png_input.is_dir():
        png_files = sorted(png_input.glob('*.png'), key=_natural_key)
    else:
        print(f"Error: {png_input} not found", file=sys.stderr)
        sys.exit(1)

    if not png_files:
        print("Error: No PNG files found", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]  # 空白版式
    total = len(png_files)
    speech_pages = None
    if speech_script:
        try:
            speech_pages = load_speech_page_entries(speech_script, expected_pages=total)
        except ValueError as exc:
            print(f"Error: {exc}", file=sys.stderr)
            sys.exit(1)

    for i, png_file in enumerate(png_files):
        slide = prs.slides.add_slide(blank)

        # 全屏放置 PNG（左上角 0,0，宽高覆盖整个幻灯片）
        slide.shapes.add_picture(
            str(png_file),
            Emu(0), Emu(0),
            Emu(SLIDE_W), Emu(SLIDE_H)
        )
        if speech_pages:
            write_slide_speaker_notes(slide, speech_pages[i]["speaker_notes"])

        print(f"  [{i+1}/{total}] {png_file.name}")
        if on_progress:
            on_progress(i + 1, total, png_file.name)

    prs.save(str(output_path))
    print(f"Saved: {output_path} ({total} slides)")


def main():
    parser = argparse.ArgumentParser(description="PNG to PPTX (full-slide images)")
    parser.add_argument('png', help='PNG file or directory')
    parser.add_argument('-o', '--output', default='presentation.pptx')
    parser.add_argument('--speech-script', default=None, help='Optional speech-script.json used to write speaker notes')
    args = parser.parse_args()
    convert(args.png, args.output, speech_script=args.speech_script)


if __name__ == '__main__':
    main()
