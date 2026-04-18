#!/usr/bin/env python3
"""Minimal end-to-end smoke test for the PPT workflow skill.

This script intentionally stays within the current markdown/code architecture.
It exercises the most failure-prone integration points:
1. Step 0 interview prompt rendering (structured/text dual templates)
2. Step 4 planning example -> planning_validator.py
3. resource_loader.py menu / resolve / images
4. prompt_harness.py for the Step 4 prompt chain
"""

from __future__ import annotations

import argparse
import base64
import json
import re
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from speech_script import extract_slide_speaker_notes
from workflow_versions import (  # noqa: E402
    PLANNING_CONTINUITY_VERSION,
    PLANNING_PACKET_VERSION,
    PLANNING_SCHEMA_VERSION,
    WORKFLOW_VERSION,
)


ROOT_DIR = Path(__file__).resolve().parents[1]
SCRIPTS_DIR = ROOT_DIR / "scripts"
REFERENCES_DIR = ROOT_DIR / "references"
PLAYBOOK_PATH = REFERENCES_DIR / "playbooks/step4/page-planning-playbook.md"
PAGE_TEMPLATE_EXPECTATIONS = {
    "cover": "# 封面页 -- 演讲的第一声呼吸",
    "toc": "# 目录页 -- 演讲的地图俯瞰",
    "section": "# 章节封面页 -- 演讲中的呼吸",
    "end": "# 结束页 -- 演讲的最后一个视觉印记",
}


@dataclass
class SmokeResult:
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    steps: list[str] = field(default_factory=list)

    def error(self, message: str) -> None:
        self.errors.append(message)

    def warn(self, message: str) -> None:
        self.warnings.append(message)

    def note(self, message: str) -> None:
        self.steps.append(message)


def run_cmd(label: str, args: list[str], result: SmokeResult, cwd: Path = ROOT_DIR) -> subprocess.CompletedProcess[str]:
    proc = subprocess.run(
        args,
        cwd=str(cwd),
        text=True,
        capture_output=True,
    )
    if proc.returncode != 0:
        result.error(
            f"{label}: exit={proc.returncode}\n"
            f"cmd={' '.join(args)}\n"
            f"stdout:\n{proc.stdout}\n"
            f"stderr:\n{proc.stderr}"
        )
    else:
        result.note(f"{label}: ok")
    return proc


def run_cmd_expect(
    label: str,
    args: list[str],
    expected_codes: set[int],
    result: SmokeResult,
    cwd: Path = ROOT_DIR,
) -> subprocess.CompletedProcess[str]:
    proc = subprocess.run(
        args,
        cwd=str(cwd),
        text=True,
        capture_output=True,
    )
    if proc.returncode not in expected_codes:
        result.error(
            f"{label}: exit={proc.returncode}, expected one of {sorted(expected_codes)}\n"
            f"cmd={' '.join(args)}\n"
            f"stdout:\n{proc.stdout}\n"
            f"stderr:\n{proc.stderr}"
        )
    else:
        result.note(f"{label}: ok(exit={proc.returncode})")
    return proc


def write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def write_svg(path: Path, width: int, height: int, fill: str) -> None:
    write_text(
        path,
        (
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">'
            f'<rect width="100%" height="100%" fill="{fill}"/>'
            f'<text x="50%" y="50%" dominant-baseline="middle" text-anchor="middle" '
            f'font-size="48" fill="#ffffff">{width}x{height}</text>'
            '</svg>'
        ),
    )


def build_content_page_fixture(image_source_hint: str | None = None) -> dict[str, object]:
    """Build a minimal content page planning fixture for smoke testing."""
    return {
        "page": {
            "slide_number": 3,
            "page_type": "content",
            "narrative_role": "evidence",
            "title": "增长判断",
            "page_goal": "证明增长成立",
            "audience_takeaway": "增长数据可信",
            "visual_weight": 7,
            "layout_hint": "hero-top",
            "layout_variation_note": "与上一页重心不同",
            "focus_zone": "右上 1/3 作为视觉锚点",
            "negative_space_target": "medium",
            "page_text_strategy": "标题强、正文短、数据做锚点",
            "rhythm_action": "推进",
            "must_avoid": ["禁止平均分栏"],
            "variation_guardrails": {
                "same_gene_as_deck": "保留统一字体和 signature_move",
                "different_from_previous": ["重心从上移到右"],
            },
            "director_command": {
                "mood": "判断感强、结论先行",
                "spatial_strategy": "主锚占据第一视线",
                "anchor_treatment": "用尺度断层强化主锚",
                "techniques": ["T1", "W3"],
                "prose": "保持证据链清晰",
            },
            "decoration_hints": {
                "background": {"feel": "轻微渐变底", "restraint": "不抢文字对比", "techniques": ["T1"]},
                "floating": {"feel": "局部辅助装饰", "restraint": "只服务锚点动线", "techniques": ["W3"]},
                "page_accent": {"feel": "强调色集中在锚点附近", "restraint": "accent 只用 1-2 种", "techniques": ["T9"]},
            },
            "source_guidance": {
                "brief_sections": ["核心发现"],
                "citation_expectation": "有数字就保留来源",
                "strictness": "不得超出 brief 结论边界",
            },
            "resources": {
                "page_template": None,
                "layout_refs": ["hero-top"],
                "block_refs": [],
                "chart_refs": ["kpi", "metric-row"],
                "principle_refs": ["visual-hierarchy", "composition"],
                "resource_rationale": "用 hero-top 放大单一结论",
            },
            "cards": [
                {
                    "card_id": "s03-anchor",
                    "role": "anchor",
                    "card_type": "data_highlight",
                    "card_style": "accent",
                    "argument_role": "claim",
                    "headline": "核心指标",
                    "body": ["一句解释它为什么重要"],
                    "data_points": [
                        {"label": "同比增长", "value": "47.3", "unit": "%", "source": "search-brief metrics[2]"}
                    ],
                    "chart": {"chart_type": "kpi"},
                    "content_budget": {"headline_max_chars": 12, "body_max_bullets": 2, "body_max_lines": 4},
                    "image": {
                        "mode": "decorate",
                        "needed": False,
                        "usage": None,
                        "placement": None,
                        "content_description": None,
                        "source_hint": None,
                        "decorate_brief": "用内联 SVG 装饰填满留白",
                    },
                    "resource_ref": {"chart": "kpi", "principle": "visual-hierarchy"},
                },
                {
                    "card_id": "s03-support-1",
                    "role": "support",
                    "card_type": "data",
                    "card_style": "outline",
                    "argument_role": "evidence",
                    "headline": "增长原因",
                    "body": ["增长主要来自高客单区域放量"],
                    "data_points": [
                        {"label": "高客单区域占比", "value": "31", "unit": "%", "source": "search-brief metrics[4]"}
                    ],
                    "chart": {"chart_type": "metric_row"},
                    "content_budget": {"headline_max_chars": 12, "body_max_bullets": 2, "body_max_lines": 4},
                    "image": {
                        "mode": "decorate",
                        "needed": False,
                        "usage": None,
                        "placement": None,
                        "content_description": None,
                        "source_hint": None,
                        "decorate_brief": "用低对比度辅助线承托信息",
                    },
                    "resource_ref": {"chart": "metric-row", "principle": "composition"},
                },
                {
                    "card_id": "s03-context-image",
                    "role": "context",
                    "card_type": "image_hero",
                    "card_style": "filled",
                    "argument_role": "evidence",
                    "headline": "场景证据",
                    "body": ["配图用于强化增长结论的真实场景感。"],
                    "data_points": [],
                    "content_budget": {"headline_max_chars": 12, "body_max_bullets": 1, "body_max_lines": 3},
                    "image": {
                        "mode": "provided",
                        "needed": True,
                        "usage": "inline-illustration",
                        "placement": "right-half",
                        "content_description": "用于支撑增长结论的场景配图",
                        "source_hint": image_source_hint,
                        "decorate_brief": None,
                    },
                    "resource_ref": {"block": "image-hero", "principle": "composition"},
                },
            ],
            "workflow_metadata": {
                "stage": "planning",
                "workflow_version": WORKFLOW_VERSION,
                "planning_schema_version": PLANNING_SCHEMA_VERSION,
                "planning_packet_version": PLANNING_PACKET_VERSION,
                "planning_continuity_version": PLANNING_CONTINUITY_VERSION,
            },
        }
    }


def assert_contains(label: str, haystack: str, needles: list[str], result: SmokeResult) -> None:
    missing = [needle for needle in needles if needle not in haystack]
    if missing:
        result.error(f"{label}: missing expected content {missing}")


def assert_no_unfilled_vars(label: str, text: str, result: SmokeResult) -> None:
    leftovers = sorted(set(re.findall(r"\{\{[A-Z_][A-Z0-9_]*\}\}", text)))
    if leftovers:
        result.error(f"{label}: unfilled template vars remain: {leftovers}")


def assert_max_bytes(label: str, text: str, max_bytes: int, result: SmokeResult) -> None:
    size = len(text.encode("utf-8"))
    if size > max_bytes:
        result.error(f"{label}: rendered prompt too large ({size} bytes > {max_bytes} bytes)")


def build_non_content_page(page_type: str) -> dict[str, object]:
    return {
        "page": {
            "slide_number": 1,
            "page_type": page_type,
            "narrative_role": "opening" if page_type == "cover" else "transition",
            "title": f"Smoke {page_type}",
            "page_goal": f"验证 {page_type} 页面模板路由",
            "audience_takeaway": f"{page_type} page template resolve",
            "visual_weight": 7,
            "focus_zone": "center",
            "negative_space_target": "medium",
            "page_text_strategy": "短句为主",
            "rhythm_action": "推进",
            "must_avoid": [],
            "variation_guardrails": {
                "same_gene_as_deck": "保留统一风格变量",
                "different_from_previous": ["验证 page template 路由"],
            },
            "director_command": {
                "mood": "测试态",
                "spatial_strategy": "居中聚焦",
                "anchor_treatment": "标题优先",
                "techniques": ["T1"],
                "prose": "用于验证非 content 页的模板消费链。",
            },
            "decoration_hints": {
                "background": {"feel": "轻量背景", "restraint": "不抢主标题", "techniques": ["T1"]},
                "floating": {"feel": "弱装饰", "restraint": "仅做陪衬", "techniques": []},
                "page_accent": {"feel": "少量强调色", "restraint": "仅一处强调", "techniques": []},
            },
            "resources": {
                "page_template": None,
                "layout_refs": [],
                "block_refs": [],
                "chart_refs": [],
                "principle_refs": [],
                "resource_rationale": "验证 page_type 自动路由到 page-templates/",
            },
            "cards": [
                {
                    "card_id": "s01-anchor",
                    "role": "anchor",
                    "card_type": "text",
                    "card_style": "accent",
                    "headline": f"{page_type} smoke",
                    "body": ["最小非 content 页冒烟样例"],
                    "content_budget": {"headline_max_chars": 12, "body_max_bullets": 1, "body_max_lines": 2},
                    "image": {
                        "mode": "decorate",
                        "needed": False,
                        "usage": None,
                        "placement": None,
                        "content_description": None,
                        "source_hint": None,
                        "decorate_brief": "只做轻量占位，不引入外部图片。",
                    },
                }
            ],
            "workflow_metadata": {
                "stage": "planning",
                "workflow_version": WORKFLOW_VERSION,
                "planning_schema_version": PLANNING_SCHEMA_VERSION,
                "planning_packet_version": PLANNING_PACKET_VERSION,
                "planning_continuity_version": PLANNING_CONTINUITY_VERSION,
            },
        }
    }


def build_fixture_tree(tmp_dir: Path) -> dict[str, Path]:
    fixtures = {
        "requirements": tmp_dir / "requirements-interview.txt",
        "outline": tmp_dir / "outline.txt",
        "brief": tmp_dir / "search-brief.txt",
        "style": tmp_dir / "style.json",
        "planning": tmp_dir / "planning/planning3.json",
        "slide": tmp_dir / "slides/slide-3.html",
        "png": tmp_dir / "png/slide-3.png",
        "images": tmp_dir / "images",
        "runtime": tmp_dir / "runtime",
        "prompt_interview_structured": tmp_dir / "runtime/prompt-interview-structured.md",
        "prompt_interview_text": tmp_dir / "runtime/prompt-interview-text.md",
        "prompt_style_phase1": tmp_dir / "runtime/prompt-style-phase1.md",
        "prompt_planning": tmp_dir / "runtime/prompt-page-planning-3.md",
        "prompt_html": tmp_dir / "runtime/prompt-page-html-3.md",
        "prompt_review": tmp_dir / "runtime/prompt-page-review-3.md",
        "prompt_orchestrator": tmp_dir / "runtime/prompt-page-orchestrator-3.md",
    }

    write_text(
        fixtures["requirements"],
        "# 需求归一化\n\n## 基本信息\n- 主题：Smoke Test\n- 项目类型：演示文稿\n- 语言：中文\n- 输入类型：示例\n- 分支：research\n",
    )
    write_text(fixtures["outline"], "# 大纲\n\n## Part 1: Demo\n\n### 第 3 页：增长判断\n- 页目标：增长成立\n")
    write_text(fixtures["brief"], "# Research Brief\n\n## 核心发现\n1. 示例发现 [来源: smoke]\n")
    write_text(
        fixtures["style"],
        json.dumps(
            {
                "style_id": "smoke",
                "style_name": "Smoke",
                "mood_keywords": ["clear", "structured", "modern"],
                "design_soul": "清晰、克制、强调论点主次。",
                "variation_strategy": "统一色彩与边角，允许每页在布局重心和装饰位置上变化。",
                "decoration_dna": {
                    "signature_move": "轻微几何线条",
                    "forbidden": ["过强噪点"],
                    "recommended_combos": ["outline + accent"],
                },
                "font_family": "Noto Sans SC",
                "css_variables": {
                    "--bg-primary": "#0f172a",
                    "--bg-secondary": "#111827",
                    "--card-bg-from": "#1f2937",
                    "--card-bg-to": "#111827",
                    "--card-border": "#334155",
                    "--card-radius": "24px",
                    "--text-primary": "#f8fafc",
                    "--text-secondary": "#cbd5e1",
                    "--accent-1": "#38bdf8",
                    "--accent-2": "#22c55e",
                    "--accent-3": "#f59e0b",
                    "--accent-4": "#a78bfa",
                    "--font-primary": "Noto Sans SC",
                },
            },
            ensure_ascii=False,
            indent=2,
        ),
    )
    fixtures["images"].mkdir(parents=True, exist_ok=True)
    write_svg(fixtures["images"] / "landscape.svg", 1600, 900, "#2563eb")
    write_svg(fixtures["images"] / "portrait.svg", 900, 1400, "#dc2626")
    write_svg(fixtures["images"] / "square.svg", 1200, 1200, "#16a34a")
    write_text(
        fixtures["planning"],
        json.dumps(
            build_content_page_fixture(str((fixtures["images"] / "landscape.svg").resolve())),
            ensure_ascii=False,
            indent=2,
        ),
    )
    return fixtures


def build_visual_qa_fixtures(tmp_dir: Path, result: SmokeResult) -> dict[str, Path] | None:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        result.warn("visual-qa smoke skipped: Pillow unavailable")
        return None

    root = tmp_dir / "visual-qa"
    good_root = root / "good"
    warn_root = root / "warn"
    footer_broken_root = root / "footer-broken"
    header_broken_root = root / "header-broken"
    for base in (good_root, warn_root, footer_broken_root, header_broken_root):
        (base / "png").mkdir(parents=True, exist_ok=True)
        (base / "planning").mkdir(parents=True, exist_ok=True)
        (base / "slides").mkdir(parents=True, exist_ok=True)
        (base / "images").mkdir(parents=True, exist_ok=True)

    font = ImageFont.load_default()

    def create_variant(base: Path, variant: str) -> tuple[Path, Path, Path]:
        image_asset = base / "images" / "scene.png"
        asset_img = Image.new("RGB", (1600, 900), "#2563eb")
        asset_draw = ImageDraw.Draw(asset_img)
        for y in range(0, 900, 12):
            color = (37 + (y % 60), 99 + (y % 40), 180 + (y % 30))
            asset_draw.rectangle((0, y, 1600, min(899, y + 11)), fill=color)
        asset_img.save(image_asset)

        slide_html = base / "slides" / "slide-3.html"
        if variant == "footer-broken":
            slide_html.write_text(
                """
<html><head><style>
body { margin: 0; width: 1280px; height: 720px; overflow: hidden; background: #0f172a; }
.slide-header { position: absolute; top: 20px; left: 40px; right: 40px; height: 50px; display: flex; gap: 16px; align-items: baseline; }
.media-shell { width: 360px; height: 240px; overflow: hidden; border-radius: 24px; }
.media-shell img { width: 100%; height: 100%; object-fit: cover; display: block; }
</style></head><body>
<header class="slide-header"><span class="overline">PART 03</span><h1 class="page-title">增长判断</h1></header>
<div class="media-shell"><img src="../images/scene.png" alt="scene"></div>
<footer class="deck-footer"><span class="section-label">GROWTH</span></footer>
</body></html>
""".strip(),
                encoding="utf-8",
            )
        elif variant == "header-broken":
            slide_html.write_text(
                """
<html><head><style>
body { margin: 0; width: 1280px; height: 720px; overflow: hidden; background: #0f172a; }
.media-shell { width: 360px; height: 240px; overflow: hidden; border-radius: 24px; }
.media-shell img { width: 100%; height: 100%; object-fit: cover; display: block; }
.slide-footer { position: absolute; left: 40px; right: 40px; bottom: 12px; height: 20px; }
</style></head><body>
<header class="top-banner"><span class="eyebrow">PART 03</span><div class="title-line">增长判断</div></header>
<div class="media-shell"><img src="../images/scene.png" alt="scene"></div>
<footer class="slide-footer"><span class="footer-section">GROWTH</span><span class="footer-page">3 / 12</span></footer>
</body></html>
""".strip(),
                encoding="utf-8",
            )
        elif variant == "warn":
            slide_html.write_text(
                """
<html><head><style>
body { margin: 0; width: 1280px; height: 720px; overflow: hidden; background: #0f172a; }
.slide-header { position: absolute; top: 20px; left: 40px; right: 40px; height: 50px; display: flex; gap: 16px; align-items: baseline; }
.media-shell { width: 360px; height: 240px; overflow: hidden; border-radius: 24px; }
.media-shell img { width: 100%; height: 100%; object-fit: cover; display: block; }
.slide-footer { position: absolute; left: 40px; right: 40px; bottom: 12px; height: 20px; }
</style></head><body>
<header class="slide-header"><span class="overline">&nbsp;</span><h1 class="page-title">增长判断</h1></header>
<div class="media-shell"><img src="../images/scene.png" alt="scene"></div>
<footer class="slide-footer"><span class="footer-section">GROWTH</span><span class="footer-page">3 / 12</span></footer>
</body></html>
""".strip(),
                encoding="utf-8",
            )
        else:
            slide_html.write_text(
                """
<html><head><style>
body { margin: 0; width: 1280px; height: 720px; overflow: hidden; background: #0f172a; }
.slide-header { position: absolute; top: 20px; left: 40px; right: 40px; height: 50px; display: flex; gap: 16px; align-items: baseline; }
.media-shell { width: 360px; height: 240px; overflow: hidden; border-radius: 24px; }
.media-shell img { width: 100%; height: 100%; object-fit: cover; display: block; }
.slide-footer { position: absolute; left: 40px; right: 40px; bottom: 12px; height: 20px; }
</style></head><body>
<header class="slide-header"><span class="overline">PART 03</span><h1 class="page-title">增长判断</h1></header>
<div class="media-shell"><img src="../images/scene.png" alt="scene"></div>
<footer class="slide-footer"><span class="footer-section">GROWTH</span><span class="footer-page">3 / 12</span></footer>
</body></html>
""".strip(),
                encoding="utf-8",
            )

        planning_path = base / "planning" / "planning3.json"
        planning_path.write_text(
            json.dumps(
                {
                    "page": {
                        "page_type": "content",
                        "cards": [
                            {
                                "card_id": "s03-image",
                                "image": {
                                    "mode": "provided",
                                    "needed": True,
                                    "placement": "right-half",
                                    "source_hint": str(image_asset.resolve()),
                                },
                            }
                        ],
                    }
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        slide_png = base / "png" / "slide-3.png"
        canvas = Image.new("RGB", (1280, 720), "#0f172a")
        draw = ImageDraw.Draw(canvas)
        draw.rounded_rectangle((56, 56, 1224, 120), radius=22, fill="#334155")
        draw.rounded_rectangle((80, 108, 710, 600), radius=24, fill="#475569")
        draw.rounded_rectangle((760, 126, 1160, 600), radius=24, fill="#1f2937")
        draw.rounded_rectangle((92, 410, 678, 584), radius=20, fill="#64748b")
        draw.rounded_rectangle((780, 430, 1138, 560), radius=20, fill="#334155")
        if variant == "warn":
            for y in range(20, 92):
                for x in range(220, 598):
                    tone = 64 + ((x * 5 + y * 3) % 112)
                    canvas.putpixel((x, y), (tone, min(255, tone + 12), min(255, tone + 28)))
        if variant != "header-broken":
            draw.text((96, 26), "PART 03 -- GROWTH", fill="#60a5fa", font=font)
            draw.text((246, 40), "增长判断", fill="#f8fafc", font=font)
        else:
            draw.text((384, 96), "header drift", fill="#e2e8f0", font=font)
        if variant == "warn":
            draw.text((240, 104), "标题串位告警", fill="#e2e8f0", font=font)
        draw.text((110, 146), "核心证据", fill="#f8fafc", font=font)
        draw.text((110, 182), "主结论与支撑正文保持在安全区内。", fill="#dbe4f0", font=font)
        draw.text((110, 210), "图文并茂页的视觉锚点应清晰稳定。", fill="#dbe4f0", font=font)
        draw.text((94, 76), "Q3 视觉 QA synthetic slide", fill="#cbd5e1", font=font)
        draw.text((920, 78), "safe export", fill="#94a3b8", font=font)

        for i in range(10):
            draw.text((110, 250 + i * 24), f"body line {i + 1}: synthetic density check", fill="#eff6ff", font=font)
        for i in range(4):
            draw.text((110, 436 + i * 22), f"support block {i + 1}", fill="#0f172a", font=font)
        for i in range(4):
            draw.text((804, 446 + i * 22), f"metric note {i + 1}", fill="#e2e8f0", font=font)

        for y in range(150, 360):
            for x in range(780, 1100):
                base_color = 72 + ((x + y) % 36)
                canvas.putpixel((x, y), (base_color, base_color + 12, base_color + 28))

        for y in range(132, 590):
            for x in range(96, 694):
                if (x + y) % 19 == 0:
                    tone = 116 + ((x * 3 + y) % 52)
                    canvas.putpixel((x, y), (tone, min(255, tone + 8), min(255, tone + 18)))

        for y in range(438, 556):
            for x in range(790, 1128):
                if (x * 2 + y) % 23 == 0:
                    tone = 82 + ((x + y * 2) % 70)
                    canvas.putpixel((x, y), (tone, min(255, tone + 20), min(255, tone + 34)))

        for y in range(390, 418):
            draw.line((790, y, 1090, y), fill="#38bdf8", width=1)

        for x in range(118, 674, 18):
            draw.line((x, 592, x + 8, 610), fill="#93c5fd", width=1)
        for x in range(790, 1120, 16):
            draw.line((x, 578, x + 10, 604), fill="#60a5fa", width=1)

        if variant != "footer-broken":
            draw.text((48, 688), "GROWTH", fill="#94a3b8", font=font)
            draw.text((1184, 688), "3 / 12", fill="#94a3b8", font=font, anchor="ra")
        else:
            draw.text((584, 688), "footer lost", fill="#cbd5e1", font=font)

        if variant == "warn":
            draw.text((818, 222), "LOW CONTRAST ON IMAGE", fill=(96, 102, 112), font=font)
            draw.text((818, 238), "SECOND LINE", fill=(99, 104, 112), font=font)
            draw.text((770, 616), "caption drift line 1", fill="#dde7f5", font=font)
            draw.text((770, 634), "source: somewhere", fill="#dde7f5", font=font)
            draw.text((780, 688), "source: camera roll", fill="#dde7f5", font=font)
        else:
            draw.text((808, 398), "场景配图", fill="#f8fafc", font=font)

        canvas.save(slide_png)
        return slide_png, planning_path, slide_html

    good_png, good_planning, good_html = create_variant(good_root, "good")
    warn_png, warn_planning, warn_html = create_variant(warn_root, "warn")
    footer_broken_png, footer_broken_planning, footer_broken_html = create_variant(footer_broken_root, "footer-broken")
    header_broken_png, header_broken_planning, header_broken_html = create_variant(header_broken_root, "header-broken")
    return {
        "good_png": good_png,
        "good_planning": good_planning,
        "good_html": good_html,
        "warn_png": warn_png,
        "warn_planning": warn_planning,
        "warn_html": warn_html,
        "footer_broken_png": footer_broken_png,
        "footer_broken_planning": footer_broken_planning,
        "footer_broken_html": footer_broken_html,
        "header_broken_png": header_broken_png,
        "header_broken_planning": header_broken_planning,
        "header_broken_html": header_broken_html,
    }


def write_png_stub(path: Path) -> None:
    png_data = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+nXK0AAAAASUVORK5CYII="
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(png_data)


def slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            value = shape.text_frame.text.strip()
            if value:
                texts.append(value)
        try:
            table = shape.table
        except Exception:
            table = None
        if table is not None:
            for row in table.rows:
                for cell in row.cells:
                    value = cell.text.strip()
                    if value:
                        texts.append(value)
        if hasattr(shape, "shapes"):
            for child_text in slide_texts(shape):
                texts.append(child_text)
    return texts


def slide_shape_names(slide) -> list[str]:
    names: list[str] = []
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if name:
            names.append(name)
        if hasattr(shape, "shapes"):
            names.extend(slide_shape_names(shape))
    return names


def slide_notes_text(slide) -> str:
    return extract_slide_speaker_notes(slide)


def run_template_update_export_smoke(tmp_dir: Path, result: SmokeResult, py: str) -> None:
    try:
        import collections.abc  # noqa: F401
        from pptx import Presentation
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.util import Inches, Pt
        from milestone_check import Checker
    except Exception as exc:
        result.error(f"template-update smoke bootstrap failed: {exc}")
        return

    output_dir = tmp_dir / "template-update-gate"
    planning_dir = output_dir / "planning"
    slides_dir = output_dir / "slides"
    svg_dir = output_dir / "svg"
    png_dir = output_dir / "png"
    runtime_dir = output_dir / "runtime"
    for path in (planning_dir, slides_dir, svg_dir, png_dir, runtime_dir):
        path.mkdir(parents=True, exist_ok=True)

    def write_json(path: Path, payload: dict[str, object]) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    speech_script_path = output_dir / "speech-script.json"
    speech_script_md_path = output_dir / "speech-script.md"
    write_json(
        speech_script_path,
        {
            "deck_title": "Template Update Smoke",
            "language": "zh-CN",
            "summary": "验证 speech-script.json 可以派生 markdown，并写入 PNG/SVG 导出的 speaker notes。",
            "pages": [
                {
                    "page": 1,
                    "slide_title": "增长总览",
                    "estimated_seconds": 45,
                    "transition_to_next": "接着切到成本构成，解释增长背后的投入结构。",
                    "speaker_notes": "先用月活跃用户数 2.4M 定住全局判断，再强调 +12.3% 说明增长不是偶发脉冲，而是可持续改善的结果。",
                },
                {
                    "page": 2,
                    "slide_title": "成本构成",
                    "estimated_seconds": 35,
                    "speaker_notes": "这一页不要只读图例，要直接讲出三类成本的占比逻辑，并把人力、云资源、市场之间的资源配置关系说清楚。",
                },
            ],
        },
    )
    speech_markdown = run_cmd(
        "speech-script-format",
        [
            py,
            str(SCRIPTS_DIR / "speech_script_formatter.py"),
            str(speech_script_path),
            "-o",
            str(speech_script_md_path),
        ],
        result,
    )
    if speech_markdown.returncode != 0:
        return
    if "# 演讲稿：Template Update Smoke" not in speech_script_md_path.read_text(encoding="utf-8"):
        result.error("speech-script formatter: markdown title missing")
    write_text(runtime_dir / "prompt-speech-orchestrator.md", "speech orchestrator smoke")

    write_text(output_dir / "preview.html", "<html><body>template update smoke</body></html>")
    for index in range(1, 3):
        write_text(
            planning_dir / f"planning{index}.json",
            json.dumps({"page": {"slide_number": index, "page_type": "content", "cards": []}}, ensure_ascii=False, indent=2),
        )
        write_png_stub(png_dir / f"slide-{index}.png")

    def write_initial_assets() -> None:
        write_text(
            slides_dir / "slide-1.html",
            """
<html><body>
  <div data-card-id="s01-title"><h1>增长总览</h1></div>
  <div data-card-id="s01-kpi" class="kpi-card metric-card">
    <p>2.4M</p>
    <p>+12.3%</p>
    <p>月活跃用户数</p>
  </div>
</body></html>
""".strip(),
        )
        write_text(
            svg_dir / "slide-1.svg",
            """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <text x="100" y="170" font-size="30" fill="#0f172a">增长总览</text>
  <text x="120" y="315" font-size="52" font-weight="700" fill="#2563eb">2.4M</text>
  <polygon points="245,292 262,320 228,320" fill="#22c55e" />
  <text x="272" y="315" font-size="18" font-weight="600" fill="#22c55e">+12.3%</text>
  <text x="120" y="350" font-size="14" fill="#64748b">月活跃用户数</text>
</svg>
""".strip(),
        )
        write_json(
            svg_dir / "slide-1.semantic.json",
            {
                "summary": {"blocks": 2, "tables": 0, "charts": 1},
                "blocks": [
                    {"block_id": "s01-title", "bbox": {"x": 80, "y": 120, "width": 420, "height": 80}, "contains_table": False, "contains_chart_like": False},
                    {"block_id": "s01-kpi", "bbox": {"x": 90, "y": 245, "width": 340, "height": 125}, "contains_table": False, "contains_chart_like": True},
                ],
                "tables": [],
                "charts": [
                    {"chart_id": "chart-s01-kpi", "block_id": "s01-kpi", "bbox": {"x": 90, "y": 245, "width": 340, "height": 125}, "chart_type_hint": "kpi"}
                ],
            },
        )
        write_text(
            slides_dir / "slide-2.html",
            """
<html><body>
  <div data-card-id="s03-stacked" class="stacked-bar-card chart-card">
    <h2>成本构成</h2>
    <p>人力 / 云资源 / 市场</p>
  </div>
</body></html>
""".strip(),
        )
        write_text(
            svg_dir / "slide-2.svg",
            """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <text x="120" y="210" font-size="28" fill="#0f172a">成本构成</text>
  <rect x="120" y="260" width="150" height="24" fill="#2563eb" />
  <rect x="270" y="260" width="95" height="24" fill="#16a34a" />
  <rect x="365" y="260" width="55" height="24" fill="#f59e0b" />
  <rect x="120" y="330" width="12" height="12" fill="#2563eb" />
  <text x="140" y="340" font-size="14" fill="#334155">人力</text>
  <rect x="220" y="330" width="12" height="12" fill="#16a34a" />
  <text x="240" y="340" font-size="14" fill="#334155">云资源</text>
  <rect x="340" y="330" width="12" height="12" fill="#f59e0b" />
  <text x="360" y="340" font-size="14" fill="#334155">市场</text>
</svg>
""".strip(),
        )
        write_json(
            svg_dir / "slide-2.semantic.json",
            {
                "summary": {"blocks": 1, "tables": 0, "charts": 1},
                "blocks": [
                    {"block_id": "s03-stacked", "bbox": {"x": 100, "y": 220, "width": 360, "height": 150}, "contains_table": False, "contains_chart_like": True},
                ],
                "tables": [],
                "charts": [
                    {"chart_id": "chart-s03-stacked", "block_id": "s03-stacked", "bbox": {"x": 100, "y": 220, "width": 360, "height": 150}, "chart_type_hint": "stacked_bar"}
                ],
            },
        )

    write_initial_assets()

    write_json(
        output_dir / "svg-export-report.json",
        {
            "generated_at": "2025-01-01T00:00:00Z",
            "html_input": "slides",
            "svg_output_dir": "svg",
            "summary": {
                "total_pages": 2,
                "editable_pages": 2,
                "raster_fallback_pages": 0,
                "pathified_pages": 0,
                "failed_pages": 0,
                "warning_pages": 0,
            },
            "pages": [
                {
                    "page_number": 1,
                    "page_name": "slide-1",
                    "source_html": "slides/slide-1.html",
                    "source_svg": "svg/slide-1.svg",
                    "semantic_path": "svg/slide-1.semantic.json",
                    "method": "dom_to_svg_editable",
                    "editable": True,
                    "success": True,
                    "text_count": 3,
                    "image_count": 0,
                    "path_count": 1,
                },
                {
                    "page_number": 2,
                    "page_name": "slide-2",
                    "source_html": "slides/slide-2.html",
                    "source_svg": "svg/slide-2.svg",
                    "semantic_path": "svg/slide-2.semantic.json",
                    "method": "dom_to_svg_editable",
                    "editable": True,
                    "success": True,
                    "text_count": 4,
                    "image_count": 0,
                    "path_count": 0,
                },
            ],
        },
    )

    template_path = output_dir / "template-source.pptx"
    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank)
    slide2 = prs.slides.add_slide(blank)
    slide3 = prs.slides.add_slide(blank)

    def add_textbox(slide, text: str, x: float, y: float, w: float, h: float, name: str | None = None) -> None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        shape.text_frame.text = text
        shape.text_frame.paragraphs[0].font.size = Pt(16)
        if name:
            shape.name = name

    def add_slot(slide, name: str, x: float, y: float, w: float, h: float) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.name = name
        shape.text_frame.text = name

    add_textbox(slide1, "KEEP HEADER 1", 0.6, 0.3, 2.8, 0.4)
    add_textbox(slide1, "PAGE 1", 11.2, 7.0, 1.0, 0.3)
    add_slot(slide1, "BlockSlot:s01-title:slot", 0.8, 1.3, 4.2, 0.8)
    add_slot(slide1, "BlockSlot:s01-kpi:slot", 0.8, 2.5, 3.8, 1.4)
    add_textbox(slide2, "UNCHANGED TEMPLATE SLIDE", 0.8, 0.9, 4.0, 0.5)
    add_textbox(slide3, "KEEP HEADER 3", 0.6, 0.3, 2.8, 0.4)
    add_textbox(slide3, "PAGE 3", 11.2, 7.0, 1.0, 0.3)
    add_slot(slide3, "BlockSlot:s03-stacked:slot", 0.8, 2.2, 4.4, 1.7)
    prs.save(template_path)

    png_convert = run_cmd(
        "template-update-png-pptx",
        [
            py,
            str(SCRIPTS_DIR / "png2pptx.py"),
            str(png_dir),
            "-o",
            str(output_dir / "presentation-png.pptx"),
            "--speech-script",
            str(speech_script_path),
        ],
        result,
    )
    if png_convert.returncode != 0:
        return

    initial_convert = run_cmd(
        "template-update-initial",
        [
            py,
            str(SCRIPTS_DIR / "svg2pptx.py"),
            str(svg_dir),
            "-o",
            str(output_dir / "presentation-svg.pptx"),
            "--html-dir",
            str(slides_dir),
            "--export-report",
            str(output_dir / "svg-export-report.json"),
            "--report-path",
            str(output_dir / "presentation-svg.report.json"),
            "--speech-script",
            str(speech_script_path),
            "--template-pptx",
            str(template_path),
            "--target-slides",
            "1,3",
        ],
        result,
    )
    if initial_convert.returncode != 0:
        return

    write_json(
        output_dir / "delivery-manifest.json",
        {
            "run_id": "template-update-smoke",
            "generated_at": "2025-01-01T00:00:00Z",
            "artifacts": {
                "preview_html": "preview.html",
                "speech_script_json": "speech-script.json",
                "speech_script_md": "speech-script.md",
                "presentation_png_pptx": "presentation-png.pptx",
                "presentation_svg_pptx": "presentation-svg.pptx",
            },
            "summary": {"total_pages": 2},
        },
    )

    png_converted = Presentation(str(output_dir / "presentation-png.pptx"))
    converted = Presentation(str(output_dir / "presentation-svg.pptx"))
    png_slide1_notes = slide_notes_text(png_converted.slides[0])
    png_slide2_notes = slide_notes_text(png_converted.slides[1])
    slide1_text = "\n".join(slide_texts(converted.slides[0]))
    slide2_text = "\n".join(slide_texts(converted.slides[1]))
    slide3_text = "\n".join(slide_texts(converted.slides[2]))
    slide1_names = slide_shape_names(converted.slides[0])
    slide3_names = slide_shape_names(converted.slides[2])
    slide1_notes = slide_notes_text(converted.slides[0])
    slide2_notes = slide_notes_text(converted.slides[1])
    slide3_notes = slide_notes_text(converted.slides[2])
    if "KEEP HEADER 1" not in slide1_text or "PAGE 1" not in slide1_text:
        result.error("template-update initial: slide 1 decorations were not preserved")
    if "UNCHANGED TEMPLATE SLIDE" not in slide2_text:
        result.error("template-update initial: untouched slide content was modified")
    if "2.4M" not in slide1_text:
        result.error("template-update initial: KPI text missing on slide 1")
    if "成本构成" not in slide3_text:
        result.error("template-update initial: stacked bar title missing on slide 3")
    if any(name.startswith("BlockSlot:") for name in slide1_names + slide3_names):
        result.error("template-update initial: block slot placeholders were not removed")
    if not any(name.startswith("NativeChart:") for name in slide1_names):
        result.error("template-update initial: slide 1 native KPI promotion missing")
    if not any(name.startswith("NativeChart:") for name in slide3_names):
        result.error("template-update initial: slide 3 native stacked bar promotion missing")
    if "2.4M" not in png_slide1_notes or "人力、云资源、市场" not in png_slide2_notes:
        result.error("template-update initial: PNG deck speaker notes missing expected content")
    if "2.4M" not in slide1_notes or "人力、云资源、市场" not in slide3_notes:
        result.error("template-update initial: SVG deck speaker notes missing expected content on target slides")
    if slide2_notes:
        result.error("template-update initial: untouched template slide unexpectedly received speaker notes")

    try:
        Checker(ROOT_DIR, output_dir, "5", quiet=True).check_step5()
        result.note("template-update-step5: ok")
    except Exception as exc:
        result.error(f"template-update-step5 failed: {exc}")

    rerun_dir = tmp_dir / "template-update-rerun"
    rerun_slides = rerun_dir / "slides"
    rerun_svg = rerun_dir / "svg"
    rerun_slides.mkdir(parents=True, exist_ok=True)
    rerun_svg.mkdir(parents=True, exist_ok=True)
    write_text(
        rerun_slides / "slide-1.html",
        """
<html><body>
  <div data-card-id="s01-title"><h1>增长二次刷新</h1></div>
  <div data-card-id="s01-kpi" class="kpi-card metric-card">
    <p>3.1M</p>
    <p>+18.0%</p>
    <p>月活跃用户数</p>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        rerun_svg / "slide-1.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <text x="100" y="170" font-size="30" fill="#0f172a">增长二次刷新</text>
  <text x="120" y="315" font-size="52" font-weight="700" fill="#2563eb">3.1M</text>
  <polygon points="245,292 262,320 228,320" fill="#22c55e" />
  <text x="272" y="315" font-size="18" font-weight="600" fill="#22c55e">+18.0%</text>
  <text x="120" y="350" font-size="14" fill="#64748b">月活跃用户数</text>
</svg>
""".strip(),
    )
    write_json(
        rerun_svg / "slide-1.semantic.json",
        {
            "summary": {"blocks": 2, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s01-title", "bbox": {"x": 80, "y": 120, "width": 420, "height": 80}, "contains_table": False, "contains_chart_like": False},
                {"block_id": "s01-kpi", "bbox": {"x": 90, "y": 245, "width": 340, "height": 125}, "contains_table": False, "contains_chart_like": True},
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s01-kpi", "block_id": "s01-kpi", "bbox": {"x": 90, "y": 245, "width": 340, "height": 125}, "chart_type_hint": "kpi"}
            ],
        },
    )
    write_json(
        rerun_dir / "svg-export-report.json",
        {
            "generated_at": "2025-01-01T00:00:00Z",
            "html_input": "slides",
            "svg_output_dir": "svg",
            "summary": {
                "total_pages": 1,
                "editable_pages": 1,
                "raster_fallback_pages": 0,
                "pathified_pages": 0,
                "failed_pages": 0,
                "warning_pages": 0,
            },
            "pages": [
                {
                    "page_number": 1,
                    "page_name": "slide-1",
                    "source_html": "slides/slide-1.html",
                    "source_svg": "svg/slide-1.svg",
                    "semantic_path": "svg/slide-1.semantic.json",
                    "method": "dom_to_svg_editable",
                    "editable": True,
                    "success": True,
                    "text_count": 3,
                    "image_count": 0,
                    "path_count": 1,
                }
            ],
        },
    )
    write_json(
        rerun_dir / "speech-script.json",
        {
            "deck_title": "Template Update Smoke Rerun",
            "language": "zh-CN",
            "pages": [
                {
                    "page": 1,
                    "slide_title": "增长二次刷新",
                    "estimated_seconds": 30,
                    "speaker_notes": "二次刷新时要明确指出核心指标已经从 2.4M 更新为 3.1M，并强调这是重跑后写入的新备注。",
                }
            ],
        },
    )
    rerun_convert = run_cmd(
        "template-update-rerun",
        [
            py,
            str(SCRIPTS_DIR / "svg2pptx.py"),
            str(rerun_svg),
            "-o",
            str(output_dir / "presentation-svg-rerun.pptx"),
            "--html-dir",
            str(rerun_slides),
            "--export-report",
            str(rerun_dir / "svg-export-report.json"),
            "--report-path",
            str(rerun_dir / "presentation-svg-rerun.report.json"),
            "--speech-script",
            str(rerun_dir / "speech-script.json"),
            "--template-pptx",
            str(output_dir / "presentation-svg.pptx"),
            "--target-slides",
            "1",
        ],
        result,
    )
    if rerun_convert.returncode != 0:
        return
    rerun_presentation = Presentation(str(output_dir / "presentation-svg-rerun.pptx"))
    rerun_slide1_text = "\n".join(slide_texts(rerun_presentation.slides[0]))
    rerun_slide2_text = "\n".join(slide_texts(rerun_presentation.slides[1]))
    rerun_slide3_names = slide_shape_names(rerun_presentation.slides[2])
    rerun_slide1_notes = slide_notes_text(rerun_presentation.slides[0])
    rerun_slide3_notes = slide_notes_text(rerun_presentation.slides[2])
    if "KEEP HEADER 1" not in rerun_slide1_text or "PAGE 1" not in rerun_slide1_text:
        result.error("template-update rerun: slide 1 decorations were not preserved")
    if "3.1M" not in rerun_slide1_text or "2.4M" in rerun_slide1_text:
        result.error("template-update rerun: old managed block content was not replaced cleanly")
    if "UNCHANGED TEMPLATE SLIDE" not in rerun_slide2_text:
        result.error("template-update rerun: untouched slide changed after rerun")
    if not any(name.startswith("NativeChart:") for name in rerun_slide3_names):
        result.error("template-update rerun: non-target slide structured chart was lost")
    if "3.1M" not in rerun_slide1_notes or "重跑后写入的新备注" not in rerun_slide1_notes:
        result.error("template-update rerun: updated slide speaker notes were not replaced cleanly")
    if slide3_notes != rerun_slide3_notes:
        result.error("template-update rerun: non-target slide speaker notes changed unexpectedly")


def run_native_chart_family_smoke(tmp_dir: Path, result: SmokeResult, py: str) -> None:
    try:
        import collections.abc  # noqa: F401
        from pptx import Presentation
    except Exception as exc:
        result.error(f"native-chart family smoke bootstrap failed: {exc}")
        return

    output_dir = tmp_dir / "native-chart-families"
    slides_dir = output_dir / "slides"
    svg_dir = output_dir / "svg"
    slides_dir.mkdir(parents=True, exist_ok=True)
    svg_dir.mkdir(parents=True, exist_ok=True)

    def write_json(path: Path, payload: dict[str, object]) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    pages: list[dict[str, object]] = []

    write_text(
        slides_dir / "slide-1.html",
        """
<html><body>
  <div data-card-id="s01-sparkline" class="sparkline chart-card">
    <svg width="280" height="120" viewBox="0 0 280 120" aria-label="用户趋势 sparkline"></svg>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        svg_dir / "slide-1.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <path d="M120 360 L170 330 L220 338 L270 286 L320 246 L370 210 L370 390 L120 390 Z" fill="#2563eb" opacity="0.12" />
  <polyline points="120,360 170,330 220,338 270,286 320,246 370,210" fill="none" stroke="#2563eb" stroke-width="3" stroke-linecap="round" />
  <circle cx="370" cy="210" r="5" fill="#2563eb" />
</svg>
""".strip(),
    )
    write_json(
        svg_dir / "slide-1.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s01-sparkline", "bbox": {"x": 110, "y": 205, "width": 285, "height": 190}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s01-sparkline", "block_id": "s01-sparkline", "bbox": {"x": 110, "y": 205, "width": 285, "height": 190}, "chart_type_hint": "sparkline"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 1,
            "page_name": "slide-1",
            "source_html": "slides/slide-1.html",
            "source_svg": "svg/slide-1.svg",
            "semantic_path": "svg/slide-1.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 0,
            "image_count": 0,
            "path_count": 1,
        }
    )

    write_text(
        slides_dir / "slide-2.html",
        """
<html><body>
  <div data-card-id="s02-rating" class="rating chart-card">
    <div aria-label="体验评分 rating"></div>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        svg_dir / "slide-2.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <circle cx="160" cy="280" r="10" fill="#2563eb" />
  <circle cx="205" cy="280" r="10" fill="#2563eb" />
  <circle cx="250" cy="280" r="10" fill="#2563eb" />
  <circle cx="295" cy="280" r="10" fill="#2563eb" />
  <circle cx="340" cy="280" r="10" fill="none" stroke="#2563eb" stroke-width="2" />
</svg>
""".strip(),
    )
    write_json(
        svg_dir / "slide-2.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s02-rating", "bbox": {"x": 140, "y": 250, "width": 220, "height": 70}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s02-rating", "block_id": "s02-rating", "bbox": {"x": 140, "y": 250, "width": 220, "height": 70}, "chart_type_hint": "rating"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 2,
            "page_name": "slide-2",
            "source_html": "slides/slide-2.html",
            "source_svg": "svg/slide-2.svg",
            "semantic_path": "svg/slide-2.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 0,
            "image_count": 0,
            "path_count": 0,
        }
    )

    write_text(
        slides_dir / "slide-3.html",
        """
<html><body>
  <div data-card-id="s03-timeline" class="timeline chart-card">
    <div aria-label="阶段里程碑 timeline"></div>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        svg_dir / "slide-3.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <line x1="140" y1="280" x2="500" y2="280" stroke="#94a3b8" stroke-width="3" />
  <circle cx="160" cy="280" r="12" fill="#2563eb" />
  <circle cx="320" cy="280" r="12" fill="#38bdf8" />
  <circle cx="480" cy="280" r="14" fill="#0f172a" stroke="#2563eb" stroke-width="3" />
  <text x="130" y="330" font-size="16" fill="#0f172a">调研</text>
  <text x="290" y="330" font-size="16" fill="#0f172a">验证</text>
  <text x="450" y="330" font-size="16" fill="#0f172a">上线</text>
</svg>
""".strip(),
    )
    write_json(
        svg_dir / "slide-3.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s03-timeline", "bbox": {"x": 120, "y": 240, "width": 420, "height": 120}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s03-timeline", "block_id": "s03-timeline", "bbox": {"x": 120, "y": 240, "width": 420, "height": 120}, "chart_type_hint": "timeline"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 3,
            "page_name": "slide-3",
            "source_html": "slides/slide-3.html",
            "source_svg": "svg/slide-3.svg",
            "semantic_path": "svg/slide-3.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 3,
            "image_count": 0,
            "path_count": 0,
        }
    )

    write_text(
        slides_dir / "slide-4.html",
        """
<html><body>
  <div data-card-id="s04-funnel" class="funnel chart-card">
    <div aria-label="转化漏斗 funnel"></div>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        svg_dir / "slide-4.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <polygon points="180,180 460,180 430,240 210,240" fill="#1d4ed8" />
  <polygon points="210,252 430,252 398,312 242,312" fill="#2563eb" />
  <polygon points="242,324 398,324 366,384 274,384" fill="#60a5fa" />
  <text x="270" y="217" font-size="15" fill="#ffffff">触达 10k</text>
  <text x="284" y="288" font-size="15" fill="#ffffff">试用 4.8k</text>
  <text x="298" y="359" font-size="15" fill="#0f172a">成交 1.6k</text>
</svg>
""".strip(),
    )
    write_json(
        svg_dir / "slide-4.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s04-funnel", "bbox": {"x": 170, "y": 165, "width": 300, "height": 230}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s04-funnel", "block_id": "s04-funnel", "bbox": {"x": 170, "y": 165, "width": 300, "height": 230}, "chart_type_hint": "funnel"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 4,
            "page_name": "slide-4",
            "source_html": "slides/slide-4.html",
            "source_svg": "svg/slide-4.svg",
            "semantic_path": "svg/slide-4.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 3,
            "image_count": 0,
            "path_count": 0,
        }
    )

    write_text(
        slides_dir / "slide-5.html",
        """
<html><body>
  <div data-card-id="s05-radar" class="radar chart-card">
    <svg width="260" height="260" viewBox="0 0 260 260" aria-label="能力维度 radar"></svg>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        svg_dir / "slide-5.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <polygon points="320,180 390,235 360,320 280,320 250,235" fill="none" stroke="#cbd5e1" stroke-width="2" />
  <polygon points="320,215 355,242 340,285 300,285 285,242" fill="none" stroke="#cbd5e1" stroke-width="1" />
  <line x1="320" y1="250" x2="320" y2="180" stroke="#cbd5e1" stroke-width="1" />
  <line x1="320" y1="250" x2="390" y2="235" stroke="#cbd5e1" stroke-width="1" />
  <line x1="320" y1="250" x2="360" y2="320" stroke="#cbd5e1" stroke-width="1" />
  <line x1="320" y1="250" x2="280" y2="320" stroke="#cbd5e1" stroke-width="1" />
  <line x1="320" y1="250" x2="250" y2="235" stroke="#cbd5e1" stroke-width="1" />
  <polygon points="320,202 368,238 344,295 296,286 276,238" fill="#2563eb" opacity="0.18" stroke="#2563eb" stroke-width="3" />
  <circle cx="320" cy="202" r="4" fill="#2563eb" />
  <circle cx="368" cy="238" r="4" fill="#2563eb" />
  <circle cx="344" cy="295" r="4" fill="#2563eb" />
  <circle cx="296" cy="286" r="4" fill="#2563eb" />
  <circle cx="276" cy="238" r="4" fill="#2563eb" />
  <text x="300" y="168" font-size="13" fill="#0f172a">战略</text>
  <text x="398" y="238" font-size="13" fill="#0f172a">产品</text>
  <text x="350" y="338" font-size="13" fill="#0f172a">运营</text>
  <text x="260" y="338" font-size="13" fill="#0f172a">交付</text>
  <text x="206" y="238" font-size="13" fill="#0f172a">品牌</text>
</svg>
""".strip(),
    )
    write_json(
        svg_dir / "slide-5.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s05-radar", "bbox": {"x": 210, "y": 155, "width": 210, "height": 200}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s05-radar", "block_id": "s05-radar", "bbox": {"x": 210, "y": 155, "width": 210, "height": 200}, "chart_type_hint": "radar"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 5,
            "page_name": "slide-5",
            "source_html": "slides/slide-5.html",
            "source_svg": "svg/slide-5.svg",
            "semantic_path": "svg/slide-5.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 5,
            "image_count": 0,
            "path_count": 0,
        }
    )

    write_json(
        output_dir / "svg-export-report.json",
        {
            "generated_at": "2025-01-01T00:00:00Z",
            "html_input": "slides",
            "svg_output_dir": "svg",
            "summary": {
                "total_pages": 7,
                "editable_pages": 7,
                "raster_fallback_pages": 0,
                "pathified_pages": 0,
                "failed_pages": 0,
                "warning_pages": 0,
            },
            "pages": pages,
        },
    )

    write_text(
        slides_dir / "slide-6.html",
        """
<html><body>
  <div data-card-id="s06-waffle" class="waffle chart-card">
    <div aria-label="占比点阵 waffle"></div>
  </div>
</body></html>
""".strip(),
    )
    waffle_cells = []
    start_x = 160
    start_y = 220
    cell = 14
    gap = 6
    for row in range(5):
        for col in range(10):
            x = start_x + col * (cell + gap)
            y = start_y + row * (cell + gap)
            fill = "#2563eb" if (row * 10 + col) < 34 else "#cbd5e1"
            waffle_cells.append(f'<rect x="{x}" y="{y}" width="{cell}" height="{cell}" rx="2" fill="{fill}" />')
    write_text(
        svg_dir / "slide-6.svg",
        (
            '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">'
            + ''.join(waffle_cells)
            + '<text x="160" y="360" font-size="18" fill="#0f172a">34%</text>'
            + '<text x="220" y="360" font-size="16" fill="#475569">完成率</text>'
            + '</svg>'
        ),
    )
    write_json(
        svg_dir / "slide-6.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s06-waffle", "bbox": {"x": 150, "y": 210, "width": 240, "height": 170}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s06-waffle", "block_id": "s06-waffle", "bbox": {"x": 150, "y": 210, "width": 240, "height": 170}, "chart_type_hint": "waffle"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 6,
            "page_name": "slide-6",
            "source_html": "slides/slide-6.html",
            "source_svg": "svg/slide-6.svg",
            "semantic_path": "svg/slide-6.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 2,
            "image_count": 0,
            "path_count": 0,
        }
    )

    write_text(
        slides_dir / "slide-7.html",
        """
<html><body>
  <div data-card-id="s07-treemap" class="treemap chart-card">
    <div aria-label="层级占比 treemap"></div>
  </div>
</body></html>
""".strip(),
    )
    write_text(
        svg_dir / "slide-7.svg",
        """
<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">
  <rect x="160" y="190" width="180" height="160" fill="#1d4ed8" rx="6" />
  <rect x="346" y="190" width="122" height="96" fill="#2563eb" rx="6" />
  <rect x="346" y="292" width="122" height="58" fill="#60a5fa" rx="6" />
  <rect x="474" y="190" width="74" height="160" fill="#93c5fd" rx="6" />
  <text x="180" y="236" font-size="18" fill="#ffffff">云资源</text>
  <text x="180" y="264" font-size="16" fill="#ffffff">46%</text>
  <text x="360" y="228" font-size="16" fill="#ffffff">人力 24%</text>
  <text x="360" y="324" font-size="14" fill="#0f172a">市场 12%</text>
  <text x="486" y="236" font-size="13" fill="#0f172a">其他 18%</text>
</svg>
""".strip(),
    )
    write_json(
        svg_dir / "slide-7.semantic.json",
        {
            "summary": {"blocks": 1, "tables": 0, "charts": 1},
            "blocks": [
                {"block_id": "s07-treemap", "bbox": {"x": 150, "y": 180, "width": 410, "height": 180}, "contains_table": False, "contains_chart_like": True}
            ],
            "tables": [],
            "charts": [
                {"chart_id": "chart-s07-treemap", "block_id": "s07-treemap", "bbox": {"x": 150, "y": 180, "width": 410, "height": 180}, "chart_type_hint": "treemap"}
            ],
        },
    )
    pages.append(
        {
            "page_number": 7,
            "page_name": "slide-7",
            "source_html": "slides/slide-7.html",
            "source_svg": "svg/slide-7.svg",
            "semantic_path": "svg/slide-7.semantic.json",
            "method": "dom_to_svg_editable",
            "editable": True,
            "success": True,
            "text_count": 5,
            "image_count": 0,
            "path_count": 0,
        }
    )

    export_proc = run_cmd(
        "native-chart-families-export",
        [
            py,
            str(SCRIPTS_DIR / "svg2pptx.py"),
            str(svg_dir),
            "-o",
            str(output_dir / "presentation-svg.pptx"),
            "--html-dir",
            str(slides_dir),
            "--export-report",
            str(output_dir / "svg-export-report.json"),
            "--report-path",
            str(output_dir / "presentation-svg.report.json"),
        ],
        result,
    )
    if export_proc.returncode != 0:
        return

    inspect_proc = run_cmd(
        "native-chart-families-inspect",
        [
            py,
            str(SCRIPTS_DIR / "inspect_pptx.py"),
            str(output_dir / "presentation-svg.pptx"),
            "--source-report",
            str(output_dir / "presentation-svg.report.json"),
            "-o",
            str(output_dir / "presentation-svg.inspect.json"),
        ],
        result,
    )
    if inspect_proc.returncode != 0:
        return

    presentation = Presentation(str(output_dir / "presentation-svg.pptx"))
    expected_native = {
        0: "sparkline",
        1: "rating",
        2: "timeline",
        3: "funnel",
        4: "radar",
        5: "waffle",
        6: "treemap",
    }
    for slide_index, chart_type in expected_native.items():
        shape_names = slide_shape_names(presentation.slides[slide_index])
        if not any(name.startswith("NativeChart:") and f":{chart_type}" in name for name in shape_names):
            result.error(f"native-chart families: {chart_type} native promotion missing on slide {slide_index + 1}")
        if any(name.startswith("ChartGroup:") for name in shape_names):
            result.error(f"native-chart families: {chart_type} degraded to ChartGroup fallback on slide {slide_index + 1}")

    try:
        inspection_payload = json.loads((output_dir / "presentation-svg.inspect.json").read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        result.error(f"native-chart families: failed to read inspection report: {exc}")
        return

    summary = inspection_payload.get("summary") if isinstance(inspection_payload, dict) else None
    if not isinstance(summary, dict):
        result.error("native-chart families: inspection summary missing")
        return
    if int(summary.get("native_chart_shapes_total") or 0) < 7:
        result.error("native-chart families: inspection did not count all promoted chart families as native charts")
    if int(summary.get("chart_group_shapes_total") or 0) != 0:
        result.error("native-chart families: inspection still reports ChartGroup fallback for promoted slides")


def run_smoke() -> SmokeResult:
    result = SmokeResult()
    with tempfile.TemporaryDirectory(prefix="ppt-skill-smoke-") as tmp:
        tmp_dir = Path(tmp)
        fx = build_fixture_tree(tmp_dir)
        py = sys.executable
        visual_fx = build_visual_qa_fixtures(tmp_dir, result)

        validator = run_cmd(
            "planning-validator",
            [
                py,
                str(SCRIPTS_DIR / "planning_validator.py"),
                str(fx["planning"].parent),
                "--refs",
                str(REFERENCES_DIR),
                "--page",
                "3",
            ],
            result,
        )
        if validator.returncode == 0:
            assert_contains("planning-validator", validator.stdout, ["OK"], result)

        menu = run_cmd(
            "resource-loader-menu",
            [py, str(SCRIPTS_DIR / "resource_loader.py"), "menu", "--refs-dir", str(REFERENCES_DIR)],
            result,
        )
        if menu.returncode == 0:
            assert_contains("resource-loader-menu", menu.stdout, ["### layouts/", "#### hero-top", "### blocks/"], result)

        resolve = run_cmd(
            "resource-loader-resolve",
            [
                py,
                str(SCRIPTS_DIR / "resource_loader.py"),
                "resolve",
                "--refs-dir",
                str(REFERENCES_DIR),
                "--planning",
                str(fx["planning"]),
            ],
            result,
        )
        if resolve.returncode == 0:
            assert_contains(
                "resource-loader-resolve",
                resolve.stdout,
                [
                    "# 顶部英雄式版式",
                    "# KPI 指标卡（数字+趋势箭头+标签）",
                    "# 指标行（数字+标签+进度条 组合）",
                    "# 视觉层级与 CRAP 原则",
                    "# 构图与留白",
                    "# Director Command Runtime Rules",
                ],
                result,
            )
            assert_no_unfilled_vars("resource-loader-resolve", resolve.stdout, result)

        images = run_cmd(
            "resource-loader-images",
            [
                py,
                str(SCRIPTS_DIR / "resource_loader.py"),
                "images",
                "--images-dir",
                str(fx["images"]),
            ],
            result,
        )
        if images.returncode == 0:
            assert_contains(
                "resource-loader-images",
                images.stdout,
                [
                    "count: 3",
                    "rel=landscape.svg",
                    "dimensions=1600x900",
                    "aspect_ratio=1.778",
                    "orientation=landscape",
                    "rel=portrait.svg",
                    "orientation=portrait",
                    "rel=square.svg",
                    "orientation=square",
                ],
                result,
            )

        if visual_fx:
            visual_good = run_cmd_expect(
                "visual-qa-good",
                [
                    py,
                    str(SCRIPTS_DIR / "visual_qa.py"),
                    str(visual_fx["good_png"]),
                    "--planning",
                    str(visual_fx["good_planning"]),
                    "--html",
                    str(visual_fx["good_html"]),
                ],
                {0},
                result,
            )
            if visual_good.returncode == 0:
                assert_contains(
                    "visual-qa-good",
                    visual_good.stdout,
                    ["CAP-01", "IMG-03", "HDR-01", "HDR-02", "FTR-01", "FTR-02", "verdict: PASS"],
                    result,
                )

            visual_warn = run_cmd_expect(
                "visual-qa-warning",
                [
                    py,
                    str(SCRIPTS_DIR / "visual_qa.py"),
                    str(visual_fx["warn_png"]),
                    "--planning",
                    str(visual_fx["warn_planning"]),
                    "--html",
                    str(visual_fx["warn_html"]),
                ],
                {2},
                result,
            )
            if visual_warn.returncode == 2:
                assert_contains(
                    "visual-qa-warning",
                    visual_warn.stdout,
                    ["CAP-01", "IMG-03", "HDR-01", "HDR-02", "FTR-01", "FTR-02", "verdict: WARN"],
                    result,
                )

            visual_broken = run_cmd_expect(
                "visual-qa-footer-broken",
                [
                    py,
                    str(SCRIPTS_DIR / "visual_qa.py"),
                    str(visual_fx["footer_broken_png"]),
                    "--planning",
                    str(visual_fx["footer_broken_planning"]),
                    "--html",
                    str(visual_fx["footer_broken_html"]),
                ],
                {1},
                result,
            )
            if visual_broken.returncode == 1:
                assert_contains(
                    "visual-qa-footer-broken",
                    visual_broken.stdout,
                    ["FTR-02", "verdict: FAIL"],
                    result,
                )

            visual_header_broken = run_cmd_expect(
                "visual-qa-header-broken",
                [
                    py,
                    str(SCRIPTS_DIR / "visual_qa.py"),
                    str(visual_fx["header_broken_png"]),
                    "--planning",
                    str(visual_fx["header_broken_planning"]),
                    "--html",
                    str(visual_fx["header_broken_html"]),
                ],
                {1},
                result,
            )
            if visual_header_broken.returncode == 1:
                assert_contains(
                    "visual-qa-header-broken",
                    visual_header_broken.stdout,
                    ["HDR-01", "verdict: FAIL"],
                    result,
                )

            run_native_chart_family_smoke(tmp_dir, result, py)
        run_template_update_export_smoke(tmp_dir, result, py)

        for page_type, expected_title in PAGE_TEMPLATE_EXPECTATIONS.items():
            planning_dir = tmp_dir / f"planning-{page_type}"
            planning_path = planning_dir / "planning1.json"
            write_text(planning_path, json.dumps(build_non_content_page(page_type), ensure_ascii=False, indent=2))
            non_content_validate = run_cmd(
                f"planning-validator-{page_type}",
                [
                    py,
                    str(SCRIPTS_DIR / "planning_validator.py"),
                    str(planning_dir),
                    "--refs",
                    str(REFERENCES_DIR),
                    "--page",
                    "1",
                ],
                result,
            )
            if non_content_validate.returncode == 0:
                assert_contains(f"planning-validator-{page_type}", non_content_validate.stdout, ["OK"], result)

            non_content_resolve = run_cmd(
                f"resource-loader-resolve-{page_type}",
                [
                    py,
                    str(SCRIPTS_DIR / "resource_loader.py"),
                    "resolve",
                    "--refs-dir",
                    str(REFERENCES_DIR),
                    "--planning",
                    str(planning_path),
                ],
                result,
            )
            if non_content_resolve.returncode == 0:
                assert_contains(f"resource-loader-resolve-{page_type}", non_content_resolve.stdout, [expected_title], result)
                assert_no_unfilled_vars(f"resource-loader-resolve-{page_type}", non_content_resolve.stdout, result)

        prompt_specs = [
            (
                "prompt-interview-structured",
                fx["prompt_interview_structured"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/tpl-interview-structured-ui.md"),
                    "--var",
                    "TOPIC=Linux.do 社区介绍",
                    "--var",
                    "USER_CONTEXT=4 页介绍型 PPT，目标是快速讲清社区定位、氛围、价值与加入理由。",
                    "--inject-file",
                    f"INTERVIEW_MODE_MODULE={REFERENCES_DIR / 'prompts/module-structured-interview-ui.md'}",
                    "--inject-file",
                    f"INTERVIEW_CORE={REFERENCES_DIR / 'prompts/tpl-interview.md'}",
                    "--output",
                    str(fx["prompt_interview_structured"]),
                ],
            ),
            (
                "prompt-interview-text",
                fx["prompt_interview_text"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/tpl-interview-text-fallback.md"),
                    "--var",
                    "TOPIC=Linux.do 社区介绍",
                    "--var",
                    "USER_CONTEXT=4 页介绍型 PPT，目标是快速讲清社区定位、氛围、价值与加入理由。",
                    "--inject-file",
                    f"INTERVIEW_MODE_MODULE={REFERENCES_DIR / 'prompts/module-text-interview-fallback.md'}",
                    "--inject-file",
                    f"INTERVIEW_CORE={REFERENCES_DIR / 'prompts/tpl-interview.md'}",
                    "--output",
                    str(fx["prompt_interview_text"]),
                ],
            ),
            (
                "prompt-style-phase1",
                fx["prompt_style_phase1"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/tpl-style-phase1.md"),
                    "--var",
                    f"REQUIREMENTS_PATH={fx['requirements']}",
                    "--var",
                    f"OUTLINE_PATH={fx['outline']}",
                    "--var",
                    f"SKILL_DIR={ROOT_DIR}",
                    "--var",
                    f"STYLE_OUTPUT={fx['style']}",
                    "--inject-file",
                    f"STYLE_RUNTIME_RULES={REFERENCES_DIR / 'styles/runtime-style-rules.md'}",
                    "--inject-file",
                    f"STYLE_PRESET_INDEX={REFERENCES_DIR / 'styles/runtime-style-palette-index.md'}",
                    "--inject-file",
                    f"PLAYBOOK={REFERENCES_DIR / 'playbooks/style-phase1-playbook.md'}",
                    "--output",
                    str(fx["prompt_style_phase1"]),
                ],
            ),
            (
                "prompt-page-planning",
                fx["prompt_planning"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/step4/tpl-page-planning.md"),
                    "--var",
                    "PAGE_NUM=3",
                    "--var",
                    "TOTAL_PAGES=8",
                    "--var",
                    f"REQUIREMENTS_PATH={fx['requirements']}",
                    "--var",
                    f"OUTLINE_PATH={fx['outline']}",
                    "--var",
                    f"BRIEF_PATH={fx['brief']}",
                    "--var",
                    f"STYLE_PATH={fx['style']}",
                    "--var",
                    f"IMAGES_DIR={fx['images']}",
                    "--var",
                    f"PLANNING_OUTPUT={fx['planning']}",
                    "--var",
                    f"SKILL_DIR={ROOT_DIR}",
                    "--var",
                    f"REFS_DIR={REFERENCES_DIR}",
                    "--inject-file",
                    f"PRINCIPLES_CHEATSHEET={REFERENCES_DIR / 'principles/design-principles-cheatsheet.md'}",
                    "--inject-file",
                    f"PLAYBOOK={REFERENCES_DIR / 'playbooks/step4/page-planning-playbook.md'}",
                    "--output",
                    str(fx["prompt_planning"]),
                ],
            ),
            (
                "prompt-page-html",
                fx["prompt_html"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/step4/tpl-page-html.md"),
                    "--var",
                    "PAGE_NUM=3",
                    "--var",
                    "TOTAL_PAGES=8",
                    "--var",
                    f"PLANNING_OUTPUT={fx['planning']}",
                    "--var",
                    f"SLIDE_OUTPUT={fx['slide']}",
                    "--var",
                    f"IMAGES_DIR={fx['images']}",
                    "--var",
                    f"STYLE_PATH={fx['style']}",
                    "--var",
                    f"SKILL_DIR={ROOT_DIR}",
                    "--var",
                    f"REFS_DIR={REFERENCES_DIR}",
                    "--inject-file",
                    f"PPTX_EXPORT_SAFETY={REFERENCES_DIR / 'design-runtime/pptx-export-safety.md'}",
                    "--inject-file",
                    f"PPTX_HTML_COMPAT={REFERENCES_DIR / 'design-runtime/pptx-html-compat.md'}",
                    "--inject-file",
                    f"PLAYBOOK={REFERENCES_DIR / 'playbooks/step4/page-html-playbook.md'}",
                    "--output",
                    str(fx["prompt_html"]),
                ],
            ),
            (
                "prompt-page-review",
                fx["prompt_review"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/step4/tpl-page-review.md"),
                    "--var",
                    "PAGE_NUM=3",
                    "--var",
                    "TOTAL_PAGES=8",
                    "--var",
                    f"PLANNING_OUTPUT={fx['planning']}",
                    "--var",
                    f"SLIDE_OUTPUT={fx['slide']}",
                    "--var",
                    f"PNG_OUTPUT={fx['png']}",
                    "--var",
                    f"STYLE_PATH={fx['style']}",
                    "--var",
                    f"SKILL_DIR={ROOT_DIR}",
                    "--var",
                    f"REVIEW_DIR={tmp_dir / 'review'}",
                    "--inject-file",
                    f"PPTX_EXPORT_SAFETY={REFERENCES_DIR / 'design-runtime/pptx-export-safety.md'}",
                    "--inject-file",
                    f"PPTX_HTML_COMPAT={REFERENCES_DIR / 'design-runtime/pptx-html-compat.md'}",
                    "--inject-file",
                    f"PLAYBOOK={REFERENCES_DIR / 'playbooks/step4/page-review-playbook.md'}",
                    "--inject-file",
                    f"FAILURE_MODES={REFERENCES_DIR / 'principles/runtime-failure-modes.md'}",
                    "--inject-file",
                    f"PRINCIPLES_CHEATSHEET={REFERENCES_DIR / 'principles/design-principles-cheatsheet.md'}",
                    "--output",
                    str(fx["prompt_review"]),
                ],
            ),
            (
                "prompt-page-orchestrator",
                fx["prompt_orchestrator"],
                [
                    py,
                    str(SCRIPTS_DIR / "prompt_harness.py"),
                    "--template",
                    str(REFERENCES_DIR / "prompts/step4/tpl-page-orchestrator.md"),
                    "--var",
                    "PAGE_NUM=3",
                    "--var",
                    "TOTAL_PAGES=8",
                    "--var",
                    f"PLANNING_PROMPT_PATH={fx['prompt_planning']}",
                    "--var",
                    f"HTML_PROMPT_PATH={fx['prompt_html']}",
                    "--var",
                    f"REVIEW_PROMPT_PATH={fx['prompt_review']}",
                    "--var",
                    f"PLANNING_OUTPUT={fx['planning']}",
                    "--var",
                    f"SLIDE_OUTPUT={fx['slide']}",
                    "--var",
                    f"PNG_OUTPUT={fx['png']}",
                    "--output",
                    str(fx["prompt_orchestrator"]),
                ],
            ),
        ]

        for label, output_path, args in prompt_specs:
            proc = run_cmd(label, args, result)
            if proc.returncode == 0:
                rendered = output_path.read_text(encoding="utf-8")
                assert_no_unfilled_vars(label, rendered, result)
                if label == "prompt-interview-structured":
                    assert_contains(
                        label,
                        rendered,
                        [
                            "# 采访问卷（Structured UI）",
                            "主题：Linux.do 社区介绍",
                            "用户背景：4 页介绍型 PPT",
                            "# Structured UI Mode -- CLI 原生结构化采访",
                            "AskUserQuestion",
                            "request_user_input",
                            "# 采访问卷共享核心",
                            "presentation_scenario",
                            "## 最终要求",
                        ],
                        result,
                    )
                    assert_max_bytes(label, rendered, 9000, result)
                if label == "prompt-interview-text":
                    assert_contains(
                        label,
                        rendered,
                        [
                            "# 采访问卷（Text Fallback）",
                            "结构化文本采访单",
                            "**A. 场景与目标**",
                            "全部按默认，用 research",
                            "# 采访问卷共享核心",
                            "presentation_scenario",
                            "## 最终要求",
                        ],
                        result,
                    )
                    assert_max_bytes(label, rendered, 9000, result)
                if label == "prompt-style-phase1":
                    assert_contains(
                        label,
                        rendered,
                        [
                            "# Runtime Style Rules",
                            "# Runtime Style Palette Index",
                            "# Style Phase 1 Playbook -- 风格合同的定义与输出",
                        ],
                        result,
                    )
                if label == "prompt-page-planning":
                    assert_contains(
                        label,
                        rendered,
                        ["# Page Planning Playbook -- 单页策划稿", "# 设计原则速查表 -- Step 4 字段级操作手册"],
                        result,
                    )
                if label == "prompt-page-html":
                    assert_contains(
                        label,
                        rendered,
                        [
                            "# Page HTML Playbook -- 单页 HTML 设计稿",
                            "## Phase 4.5：图片 + 正文 + 图表共存排版合同（图文并茂页强制）",
                            "### 图片容器合同",
                        ],
                        result,
                    )
                if label == "prompt-page-review":
                    assert_contains(
                        label,
                        rendered,
                        [
                            "# Page Visual Review & Fix Playbook -- 单页图审与 HTML 修复",
                            "### 图片共存专项扫描（必须逐条过）",
                            "当前自动断言已覆盖：",
                            "# Runtime Failure Modes",
                        ],
                        result,
                    )

    return result


def print_messages(title: str, messages: list[str]) -> None:
    if not messages:
        return
    print(title)
    for item in messages:
        print(f"- {item}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Minimal end-to-end smoke test for the PPT skill")
    parser.add_argument(
        "--strict-warnings",
        action="store_true",
        help="treat warnings as failures",
    )
    args = parser.parse_args()

    result = run_smoke()
    print("PPT skill smoke test")
    print(f"errors: {len(result.errors)}")
    print(f"warnings: {len(result.warnings)}")
    print_messages("Steps", result.steps)
    print_messages("Errors", result.errors)
    print_messages("Warnings", result.warnings)

    if result.errors:
        return 1
    if args.strict_warnings and result.warnings:
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
