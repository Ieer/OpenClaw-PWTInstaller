#!/usr/bin/env python3
"""Inspect PPTX structure and report per-slide editability signals.

Usage:
    python3 scripts/inspect_pptx.py presentation-svg.pptx
    python3 scripts/inspect_pptx.py presentation-svg.pptx -o presentation-svg.inspect.json
"""

from __future__ import annotations

import argparse
import collections.abc
import json
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from speech_script import extract_slide_speaker_notes


CHART_GROUP_NAME_PREFIX = "ChartGroup:"
NATIVE_CHART_NAME_PREFIX = "NativeChart:"


def iso_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def load_source_report(path: Path | None) -> dict | None:
    if path is None or not path.exists():
        return None
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    return payload if isinstance(payload, dict) else None


def shape_type_name(shape) -> str:
    shape_type = getattr(shape, "shape_type", None)
    return getattr(shape_type, "name", str(shape_type))


def shape_name(shape) -> str:
    value = getattr(shape, "name", None)
    return value if isinstance(value, str) else ""


def optional_non_negative_int(source_page: dict | None, field_name: str) -> int | None:
    if not isinstance(source_page, dict):
        return None
    value = source_page.get(field_name)
    return value if isinstance(value, int) and value >= 0 else None


def ratio_or_none(actual: int, expected: int | None) -> float | None:
    if expected is None or expected <= 0:
        return None
    return round(actual / expected, 4)


def is_managed_slide(payload: dict[str, object]) -> bool:
    return payload.get("expected_rendered_charts") is not None or payload.get("expected_structured_charts") is not None


def is_full_slide_picture(shape, slide_width: int, slide_height: int) -> bool:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return False
    tolerance_x = max(int(slide_width * 0.01), 12700)
    tolerance_y = max(int(slide_height * 0.01), 12700)
    return (
        abs(int(shape.left)) <= tolerance_x
        and abs(int(shape.top)) <= tolerance_y
        and abs(int(shape.width) - slide_width) <= tolerance_x
        and abs(int(shape.height) - slide_height) <= tolerance_y
    )


def inspect_slide(slide, slide_number: int, slide_width: int, slide_height: int, source_page: dict | None = None) -> dict[str, object]:
    shape_count = 0
    picture_count = 0
    text_shape_count = 0
    auto_shape_count = 0
    table_count = 0
    chart_count = 0
    group_count = 0
    chart_group_count = 0
    native_chart_count = 0
    connector_count = 0
    other_count = 0
    full_slide_picture_count = 0
    shape_type_breakdown: dict[str, int] = {}

    for shape in slide.shapes:
        shape_count += 1
        type_name = shape_type_name(shape)
        current_name = shape_name(shape)
        shape_type_breakdown[type_name] = shape_type_breakdown.get(type_name, 0) + 1

        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
            if is_full_slide_picture(shape, slide_width, slide_height):
                full_slide_picture_count += 1
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TEXT_BOX:
            text_shape_count += 1
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shape_count += 1
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            group_count += 1
            if current_name.startswith(CHART_GROUP_NAME_PREFIX):
                chart_group_count += 1
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE:
            connector_count += 1
        else:
            other_count += 1

        if getattr(shape, "has_text_frame", False):
            text_shape_count += 1 if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.TEXT_BOX else 0
        if getattr(shape, "has_table", False):
            table_count += 1
        if current_name.startswith(NATIVE_CHART_NAME_PREFIX):
            native_chart_count += 1
        if getattr(shape, "has_chart", False):
            chart_count += 1

    non_picture_shape_count = shape_count - picture_count
    degraded_full_slide_picture = full_slide_picture_count > 0 and non_picture_shape_count == 0 and text_shape_count == 0 and table_count == 0 and chart_count == 0
    structured_chart_count = chart_group_count + native_chart_count
    speaker_notes = extract_slide_speaker_notes(slide)
    has_notes = bool(speaker_notes)
    notes_char_count = len(speaker_notes)
    expected_rendered_charts = optional_non_negative_int(source_page, "rendered_charts")
    expected_chart_groups = optional_non_negative_int(source_page, "structured_chart_groups")
    expected_native_charts = optional_non_negative_int(source_page, "native_charts")
    expected_structured_charts = None
    if expected_chart_groups is not None or expected_native_charts is not None:
        expected_structured_charts = (expected_chart_groups or 0) + (expected_native_charts or 0)

    payload: dict[str, object] = {
        "slide_number": slide_number,
        "shape_count": shape_count,
        "picture_count": picture_count,
        "text_shape_count": text_shape_count,
        "auto_shape_count": auto_shape_count,
        "table_count": table_count,
        "chart_count": chart_count,
        "group_count": group_count,
        "chart_group_count": chart_group_count,
        "native_chart_count": native_chart_count,
        "structured_chart_count": structured_chart_count,
        "connector_count": connector_count,
        "other_count": other_count,
        "full_slide_picture_count": full_slide_picture_count,
        "degraded_full_slide_picture": degraded_full_slide_picture,
        "has_notes": has_notes,
        "notes_char_count": notes_char_count,
        "shape_type_breakdown": shape_type_breakdown,
        "expected_rendered_charts": expected_rendered_charts,
        "expected_chart_groups": expected_chart_groups,
        "expected_native_charts": expected_native_charts,
        "expected_structured_charts": expected_structured_charts,
        "chart_group_hit_rate": ratio_or_none(chart_group_count, expected_rendered_charts),
        "structured_chart_hit_rate": ratio_or_none(structured_chart_count, expected_rendered_charts),
    }
    if isinstance(source_page, dict):
        payload["source_method"] = source_page.get("source_method")
        payload["source_editable"] = source_page.get("source_editable")
    return payload


def build_report(pptx_path: Path, slides_payload: list[dict[str, object]], slide_width: int, slide_height: int, source_report_path: Path | None) -> dict[str, object]:
    total_shapes = sum(int(item["shape_count"]) for item in slides_payload)
    slides_with_notes = sum(1 for item in slides_payload if bool(item.get("has_notes")))
    notes_char_total = sum(int(item.get("notes_char_count") or 0) for item in slides_payload)
    picture_only_slides = sum(1 for item in slides_payload if int(item["shape_count"]) > 0 and int(item["picture_count"]) == int(item["shape_count"]))
    full_slide_picture_slides = sum(1 for item in slides_payload if int(item["full_slide_picture_count"]) > 0)
    degraded_slides = sum(1 for item in slides_payload if bool(item["degraded_full_slide_picture"]))
    table_slides = sum(1 for item in slides_payload if int(item["table_count"]) > 0)
    chart_slides = sum(1 for item in slides_payload if int(item["chart_count"]) > 0)
    managed_slides = [item for item in slides_payload if is_managed_slide(item)]
    chart_scope = managed_slides if managed_slides else slides_payload
    chart_group_slides = sum(1 for item in chart_scope if int(item.get("chart_group_count") or 0) > 0)
    native_chart_slides = sum(1 for item in chart_scope if int(item.get("native_chart_count") or 0) > 0)
    chart_group_shapes_total = sum(int(item.get("chart_group_count") or 0) for item in chart_scope)
    native_chart_shapes_total = sum(int(item.get("native_chart_count") or 0) for item in chart_scope)
    structured_chart_shapes_total = sum(int(item.get("structured_chart_count") or 0) for item in chart_scope)
    expected_rendered_chart_regions_total = sum(int(item.get("expected_rendered_charts") or 0) for item in chart_scope)
    expected_chart_groups_total = sum(int(item.get("expected_chart_groups") or 0) for item in chart_scope)
    expected_native_charts_total = sum(int(item.get("expected_native_charts") or 0) for item in chart_scope)
    expected_structured_chart_shapes_total = sum(int(item.get("expected_structured_charts") or 0) for item in chart_scope)
    return {
        "generated_at": iso_now(),
        "pptx_path": str(pptx_path.resolve()),
        "source_report_path": str(source_report_path.resolve()) if source_report_path else None,
        "slide_size": {
            "width": slide_width,
            "height": slide_height,
        },
        "summary": {
            "total_slides": len(slides_payload),
            "total_shapes": total_shapes,
            "slides_with_notes": slides_with_notes,
            "notes_char_total": notes_char_total,
            "picture_only_slides": picture_only_slides,
            "full_slide_picture_slides": full_slide_picture_slides,
            "degraded_full_slide_picture_slides": degraded_slides,
            "table_slides": table_slides,
            "chart_slides": chart_slides,
            "managed_structured_chart_slides": len(chart_scope),
            "chart_group_slides": chart_group_slides,
            "native_chart_slides": native_chart_slides,
            "chart_group_shapes_total": chart_group_shapes_total,
            "native_chart_shapes_total": native_chart_shapes_total,
            "structured_chart_shapes_total": structured_chart_shapes_total,
            "expected_rendered_chart_regions_total": expected_rendered_chart_regions_total,
            "expected_chart_groups_total": expected_chart_groups_total,
            "expected_native_charts_total": expected_native_charts_total,
            "expected_structured_chart_shapes_total": expected_structured_chart_shapes_total,
            "chart_group_hit_rate": ratio_or_none(chart_group_shapes_total, expected_rendered_chart_regions_total),
            "structured_chart_hit_rate": ratio_or_none(structured_chart_shapes_total, expected_rendered_chart_regions_total),
        },
        "slides": slides_payload,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect PPTX shape structure")
    parser.add_argument("pptx", help="Path to the PPTX file")
    parser.add_argument("-o", "--output", default=None, help="Path to write inspection JSON")
    parser.add_argument("--source-report", default=None, help="Optional PPTX export report JSON for source method annotation")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).resolve()
    if not pptx_path.exists():
        raise SystemExit(f"ERROR: path not found: {pptx_path}")

    prs = Presentation(str(pptx_path))
    source_report_path = Path(args.source_report).resolve() if args.source_report else None
    source_report = load_source_report(source_report_path)
    source_pages = source_report.get("pages") if isinstance(source_report, dict) and isinstance(source_report.get("pages"), list) else []
    source_page_lookup = {
        int(page.get("slide_number")): page
        for page in source_pages
        if isinstance(page, dict) and isinstance(page.get("slide_number"), int) and int(page.get("slide_number")) > 0
    }

    slides_payload = []
    for index, slide in enumerate(prs.slides, start=1):
        source_page = source_page_lookup.get(index)
        slides_payload.append(inspect_slide(slide, index, prs.slide_width, prs.slide_height, source_page))

    report = build_report(pptx_path, slides_payload, prs.slide_width, prs.slide_height, source_report_path)
    output_path = Path(args.output).resolve() if args.output else pptx_path.with_suffix(".inspect.json")
    output_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Saved: {output_path}")
    print(
        "Slides="
        f"{report['summary']['total_slides']}, "
        f"degraded_full_slide_picture_slides={report['summary']['degraded_full_slide_picture_slides']}, "
        f"slides_with_notes={report['summary']['slides_with_notes']}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())