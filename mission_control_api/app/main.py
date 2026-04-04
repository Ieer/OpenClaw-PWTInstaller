from __future__ import annotations

import json
import hashlib
import math
import mimetypes
import re
import time
import urllib.parse
from datetime import datetime, timedelta, timezone
from pathlib import Path
from urllib.parse import urlparse
from uuid import UUID, uuid4

import sqlalchemy as sa
from sqlalchemy.dialects.postgresql import insert as pg_insert
from fastapi import Depends, FastAPI, Header, HTTPException, WebSocket, WebSocketDisconnect, Request
from fastapi.responses import HTMLResponse, Response
import httpx
import asyncio
from redis.asyncio import Redis

from .config import Settings, load_settings
from .db import create_engine, create_session_factory
from .models import (
    agent_skill_mappings,
    comments,
    events,
    knowledge_feedback_events,
    knowledge_unit_lifecycle_events,
    knowledge_resolve_ranking_profiles,
    knowledge_resolve_audits,
    knowledge_sources,
    knowledge_unit_embeddings,
    knowledge_units,
    knowledge_validation_policy_bundles,
    knowledge_validation_policy_change_events,
    knowledge_validation_policies,
    knowledge_validation_policy_rollouts,
    knowledge_validation_policy_rules,
    knowledge_validations,
    tasks,
)
from .schemas import (
    AgentControlActionIn,
    AgentControlActionOut,
    AgentSkillMappingOut,
    AgentUsageSnapshotOut,
    BoardColumn,
    BoardOut,
    CommentCreate,
    CommentOut,
    EventIn,
    EventLiteOut,
    EventOut,
    ContainerHealthSummaryOut,
    Health,
    KnowledgeFeedbackIn,
    KnowledgeFeedbackOut,
    KnowledgeFeedbackSummaryOut,
    KnowledgeLifecycleActionIn,
    KnowledgeLifecycleEventOut,
    KnowledgeResolveRankingProfileOut,
    KnowledgeResolveRankingProfileUpsertIn,
    KnowledgeResolveAuditOut,
    KnowledgeResolveMetricsOut,
    KnowledgeResolveRejectSummaryOut,
    KnowledgeResolveRiskMetricsOut,
    KnowledgeSearchIn,
    KnowledgeSearchItemOut,
    KnowledgeSearchOut,
    KnowledgeSourceImportIn,
    KnowledgeSourceChunkIn,
    KnowledgeSourceChunkOut,
    KnowledgeSourceOut,
    KnowledgeSourceScanIn,
    KnowledgeSourceScanOut,
    KnowledgeResolveIn,
    KnowledgeResolveItemOut,
    KnowledgeResolveRejectedOut,
    KnowledgeResolveOut,
    KnowledgeUnitCreateIn,
    KnowledgeUnitOut,
    KnowledgeValidationCreateIn,
    KnowledgeValidationPolicyBundleOut,
    KnowledgeValidationPolicyBundleUpsertIn,
    KnowledgeValidationPolicyChangeEventOut,
    KnowledgeValidationPolicyObservabilityItemOut,
    KnowledgeValidationPolicyObservabilityOut,
    KnowledgeValidationPolicyRollbackIn,
    KnowledgeValidationPolicyRollbackOut,
    KnowledgeValidationPolicyOut,
    KnowledgeValidationPolicyRolloutOut,
    KnowledgeValidationPolicyRolloutUpsertIn,
    KnowledgeValidationPolicyRuleOut,
    KnowledgeValidationPolicyRuleUpsertIn,
    KnowledgeValidationPolicyUpsertIn,
    KnowledgeValidationOut,
    HealthSignalOut,
    ObservabilitySummaryOut,
    SkillItem,
    SkillsMappingFailureItem,
    SkillsMappingPatchIn,
    SkillsMappingPatchOut,
    TaskCreate,
    TaskOut,
    WorkspaceSkillGroup,
)


ALLOWED_TASK_STATUSES = {"INBOX", "ASSIGNED", "IN PROGRESS", "REVIEW", "DONE"}
TASK_STATUS_TRANSITIONS = {
    "INBOX": {"ASSIGNED"},
    "ASSIGNED": {"IN PROGRESS", "REVIEW"},
    "IN PROGRESS": {"REVIEW", "DONE"},
    "REVIEW": {"IN PROGRESS", "DONE"},
    "DONE": set(),
}

USAGE_CACHE_TTL_SECONDS = 15.0
_AGENT_USAGE_CACHE: dict[str, object] = {
    "days": None,
    "generated_at": 0.0,
    "data": [],
}

ALLOWED_AGENT_CONTROL_ACTIONS = {"start", "stop", "restart"}
KNOWLEDGE_FEEDBACK_TYPES = {"usage", "conflict", "invalidation", "promotion"}
KNOWLEDGE_RISK_ORDER = {"low": 0, "normal": 1, "high": 2, "critical": 3}
KNOWLEDGE_LIFECYCLE_STAGES = {"active", "preferred", "deprecated", "inactive", "archived", "superseded"}
KNOWLEDGE_LIFECYCLE_ACTIONS = {"promote", "demote", "invalidate", "reactivate", "archive", "supersede"}
KNOWLEDGE_POLICY_ROLLOUT_MODES = {"full", "percentage"}
KNOWLEDGE_TEXT_EXTENSIONS = {"txt", "md", "rst", "csv", "json", "yaml", "yml", "log"}
KNOWLEDGE_DOCUMENT_EXTENSIONS = {"pdf", "docx", "pptx", "xlsx"}
KNOWLEDGE_OCR_CANDIDATE_EXTENSIONS = {"png", "jpg", "jpeg", "webp", "pdf"}
KNOWLEDGE_OCR_LANGUAGES = "eng+chi_sim"
KNOWLEDGE_OCR_MAX_PDF_PAGES = 5
KNOWLEDGE_PDF_TEXT_MIN_CHARS = 32
KNOWLEDGE_EMBEDDING_PROTOCOLS = {"openai-embeddings", "openai-completions"}
DEFAULT_VALIDATION_POLICIES: dict[str, dict] = {
    "low": {
        "strict_mode": False,
        "require_validation": False,
        "require_approved": False,
        "require_not_expired": False,
        "min_confidence": None,
        "max_validation_age_days": None,
    },
    "normal": {
        "strict_mode": False,
        "require_validation": True,
        "require_approved": True,
        "require_not_expired": False,
        "min_confidence": None,
        "max_validation_age_days": None,
    },
    "high": {
        "strict_mode": False,
        "require_validation": True,
        "require_approved": True,
        "require_not_expired": True,
        "min_confidence": 0.7,
        "max_validation_age_days": 30,
    },
    "critical": {
        "strict_mode": False,
        "require_validation": True,
        "require_approved": True,
        "require_not_expired": True,
        "min_confidence": 0.85,
        "max_validation_age_days": 14,
    },
}
DEFAULT_RANKING_PROFILES: dict[str, dict[str, float | str | bool | None]] = {
    "balanced": {
        "description": "Balanced default profile for mixed lexical and semantic ranking.",
        "enabled": True,
        "base_score": 1.0,
        "lexical_weight": 0.8,
        "semantic_weight": 1.0,
        "tag_weight": 0.3,
        "validation_confidence_weight": 0.5,
        "approved_bonus": 0.5,
        "preferred_bonus": 0.35,
        "deprecated_penalty": 0.25,
    },
    "precision": {
        "description": "Precision-focused profile for high-risk retrieval.",
        "enabled": True,
        "base_score": 1.0,
        "lexical_weight": 0.8,
        "semantic_weight": 1.25,
        "tag_weight": 0.2,
        "validation_confidence_weight": 0.65,
        "approved_bonus": 0.5,
        "preferred_bonus": 0.35,
        "deprecated_penalty": 0.25,
    },
    "recall": {
        "description": "Recall-focused profile with stronger lexical and tag emphasis.",
        "enabled": True,
        "base_score": 1.0,
        "lexical_weight": 1.0,
        "semantic_weight": 0.9,
        "tag_weight": 0.4,
        "validation_confidence_weight": 0.4,
        "approved_bonus": 0.4,
        "preferred_bonus": 0.25,
        "deprecated_penalty": 0.15,
    },
}


def _safe_resolve_under(root: Path, value: str) -> Path:
    raw = str(value or "").strip().replace("\\", "/")
    if not raw:
        raise HTTPException(status_code=422, detail="path must be non-empty")
    if raw.startswith("/"):
        raise HTTPException(status_code=422, detail="absolute path is not allowed")

    candidate = (root / raw).resolve()
    try:
        candidate.relative_to(root)
    except ValueError as exc:
        raise HTTPException(status_code=422, detail="path must stay under configured knowledge root") from exc
    return candidate


def _relative_storage_path(root: Path, absolute_path: Path) -> str:
    try:
        return str(absolute_path.resolve().relative_to(root.resolve()))
    except ValueError:
        return str(absolute_path)


def _guess_mime_type(path: Path) -> str | None:
    guessed, _ = mimetypes.guess_type(str(path))
    return guessed


def _sha256sum(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _knowledge_source_row_to_out(row) -> KnowledgeSourceOut:
    return KnowledgeSourceOut(
        id=row.id,
        source_type=row.source_type,
        title=row.title,
        external_uri=row.external_uri,
        storage_path=row.storage_path,
        checksum_sha256=row.checksum_sha256,
        mime_type=row.mime_type,
        owner=row.owner,
        version_label=row.version_label,
        status=row.status,
        meta=row.meta or {},
        collected_at=row.collected_at,
        updated_at=row.updated_at,
    )


def _knowledge_unit_row_to_out(row) -> KnowledgeUnitOut:
    return KnowledgeUnitOut(
        id=row.id,
        source_id=row.source_id,
        unit_key=row.unit_key,
        title=row.title,
        content=row.content,
        content_sha256=row.content_sha256,
        tags=list(row.tags or []),
        agent_scope=list(row.agent_scope or []),
        risk_level=row.risk_level,
        status=row.status,
        lifecycle_stage=getattr(row, "lifecycle_stage", "active"),
        superseded_by_unit_id=getattr(row, "superseded_by_unit_id", None),
        retired_at=getattr(row, "retired_at", None),
        meta=row.meta or {},
        created_at=row.created_at,
        updated_at=row.updated_at,
    )


def _knowledge_validation_row_to_out(row) -> KnowledgeValidationOut:
    return KnowledgeValidationOut(
        id=row.id,
        unit_id=row.unit_id,
        validator=row.validator,
        validation_status=row.validation_status,
        validated_at=row.validated_at,
        expires_at=row.expires_at,
        confidence=row.confidence,
        notes=row.notes,
        meta=row.meta or {},
        created_at=row.created_at,
    )


def _knowledge_feedback_row_to_out(row) -> KnowledgeFeedbackOut:
    return KnowledgeFeedbackOut(
        id=row.id,
        unit_id=row.unit_id,
        feedback_type=row.feedback_type,
        agent=row.agent,
        severity=row.severity,
        payload=row.payload or {},
        created_at=row.created_at,
    )


def _knowledge_validation_policy_row_to_out(row) -> KnowledgeValidationPolicyOut:
    return KnowledgeValidationPolicyOut(
        risk_level=row.risk_level,
        strict_mode=bool(row.strict_mode),
        require_validation=bool(row.require_validation),
        require_approved=bool(row.require_approved),
        require_not_expired=bool(row.require_not_expired),
        min_confidence=row.min_confidence,
        max_validation_age_days=row.max_validation_age_days,
        created_at=row.created_at,
        updated_at=row.updated_at,
    )


def _knowledge_validation_policy_bundle_row_to_out(row) -> KnowledgeValidationPolicyBundleOut:
    return KnowledgeValidationPolicyBundleOut(
        id=row.id,
        bundle_key=row.bundle_key,
        description=row.description,
        is_default=bool(row.is_default),
        enabled=bool(row.enabled),
        created_at=row.created_at,
        updated_at=row.updated_at,
    )


def _knowledge_validation_policy_rule_row_to_out(row) -> KnowledgeValidationPolicyRuleOut:
    return KnowledgeValidationPolicyRuleOut(
        id=row.id,
        bundle_id=row.bundle_id,
        rule_key=row.rule_key,
        risk_level=row.risk_level,
        task_pattern=row.task_pattern,
        agent_slug=row.agent_slug,
        source_type=row.source_type,
        priority=int(row.priority or 0),
        enabled=bool(row.enabled),
        strict_mode=bool(row.strict_mode),
        require_validation=bool(row.require_validation),
        require_approved=bool(row.require_approved),
        require_not_expired=bool(row.require_not_expired),
        min_confidence=row.min_confidence,
        max_validation_age_days=row.max_validation_age_days,
        description=row.description,
        created_at=row.created_at,
        updated_at=row.updated_at,
    )


def _knowledge_validation_policy_rollout_row_to_out(row) -> KnowledgeValidationPolicyRolloutOut:
    return KnowledgeValidationPolicyRolloutOut(
        id=row.id,
        bundle_id=row.bundle_id,
        rollout_key=row.rollout_key,
        target_agent_slug=row.target_agent_slug,
        target_source_type=row.target_source_type,
        task_pattern=row.task_pattern,
        priority=int(row.priority or 0),
        enabled=bool(row.enabled),
        rollout_mode=row.rollout_mode,
        rollout_percentage=row.rollout_percentage,
        created_at=row.created_at,
        updated_at=row.updated_at,
    )


def _knowledge_validation_policy_change_event_row_to_out(row) -> KnowledgeValidationPolicyChangeEventOut:
    return KnowledgeValidationPolicyChangeEventOut(
        id=row.id,
        entity_type=row.entity_type,
        entity_id=row.entity_id,
        entity_key=row.entity_key,
        action=row.action,
        actor=row.actor,
        before_state=row.before_state or {},
        after_state=row.after_state or {},
        created_at=row.created_at,
    )


def _knowledge_resolve_ranking_profile_row_to_out(row) -> KnowledgeResolveRankingProfileOut:
    return KnowledgeResolveRankingProfileOut(
        id=row.id,
        profile_key=row.profile_key,
        description=row.description,
        enabled=bool(row.enabled),
        base_score=float(row.base_score),
        lexical_weight=float(row.lexical_weight),
        semantic_weight=float(row.semantic_weight),
        tag_weight=float(row.tag_weight),
        validation_confidence_weight=float(row.validation_confidence_weight),
        approved_bonus=float(row.approved_bonus),
        preferred_bonus=float(row.preferred_bonus),
        deprecated_penalty=float(row.deprecated_penalty),
        created_at=row.created_at,
        updated_at=row.updated_at,
    )


def _knowledge_lifecycle_event_row_to_out(row) -> KnowledgeLifecycleEventOut:
    return KnowledgeLifecycleEventOut(
        id=row.id,
        unit_id=row.unit_id,
        action=row.action,
        actor=row.actor,
        payload=row.payload or {},
        created_at=row.created_at,
    )


def _knowledge_resolve_audit_row_to_out(row) -> KnowledgeResolveAuditOut:
    return KnowledgeResolveAuditOut(
        id=row.id,
        task=row.task,
        agent_slug=row.agent_slug,
        requested_risk_level=row.requested_risk_level,
        tags=list(row.tags or []),
        selected_count=int(row.selected_count or 0),
        rejected_count=int(row.rejected_count or 0),
        payload=row.payload or {},
        created_at=row.created_at,
    )


def _default_validation_policy(risk_level: str) -> dict:
    return dict(DEFAULT_VALIDATION_POLICIES.get(risk_level, DEFAULT_VALIDATION_POLICIES["normal"]))


def _normalize_task_pattern(value: str | None) -> str | None:
    pattern = str(value or "").strip()
    return pattern or None


def _task_pattern_matches(pattern: str | None, task: str) -> bool:
    normalized = _normalize_task_pattern(pattern)
    if not normalized:
        return True
    escaped = re.escape(normalized).replace(r"\*", ".*")
    return re.search(escaped, task, flags=re.IGNORECASE) is not None


def _rollout_bucket(task: str, agent_slug: str | None, source_type: str | None) -> int:
    raw = "|".join([task.strip().lower(), str(agent_slug or "").strip().lower(), str(source_type or "").strip().lower()])
    digest = hashlib.sha256(raw.encode("utf-8")).hexdigest()
    return int(digest[:8], 16) % 100


def _policy_specificity_score(item: dict | object) -> int:
    mapping = item if isinstance(item, dict) else getattr(item, "_mapping", None) or item
    score = 0
    if getattr(mapping, "task_pattern", None) if not isinstance(mapping, dict) else mapping.get("task_pattern"):
        score += 4
    if getattr(mapping, "agent_slug", None) if not isinstance(mapping, dict) else mapping.get("agent_slug"):
        score += 2
    if getattr(mapping, "source_type", None) if not isinstance(mapping, dict) else mapping.get("source_type"):
        score += 1
    return score


def _tokenize_text(value: str) -> set[str]:
    return {token for token in re.findall(r"[a-zA-Z0-9_\-]{2,}", str(value or "").lower())}


def _compute_lexical_overlap(task: str, title: str, content: str, tags: list[str]) -> float:
    query_tokens = _tokenize_text(task)
    if not query_tokens:
        return 0.0
    doc_tokens = _tokenize_text(title) | _tokenize_text(content) | {str(tag).lower() for tag in (tags or [])}
    if not doc_tokens:
        return 0.0
    return float(len(query_tokens & doc_tokens)) / float(len(query_tokens))


def _lifecycle_score_adjustment(stage: str | None) -> float:
    normalized = str(stage or "active").strip().lower()
    if normalized == "preferred":
        return 0.35
    if normalized == "deprecated":
        return -0.25
    return 0.0


def _serialize_row_state(row: object) -> dict:
    def _normalize_value(value):
        if isinstance(value, UUID):
            return str(value)
        if isinstance(value, datetime):
            return value.isoformat()
        if isinstance(value, list):
            return [_normalize_value(item) for item in value]
        if isinstance(value, dict):
            return {str(key): _normalize_value(item) for key, item in value.items()}
        return value

    if row is None:
        return {}
    if isinstance(row, dict):
        return {key: _normalize_value(value) for key, value in row.items()}
    mapping = getattr(row, "_mapping", None)
    if mapping is not None:
        return {key: _normalize_value(value) for key, value in mapping.items()}
    if hasattr(row, "keys"):
        return {key: _normalize_value(row[key]) for key in row.keys()}
    return {key: _normalize_value(value) for key, value in dict(row).items()}


async def _ensure_default_policy_bundle(session) -> dict:
    existing = (
        await session.execute(
            sa.select(
                knowledge_validation_policy_bundles.c.id,
                knowledge_validation_policy_bundles.c.bundle_key,
                knowledge_validation_policy_bundles.c.description,
                knowledge_validation_policy_bundles.c.is_default,
                knowledge_validation_policy_bundles.c.enabled,
                knowledge_validation_policy_bundles.c.created_at,
                knowledge_validation_policy_bundles.c.updated_at,
            ).where(knowledge_validation_policy_bundles.c.bundle_key == "default-v1")
        )
    ).first()
    if existing:
        return dict(existing._mapping)

    now = datetime.utcnow()
    row = (
        await session.execute(
            knowledge_validation_policy_bundles.insert()
            .values(
                id=uuid4(),
                bundle_key="default-v1",
                description="default fallback validation policy bundle",
                is_default=True,
                enabled=True,
                created_at=now,
                updated_at=now,
            )
            .returning(
                knowledge_validation_policy_bundles.c.id,
                knowledge_validation_policy_bundles.c.bundle_key,
                knowledge_validation_policy_bundles.c.description,
                knowledge_validation_policy_bundles.c.is_default,
                knowledge_validation_policy_bundles.c.enabled,
                knowledge_validation_policy_bundles.c.created_at,
                knowledge_validation_policy_bundles.c.updated_at,
            )
        )
    ).one()
    return dict(row._mapping)


async def _ensure_default_ranking_profiles(session) -> None:
    now = datetime.utcnow()
    for profile_key, config in DEFAULT_RANKING_PROFILES.items():
        await session.execute(
            pg_insert(knowledge_resolve_ranking_profiles)
            .values(
                id=uuid4(),
                profile_key=profile_key,
                description=config.get("description"),
                enabled=bool(config.get("enabled", True)),
                base_score=float(config.get("base_score", 1.0)),
                lexical_weight=float(config.get("lexical_weight", 0.8)),
                semantic_weight=float(config.get("semantic_weight", 1.0)),
                tag_weight=float(config.get("tag_weight", 0.3)),
                validation_confidence_weight=float(config.get("validation_confidence_weight", 0.5)),
                approved_bonus=float(config.get("approved_bonus", 0.5)),
                preferred_bonus=float(config.get("preferred_bonus", 0.35)),
                deprecated_penalty=float(config.get("deprecated_penalty", 0.25)),
                created_at=now,
                updated_at=now,
            )
            .on_conflict_do_nothing(index_elements=[knowledge_resolve_ranking_profiles.c.profile_key])
        )


async def _load_ranking_profile(session, profile_key: str) -> dict:
    await _ensure_default_ranking_profiles(session)
    normalized = str(profile_key or "balanced").strip().lower() or "balanced"
    row = (
        await session.execute(
            sa.select(
                knowledge_resolve_ranking_profiles.c.id,
                knowledge_resolve_ranking_profiles.c.profile_key,
                knowledge_resolve_ranking_profiles.c.description,
                knowledge_resolve_ranking_profiles.c.enabled,
                knowledge_resolve_ranking_profiles.c.base_score,
                knowledge_resolve_ranking_profiles.c.lexical_weight,
                knowledge_resolve_ranking_profiles.c.semantic_weight,
                knowledge_resolve_ranking_profiles.c.tag_weight,
                knowledge_resolve_ranking_profiles.c.validation_confidence_weight,
                knowledge_resolve_ranking_profiles.c.approved_bonus,
                knowledge_resolve_ranking_profiles.c.preferred_bonus,
                knowledge_resolve_ranking_profiles.c.deprecated_penalty,
                knowledge_resolve_ranking_profiles.c.created_at,
                knowledge_resolve_ranking_profiles.c.updated_at,
            ).where(knowledge_resolve_ranking_profiles.c.profile_key == normalized)
        )
    ).first()
    if not row:
        raise HTTPException(status_code=422, detail=f"unknown ranking_profile: {normalized}")
    if not bool(row.enabled):
        raise HTTPException(status_code=422, detail=f"ranking_profile is disabled: {normalized}")
    return dict(row._mapping)


async def _record_policy_change_event(
    session,
    *,
    entity_type: str,
    entity_id: UUID | None,
    entity_key: str,
    action: str,
    actor: str | None,
    before_state: dict | None,
    after_state: dict | None,
    now: datetime,
) -> None:
    await session.execute(
        knowledge_validation_policy_change_events.insert().values(
            id=uuid4(),
            entity_type=entity_type,
            entity_id=entity_id,
            entity_key=entity_key,
            action=action,
            actor=actor,
            before_state=before_state or {},
            after_state=after_state or {},
            created_at=now,
        )
    )


def _observability_item(key: str, rows: list[dict]) -> KnowledgeValidationPolicyObservabilityItemOut:
    total_resolve_requests = len(rows)
    requests_with_hits = sum(1 for row in rows if int(row.get("selected_count") or 0) > 0)
    requests_with_rejections = sum(1 for row in rows if int(row.get("rejected_count") or 0) > 0)
    total_selected_units = sum(int(row.get("selected_count") or 0) for row in rows)
    total_rejected_units = sum(int(row.get("rejected_count") or 0) for row in rows)
    denominator = max(total_resolve_requests, 1)
    return KnowledgeValidationPolicyObservabilityItemOut(
        key=key,
        total_resolve_requests=total_resolve_requests,
        requests_with_hits=requests_with_hits,
        requests_with_rejections=requests_with_rejections,
        total_selected_units=total_selected_units,
        total_rejected_units=total_rejected_units,
        resolve_hit_rate=float(requests_with_hits) / float(denominator),
        resolve_rejection_rate=float(requests_with_rejections) / float(denominator),
    )


async def _rollback_policy_entity(
    session,
    *,
    entity_type: str,
    entity_id: UUID,
    actor: str | None,
) -> KnowledgeValidationPolicyRollbackOut:
    rows = (
        await session.execute(
            sa.select(
                knowledge_validation_policy_change_events.c.id,
                knowledge_validation_policy_change_events.c.entity_type,
                knowledge_validation_policy_change_events.c.entity_id,
                knowledge_validation_policy_change_events.c.entity_key,
                knowledge_validation_policy_change_events.c.action,
                knowledge_validation_policy_change_events.c.actor,
                knowledge_validation_policy_change_events.c.before_state,
                knowledge_validation_policy_change_events.c.after_state,
                knowledge_validation_policy_change_events.c.created_at,
            )
            .where(knowledge_validation_policy_change_events.c.entity_type == entity_type)
            .where(knowledge_validation_policy_change_events.c.entity_id == entity_id)
            .order_by(knowledge_validation_policy_change_events.c.created_at.desc())
            .limit(1)
        )
    ).all()
    if not rows:
        raise HTTPException(status_code=404, detail=f"no change history found for {entity_type}")
    latest = rows[0]
    restore_state = dict(latest.before_state or {})
    if entity_type == "bundle":
        if not restore_state:
            restore_state = {"enabled": False, "is_default": False}
        await session.execute(
            knowledge_validation_policy_bundles.update()
            .where(knowledge_validation_policy_bundles.c.id == entity_id)
            .values(**restore_state, updated_at=datetime.utcnow())
        )
        current = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_bundles.c.id,
                    knowledge_validation_policy_bundles.c.bundle_key,
                    knowledge_validation_policy_bundles.c.description,
                    knowledge_validation_policy_bundles.c.is_default,
                    knowledge_validation_policy_bundles.c.enabled,
                    knowledge_validation_policy_bundles.c.created_at,
                    knowledge_validation_policy_bundles.c.updated_at,
                ).where(knowledge_validation_policy_bundles.c.id == entity_id)
            )
        ).one()
        current_state = _serialize_row_state(current)
        await _record_policy_change_event(
            session,
            entity_type=entity_type,
            entity_id=entity_id,
            entity_key=current.bundle_key,
            action="rollback",
            actor=actor,
            before_state=dict(latest.after_state or {}),
            after_state=current_state,
            now=datetime.utcnow(),
        )
        return KnowledgeValidationPolicyRollbackOut(
            entity_type=entity_type,
            entity_id=entity_id,
            entity_key=current.bundle_key,
            rollback_action="rollback",
            restored_from_event_id=latest.id,
            current_state=current_state,
        )
    if entity_type == "rollout":
        if not restore_state:
            restore_state = {"enabled": False}
        await session.execute(
            knowledge_validation_policy_rollouts.update()
            .where(knowledge_validation_policy_rollouts.c.id == entity_id)
            .values(**restore_state, updated_at=datetime.utcnow())
        )
        current = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rollouts.c.id,
                    knowledge_validation_policy_rollouts.c.bundle_id,
                    knowledge_validation_policy_rollouts.c.rollout_key,
                    knowledge_validation_policy_rollouts.c.target_agent_slug,
                    knowledge_validation_policy_rollouts.c.target_source_type,
                    knowledge_validation_policy_rollouts.c.task_pattern,
                    knowledge_validation_policy_rollouts.c.priority,
                    knowledge_validation_policy_rollouts.c.enabled,
                    knowledge_validation_policy_rollouts.c.rollout_mode,
                    knowledge_validation_policy_rollouts.c.rollout_percentage,
                    knowledge_validation_policy_rollouts.c.created_at,
                    knowledge_validation_policy_rollouts.c.updated_at,
                ).where(knowledge_validation_policy_rollouts.c.id == entity_id)
            )
        ).one()
        current_state = _serialize_row_state(current)
        await _record_policy_change_event(
            session,
            entity_type=entity_type,
            entity_id=entity_id,
            entity_key=current.rollout_key,
            action="rollback",
            actor=actor,
            before_state=dict(latest.after_state or {}),
            after_state=current_state,
            now=datetime.utcnow(),
        )
        return KnowledgeValidationPolicyRollbackOut(
            entity_type=entity_type,
            entity_id=entity_id,
            entity_key=current.rollout_key,
            rollback_action="rollback",
            restored_from_event_id=latest.id,
            current_state=current_state,
        )
    raise HTTPException(status_code=422, detail=f"unsupported rollback entity_type: {entity_type}")


async def _select_policy_bundle(session, *, task: str, agent_slug: str | None, source_type: str | None) -> tuple[dict, dict | None]:
    default_bundle = await _ensure_default_policy_bundle(session)
    rollout_rows = (
        await session.execute(
            sa.select(
                knowledge_validation_policy_rollouts.c.id,
                knowledge_validation_policy_rollouts.c.bundle_id,
                knowledge_validation_policy_rollouts.c.rollout_key,
                knowledge_validation_policy_rollouts.c.target_agent_slug,
                knowledge_validation_policy_rollouts.c.target_source_type,
                knowledge_validation_policy_rollouts.c.task_pattern,
                knowledge_validation_policy_rollouts.c.priority,
                knowledge_validation_policy_rollouts.c.enabled,
                knowledge_validation_policy_rollouts.c.rollout_mode,
                knowledge_validation_policy_rollouts.c.rollout_percentage,
                knowledge_validation_policy_bundles.c.bundle_key,
                knowledge_validation_policy_bundles.c.description,
                knowledge_validation_policy_bundles.c.is_default,
                knowledge_validation_policy_bundles.c.enabled.label("bundle_enabled"),
            )
            .select_from(
                knowledge_validation_policy_rollouts.join(
                    knowledge_validation_policy_bundles,
                    knowledge_validation_policy_bundles.c.id == knowledge_validation_policy_rollouts.c.bundle_id,
                )
            )
            .where(knowledge_validation_policy_rollouts.c.enabled == sa.true())
            .order_by(
                knowledge_validation_policy_rollouts.c.priority.desc(),
                knowledge_validation_policy_rollouts.c.updated_at.desc(),
                knowledge_validation_policy_rollouts.c.created_at.desc(),
            )
        )
    ).mappings().all()

    bucket = _rollout_bucket(task, agent_slug, source_type)
    for row in rollout_rows:
        if not bool(row["bundle_enabled"]):
            continue
        if row["target_agent_slug"] and row["target_agent_slug"] != agent_slug:
            continue
        if row["target_source_type"] and row["target_source_type"] != source_type:
            continue
        if not _task_pattern_matches(row["task_pattern"], task):
            continue
        mode = str(row["rollout_mode"] or "full").strip().lower()
        if mode not in KNOWLEDGE_POLICY_ROLLOUT_MODES:
            continue
        if mode == "percentage":
            percentage = int(row["rollout_percentage"] or 0)
            if percentage <= 0 or bucket >= percentage:
                continue
        bundle = {
            "id": row["bundle_id"],
            "bundle_key": row["bundle_key"],
            "description": row["description"],
            "is_default": bool(row["is_default"]),
            "enabled": bool(row["bundle_enabled"]),
        }
        return bundle, dict(row)

    return default_bundle, None


async def _resolve_validation_policy(
    session,
    *,
    task: str,
    agent_slug: str | None,
    source_type: str | None,
    risk_level: str,
) -> tuple[dict, dict]:
    normalized_risk = str(risk_level or "normal").strip().lower() or "normal"
    bundle, rollout = await _select_policy_bundle(
        session,
        task=task,
        agent_slug=agent_slug,
        source_type=source_type,
    )

    rule_rows = (
        await session.execute(
            sa.select(
                knowledge_validation_policy_rules.c.id,
                knowledge_validation_policy_rules.c.bundle_id,
                knowledge_validation_policy_rules.c.rule_key,
                knowledge_validation_policy_rules.c.risk_level,
                knowledge_validation_policy_rules.c.task_pattern,
                knowledge_validation_policy_rules.c.agent_slug,
                knowledge_validation_policy_rules.c.source_type,
                knowledge_validation_policy_rules.c.priority,
                knowledge_validation_policy_rules.c.enabled,
                knowledge_validation_policy_rules.c.strict_mode,
                knowledge_validation_policy_rules.c.require_validation,
                knowledge_validation_policy_rules.c.require_approved,
                knowledge_validation_policy_rules.c.require_not_expired,
                knowledge_validation_policy_rules.c.min_confidence,
                knowledge_validation_policy_rules.c.max_validation_age_days,
                knowledge_validation_policy_rules.c.description,
                knowledge_validation_policy_rules.c.created_at,
                knowledge_validation_policy_rules.c.updated_at,
            )
            .where(knowledge_validation_policy_rules.c.bundle_id == bundle["id"])
            .where(knowledge_validation_policy_rules.c.risk_level == normalized_risk)
            .where(knowledge_validation_policy_rules.c.enabled == sa.true())
            .order_by(
                knowledge_validation_policy_rules.c.priority.desc(),
                knowledge_validation_policy_rules.c.updated_at.desc(),
                knowledge_validation_policy_rules.c.created_at.desc(),
            )
        )
    ).mappings().all()

    matching_rules: list[dict] = []
    for row in rule_rows:
        if row["agent_slug"] and row["agent_slug"] != agent_slug:
            continue
        if row["source_type"] and row["source_type"] != source_type:
            continue
        if not _task_pattern_matches(row["task_pattern"], task):
            continue
        matching_rules.append(dict(row))

    matching_rules.sort(
        key=lambda item: (
            int(item["priority"] or 0),
            _policy_specificity_score(item),
            str(item.get("updated_at") or ""),
            str(item.get("created_at") or ""),
        ),
        reverse=True,
    )
    if matching_rules:
        matched_rule = matching_rules[0]
        policy = {
            "strict_mode": bool(matched_rule["strict_mode"]),
            "require_validation": bool(matched_rule["require_validation"]),
            "require_approved": bool(matched_rule["require_approved"]),
            "require_not_expired": bool(matched_rule["require_not_expired"]),
            "min_confidence": matched_rule["min_confidence"],
            "max_validation_age_days": matched_rule["max_validation_age_days"],
        }
        metadata = {
            "policy_source": "dynamic_rule",
            "bundle_id": str(bundle["id"]),
            "bundle_key": bundle["bundle_key"],
            "rollout_id": str(rollout["id"]) if rollout else None,
            "rollout_key": rollout["rollout_key"] if rollout else None,
            "rule_id": str(matched_rule["id"]),
            "rule_key": matched_rule["rule_key"],
        }
        return policy, metadata

    policy_row = (
        await session.execute(
            sa.select(
                knowledge_validation_policies.c.strict_mode,
                knowledge_validation_policies.c.require_validation,
                knowledge_validation_policies.c.require_approved,
                knowledge_validation_policies.c.require_not_expired,
                knowledge_validation_policies.c.min_confidence,
                knowledge_validation_policies.c.max_validation_age_days,
            ).where(knowledge_validation_policies.c.risk_level == normalized_risk)
        )
    ).first()
    policy = dict(policy_row._mapping) if policy_row else _default_validation_policy(normalized_risk)
    metadata = {
        "policy_source": "risk_level_default",
        "bundle_id": str(bundle["id"]),
        "bundle_key": bundle["bundle_key"],
        "rollout_id": str(rollout["id"]) if rollout else None,
        "rollout_key": rollout["rollout_key"] if rollout else None,
        "rule_id": None,
        "rule_key": None,
    }
    return policy, metadata


async def _record_lifecycle_event(session, *, unit_id: UUID, action: str, actor: str | None, payload: dict, now: datetime) -> None:
    await session.execute(
        knowledge_unit_lifecycle_events.insert().values(
            id=uuid4(),
            unit_id=unit_id,
            action=action,
            actor=actor,
            payload=payload,
            created_at=now,
        )
    )


def _split_text_chunks(content: str, *, chunk_chars: int, overlap: int, max_chunks: int) -> list[str]:
    text = str(content or "").strip()
    if not text:
        return []

    size = min(max(int(chunk_chars or 1200), 200), 8000)
    step = max(100, size - min(max(int(overlap or 120), 0), size - 50))
    chunks: list[str] = []
    idx = 0

    while idx < len(text) and len(chunks) < max_chunks:
        end = min(idx + size, len(text))
        window = text[idx:end]
        if end < len(text):
            last_break = max(window.rfind("\n\n"), window.rfind("\n"), window.rfind("。"), window.rfind("."))
            if last_break > int(size * 0.6):
                end = idx + last_break + 1
                window = text[idx:end]
        chunks.append(window.strip())
        idx += step if end < len(text) else len(text)

    return [chunk for chunk in chunks if chunk]


def _risk_rank(level: str | None) -> int:
    return KNOWLEDGE_RISK_ORDER.get(str(level or "normal").strip().lower(), KNOWLEDGE_RISK_ORDER["normal"])


def _normalize_ocr_languages(value: str | None) -> str:
    normalized = str(value or "").strip()
    return normalized or KNOWLEDGE_OCR_LANGUAGES


def _normalize_max_pdf_pages(value: int | None) -> int:
    if value is None:
        return KNOWLEDGE_OCR_MAX_PDF_PAGES
    return min(max(int(value), 1), 20)


def _knowledge_embedding_runtime_config(settings: Settings) -> dict | None:
    if not settings.knowledge_embedding_enabled:
        return None

    protocol = str(settings.knowledge_embedding_api_protocol or "openai-embeddings").strip().lower()
    if protocol not in KNOWLEDGE_EMBEDDING_PROTOCOLS:
        raise ValueError(f"unsupported embedding protocol: {protocol}")

    model = str(settings.knowledge_embedding_model or "").strip()
    base_url = str(settings.knowledge_embedding_base_url or "").strip().rstrip("/")
    if not model:
        raise ValueError("MC_KNOWLEDGE_EMBEDDING_MODEL is required when embedding is enabled")
    if not base_url:
        raise ValueError("MC_KNOWLEDGE_EMBEDDING_BASE_URL is required when embedding is enabled")

    dimensions = settings.knowledge_embedding_dimensions
    if dimensions is not None:
        dimensions = int(dimensions)
        if dimensions <= 0:
            raise ValueError("MC_KNOWLEDGE_EMBEDDING_DIMENSIONS must be a positive integer when provided")

    timeout = max(float(settings.knowledge_embedding_timeout_seconds or 20.0), 1.0)
    return {
        "enabled": True,
        "model": model,
        "base_url": base_url,
        "api_key": str(settings.knowledge_embedding_api_key or "").strip(),
        "protocol": protocol,
        "dimensions": dimensions,
        "timeout_seconds": timeout,
    }


def _vector_literal(values: list[float]) -> str:
    return "[" + ",".join(f"{float(value):.12g}" for value in values) + "]"


def _vector_cast_type(dimensions: int) -> str:
    normalized = int(dimensions)
    if normalized <= 0:
        raise ValueError(f"vector dimensions must be positive, got {normalized}")
    return f"vector({normalized})"


def _normalize_embedding(values: object, *, expected_dimensions: int | None) -> list[float]:
    if not isinstance(values, list):
        raise ValueError("embedding response item is not a list")
    normalized = [float(value) for value in values]
    if not normalized:
        raise ValueError("embedding response item is empty")
    if expected_dimensions is not None and len(normalized) != expected_dimensions:
        raise ValueError(f"embedding dimensions mismatch: expected {expected_dimensions}, got {len(normalized)}")
    if any(not math.isfinite(value) for value in normalized):
        raise ValueError("embedding contains non-finite values")
    return normalized


async def _request_knowledge_embeddings(
    texts: list[str],
    *,
    config: dict,
) -> tuple[list[list[float]], str, int]:
    headers = {"Content-Type": "application/json"}
    api_key = str(config.get("api_key") or "").strip()
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    payload = {
        "model": config["model"],
        "input": texts,
    }
    url = f"{config['base_url']}/embeddings"

    async with httpx.AsyncClient(timeout=config["timeout_seconds"]) as client:
        response = await client.post(url, headers=headers, json=payload)

    if response.status_code >= 400:
        detail = response.text.strip()
        raise ValueError(f"embedding request failed: status={response.status_code} detail={detail[:400]}")

    body = response.json()
    data = body.get("data")
    if not isinstance(data, list) or len(data) != len(texts):
        raise ValueError("embedding response data length mismatch")

    vectors: list[list[float]] = []
    actual_dimensions: int | None = None
    for item in data:
        if not isinstance(item, dict):
            raise ValueError("embedding response item is not an object")
        vector = _normalize_embedding(item.get("embedding"), expected_dimensions=config.get("dimensions"))
        if actual_dimensions is None:
            actual_dimensions = len(vector)
        elif len(vector) != actual_dimensions:
            raise ValueError("embedding response dimensions are inconsistent across items")
        vectors.append(vector)

    if actual_dimensions is None:
        raise ValueError("embedding response returned no vectors")

    model_name = str(body.get("model") or config["model"]).strip() or config["model"]
    return vectors, model_name, actual_dimensions


def _similarity_from_cosine_distance(distance: float) -> float:
    return max(min(1.0 - float(distance), 1.0), -1.0)


async def _search_knowledge_units_by_embedding(
    session,
    *,
    query_embedding: list[float],
    embedding_model: str,
    embedding_dimensions: int,
    limit: int,
    source_id: UUID | None,
    source_type: str | None,
    agent_slug: str | None,
    risk_level: str | None,
    require_approved_validation: bool,
    tags: list[str],
) -> list[dict]:
    cast_type = _vector_cast_type(embedding_dimensions)
    allowed_levels = [
        name for name, rank in KNOWLEDGE_RISK_ORDER.items()
        if rank <= _risk_rank(risk_level)
    ]

    sql_parts = [
        """
        SELECT
          ku.id,
          ku.source_id,
          ku.unit_key,
          ku.title,
          ku.content,
          ku.content_sha256,
          ku.tags,
          ku.agent_scope,
          ku.risk_level,
          ku.status,
          ku.lifecycle_stage,
          ku.superseded_by_unit_id,
          ku.retired_at,
          ku.meta,
          ku.created_at,
          ku.updated_at,
                    ks.source_type,
          kue.embedding_model,
          kue.embedding_dimensions,
          ((kue.embedding::""" + cast_type + """) <=> CAST(:query_embedding AS """ + cast_type + """)) AS distance
        FROM knowledge_unit_embeddings kue
        JOIN knowledge_units ku ON ku.id = kue.unit_id
                LEFT JOIN knowledge_sources ks ON ks.id = ku.source_id
        WHERE ku.status = 'active'
          AND kue.embedding_model = :embedding_model
          AND kue.embedding_dimensions = :embedding_dimensions
          AND ku.risk_level = ANY(:allowed_levels)
        """
    ]
    params: dict[str, object] = {
        "query_embedding": _vector_literal(query_embedding),
        "embedding_model": embedding_model,
        "embedding_dimensions": int(embedding_dimensions),
        "allowed_levels": allowed_levels,
        "limit": int(limit),
    }

    if source_id is not None:
        sql_parts.append("AND ku.source_id = :source_id")
        params["source_id"] = source_id

    if source_type:
        sql_parts.append("AND ks.source_type = :source_type")
        params["source_type"] = source_type

    if agent_slug:
        sql_parts.append(
            "AND (cardinality(ku.agent_scope) = 0 OR :agent_slug = ANY(ku.agent_scope))"
        )
        params["agent_slug"] = agent_slug

    if require_approved_validation:
        sql_parts.append(
            """
            AND EXISTS (
              SELECT 1
              FROM (
                SELECT kv.validation_status
                FROM knowledge_validations kv
                WHERE kv.unit_id = ku.id
                ORDER BY kv.validated_at DESC
                LIMIT 1
              ) latest_validation
              WHERE latest_validation.validation_status = 'approved'
            )
            """
        )

    if tags:
        sql_parts.append("AND ku.tags && CAST(:tags AS text[])")
        params["tags"] = tags

    sql_parts.append("ORDER BY distance ASC, ku.updated_at DESC")
    sql_parts.append("LIMIT :limit")
    stmt = sa.text("\n".join(sql_parts))
    rows = (await session.execute(stmt, params)).mappings().all()
    return [dict(row) for row in rows]


def _coerce_row_mapping(row: object) -> dict:
    if isinstance(row, dict):
        return dict(row)
    mapping = getattr(row, "_mapping", None)
    if mapping is not None:
        return dict(mapping)
    raise TypeError("unsupported row type")


async def _upsert_knowledge_unit_embedding(
    session,
    *,
    unit_id: UUID,
    embedding_model: str,
    embedding_dimensions: int,
    content_sha256: str,
    embedding_values: list[float],
    meta: dict,
    now: datetime,
) -> None:
    await session.execute(
        sa.text(
            """
            INSERT INTO knowledge_unit_embeddings (
              id,
              unit_id,
              embedding_model,
              embedding_dimensions,
              content_sha256,
              embedding,
              meta,
              created_at,
              updated_at
            )
            VALUES (
              :id,
              :unit_id,
              :embedding_model,
              :embedding_dimensions,
              :content_sha256,
              CAST(:embedding AS vector),
              CAST(:meta AS jsonb),
              :created_at,
              :updated_at
            )
            ON CONFLICT (unit_id, embedding_model)
            DO UPDATE SET
              embedding_dimensions = EXCLUDED.embedding_dimensions,
              content_sha256 = EXCLUDED.content_sha256,
              embedding = EXCLUDED.embedding,
              meta = EXCLUDED.meta,
              updated_at = EXCLUDED.updated_at
            """
        ),
        {
            "id": uuid4(),
            "unit_id": unit_id,
            "embedding_model": embedding_model,
            "embedding_dimensions": embedding_dimensions,
            "content_sha256": content_sha256,
            "embedding": _vector_literal(embedding_values),
            "meta": json.dumps(meta, ensure_ascii=True),
            "created_at": now,
            "updated_at": now,
        },
    )


def _extract_plain_text(path: Path) -> tuple[str, str, list[str]]:
    return path.read_text(encoding="utf-8", errors="replace"), "plain_text", []


def _ocr_image_to_text(image, *, ocr_languages: str) -> str:
    import pytesseract
    from PIL import ImageOps

    normalized = ImageOps.exif_transpose(image)
    grayscale = ImageOps.grayscale(normalized)
    return str(pytesseract.image_to_string(grayscale, lang=ocr_languages) or "").strip()


def _extract_image_ocr_text(path: Path, *, ocr_languages: str) -> tuple[str, str, list[str], dict]:
    from PIL import Image

    with Image.open(path) as image:
        text = _ocr_image_to_text(image, ocr_languages=ocr_languages)
    if not text:
        raise ValueError("image OCR produced no text")
    return text, "tesseract_ocr", [], {
        "ocr_used": True,
        "ocr_fallback_used": False,
        "ocr_languages": ocr_languages,
        "ocr_pages_processed": 1,
        "ocr_pages_truncated": 0,
    }


def _extract_pdf_ocr_text(path: Path, *, ocr_languages: str, max_pdf_pages: int) -> tuple[str, str, list[str], dict]:
    from pdf2image import convert_from_path
    from pypdf import PdfReader

    page_total = len(PdfReader(str(path)).pages)
    pages_to_process = min(page_total, max_pdf_pages)
    images = convert_from_path(str(path), dpi=150, first_page=1, last_page=pages_to_process)
    parts: list[str] = []
    warnings: list[str] = []

    for index, image in enumerate(images, start=1):
        page_text = _ocr_image_to_text(image, ocr_languages=ocr_languages)
        if page_text:
            parts.append(page_text)
        else:
            warnings.append(f"ocr_page_{index}_no_text")

    pages_truncated = max(page_total - pages_to_process, 0)
    if pages_truncated > 0:
        warnings.append("ocr_pdf_page_limit_applied")

    text = "\n\n".join(parts).strip()
    if not text:
        raise ValueError("scanned pdf OCR produced no text")
    return text, "pypdf+tesseract_ocr", warnings, {
        "ocr_used": True,
        "ocr_fallback_used": True,
        "ocr_languages": ocr_languages,
        "ocr_pages_processed": pages_to_process,
        "ocr_pages_truncated": pages_truncated,
    }


def _extract_pdf_text(path: Path, *, ocr_enabled: bool, ocr_languages: str, max_pdf_pages: int) -> tuple[str, str, list[str], dict]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    warnings: list[str] = []

    for index, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            parts.append(page_text)
        else:
            warnings.append(f"page_{index}_no_text")

    text = "\n\n".join(parts).strip()
    if len(text) >= KNOWLEDGE_PDF_TEXT_MIN_CHARS:
        return text, "pypdf", warnings, {
            "ocr_used": False,
            "ocr_fallback_used": False,
            "ocr_languages": ocr_languages,
            "ocr_pages_processed": 0,
            "ocr_pages_truncated": 0,
        }

    if not ocr_enabled:
        raise ValueError("pdf text layer is insufficient and OCR is disabled for this request")

    ocr_text, parser, ocr_warnings, ocr_meta = _extract_pdf_ocr_text(
        path,
        ocr_languages=ocr_languages,
        max_pdf_pages=max_pdf_pages,
    )
    combined_warnings = list(warnings)
    combined_warnings.append("pdf_text_layer_insufficient_fallback_to_ocr")
    combined_warnings.extend(ocr_warnings)
    return ocr_text, parser, combined_warnings, ocr_meta


def _extract_docx_text(path: Path) -> tuple[str, str, list[str]]:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = str(paragraph.text or "").strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [str(cell.text or "").strip() for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if line:
                parts.append(line)

    text = "\n".join(parts).strip()
    if not text:
        raise ValueError("docx has no extractable text")
    return text, "python-docx", []


def _extract_pptx_text(path: Path) -> tuple[str, str, list[str]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"[slide {slide_index}]\n" + "\n".join(slide_parts))

    text = "\n\n".join(parts).strip()
    if not text:
        raise ValueError("pptx has no extractable text")
    return text, "python-pptx", []


def _extract_xlsx_text(path: Path) -> tuple[str, str, list[str]]:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    warnings: list[str] = []

    for sheet in workbook.worksheets:
        sheet_lines: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value not in (None, "")]
            if values:
                sheet_lines.append("\t".join(values))
        if sheet_lines:
            parts.append(f"[sheet {sheet.title}]\n" + "\n".join(sheet_lines))
        else:
            warnings.append(f"sheet_{sheet.title}_empty")

    text = "\n\n".join(parts).strip()
    if not text:
        raise ValueError("xlsx has no extractable text")
    return text, "openpyxl", warnings


def _extract_source_text(
    path: Path,
    *,
    ocr_enabled: bool,
    ocr_languages: str,
    max_pdf_pages: int,
) -> tuple[str, str, list[str], dict]:
    ext = path.suffix.lower().lstrip(".")
    if ext in KNOWLEDGE_TEXT_EXTENSIONS:
        text, parser, warnings = _extract_plain_text(path)
        return text, parser, warnings, {
            "ocr_used": False,
            "ocr_fallback_used": False,
            "ocr_languages": ocr_languages,
            "ocr_pages_processed": 0,
            "ocr_pages_truncated": 0,
        }
    if ext == "pdf":
        return _extract_pdf_text(
            path,
            ocr_enabled=ocr_enabled,
            ocr_languages=ocr_languages,
            max_pdf_pages=max_pdf_pages,
        )
    if ext == "docx":
        text, parser, warnings = _extract_docx_text(path)
        return text, parser, warnings, {
            "ocr_used": False,
            "ocr_fallback_used": False,
            "ocr_languages": ocr_languages,
            "ocr_pages_processed": 0,
            "ocr_pages_truncated": 0,
        }
    if ext == "pptx":
        text, parser, warnings = _extract_pptx_text(path)
        return text, parser, warnings, {
            "ocr_used": False,
            "ocr_fallback_used": False,
            "ocr_languages": ocr_languages,
            "ocr_pages_processed": 0,
            "ocr_pages_truncated": 0,
        }
    if ext == "xlsx":
        text, parser, warnings = _extract_xlsx_text(path)
        return text, parser, warnings, {
            "ocr_used": False,
            "ocr_fallback_used": False,
            "ocr_languages": ocr_languages,
            "ocr_pages_processed": 0,
            "ocr_pages_truncated": 0,
        }
    if ext in {"png", "jpg", "jpeg", "webp"}:
        if not ocr_enabled:
            raise ValueError(f"{ext} requires OCR but OCR is disabled for this request")
        return _extract_image_ocr_text(path, ocr_languages=ocr_languages)
    if ext in KNOWLEDGE_OCR_CANDIDATE_EXTENSIONS:
        raise ValueError(f"unsupported OCR candidate extension for current implementation: {ext}")
    raise ValueError(f"unsupported source extension for chunking: {ext}")


def _rewrite_avatar_paths(obj, agent: str):
    if isinstance(obj, str):
        if obj.startswith("/avatar/"):
            return f"/chat/{agent}{obj}"
        return obj
    if isinstance(obj, list):
        return [_rewrite_avatar_paths(item, agent) for item in obj]
    if isinstance(obj, dict):
        return {k: _rewrite_avatar_paths(v, agent) for k, v in obj.items()}
    return obj


def _build_chat_inject_script(
    agent: str,
    token: str | None,
    *,
    clear_device_auth_storage: bool = True,
    inject_gateway_settings: bool = True,
    dom_avatar_rewrite: bool = True,
) -> str:
    base_path = f"/chat/{agent}"
    token_json = json.dumps(token or "", ensure_ascii=False)
    parts = [f'window.__OPENCLAW_CONTROL_UI_BASE_PATH__="{base_path}";', "(function(){"]
    if clear_device_auth_storage or inject_gateway_settings:
        parts.append("try{")
        if clear_device_auth_storage:
            parts.append('localStorage.removeItem("openclaw.device.auth.v1");')
            parts.append('sessionStorage.removeItem("openclaw.device.auth.v1");')
        if inject_gateway_settings:
            parts.extend(
                [
                    'const k="openclaw.control.settings.v1";',
                    'const raw=localStorage.getItem(k);',
                    'let v={};',
                    'try{v=raw?JSON.parse(raw):{};}catch{}',
                    f'v.gatewayUrl=(location.protocol==="https:"?"wss":"ws")+"://"+location.host+"{base_path}/";',
                    f'v.token={token_json};',
                    'localStorage.setItem(k,JSON.stringify(v));',
                ]
            )
        parts.append("}catch(e){}")
    if dom_avatar_rewrite:
        parts.extend(
            [
                "try{",
                f'const p="{base_path}";',
                "const scan=()=>{",
                "document.querySelectorAll('img.chat-avatar.assistant[src^=\"/avatar/\"]').forEach((img)=>{",
                'const s=img.getAttribute("src")||"";',
                'if(s.startsWith("/avatar/"))img.setAttribute("src",p+s);',
                "});",
                "};",
                "scan();",
                "const mo=new MutationObserver(()=>scan());",
                'mo.observe(document.documentElement,{subtree:true,childList:true,attributes:true,attributeFilter:["src"]});',
                'window.addEventListener("beforeunload",()=>mo.disconnect(),{once:true});',
                "}catch(e){}",
            ]
        )
    parts.append("})();")
    return "".join(parts)


def _normalize_control_ui_origin(origin: str | None) -> str | None:
    if not origin:
        return None
    if origin.startswith("http://127.0.0.1:"):
        return origin.replace("http://127.0.0.1", "http://localhost", 1)
    if origin.startswith("https://127.0.0.1:"):
        return origin.replace("https://127.0.0.1", "https://localhost", 1)
    return origin


def _sanitize_connect_auth(
    message: str,
    token: str | None,
    *,
    strip_stale_device_fields: bool = True,
    force_token_in_connect: bool = True,
) -> str:
    try:
        req = json.loads(message)
    except Exception:
        return message

    if not isinstance(req, dict):
        return message
    if req.get("type") != "req" or req.get("method") != "connect":
        return message
    if not isinstance(req.get("params"), dict):
        return message

    params = req["params"]

    stale_device_keys = {
        "device",
        "deviceId",
        "deviceID",
        "device_id",
        "deviceKey",
        "device_key",
        "deviceSecret",
        "device_secret",
        "deviceToken",
        "device_token",
        "clientId",
        "client_id",
    }

    auth = params.get("auth")
    if not isinstance(auth, dict):
        auth = {}

    if strip_stale_device_fields:
        for key in stale_device_keys:
            auth.pop(key, None)
            params.pop(key, None)

    if token and force_token_in_connect:
        params["auth"] = {"token": token}
    elif token and not auth.get("token"):
        auth["token"] = token
        params["auth"] = auth
    elif auth:
        params["auth"] = auth
    else:
        params.pop("auth", None)

    req["params"] = params
    return json.dumps(req, ensure_ascii=False)


def _rewrite_avatar_meta(content: bytes, agent: str, query: str) -> bytes:
    try:
        payload = json.loads(content.decode("utf-8", errors="replace"))
    except Exception:
        return content

    avatar_url = payload.get("avatarUrl") if isinstance(payload, dict) else None
    if isinstance(avatar_url, str) and avatar_url.startswith("/avatar/"):
        rewritten = f"/chat/{agent}{avatar_url}"
        if query:
            rewritten = f"{rewritten}?{query}"
        payload["avatarUrl"] = rewritten
        return json.dumps(payload, ensure_ascii=False).encode("utf-8")
    return content


def _rewrite_control_ui_config(content: bytes, agent: str, *, rewrite_avatar: bool = True) -> bytes:
    try:
        payload = json.loads(content.decode("utf-8", errors="replace"))
    except Exception:
        return content

    if not isinstance(payload, dict):
        return content

    base_path = f"/chat/{agent}"
    payload["basePath"] = base_path

    if rewrite_avatar:
        avatar = payload.get("assistantAvatar")
        if isinstance(avatar, str) and avatar.startswith("/avatar/"):
            payload["assistantAvatar"] = f"{base_path}{avatar}"

    return json.dumps(payload, ensure_ascii=False).encode("utf-8")


def _avatar_fallback_svg(agent: str) -> bytes:
    label = (agent[:1] or "A").upper()
    return (
        "<svg xmlns='http://www.w3.org/2000/svg' width='96' height='96' viewBox='0 0 96 96'>"
        "<rect width='96' height='96' rx='48' fill='#2f3747'/>"
        "<text x='50%' y='54%' dominant-baseline='middle' text-anchor='middle' "
        "font-family='system-ui, -apple-system, Segoe UI, Roboto, sans-serif' "
        "font-size='42' fill='#ffffff'>"
        f"{label}</text></svg>"
    ).encode("utf-8")


def require_auth(settings: Settings, authorization: str | None = Header(default=None)) -> None:
    if not settings.auth_token:
        return
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(status_code=401, detail="missing bearer token")
    token = authorization.split(" ", 1)[1].strip()
    if token != settings.auth_token:
        raise HTTPException(status_code=403, detail="invalid token")


def _parse_skill_frontmatter(skill_file: Path, fallback_slug: str) -> tuple[str, str | None]:
    name = fallback_slug
    description: str | None = None

    try:
        text = skill_file.read_text(encoding="utf-8")
    except OSError:
        return name, description

    lines = text.splitlines()
    if not lines or lines[0].strip() != "---":
        return name, description

    for line in lines[1:]:
        if line.strip() == "---":
            break
        if line.startswith("name:"):
            parsed = line.split(":", 1)[1].strip().strip('"').strip("'")
            if parsed:
                name = parsed
        elif line.startswith("description:"):
            parsed = line.split(":", 1)[1].strip().strip('"').strip("'")
            if parsed:
                description = parsed
    return name, description


def _scan_global_skills(settings: Settings) -> list[SkillItem]:
    root = Path(settings.global_skills_dir)
    if not root.exists() or not root.is_dir():
        return []

    out: list[SkillItem] = []
    for skill_file in sorted(root.glob("*/SKILL.md")):
        slug = skill_file.parent.name
        name, description = _parse_skill_frontmatter(skill_file, slug)
        out.append(
            SkillItem(
                slug=slug,
                name=name,
                description=description,
                path=str(skill_file),
                scope="global",
            )
        )
    return out


def _scan_workspace_skills(settings: Settings, agent_slug: str | None = None) -> list[WorkspaceSkillGroup]:
    root = Path(settings.agent_homes_dir)
    if not root.exists() or not root.is_dir():
        return []

    groups: list[WorkspaceSkillGroup] = []
    agent_dirs = sorted(path for path in root.iterdir() if path.is_dir())

    for agent_dir in agent_dirs:
        slug = agent_dir.name
        if agent_slug and slug != agent_slug:
            continue

        skills_dir = agent_dir / "skills"
        skill_items: list[SkillItem] = []
        if skills_dir.exists() and skills_dir.is_dir():
            for skill_file in sorted(skills_dir.glob("*/SKILL.md")):
                s_slug = skill_file.parent.name
                s_name, s_desc = _parse_skill_frontmatter(skill_file, s_slug)
                skill_items.append(
                    SkillItem(
                        slug=s_slug,
                        name=s_name,
                        description=s_desc,
                        path=str(skill_file),
                        scope="workspace",
                    )
                )

        groups.append(WorkspaceSkillGroup(agent_slug=slug, skills=skill_items))
    return groups


def _scan_agent_slugs(settings: Settings) -> list[str]:
    root = Path(settings.agent_homes_dir)
    if not root.exists() or not root.is_dir():
        return []
    return sorted(path.name for path in root.iterdir() if path.is_dir())


def _known_agent_slugs(settings: Settings) -> set[str]:
    root = Path(settings.agent_homes_dir)
    if not root.exists() or not root.is_dir():
        return set()
    return {path.name for path in root.iterdir() if path.is_dir()}


async def _forward_agent_control(
    settings: Settings,
    *,
    agent: str,
    action: str,
) -> AgentControlActionOut:
    base_url = (settings.agent_controller_url or "").strip().rstrip("/")
    if not base_url:
        raise HTTPException(status_code=503, detail="agent controller is not configured")

    target_url = f"{base_url}/v1/containers/{urllib.parse.quote(agent, safe='')}/control"
    timeout = max(1.0, float(settings.agent_controller_timeout_seconds or 5.0))

    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            resp = await client.post(target_url, json={"action": action})
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"agent controller unavailable: {exc.__class__.__name__}")

    try:
        payload = resp.json()
    except Exception:
        payload = {}

    if resp.status_code >= 400:
        detail = payload.get("detail") if isinstance(payload, dict) else "controller returned error"
        raise HTTPException(status_code=resp.status_code, detail=str(detail or "controller returned error"))

    if not isinstance(payload, dict):
        raise HTTPException(status_code=502, detail="invalid controller response")

    return AgentControlActionOut(
        ok=bool(payload.get("ok", False)),
        agent=str(payload.get("agent") or agent),
        action=str(payload.get("action") or action),
        container=str(payload.get("container") or f"openclaw-{agent}"),
        status=str(payload.get("status") or "unknown"),
        detail=str(payload.get("detail") or ""),
    )


def _validate_handoff_payload(payload: dict, known_agents: set[str]) -> list[str]:
    errors: list[str] = []

    target_agent = str(payload.get("to") or "").strip()
    if not target_agent:
        errors.append("payload.to is required")
    elif known_agents and target_agent not in known_agents:
        errors.append(f"payload.to agent not found: {target_agent}")

    for field in ("problem", "context", "expected_output"):
        value = str(payload.get(field) or "").strip()
        if not value:
            errors.append(f"payload.{field} is required")

    artifact_refs = payload.get("artifact_refs")
    if not isinstance(artifact_refs, list) or not artifact_refs:
        errors.append("payload.artifact_refs must be a non-empty list")
    elif not all(isinstance(item, str) and item.strip() for item in artifact_refs):
        errors.append("payload.artifact_refs must contain non-empty strings")

    if not isinstance(payload.get("review_gate"), bool):
        errors.append("payload.review_gate must be boolean")

    return errors


def _parse_iso_datetime(value: str | None) -> datetime | None:
    text = str(value or "").strip()
    if not text:
        return None
    if text.endswith("Z"):
        text = text[:-1] + "+00:00"
    try:
        return datetime.fromisoformat(text)
    except ValueError:
        return None


def _host_port_from_url(raw_url: str, default_port: int) -> tuple[str | None, int]:
    parsed = urlparse(raw_url)
    host = parsed.hostname
    port = parsed.port or default_port
    return host, int(port)


async def _probe_tcp(host: str, port: int, *, timeout_seconds: float = 1.2) -> tuple[bool, int | None, str]:
    started = time.perf_counter()
    try:
        _, writer = await asyncio.wait_for(asyncio.open_connection(host, int(port)), timeout=timeout_seconds)
        writer.close()
        await writer.wait_closed()
        latency_ms = int((time.perf_counter() - started) * 1000)
        return True, latency_ms, "connected"
    except Exception as exc:
        return False, None, f"{exc.__class__.__name__}: {exc}"


async def _probe_http(url: str, *, timeout_seconds: float = 1.5) -> tuple[bool, int | None, str]:
    started = time.perf_counter()
    try:
        async with httpx.AsyncClient(timeout=timeout_seconds, follow_redirects=True) as client:
            response = await client.get(url)
        latency_ms = int((time.perf_counter() - started) * 1000)
        ok = int(response.status_code) < 500
        return ok, latency_ms, f"HTTP {response.status_code}"
    except Exception as exc:
        return False, None, f"{exc.__class__.__name__}: {exc}"


def _line_usage_from_session_event(line_obj: dict) -> tuple[datetime | None, dict | None]:
    message = line_obj.get("message") if isinstance(line_obj.get("message"), dict) else {}
    payload = line_obj.get("payload") if isinstance(line_obj.get("payload"), dict) else {}

    usage = None
    for candidate in (message.get("usage"), line_obj.get("usage"), payload.get("usage")):
        if isinstance(candidate, dict):
            usage = candidate
            break

    if not isinstance(usage, dict):
        return None, None

    mapped_usage = {
        "input": _safe_int(usage.get("input") or usage.get("input_tokens")),
        "output": _safe_int(usage.get("output") or usage.get("output_tokens")),
        "cacheRead": _safe_int(usage.get("cacheRead") or usage.get("cache_read_tokens")),
        "cacheWrite": _safe_int(usage.get("cacheWrite") or usage.get("cache_write_tokens")),
        "totalTokens": _safe_int(usage.get("totalTokens") or usage.get("total_tokens")),
        "cost": {
            "total": _safe_float(
                (usage.get("cost") or {}).get("total") if isinstance(usage.get("cost"), dict) else usage.get("total_cost")
            ),
        },
    }

    if mapped_usage["totalTokens"] <= 0:
        mapped_usage["totalTokens"] = (
            mapped_usage["input"]
            + mapped_usage["output"]
            + mapped_usage["cacheRead"]
            + mapped_usage["cacheWrite"]
        )

    event_ts = _parse_iso_datetime(
        line_obj.get("timestamp")
        or line_obj.get("created_at")
        or line_obj.get("createdAt")
        or line_obj.get("updated_at")
        or line_obj.get("updatedAt")
    )
    if event_ts is None:
        ts_value = line_obj.get("ts")
        if isinstance(ts_value, (int, float)):
            try:
                epoch = float(ts_value)
                if epoch > 10_000_000_000:
                    epoch = epoch / 1000.0
                event_ts = datetime.fromtimestamp(epoch)
            except (ValueError, OSError):
                event_ts = None

    return event_ts, mapped_usage


def _line_usage_from_cron_event(line_obj: dict) -> tuple[datetime | None, dict | None]:
    usage = line_obj.get("usage")
    if not isinstance(usage, dict):
        payload = line_obj.get("payload") if isinstance(line_obj.get("payload"), dict) else {}
        usage = payload.get("usage")
    if not isinstance(usage, dict):
        return None, None

    ts_ms_raw = line_obj.get("ts")
    event_ts = None
    if isinstance(ts_ms_raw, (int, float)):
        try:
            event_ts = datetime.fromtimestamp(float(ts_ms_raw) / 1000.0)
        except (ValueError, OSError):
            event_ts = None

    mapped_usage = {
        "input": int(usage.get("input_tokens") or 0),
        "output": int(usage.get("output_tokens") or 0),
        "cacheRead": int(usage.get("cache_read_tokens") or 0),
        "cacheWrite": int(usage.get("cache_write_tokens") or 0),
        "totalTokens": int(usage.get("total_tokens") or 0),
        "cost": {
            "total": float(usage.get("total_cost") or 0.0),
        },
    }
    return event_ts, mapped_usage


def _safe_int(value) -> int:
    try:
        return int(value or 0)
    except (TypeError, ValueError):
        return 0


def _safe_float(value) -> float:
    try:
        return float(value or 0.0)
    except (TypeError, ValueError):
        return 0.0


def _accumulate_usage(aggregate: dict, usage: dict, *, window_key: str) -> None:
    aggregate[f"input_tokens_{window_key}"] += _safe_int(usage.get("input"))
    aggregate[f"output_tokens_{window_key}"] += _safe_int(usage.get("output"))
    aggregate[f"cache_read_tokens_{window_key}"] += _safe_int(usage.get("cacheRead"))
    aggregate[f"cache_write_tokens_{window_key}"] += _safe_int(usage.get("cacheWrite"))

    total_tokens = _safe_int(usage.get("totalTokens"))
    if total_tokens <= 0:
        total_tokens = (
            _safe_int(usage.get("input"))
            + _safe_int(usage.get("output"))
            + _safe_int(usage.get("cacheRead"))
            + _safe_int(usage.get("cacheWrite"))
        )
    aggregate[f"total_tokens_{window_key}"] += total_tokens

    usage_cost = usage.get("cost") if isinstance(usage.get("cost"), dict) else {}
    if isinstance(usage_cost, dict):
        total_cost = _safe_float(usage_cost.get("total"))
        aggregate[f"total_cost_{window_key}"] += total_cost
        if total_cost <= 0 and total_tokens > 0:
            aggregate["missing_cost_entries_window"] += 1


def _iter_agent_usage_lines(agent_dir: Path):
    yielded: set[str] = set()

    session_patterns = [
        "agents/*/sessions/**/*.jsonl",
        "agents/*/sessions/**/*.jsonl.reset.*",
        "agents/*/sessions/**/*.jsonl.deleted.*",
        "sessions/**/*.jsonl",
        "sessions/**/*.jsonl.reset.*",
        "sessions/**/*.jsonl.deleted.*",
    ]
    cron_patterns = [
        "cron/runs/*.jsonl",
        "cron/runs/**/*.jsonl",
    ]

    for pattern in session_patterns:
        for file_path in agent_dir.glob(pattern):
            if file_path.is_file():
                key = str(file_path)
                if key not in yielded:
                    yielded.add(key)
                    yield file_path, "session"

    for pattern in cron_patterns:
        for file_path in agent_dir.glob(pattern):
            if file_path.is_file():
                key = str(file_path)
                if key not in yielded:
                    yielded.add(key)
                    yield file_path, "cron"


def _collect_agent_usage_snapshot(settings: Settings, days: int) -> list[AgentUsageSnapshotOut]:
    root = Path(settings.agent_homes_dir)
    if not root.exists() or not root.is_dir():
        return []

    window_seconds = max(1, min(days, 90)) * 24 * 60 * 60
    window_start_epoch = time.time() - window_seconds
    day_24h_epoch = time.time() - 24 * 60 * 60

    out: list[AgentUsageSnapshotOut] = []

    for agent_dir in sorted(path for path in root.iterdir() if path.is_dir()):
        aggregate = {
            "agent": agent_dir.name,
            "input_tokens_24h": 0,
            "output_tokens_24h": 0,
            "cache_read_tokens_24h": 0,
            "cache_write_tokens_24h": 0,
            "total_tokens_24h": 0,
            "total_cost_24h": 0.0,
            "input_tokens_window": 0,
            "output_tokens_window": 0,
            "cache_read_tokens_window": 0,
            "cache_write_tokens_window": 0,
            "total_tokens_window": 0,
            "total_cost_window": 0.0,
            "missing_cost_entries_window": 0,
            "days": max(1, min(days, 90)),
        }

        for file_path, source_type in _iter_agent_usage_lines(agent_dir):
            try:
                with file_path.open("r", encoding="utf-8") as handle:
                    for raw_line in handle:
                        line = raw_line.strip()
                        if not line:
                            continue
                        try:
                            obj = json.loads(line)
                        except json.JSONDecodeError:
                            continue

                        if source_type == "session":
                            event_dt, usage = _line_usage_from_session_event(obj)
                        else:
                            event_dt, usage = _line_usage_from_cron_event(obj)

                        if usage is None:
                            continue

                        include_24h = False
                        include_window = True
                        if event_dt is not None:
                            event_epoch = event_dt.timestamp()
                            include_window = event_epoch >= window_start_epoch
                            include_24h = event_epoch >= day_24h_epoch

                        if include_window:
                            _accumulate_usage(aggregate, usage, window_key="window")
                        if include_24h:
                            _accumulate_usage(aggregate, usage, window_key="24h")
            except OSError:
                continue

        out.append(AgentUsageSnapshotOut(**aggregate))

    return out


def _get_agent_usage_snapshot(settings: Settings, days: int) -> list[AgentUsageSnapshotOut]:
    clamped_days = max(1, min(int(days), 90))
    now_ts = time.time()

    cached_days = _AGENT_USAGE_CACHE.get("days")
    cached_at = _safe_float(_AGENT_USAGE_CACHE.get("generated_at"))
    cached_data = _AGENT_USAGE_CACHE.get("data")

    if (
        cached_days == clamped_days
        and isinstance(cached_data, list)
        and now_ts - cached_at <= USAGE_CACHE_TTL_SECONDS
    ):
        return cached_data  # type: ignore[return-value]

    fresh = _collect_agent_usage_snapshot(settings, clamped_days)
    _AGENT_USAGE_CACHE["days"] = clamped_days
    _AGENT_USAGE_CACHE["generated_at"] = now_ts
    _AGENT_USAGE_CACHE["data"] = fresh
    return fresh


def create_app() -> FastAPI:
    settings = load_settings()

    engine = create_engine(settings.database_url)
    session_factory = create_session_factory(engine)
    redis = Redis.from_url(settings.redis_url, decode_responses=True)

    app = FastAPI(title="Mission Control API", version="0.1.0")

    async def get_session():
        async with session_factory() as session:
            yield session

    async def publish_validation_result(
        *,
        accepted: bool,
        body: EventIn,
        errors: list[str],
        details: dict | None = None,
    ) -> None:
        await publish_event(
            redis,
            stream_key=settings.redis_stream_key,
            event={
                "id": str(uuid4()),
                "type": "event.validation",
                "agent": body.agent,
                "task_id": str(body.task_id) if body.task_id else None,
                "payload": {
                    "event_type": body.type,
                    "accepted": accepted,
                    "errors": errors,
                    "details": details or {},
                },
                "created_at": datetime.utcnow().isoformat() + "Z",
            },
        )

    @app.get("/health", response_model=Health)
    async def healthcheck() -> Health:
        return Health(ok=True)

    @app.get("/v1/observability/summary", response_model=ObservabilitySummaryOut)
    async def get_observability_summary(
        window_minutes: int = 5,
        heartbeat_stale_seconds: int = 180,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> ObservabilitySummaryOut:
        now = datetime.now(timezone.utc)
        window = max(1, min(int(window_minutes or 5), 60))
        stale = max(30, min(int(heartbeat_stale_seconds or 180), 3600))
        since = now - timedelta(minutes=window)

        status_code_expr = sa.cast(
            sa.func.nullif(sa.func.jsonb_extract_path_text(events.c.payload, "status_code"), ""),
            sa.Integer,
        )
        error_type_expr = sa.func.jsonb_extract_path_text(events.c.payload, "error_type")
        accepted_expr = sa.func.jsonb_extract_path_text(events.c.payload, "accepted")

        request_total_stmt = sa.select(sa.func.count()).where(
            events.c.created_at >= since,
            events.c.type == "chat.gateway.access",
        )
        request_total = int((await session.execute(request_total_stmt)).scalar_one() or 0)

        error_total_stmt = sa.select(sa.func.count()).where(
            events.c.created_at >= since,
            sa.or_(
                sa.and_(events.c.type == "event.validation", accepted_expr == "false"),
                events.c.type.like("%.error"),
                status_code_expr >= 400,
                sa.and_(error_type_expr.is_not(None), error_type_expr != ""),
            ),
        )
        error_total = int((await session.execute(error_total_stmt)).scalar_one() or 0)

        events_total_stmt = sa.select(sa.func.count()).where(events.c.created_at >= since)
        events_total = int((await session.execute(events_total_stmt)).scalar_one() or 0)

        tasks_done_stmt = sa.select(sa.func.count()).where(
            tasks.c.updated_at >= since,
            tasks.c.status == "DONE",
        )
        tasks_done_total = int((await session.execute(tasks_done_stmt)).scalar_one() or 0)

        task_backlog_stmt = sa.select(sa.func.count()).where(tasks.c.status != "DONE")
        task_backlog_total = int((await session.execute(task_backlog_stmt)).scalar_one() or 0)

        try:
            event_backlog_total = int(await redis.xlen(settings.redis_stream_key))
        except Exception:
            event_backlog_total = 0

        known_agents = sorted(_known_agent_slugs(settings))
        total_agents = len(known_agents)
        healthy_agents = 0

        if known_agents:
            heartbeats_stmt = (
                sa.select(events.c.agent, sa.func.max(events.c.created_at).label("last_seen"))
                .where(
                    events.c.type == "agent.heartbeat",
                    events.c.agent.is_not(None),
                    events.c.agent.in_(known_agents),
                )
                .group_by(events.c.agent)
            )
            rows = (await session.execute(heartbeats_stmt)).all()
            cutoff = now - timedelta(seconds=stale)
            for row in rows:
                last_seen = row.last_seen
                if isinstance(last_seen, datetime) and last_seen >= cutoff:
                    healthy_agents += 1

        denominator = request_total if request_total > 0 else max(events_total, 1)
        error_rate = float(error_total) / float(denominator)

        return ObservabilitySummaryOut(
            generated_at=now,
            window_minutes=window,
            request_total=request_total,
            error_total=error_total,
            error_rate=error_rate,
            event_throughput_per_min=float(events_total) / float(window),
            task_throughput_per_min=float(tasks_done_total) / float(window),
            event_backlog_total=event_backlog_total,
            task_backlog_total=task_backlog_total,
            healthy_agents=healthy_agents,
            total_agents=total_agents,
            agent_health_ratio=(float(healthy_agents) / float(total_agents)) if total_agents > 0 else 0.0,
            heartbeat_stale_seconds=stale,
        )

    @app.get("/v1/observability/container-health", response_model=ContainerHealthSummaryOut)
    async def get_container_health_summary(
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> ContainerHealthSummaryOut:
        now = datetime.now(timezone.utc)
        signals: list[HealthSignalOut] = []

        redis_host, redis_port = _host_port_from_url(settings.redis_url, 6379)
        db_host, db_port = _host_port_from_url(settings.database_url, 5432)

        compose_ok = 0
        compose_total = 0

        compose_total += 1
        redis_started = time.perf_counter()
        try:
            redis_ok = bool(await redis.ping())
            redis_latency = int((time.perf_counter() - redis_started) * 1000)
            if redis_ok:
                compose_ok += 1
            signals.append(
                HealthSignalOut(
                    name="redis.ping",
                    source="compose",
                    target=f"{redis_host or 'redis'}:{redis_port}",
                    ok=redis_ok,
                    latency_ms=redis_latency,
                    detail="PONG" if redis_ok else "No PONG",
                )
            )
        except Exception as exc:
            signals.append(
                HealthSignalOut(
                    name="redis.ping",
                    source="compose",
                    target=f"{redis_host or 'redis'}:{redis_port}",
                    ok=False,
                    detail=f"{exc.__class__.__name__}: {exc}",
                )
            )

        compose_total += 1
        db_started = time.perf_counter()
        try:
            db_ok = bool((await session.execute(sa.select(sa.literal(1)))).scalar_one() == 1)
            db_latency = int((time.perf_counter() - db_started) * 1000)
            if db_ok:
                compose_ok += 1
            signals.append(
                HealthSignalOut(
                    name="postgres.select_1",
                    source="compose",
                    target=f"{db_host or 'postgres'}:{db_port}",
                    ok=db_ok,
                    latency_ms=db_latency,
                    detail="SELECT 1",
                )
            )
        except Exception as exc:
            signals.append(
                HealthSignalOut(
                    name="postgres.select_1",
                    source="compose",
                    target=f"{db_host or 'postgres'}:{db_port}",
                    ok=False,
                    detail=f"{exc.__class__.__name__}: {exc}",
                )
            )

        port_targets: list[tuple[str, str, int]] = [
            ("mission-control-api", "mission-control-api", 9090),
            ("mission-control-ui", "mission-control-ui", 9090),
            ("mission-control-gateway", "mission-control-gateway", 80),
        ]
        for slug in sorted(_known_agent_slugs(settings)):
            port_targets.append((f"openclaw-{slug}", f"openclaw-{slug}", settings.chat_upstream_port))

        port_ok = 0
        port_total = len(port_targets)
        for name, host, port in port_targets:
            ok, latency_ms, detail = await _probe_tcp(host, port)
            if ok:
                port_ok += 1
            signals.append(
                HealthSignalOut(
                    name=name,
                    source="port",
                    target=f"{host}:{port}",
                    ok=ok,
                    latency_ms=latency_ms,
                    detail=detail,
                )
            )

        http_targets: list[tuple[str, str]] = [
            ("mission-control-api.health", "http://mission-control-api:9090/health"),
            ("mission-control-gateway.root", "http://mission-control-gateway/"),
            ("mission-control-ui.root", "http://mission-control-ui:9090/"),
        ]

        http_ok = 0
        http_total = len(http_targets)
        for name, url in http_targets:
            ok, latency_ms, detail = await _probe_http(url)
            if ok:
                http_ok += 1
            signals.append(
                HealthSignalOut(
                    name=name,
                    source="http",
                    target=url,
                    ok=ok,
                    latency_ms=latency_ms,
                    detail=detail,
                )
            )

        overall_ok = compose_ok + port_ok + http_ok
        overall_total = compose_total + port_total + http_total

        return ContainerHealthSummaryOut(
            generated_at=now,
            compose_ok=compose_ok,
            compose_total=compose_total,
            port_ok=port_ok,
            port_total=port_total,
            http_ok=http_ok,
            http_total=http_total,
            overall_ok=overall_ok,
            overall_total=overall_total,
            overall_ratio=(float(overall_ok) / float(overall_total)) if overall_total > 0 else 0.0,
            signals=signals,
        )

    @app.get("/v1/skills/global", response_model=list[SkillItem])
    async def get_global_skills(
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
    ) -> list[SkillItem]:
        return _scan_global_skills(settings)

    @app.get("/v1/skills/workspace", response_model=list[WorkspaceSkillGroup])
    async def get_workspace_skills(
        agent_slug: str | None = None,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
    ) -> list[WorkspaceSkillGroup]:
        return _scan_workspace_skills(settings, agent_slug=agent_slug)

    @app.get("/v1/usage/agents", response_model=list[AgentUsageSnapshotOut])
    async def get_agent_usage_snapshot(
        days: int = 7,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
    ) -> list[AgentUsageSnapshotOut]:
        return _get_agent_usage_snapshot(settings, days)

    @app.post("/v1/knowledge/sources/import", response_model=KnowledgeSourceOut)
    async def import_knowledge_source(
        body: KnowledgeSourceImportIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeSourceOut:
        root = Path(settings.knowledge_raw_sources_dir).resolve()
        if not root.exists() or not root.is_dir():
            raise HTTPException(status_code=503, detail=f"knowledge source root unavailable: {root}")

        target = _safe_resolve_under(root, body.relative_path)
        if not target.exists() or not target.is_file():
            raise HTTPException(status_code=422, detail=f"file not found under knowledge root: {body.relative_path}")

        storage_path = _relative_storage_path(root, target)
        existing = (
            await session.execute(
                sa.select(
                    knowledge_sources.c.id,
                    knowledge_sources.c.source_type,
                    knowledge_sources.c.title,
                    knowledge_sources.c.external_uri,
                    knowledge_sources.c.storage_path,
                    knowledge_sources.c.checksum_sha256,
                    knowledge_sources.c.mime_type,
                    knowledge_sources.c.owner,
                    knowledge_sources.c.version_label,
                    knowledge_sources.c.status,
                    knowledge_sources.c.meta,
                    knowledge_sources.c.collected_at,
                    knowledge_sources.c.updated_at,
                ).where(knowledge_sources.c.storage_path == storage_path)
            )
        ).first()
        if existing:
            return _knowledge_source_row_to_out(existing)

        checksum = _sha256sum(target)
        mime_type = _guess_mime_type(target)
        source_id = uuid4()
        now = datetime.utcnow()

        stmt = (
            knowledge_sources.insert()
            .values(
                id=source_id,
                source_type=(body.source_type or "file").strip() or "file",
                title=(body.title or target.stem).strip() or target.name,
                external_uri=body.external_uri,
                storage_path=storage_path,
                checksum_sha256=checksum,
                mime_type=mime_type,
                owner=body.owner,
                version_label=body.version_label,
                status="active",
                meta=body.meta,
                collected_at=now,
                updated_at=now,
            )
            .returning(
                knowledge_sources.c.id,
                knowledge_sources.c.source_type,
                knowledge_sources.c.title,
                knowledge_sources.c.external_uri,
                knowledge_sources.c.storage_path,
                knowledge_sources.c.checksum_sha256,
                knowledge_sources.c.mime_type,
                knowledge_sources.c.owner,
                knowledge_sources.c.version_label,
                knowledge_sources.c.status,
                knowledge_sources.c.meta,
                knowledge_sources.c.collected_at,
                knowledge_sources.c.updated_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()
        return _knowledge_source_row_to_out(row)

    @app.get("/v1/knowledge/sources", response_model=list[KnowledgeSourceOut])
    async def list_knowledge_sources(
        source_type: str | None = None,
        status: str | None = None,
        limit: int = 100,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeSourceOut]:
        stmt = sa.select(
            knowledge_sources.c.id,
            knowledge_sources.c.source_type,
            knowledge_sources.c.title,
            knowledge_sources.c.external_uri,
            knowledge_sources.c.storage_path,
            knowledge_sources.c.checksum_sha256,
            knowledge_sources.c.mime_type,
            knowledge_sources.c.owner,
            knowledge_sources.c.version_label,
            knowledge_sources.c.status,
            knowledge_sources.c.meta,
            knowledge_sources.c.collected_at,
            knowledge_sources.c.updated_at,
        )
        if source_type:
            stmt = stmt.where(knowledge_sources.c.source_type == source_type)
        if status:
            stmt = stmt.where(knowledge_sources.c.status == status)
        stmt = stmt.order_by(knowledge_sources.c.updated_at.desc()).limit(min(max(limit, 1), 500))
        rows = (await session.execute(stmt)).all()
        return [_knowledge_source_row_to_out(row) for row in rows]

    @app.get("/v1/knowledge/sources/{source_id}", response_model=KnowledgeSourceOut)
    async def get_knowledge_source(
        source_id: UUID,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeSourceOut:
        row = (
            await session.execute(
                sa.select(
                    knowledge_sources.c.id,
                    knowledge_sources.c.source_type,
                    knowledge_sources.c.title,
                    knowledge_sources.c.external_uri,
                    knowledge_sources.c.storage_path,
                    knowledge_sources.c.checksum_sha256,
                    knowledge_sources.c.mime_type,
                    knowledge_sources.c.owner,
                    knowledge_sources.c.version_label,
                    knowledge_sources.c.status,
                    knowledge_sources.c.meta,
                    knowledge_sources.c.collected_at,
                    knowledge_sources.c.updated_at,
                ).where(knowledge_sources.c.id == source_id)
            )
        ).first()
        if not row:
            raise HTTPException(status_code=404, detail="knowledge source not found")
        return _knowledge_source_row_to_out(row)

    @app.post("/v1/knowledge/sources/scan", response_model=KnowledgeSourceScanOut)
    async def scan_knowledge_sources(
        body: KnowledgeSourceScanIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeSourceScanOut:
        root = Path(settings.knowledge_raw_sources_dir).resolve()
        if not root.exists() or not root.is_dir():
            raise HTTPException(status_code=503, detail=f"knowledge source root unavailable: {root}")

        scan_root = root if not body.subdir.strip() else _safe_resolve_under(root, body.subdir)
        if not scan_root.exists() or not scan_root.is_dir():
            raise HTTPException(status_code=422, detail=f"scan directory not found: {scan_root}")

        max_files = min(max(int(body.max_files or 200), 1), 2000)
        extension_filter = {str(ext).lower().lstrip(".") for ext in body.include_extensions if str(ext).strip()}

        scanned = 0
        imported = 0
        skipped_existing = 0
        skipped_invalid = 0
        items: list[KnowledgeSourceOut] = []

        now = datetime.utcnow()
        for file_path in sorted(path for path in scan_root.rglob("*") if path.is_file()):
            if scanned >= max_files:
                break
            scanned += 1

            if extension_filter:
                ext = file_path.suffix.lower().lstrip(".")
                if ext not in extension_filter:
                    skipped_invalid += 1
                    continue

            storage_path = _relative_storage_path(root, file_path)
            exists = (
                await session.execute(
                    sa.select(knowledge_sources.c.id).where(knowledge_sources.c.storage_path == storage_path)
                )
            ).first()
            if exists:
                skipped_existing += 1
                continue

            checksum = _sha256sum(file_path)
            mime_type = _guess_mime_type(file_path)
            out = KnowledgeSourceOut(
                id=uuid4(),
                source_type=(body.source_type or "file").strip() or "file",
                title=file_path.stem,
                external_uri=None,
                storage_path=storage_path,
                checksum_sha256=checksum,
                mime_type=mime_type,
                owner=body.owner,
                version_label=body.version_label,
                status="active",
                meta={"scanned_from": str(scan_root.relative_to(root)) if scan_root != root else "."},
                collected_at=now,
                updated_at=now,
            )
            items.append(out)

            if body.dry_run:
                imported += 1
                continue

            await session.execute(
                knowledge_sources.insert().values(
                    id=out.id,
                    source_type=out.source_type,
                    title=out.title,
                    external_uri=out.external_uri,
                    storage_path=out.storage_path,
                    checksum_sha256=out.checksum_sha256,
                    mime_type=out.mime_type,
                    owner=out.owner,
                    version_label=out.version_label,
                    status=out.status,
                    meta=out.meta,
                    collected_at=out.collected_at,
                    updated_at=out.updated_at,
                )
            )
            imported += 1

        if not body.dry_run:
            await session.commit()

        return KnowledgeSourceScanOut(
            scanned=scanned,
            imported=imported,
            skipped_existing=skipped_existing,
            skipped_invalid=skipped_invalid,
            dry_run=body.dry_run,
            items=items[: min(len(items), 200)],
        )

    @app.post("/v1/knowledge/sources/{source_id}/chunk", response_model=KnowledgeSourceChunkOut)
    async def chunk_knowledge_source(
        source_id: UUID,
        body: KnowledgeSourceChunkIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeSourceChunkOut:
        try:
            embedding_config = _knowledge_embedding_runtime_config(settings)
        except ValueError as exc:
            raise HTTPException(status_code=503, detail=str(exc)) from exc

        source_row = (
            await session.execute(
                sa.select(
                    knowledge_sources.c.id,
                    knowledge_sources.c.storage_path,
                    knowledge_sources.c.title,
                    knowledge_sources.c.source_type,
                    knowledge_sources.c.meta,
                ).where(knowledge_sources.c.id == source_id)
            )
        ).first()
        if not source_row:
            raise HTTPException(status_code=404, detail="knowledge source not found")

        root = Path(settings.knowledge_raw_sources_dir).resolve()
        source_path = _safe_resolve_under(root, source_row.storage_path)
        if not source_path.exists() or not source_path.is_file():
            raise HTTPException(status_code=422, detail="knowledge source file not found")

        ocr_languages = _normalize_ocr_languages(body.ocr_languages)
        max_pdf_pages = _normalize_max_pdf_pages(body.max_pdf_pages)

        try:
            raw_text, extraction_method, extraction_warnings, extraction_meta = _extract_source_text(
                source_path,
                ocr_enabled=body.ocr_enabled,
                ocr_languages=ocr_languages,
                max_pdf_pages=max_pdf_pages,
            )
        except Exception as exc:
            now = datetime.utcnow()
            source_meta = dict(source_row.meta or {})
            source_meta.update(
                {
                    "parse_status": "failed",
                    "parse_error": str(exc),
                    "ocr_enabled": bool(body.ocr_enabled),
                    "ocr_languages": ocr_languages,
                    "max_pdf_pages": max_pdf_pages,
                    "ocr_used": source_path.suffix.lower().lstrip(".") in KNOWLEDGE_OCR_CANDIDATE_EXTENSIONS,
                    "ocr_fallback_used": False,
                }
            )
            await session.execute(
                knowledge_sources.update()
                .where(knowledge_sources.c.id == source_id)
                .values(
                    meta=source_meta,
                    updated_at=now,
                )
            )
            await session.commit()
            raise HTTPException(status_code=422, detail=str(exc)) from exc

        chunk_list = _split_text_chunks(
            raw_text,
            chunk_chars=body.chunk_chars,
            overlap=body.chunk_overlap,
            max_chunks=min(max(int(body.max_chunks or 200), 1), 1000),
        )
        now = datetime.utcnow()
        created = 0
        skipped_existing = 0
        items: list[KnowledgeUnitOut] = []
        embedding_failed_count = 0
        embedding_status = "disabled"
        embedding_error: str | None = None
        embedding_model_name: str | None = None
        embedding_dimensions: int | None = None

        for index, chunk in enumerate(chunk_list, start=1):
            unit_key = f"source:{source_id}:chunk:{index:04d}"
            exists = (
                await session.execute(sa.select(knowledge_units.c.id).where(knowledge_units.c.unit_key == unit_key))
            ).first()
            if exists:
                skipped_existing += 1
                continue

            unit_out = KnowledgeUnitOut(
                id=uuid4(),
                source_id=source_id,
                unit_key=unit_key,
                title=f"{source_row.title} #{index}",
                content=chunk,
                content_sha256=hashlib.sha256(chunk.encode("utf-8")).hexdigest(),
                tags=list({*body.tags, "auto-chunk"}),
                agent_scope=body.agent_scope,
                risk_level=(body.risk_level or "normal").strip().lower() or "normal",
                status="active",
                lifecycle_stage="active",
                meta={
                    "generated_by": "source_chunker",
                    "chunk_index": index,
                    "source_storage_path": source_row.storage_path,
                    "owner": body.owner,
                },
                created_at=now,
                updated_at=now,
            )
            items.append(unit_out)
            created += 1

        embedding_vectors: list[list[float]] = []
        if body.dry_run:
            embedding_status = "skipped_dry_run"
        elif not embedding_config:
            embedding_status = "disabled"
        elif not items:
            embedding_status = "skipped_no_new_units"
            embedding_model_name = str(embedding_config["model"])
            if embedding_config.get("dimensions") is not None:
                embedding_dimensions = int(embedding_config["dimensions"])
        else:
            try:
                embedding_vectors, embedding_model_name, embedding_dimensions = await _request_knowledge_embeddings(
                    [item.content for item in items],
                    config=embedding_config,
                )
                embedding_status = "success"
            except Exception as exc:
                embedding_status = "failed"
                embedding_error = str(exc)
                embedding_failed_count = len(items)
                embedding_model_name = str(embedding_config["model"])

        for index, unit_out in enumerate(items):
            unit_meta = dict(unit_out.meta or {})
            unit_meta.update(
                {
                    "embedding_enabled": bool(embedding_config),
                    "embedding_status": embedding_status,
                }
            )
            if embedding_model_name:
                unit_meta["embedding_model"] = embedding_model_name
            if embedding_dimensions is not None:
                unit_meta["embedding_dimensions"] = embedding_dimensions
            if embedding_error:
                unit_meta["embedding_error"] = embedding_error
            if embedding_vectors:
                unit_meta["embedding_updated_at"] = now.isoformat() + "Z"
            unit_out.meta = unit_meta

            if not body.dry_run:
                await session.execute(
                    knowledge_units.insert().values(
                        id=unit_out.id,
                        source_id=unit_out.source_id,
                        unit_key=unit_out.unit_key,
                        title=unit_out.title,
                        content=unit_out.content,
                        content_sha256=unit_out.content_sha256,
                        tags=unit_out.tags,
                        agent_scope=unit_out.agent_scope,
                        risk_level=unit_out.risk_level,
                        status=unit_out.status,
                        lifecycle_stage=unit_out.lifecycle_stage,
                        superseded_by_unit_id=unit_out.superseded_by_unit_id,
                        retired_at=unit_out.retired_at,
                        meta=unit_out.meta,
                        created_at=unit_out.created_at,
                        updated_at=unit_out.updated_at,
                    )
                )

                if embedding_vectors:
                    await _upsert_knowledge_unit_embedding(
                        session,
                        unit_id=unit_out.id,
                        embedding_model=str(embedding_model_name or embedding_config["model"]),
                        embedding_dimensions=int(embedding_dimensions or len(embedding_vectors[index])),
                        content_sha256=unit_out.content_sha256 or "",
                        embedding_values=embedding_vectors[index],
                        meta={
                            "source_id": str(source_id),
                            "unit_key": unit_out.unit_key,
                            "generated_by": "knowledge_chunk_embedding",
                        },
                        now=now,
                    )

        source_meta_update = dict(source_row.meta or {})
        source_meta_update.update(
            {
                "parse_status": "success",
                "parse_error": None,
                "parser": extraction_method,
                "text_chars": len(raw_text),
                "extraction_warnings": extraction_warnings,
                "ocr_enabled": bool(body.ocr_enabled),
                "ocr_languages": ocr_languages,
                "max_pdf_pages": max_pdf_pages,
                "embedding_enabled": bool(embedding_config),
                "embedding_status": embedding_status,
                "embedding_model": embedding_model_name,
                "embedding_dimensions": embedding_dimensions,
                "embedding_created_count": len(embedding_vectors),
                "embedding_failed_count": embedding_failed_count,
                "embedding_updated_at": now.isoformat() + "Z" if embedding_vectors else None,
                "embedding_error": embedding_error,
                **extraction_meta,
            }
        )
        await session.execute(
            knowledge_sources.update()
            .where(knowledge_sources.c.id == source_id)
            .values(
                meta=source_meta_update,
                updated_at=now,
            )
        )

        if not body.dry_run:
            await session.commit()

        return KnowledgeSourceChunkOut(
            source_id=source_id,
            scanned_chars=len(raw_text),
            chunks_total=len(chunk_list),
            created=created,
            skipped_existing=skipped_existing,
            dry_run=body.dry_run,
            items=items[: min(len(items), 200)],
        )

    @app.post("/v1/knowledge/units", response_model=KnowledgeUnitOut)
    async def create_knowledge_unit(
        body: KnowledgeUnitCreateIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeUnitOut:
        unit_key = str(body.unit_key or "").strip()
        if not unit_key:
            raise HTTPException(status_code=422, detail="unit_key is required")

        existing = (
            await session.execute(
                sa.select(
                    knowledge_units.c.id,
                    knowledge_units.c.source_id,
                    knowledge_units.c.unit_key,
                    knowledge_units.c.title,
                    knowledge_units.c.content,
                    knowledge_units.c.content_sha256,
                    knowledge_units.c.tags,
                    knowledge_units.c.agent_scope,
                    knowledge_units.c.risk_level,
                    knowledge_units.c.status,
                    knowledge_units.c.lifecycle_stage,
                    knowledge_units.c.superseded_by_unit_id,
                    knowledge_units.c.retired_at,
                    knowledge_units.c.meta,
                    knowledge_units.c.created_at,
                    knowledge_units.c.updated_at,
                ).where(knowledge_units.c.unit_key == unit_key)
            )
        ).first()
        if existing:
            return _knowledge_unit_row_to_out(existing)

        if body.source_id is not None:
            source_exists = (
                await session.execute(
                    sa.select(knowledge_sources.c.id).where(knowledge_sources.c.id == body.source_id)
                )
            ).first()
            if not source_exists:
                raise HTTPException(status_code=404, detail="knowledge source not found")

        now = datetime.utcnow()
        content_sha256 = hashlib.sha256((body.content or "").encode("utf-8")).hexdigest()
        stmt = (
            knowledge_units.insert()
            .values(
                id=uuid4(),
                source_id=body.source_id,
                unit_key=unit_key,
                title=(body.title or "").strip() or unit_key,
                content=body.content,
                content_sha256=content_sha256,
                tags=body.tags,
                agent_scope=body.agent_scope,
                risk_level=(body.risk_level or "normal").strip().lower() or "normal",
                status=(body.status or "active").strip().lower() or "active",
                lifecycle_stage="active",
                meta=body.meta,
                created_at=now,
                updated_at=now,
            )
            .returning(
                knowledge_units.c.id,
                knowledge_units.c.source_id,
                knowledge_units.c.unit_key,
                knowledge_units.c.title,
                knowledge_units.c.content,
                knowledge_units.c.content_sha256,
                knowledge_units.c.tags,
                knowledge_units.c.agent_scope,
                knowledge_units.c.risk_level,
                knowledge_units.c.status,
                knowledge_units.c.lifecycle_stage,
                knowledge_units.c.superseded_by_unit_id,
                knowledge_units.c.retired_at,
                knowledge_units.c.meta,
                knowledge_units.c.created_at,
                knowledge_units.c.updated_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()
        return _knowledge_unit_row_to_out(row)

    @app.get("/v1/knowledge/units", response_model=list[KnowledgeUnitOut])
    async def list_knowledge_units(
        source_id: UUID | None = None,
        status: str | None = None,
        agent_slug: str | None = None,
        risk_level: str | None = None,
        limit: int = 100,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeUnitOut]:
        stmt = sa.select(
            knowledge_units.c.id,
            knowledge_units.c.source_id,
            knowledge_units.c.unit_key,
            knowledge_units.c.title,
            knowledge_units.c.content,
            knowledge_units.c.content_sha256,
            knowledge_units.c.tags,
            knowledge_units.c.agent_scope,
            knowledge_units.c.risk_level,
            knowledge_units.c.status,
            knowledge_units.c.lifecycle_stage,
            knowledge_units.c.superseded_by_unit_id,
            knowledge_units.c.retired_at,
            knowledge_units.c.meta,
            knowledge_units.c.created_at,
            knowledge_units.c.updated_at,
        )
        if source_id is not None:
            stmt = stmt.where(knowledge_units.c.source_id == source_id)
        if status:
            stmt = stmt.where(knowledge_units.c.status == status)
        if agent_slug:
            stmt = stmt.where(
                sa.or_(
                    sa.func.cardinality(knowledge_units.c.agent_scope) == 0,
                    knowledge_units.c.agent_scope.any(agent_slug),
                )
            )
        if risk_level:
            requested_rank = _risk_rank(risk_level)
            allowed_levels = [name for name, rank in KNOWLEDGE_RISK_ORDER.items() if rank <= requested_rank]
            stmt = stmt.where(knowledge_units.c.risk_level.in_(allowed_levels))
        stmt = stmt.order_by(knowledge_units.c.updated_at.desc()).limit(min(max(limit, 1), 500))
        rows = (await session.execute(stmt)).all()
        return [_knowledge_unit_row_to_out(row) for row in rows]

    @app.post("/v1/knowledge/validations", response_model=KnowledgeValidationOut)
    async def create_knowledge_validation(
        body: KnowledgeValidationCreateIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationOut:
        unit_exists = (
            await session.execute(sa.select(knowledge_units.c.id).where(knowledge_units.c.id == body.unit_id))
        ).first()
        if not unit_exists:
            raise HTTPException(status_code=404, detail="knowledge unit not found")

        now = datetime.utcnow()
        stmt = (
            knowledge_validations.insert()
            .values(
                id=uuid4(),
                unit_id=body.unit_id,
                validator=(body.validator or "").strip(),
                validation_status=(body.validation_status or "approved").strip().lower() or "approved",
                validated_at=now,
                expires_at=body.expires_at,
                confidence=body.confidence,
                notes=body.notes,
                meta=body.meta,
                created_at=now,
            )
            .returning(
                knowledge_validations.c.id,
                knowledge_validations.c.unit_id,
                knowledge_validations.c.validator,
                knowledge_validations.c.validation_status,
                knowledge_validations.c.validated_at,
                knowledge_validations.c.expires_at,
                knowledge_validations.c.confidence,
                knowledge_validations.c.notes,
                knowledge_validations.c.meta,
                knowledge_validations.c.created_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()
        return _knowledge_validation_row_to_out(row)

    @app.get("/v1/knowledge/validations", response_model=list[KnowledgeValidationOut])
    async def list_knowledge_validations(
        unit_id: UUID | None = None,
        limit: int = 100,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeValidationOut]:
        stmt = sa.select(
            knowledge_validations.c.id,
            knowledge_validations.c.unit_id,
            knowledge_validations.c.validator,
            knowledge_validations.c.validation_status,
            knowledge_validations.c.validated_at,
            knowledge_validations.c.expires_at,
            knowledge_validations.c.confidence,
            knowledge_validations.c.notes,
            knowledge_validations.c.meta,
            knowledge_validations.c.created_at,
        )
        if unit_id is not None:
            stmt = stmt.where(knowledge_validations.c.unit_id == unit_id)
        stmt = stmt.order_by(knowledge_validations.c.validated_at.desc()).limit(min(max(limit, 1), 500))
        rows = (await session.execute(stmt)).all()
        return [_knowledge_validation_row_to_out(row) for row in rows]

    @app.get("/v1/knowledge/validation-policy", response_model=list[KnowledgeValidationPolicyOut])
    async def list_knowledge_validation_policy(
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeValidationPolicyOut]:
        now = datetime.utcnow()
        out: list[KnowledgeValidationPolicyOut] = []
        for risk_level in sorted(KNOWLEDGE_RISK_ORDER.keys(), key=lambda name: KNOWLEDGE_RISK_ORDER[name]):
            row = (
                await session.execute(
                    sa.select(
                        knowledge_validation_policies.c.risk_level,
                        knowledge_validation_policies.c.strict_mode,
                        knowledge_validation_policies.c.require_validation,
                        knowledge_validation_policies.c.require_approved,
                        knowledge_validation_policies.c.require_not_expired,
                        knowledge_validation_policies.c.min_confidence,
                        knowledge_validation_policies.c.max_validation_age_days,
                        knowledge_validation_policies.c.created_at,
                        knowledge_validation_policies.c.updated_at,
                    ).where(knowledge_validation_policies.c.risk_level == risk_level)
                )
            ).first()

            if row:
                out.append(_knowledge_validation_policy_row_to_out(row))
                continue

            default_policy = _default_validation_policy(risk_level)
            inserted = (
                await session.execute(
                    knowledge_validation_policies.insert()
                    .values(
                        id=uuid4(),
                        risk_level=risk_level,
                        strict_mode=default_policy["strict_mode"],
                        require_validation=default_policy["require_validation"],
                        require_approved=default_policy["require_approved"],
                        require_not_expired=default_policy["require_not_expired"],
                        min_confidence=default_policy["min_confidence"],
                        max_validation_age_days=default_policy["max_validation_age_days"],
                        created_at=now,
                        updated_at=now,
                    )
                    .returning(
                        knowledge_validation_policies.c.risk_level,
                        knowledge_validation_policies.c.strict_mode,
                        knowledge_validation_policies.c.require_validation,
                        knowledge_validation_policies.c.require_approved,
                        knowledge_validation_policies.c.require_not_expired,
                        knowledge_validation_policies.c.min_confidence,
                        knowledge_validation_policies.c.max_validation_age_days,
                        knowledge_validation_policies.c.created_at,
                        knowledge_validation_policies.c.updated_at,
                    )
                )
            ).one()
            out.append(_knowledge_validation_policy_row_to_out(inserted))

        await session.commit()
        return out

    @app.get("/v1/knowledge/validation-policy/bundles", response_model=list[KnowledgeValidationPolicyBundleOut])
    async def list_knowledge_validation_policy_bundles(
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeValidationPolicyBundleOut]:
        await _ensure_default_policy_bundle(session)
        rows = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_bundles.c.id,
                    knowledge_validation_policy_bundles.c.bundle_key,
                    knowledge_validation_policy_bundles.c.description,
                    knowledge_validation_policy_bundles.c.is_default,
                    knowledge_validation_policy_bundles.c.enabled,
                    knowledge_validation_policy_bundles.c.created_at,
                    knowledge_validation_policy_bundles.c.updated_at,
                ).order_by(knowledge_validation_policy_bundles.c.is_default.desc(), knowledge_validation_policy_bundles.c.bundle_key.asc())
            )
        ).all()
        await session.commit()
        return [_knowledge_validation_policy_bundle_row_to_out(row) for row in rows]

    @app.post("/v1/knowledge/validation-policy/bundles", response_model=KnowledgeValidationPolicyBundleOut)
    async def create_knowledge_validation_policy_bundle(
        body: KnowledgeValidationPolicyBundleUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyBundleOut:
        bundle_key = str(body.bundle_key or "").strip()
        if not bundle_key:
            raise HTTPException(status_code=422, detail="bundle_key is required")

        before_row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_bundles.c.id,
                    knowledge_validation_policy_bundles.c.bundle_key,
                    knowledge_validation_policy_bundles.c.description,
                    knowledge_validation_policy_bundles.c.is_default,
                    knowledge_validation_policy_bundles.c.enabled,
                    knowledge_validation_policy_bundles.c.created_at,
                    knowledge_validation_policy_bundles.c.updated_at,
                ).where(knowledge_validation_policy_bundles.c.bundle_key == bundle_key)
            )
        ).first()
        now = datetime.utcnow()
        row = (
            await session.execute(
                pg_insert(knowledge_validation_policy_bundles)
                .values(
                    id=uuid4(),
                    bundle_key=bundle_key,
                    description=body.description,
                    is_default=body.is_default,
                    enabled=body.enabled,
                    created_at=now,
                    updated_at=now,
                )
                .on_conflict_do_update(
                    index_elements=[knowledge_validation_policy_bundles.c.bundle_key],
                    set_={
                        "description": body.description,
                        "is_default": body.is_default,
                        "enabled": body.enabled,
                        "updated_at": now,
                    },
                )
                .returning(
                    knowledge_validation_policy_bundles.c.id,
                    knowledge_validation_policy_bundles.c.bundle_key,
                    knowledge_validation_policy_bundles.c.description,
                    knowledge_validation_policy_bundles.c.is_default,
                    knowledge_validation_policy_bundles.c.enabled,
                    knowledge_validation_policy_bundles.c.created_at,
                    knowledge_validation_policy_bundles.c.updated_at,
                )
            )
        ).one()
        if body.is_default:
            await session.execute(
                knowledge_validation_policy_bundles.update()
                .where(knowledge_validation_policy_bundles.c.id != row.id)
                .values(is_default=False, updated_at=now)
            )
        await _record_policy_change_event(
            session,
            entity_type="bundle",
            entity_id=row.id,
            entity_key=row.bundle_key,
            action="upsert",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_validation_policy_bundle_row_to_out(row)

    @app.put("/v1/knowledge/validation-policy/bundles/{bundle_id}", response_model=KnowledgeValidationPolicyBundleOut)
    async def update_knowledge_validation_policy_bundle(
        bundle_id: UUID,
        body: KnowledgeValidationPolicyBundleUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyBundleOut:
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_bundles.c.id,
                    knowledge_validation_policy_bundles.c.bundle_key,
                    knowledge_validation_policy_bundles.c.description,
                    knowledge_validation_policy_bundles.c.is_default,
                    knowledge_validation_policy_bundles.c.enabled,
                    knowledge_validation_policy_bundles.c.created_at,
                    knowledge_validation_policy_bundles.c.updated_at,
                ).where(knowledge_validation_policy_bundles.c.id == bundle_id)
            )
        ).first()
        if not before_row:
            raise HTTPException(status_code=404, detail="policy bundle not found")
        now = datetime.utcnow()
        await session.execute(
            knowledge_validation_policy_bundles.update()
            .where(knowledge_validation_policy_bundles.c.id == bundle_id)
            .values(
                bundle_key=str(body.bundle_key or "").strip(),
                description=body.description,
                is_default=body.is_default,
                enabled=body.enabled,
                updated_at=now,
            )
        )
        if body.is_default:
            await session.execute(
                knowledge_validation_policy_bundles.update()
                .where(knowledge_validation_policy_bundles.c.id != bundle_id)
                .values(is_default=False, updated_at=now)
            )
        row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_bundles.c.id,
                    knowledge_validation_policy_bundles.c.bundle_key,
                    knowledge_validation_policy_bundles.c.description,
                    knowledge_validation_policy_bundles.c.is_default,
                    knowledge_validation_policy_bundles.c.enabled,
                    knowledge_validation_policy_bundles.c.created_at,
                    knowledge_validation_policy_bundles.c.updated_at,
                ).where(knowledge_validation_policy_bundles.c.id == bundle_id)
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="bundle",
            entity_id=row.id,
            entity_key=row.bundle_key,
            action="update",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_validation_policy_bundle_row_to_out(row)

    @app.get("/v1/knowledge/validation-policy/rules", response_model=list[KnowledgeValidationPolicyRuleOut])
    async def list_knowledge_validation_policy_rules(
        bundle_id: UUID | None = None,
        risk_level: str | None = None,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeValidationPolicyRuleOut]:
        stmt = sa.select(
            knowledge_validation_policy_rules.c.id,
            knowledge_validation_policy_rules.c.bundle_id,
            knowledge_validation_policy_rules.c.rule_key,
            knowledge_validation_policy_rules.c.risk_level,
            knowledge_validation_policy_rules.c.task_pattern,
            knowledge_validation_policy_rules.c.agent_slug,
            knowledge_validation_policy_rules.c.source_type,
            knowledge_validation_policy_rules.c.priority,
            knowledge_validation_policy_rules.c.enabled,
            knowledge_validation_policy_rules.c.strict_mode,
            knowledge_validation_policy_rules.c.require_validation,
            knowledge_validation_policy_rules.c.require_approved,
            knowledge_validation_policy_rules.c.require_not_expired,
            knowledge_validation_policy_rules.c.min_confidence,
            knowledge_validation_policy_rules.c.max_validation_age_days,
            knowledge_validation_policy_rules.c.description,
            knowledge_validation_policy_rules.c.created_at,
            knowledge_validation_policy_rules.c.updated_at,
        )
        if bundle_id is not None:
            stmt = stmt.where(knowledge_validation_policy_rules.c.bundle_id == bundle_id)
        if risk_level:
            stmt = stmt.where(knowledge_validation_policy_rules.c.risk_level == risk_level.strip().lower())
        stmt = stmt.order_by(knowledge_validation_policy_rules.c.priority.desc(), knowledge_validation_policy_rules.c.rule_key.asc())
        rows = (await session.execute(stmt)).all()
        return [_knowledge_validation_policy_rule_row_to_out(row) for row in rows]

    @app.post("/v1/knowledge/validation-policy/rules", response_model=KnowledgeValidationPolicyRuleOut)
    async def create_knowledge_validation_policy_rule(
        body: KnowledgeValidationPolicyRuleUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyRuleOut:
        normalized_risk = str(body.risk_level or "").strip().lower()
        if normalized_risk not in KNOWLEDGE_RISK_ORDER:
            raise HTTPException(status_code=422, detail=f"risk_level must be one of {sorted(KNOWLEDGE_RISK_ORDER.keys())}")
        bundle_exists = (await session.execute(sa.select(knowledge_validation_policy_bundles.c.id).where(knowledge_validation_policy_bundles.c.id == body.bundle_id))).first()
        if not bundle_exists:
            raise HTTPException(status_code=404, detail="policy bundle not found")
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rules.c.id,
                    knowledge_validation_policy_rules.c.bundle_id,
                    knowledge_validation_policy_rules.c.rule_key,
                    knowledge_validation_policy_rules.c.risk_level,
                    knowledge_validation_policy_rules.c.task_pattern,
                    knowledge_validation_policy_rules.c.agent_slug,
                    knowledge_validation_policy_rules.c.source_type,
                    knowledge_validation_policy_rules.c.priority,
                    knowledge_validation_policy_rules.c.enabled,
                    knowledge_validation_policy_rules.c.strict_mode,
                    knowledge_validation_policy_rules.c.require_validation,
                    knowledge_validation_policy_rules.c.require_approved,
                    knowledge_validation_policy_rules.c.require_not_expired,
                    knowledge_validation_policy_rules.c.min_confidence,
                    knowledge_validation_policy_rules.c.max_validation_age_days,
                    knowledge_validation_policy_rules.c.description,
                    knowledge_validation_policy_rules.c.created_at,
                    knowledge_validation_policy_rules.c.updated_at,
                )
                .where(knowledge_validation_policy_rules.c.bundle_id == body.bundle_id)
                .where(knowledge_validation_policy_rules.c.rule_key == str(body.rule_key or "").strip())
            )
        ).first()
        now = datetime.utcnow()
        row = (
            await session.execute(
                pg_insert(knowledge_validation_policy_rules)
                .values(
                    id=uuid4(),
                    bundle_id=body.bundle_id,
                    rule_key=str(body.rule_key or "").strip(),
                    risk_level=normalized_risk,
                    task_pattern=_normalize_task_pattern(body.task_pattern),
                    agent_slug=(str(body.agent_slug or "").strip() or None),
                    source_type=(str(body.source_type or "").strip() or None),
                    priority=int(body.priority),
                    enabled=body.enabled,
                    strict_mode=body.strict_mode,
                    require_validation=body.require_validation,
                    require_approved=body.require_approved,
                    require_not_expired=body.require_not_expired,
                    min_confidence=body.min_confidence,
                    max_validation_age_days=body.max_validation_age_days,
                    description=body.description,
                    created_at=now,
                    updated_at=now,
                )
                .on_conflict_do_update(
                    index_elements=[knowledge_validation_policy_rules.c.bundle_id, knowledge_validation_policy_rules.c.rule_key],
                    set_={
                        "risk_level": normalized_risk,
                        "task_pattern": _normalize_task_pattern(body.task_pattern),
                        "agent_slug": (str(body.agent_slug or "").strip() or None),
                        "source_type": (str(body.source_type or "").strip() or None),
                        "priority": int(body.priority),
                        "enabled": body.enabled,
                        "strict_mode": body.strict_mode,
                        "require_validation": body.require_validation,
                        "require_approved": body.require_approved,
                        "require_not_expired": body.require_not_expired,
                        "min_confidence": body.min_confidence,
                        "max_validation_age_days": body.max_validation_age_days,
                        "description": body.description,
                        "updated_at": now,
                    },
                )
                .returning(
                    knowledge_validation_policy_rules.c.id,
                    knowledge_validation_policy_rules.c.bundle_id,
                    knowledge_validation_policy_rules.c.rule_key,
                    knowledge_validation_policy_rules.c.risk_level,
                    knowledge_validation_policy_rules.c.task_pattern,
                    knowledge_validation_policy_rules.c.agent_slug,
                    knowledge_validation_policy_rules.c.source_type,
                    knowledge_validation_policy_rules.c.priority,
                    knowledge_validation_policy_rules.c.enabled,
                    knowledge_validation_policy_rules.c.strict_mode,
                    knowledge_validation_policy_rules.c.require_validation,
                    knowledge_validation_policy_rules.c.require_approved,
                    knowledge_validation_policy_rules.c.require_not_expired,
                    knowledge_validation_policy_rules.c.min_confidence,
                    knowledge_validation_policy_rules.c.max_validation_age_days,
                    knowledge_validation_policy_rules.c.description,
                    knowledge_validation_policy_rules.c.created_at,
                    knowledge_validation_policy_rules.c.updated_at,
                )
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="rule",
            entity_id=row.id,
            entity_key=row.rule_key,
            action="upsert",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_validation_policy_rule_row_to_out(row)

    @app.put("/v1/knowledge/validation-policy/rules/{rule_id}", response_model=KnowledgeValidationPolicyRuleOut)
    async def update_knowledge_validation_policy_rule(
        rule_id: UUID,
        body: KnowledgeValidationPolicyRuleUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyRuleOut:
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rules.c.id,
                    knowledge_validation_policy_rules.c.bundle_id,
                    knowledge_validation_policy_rules.c.rule_key,
                    knowledge_validation_policy_rules.c.risk_level,
                    knowledge_validation_policy_rules.c.task_pattern,
                    knowledge_validation_policy_rules.c.agent_slug,
                    knowledge_validation_policy_rules.c.source_type,
                    knowledge_validation_policy_rules.c.priority,
                    knowledge_validation_policy_rules.c.enabled,
                    knowledge_validation_policy_rules.c.strict_mode,
                    knowledge_validation_policy_rules.c.require_validation,
                    knowledge_validation_policy_rules.c.require_approved,
                    knowledge_validation_policy_rules.c.require_not_expired,
                    knowledge_validation_policy_rules.c.min_confidence,
                    knowledge_validation_policy_rules.c.max_validation_age_days,
                    knowledge_validation_policy_rules.c.description,
                    knowledge_validation_policy_rules.c.created_at,
                    knowledge_validation_policy_rules.c.updated_at,
                ).where(knowledge_validation_policy_rules.c.id == rule_id)
            )
        ).first()
        if not before_row:
            raise HTTPException(status_code=404, detail="policy rule not found")
        normalized_risk = str(body.risk_level or "").strip().lower()
        if normalized_risk not in KNOWLEDGE_RISK_ORDER:
            raise HTTPException(status_code=422, detail=f"risk_level must be one of {sorted(KNOWLEDGE_RISK_ORDER.keys())}")
        now = datetime.utcnow()
        await session.execute(
            knowledge_validation_policy_rules.update()
            .where(knowledge_validation_policy_rules.c.id == rule_id)
            .values(
                bundle_id=body.bundle_id,
                rule_key=str(body.rule_key or "").strip(),
                risk_level=normalized_risk,
                task_pattern=_normalize_task_pattern(body.task_pattern),
                agent_slug=(str(body.agent_slug or "").strip() or None),
                source_type=(str(body.source_type or "").strip() or None),
                priority=int(body.priority),
                enabled=body.enabled,
                strict_mode=body.strict_mode,
                require_validation=body.require_validation,
                require_approved=body.require_approved,
                require_not_expired=body.require_not_expired,
                min_confidence=body.min_confidence,
                max_validation_age_days=body.max_validation_age_days,
                description=body.description,
                updated_at=now,
            )
        )
        row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rules.c.id,
                    knowledge_validation_policy_rules.c.bundle_id,
                    knowledge_validation_policy_rules.c.rule_key,
                    knowledge_validation_policy_rules.c.risk_level,
                    knowledge_validation_policy_rules.c.task_pattern,
                    knowledge_validation_policy_rules.c.agent_slug,
                    knowledge_validation_policy_rules.c.source_type,
                    knowledge_validation_policy_rules.c.priority,
                    knowledge_validation_policy_rules.c.enabled,
                    knowledge_validation_policy_rules.c.strict_mode,
                    knowledge_validation_policy_rules.c.require_validation,
                    knowledge_validation_policy_rules.c.require_approved,
                    knowledge_validation_policy_rules.c.require_not_expired,
                    knowledge_validation_policy_rules.c.min_confidence,
                    knowledge_validation_policy_rules.c.max_validation_age_days,
                    knowledge_validation_policy_rules.c.description,
                    knowledge_validation_policy_rules.c.created_at,
                    knowledge_validation_policy_rules.c.updated_at,
                ).where(knowledge_validation_policy_rules.c.id == rule_id)
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="rule",
            entity_id=row.id,
            entity_key=row.rule_key,
            action="update",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_validation_policy_rule_row_to_out(row)

    @app.get("/v1/knowledge/validation-policy/rollouts", response_model=list[KnowledgeValidationPolicyRolloutOut])
    async def list_knowledge_validation_policy_rollouts(
        bundle_id: UUID | None = None,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeValidationPolicyRolloutOut]:
        stmt = sa.select(
            knowledge_validation_policy_rollouts.c.id,
            knowledge_validation_policy_rollouts.c.bundle_id,
            knowledge_validation_policy_rollouts.c.rollout_key,
            knowledge_validation_policy_rollouts.c.target_agent_slug,
            knowledge_validation_policy_rollouts.c.target_source_type,
            knowledge_validation_policy_rollouts.c.task_pattern,
            knowledge_validation_policy_rollouts.c.priority,
            knowledge_validation_policy_rollouts.c.enabled,
            knowledge_validation_policy_rollouts.c.rollout_mode,
            knowledge_validation_policy_rollouts.c.rollout_percentage,
            knowledge_validation_policy_rollouts.c.created_at,
            knowledge_validation_policy_rollouts.c.updated_at,
        )
        if bundle_id is not None:
            stmt = stmt.where(knowledge_validation_policy_rollouts.c.bundle_id == bundle_id)
        stmt = stmt.order_by(knowledge_validation_policy_rollouts.c.priority.desc(), knowledge_validation_policy_rollouts.c.rollout_key.asc())
        rows = (await session.execute(stmt)).all()
        return [_knowledge_validation_policy_rollout_row_to_out(row) for row in rows]

    @app.post("/v1/knowledge/validation-policy/rollouts", response_model=KnowledgeValidationPolicyRolloutOut)
    async def create_knowledge_validation_policy_rollout(
        body: KnowledgeValidationPolicyRolloutUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyRolloutOut:
        mode = str(body.rollout_mode or "full").strip().lower()
        if mode not in KNOWLEDGE_POLICY_ROLLOUT_MODES:
            raise HTTPException(status_code=422, detail=f"rollout_mode must be one of {sorted(KNOWLEDGE_POLICY_ROLLOUT_MODES)}")
        percentage = body.rollout_percentage
        if mode == "percentage":
            if percentage is None or int(percentage) < 1 or int(percentage) > 100:
                raise HTTPException(status_code=422, detail="rollout_percentage must be between 1 and 100 for percentage rollout")
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rollouts.c.id,
                    knowledge_validation_policy_rollouts.c.bundle_id,
                    knowledge_validation_policy_rollouts.c.rollout_key,
                    knowledge_validation_policy_rollouts.c.target_agent_slug,
                    knowledge_validation_policy_rollouts.c.target_source_type,
                    knowledge_validation_policy_rollouts.c.task_pattern,
                    knowledge_validation_policy_rollouts.c.priority,
                    knowledge_validation_policy_rollouts.c.enabled,
                    knowledge_validation_policy_rollouts.c.rollout_mode,
                    knowledge_validation_policy_rollouts.c.rollout_percentage,
                    knowledge_validation_policy_rollouts.c.created_at,
                    knowledge_validation_policy_rollouts.c.updated_at,
                )
                .where(knowledge_validation_policy_rollouts.c.bundle_id == body.bundle_id)
                .where(knowledge_validation_policy_rollouts.c.rollout_key == str(body.rollout_key or "").strip())
            )
        ).first()
        now = datetime.utcnow()
        row = (
            await session.execute(
                pg_insert(knowledge_validation_policy_rollouts)
                .values(
                    id=uuid4(),
                    bundle_id=body.bundle_id,
                    rollout_key=str(body.rollout_key or "").strip(),
                    target_agent_slug=(str(body.target_agent_slug or "").strip() or None),
                    target_source_type=(str(body.target_source_type or "").strip() or None),
                    task_pattern=_normalize_task_pattern(body.task_pattern),
                    priority=int(body.priority),
                    enabled=body.enabled,
                    rollout_mode=mode,
                    rollout_percentage=int(percentage) if percentage is not None else None,
                    created_at=now,
                    updated_at=now,
                )
                .on_conflict_do_update(
                    index_elements=[knowledge_validation_policy_rollouts.c.bundle_id, knowledge_validation_policy_rollouts.c.rollout_key],
                    set_={
                        "target_agent_slug": (str(body.target_agent_slug or "").strip() or None),
                        "target_source_type": (str(body.target_source_type or "").strip() or None),
                        "task_pattern": _normalize_task_pattern(body.task_pattern),
                        "priority": int(body.priority),
                        "enabled": body.enabled,
                        "rollout_mode": mode,
                        "rollout_percentage": int(percentage) if percentage is not None else None,
                        "updated_at": now,
                    },
                )
                .returning(
                    knowledge_validation_policy_rollouts.c.id,
                    knowledge_validation_policy_rollouts.c.bundle_id,
                    knowledge_validation_policy_rollouts.c.rollout_key,
                    knowledge_validation_policy_rollouts.c.target_agent_slug,
                    knowledge_validation_policy_rollouts.c.target_source_type,
                    knowledge_validation_policy_rollouts.c.task_pattern,
                    knowledge_validation_policy_rollouts.c.priority,
                    knowledge_validation_policy_rollouts.c.enabled,
                    knowledge_validation_policy_rollouts.c.rollout_mode,
                    knowledge_validation_policy_rollouts.c.rollout_percentage,
                    knowledge_validation_policy_rollouts.c.created_at,
                    knowledge_validation_policy_rollouts.c.updated_at,
                )
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="rollout",
            entity_id=row.id,
            entity_key=row.rollout_key,
            action="upsert",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_validation_policy_rollout_row_to_out(row)

    @app.put("/v1/knowledge/validation-policy/rollouts/{rollout_id}", response_model=KnowledgeValidationPolicyRolloutOut)
    async def update_knowledge_validation_policy_rollout(
        rollout_id: UUID,
        body: KnowledgeValidationPolicyRolloutUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyRolloutOut:
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rollouts.c.id,
                    knowledge_validation_policy_rollouts.c.bundle_id,
                    knowledge_validation_policy_rollouts.c.rollout_key,
                    knowledge_validation_policy_rollouts.c.target_agent_slug,
                    knowledge_validation_policy_rollouts.c.target_source_type,
                    knowledge_validation_policy_rollouts.c.task_pattern,
                    knowledge_validation_policy_rollouts.c.priority,
                    knowledge_validation_policy_rollouts.c.enabled,
                    knowledge_validation_policy_rollouts.c.rollout_mode,
                    knowledge_validation_policy_rollouts.c.rollout_percentage,
                    knowledge_validation_policy_rollouts.c.created_at,
                    knowledge_validation_policy_rollouts.c.updated_at,
                ).where(knowledge_validation_policy_rollouts.c.id == rollout_id)
            )
        ).first()
        if not before_row:
            raise HTTPException(status_code=404, detail="policy rollout not found")
        mode = str(body.rollout_mode or "full").strip().lower()
        if mode not in KNOWLEDGE_POLICY_ROLLOUT_MODES:
            raise HTTPException(status_code=422, detail=f"rollout_mode must be one of {sorted(KNOWLEDGE_POLICY_ROLLOUT_MODES)}")
        percentage = body.rollout_percentage
        if mode == "percentage":
            if percentage is None or int(percentage) < 1 or int(percentage) > 100:
                raise HTTPException(status_code=422, detail="rollout_percentage must be between 1 and 100 for percentage rollout")
        now = datetime.utcnow()
        await session.execute(
            knowledge_validation_policy_rollouts.update()
            .where(knowledge_validation_policy_rollouts.c.id == rollout_id)
            .values(
                bundle_id=body.bundle_id,
                rollout_key=str(body.rollout_key or "").strip(),
                target_agent_slug=(str(body.target_agent_slug or "").strip() or None),
                target_source_type=(str(body.target_source_type or "").strip() or None),
                task_pattern=_normalize_task_pattern(body.task_pattern),
                priority=int(body.priority),
                enabled=body.enabled,
                rollout_mode=mode,
                rollout_percentage=int(percentage) if percentage is not None else None,
                updated_at=now,
            )
        )
        row = (
            await session.execute(
                sa.select(
                    knowledge_validation_policy_rollouts.c.id,
                    knowledge_validation_policy_rollouts.c.bundle_id,
                    knowledge_validation_policy_rollouts.c.rollout_key,
                    knowledge_validation_policy_rollouts.c.target_agent_slug,
                    knowledge_validation_policy_rollouts.c.target_source_type,
                    knowledge_validation_policy_rollouts.c.task_pattern,
                    knowledge_validation_policy_rollouts.c.priority,
                    knowledge_validation_policy_rollouts.c.enabled,
                    knowledge_validation_policy_rollouts.c.rollout_mode,
                    knowledge_validation_policy_rollouts.c.rollout_percentage,
                    knowledge_validation_policy_rollouts.c.created_at,
                    knowledge_validation_policy_rollouts.c.updated_at,
                ).where(knowledge_validation_policy_rollouts.c.id == rollout_id)
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="rollout",
            entity_id=row.id,
            entity_key=row.rollout_key,
            action="update",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_validation_policy_rollout_row_to_out(row)

    @app.get("/v1/knowledge/resolve/ranking-profiles", response_model=list[KnowledgeResolveRankingProfileOut])
    async def list_knowledge_resolve_ranking_profiles(
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeResolveRankingProfileOut]:
        await _ensure_default_ranking_profiles(session)
        rows = (
            await session.execute(
                sa.select(
                    knowledge_resolve_ranking_profiles.c.id,
                    knowledge_resolve_ranking_profiles.c.profile_key,
                    knowledge_resolve_ranking_profiles.c.description,
                    knowledge_resolve_ranking_profiles.c.enabled,
                    knowledge_resolve_ranking_profiles.c.base_score,
                    knowledge_resolve_ranking_profiles.c.lexical_weight,
                    knowledge_resolve_ranking_profiles.c.semantic_weight,
                    knowledge_resolve_ranking_profiles.c.tag_weight,
                    knowledge_resolve_ranking_profiles.c.validation_confidence_weight,
                    knowledge_resolve_ranking_profiles.c.approved_bonus,
                    knowledge_resolve_ranking_profiles.c.preferred_bonus,
                    knowledge_resolve_ranking_profiles.c.deprecated_penalty,
                    knowledge_resolve_ranking_profiles.c.created_at,
                    knowledge_resolve_ranking_profiles.c.updated_at,
                ).order_by(knowledge_resolve_ranking_profiles.c.profile_key.asc())
            )
        ).all()
        await session.commit()
        return [_knowledge_resolve_ranking_profile_row_to_out(row) for row in rows]

    @app.post("/v1/knowledge/resolve/ranking-profiles", response_model=KnowledgeResolveRankingProfileOut)
    async def create_knowledge_resolve_ranking_profile(
        body: KnowledgeResolveRankingProfileUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeResolveRankingProfileOut:
        profile_key = str(body.profile_key or "").strip().lower()
        if not profile_key:
            raise HTTPException(status_code=422, detail="profile_key is required")
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_resolve_ranking_profiles.c.id,
                    knowledge_resolve_ranking_profiles.c.profile_key,
                    knowledge_resolve_ranking_profiles.c.description,
                    knowledge_resolve_ranking_profiles.c.enabled,
                    knowledge_resolve_ranking_profiles.c.base_score,
                    knowledge_resolve_ranking_profiles.c.lexical_weight,
                    knowledge_resolve_ranking_profiles.c.semantic_weight,
                    knowledge_resolve_ranking_profiles.c.tag_weight,
                    knowledge_resolve_ranking_profiles.c.validation_confidence_weight,
                    knowledge_resolve_ranking_profiles.c.approved_bonus,
                    knowledge_resolve_ranking_profiles.c.preferred_bonus,
                    knowledge_resolve_ranking_profiles.c.deprecated_penalty,
                    knowledge_resolve_ranking_profiles.c.created_at,
                    knowledge_resolve_ranking_profiles.c.updated_at,
                ).where(knowledge_resolve_ranking_profiles.c.profile_key == profile_key)
            )
        ).first()
        now = datetime.utcnow()
        row = (
            await session.execute(
                pg_insert(knowledge_resolve_ranking_profiles)
                .values(
                    id=uuid4(),
                    profile_key=profile_key,
                    description=body.description,
                    enabled=body.enabled,
                    base_score=body.base_score,
                    lexical_weight=body.lexical_weight,
                    semantic_weight=body.semantic_weight,
                    tag_weight=body.tag_weight,
                    validation_confidence_weight=body.validation_confidence_weight,
                    approved_bonus=body.approved_bonus,
                    preferred_bonus=body.preferred_bonus,
                    deprecated_penalty=body.deprecated_penalty,
                    created_at=now,
                    updated_at=now,
                )
                .on_conflict_do_update(
                    index_elements=[knowledge_resolve_ranking_profiles.c.profile_key],
                    set_={
                        "description": body.description,
                        "enabled": body.enabled,
                        "base_score": body.base_score,
                        "lexical_weight": body.lexical_weight,
                        "semantic_weight": body.semantic_weight,
                        "tag_weight": body.tag_weight,
                        "validation_confidence_weight": body.validation_confidence_weight,
                        "approved_bonus": body.approved_bonus,
                        "preferred_bonus": body.preferred_bonus,
                        "deprecated_penalty": body.deprecated_penalty,
                        "updated_at": now,
                    },
                )
                .returning(
                    knowledge_resolve_ranking_profiles.c.id,
                    knowledge_resolve_ranking_profiles.c.profile_key,
                    knowledge_resolve_ranking_profiles.c.description,
                    knowledge_resolve_ranking_profiles.c.enabled,
                    knowledge_resolve_ranking_profiles.c.base_score,
                    knowledge_resolve_ranking_profiles.c.lexical_weight,
                    knowledge_resolve_ranking_profiles.c.semantic_weight,
                    knowledge_resolve_ranking_profiles.c.tag_weight,
                    knowledge_resolve_ranking_profiles.c.validation_confidence_weight,
                    knowledge_resolve_ranking_profiles.c.approved_bonus,
                    knowledge_resolve_ranking_profiles.c.preferred_bonus,
                    knowledge_resolve_ranking_profiles.c.deprecated_penalty,
                    knowledge_resolve_ranking_profiles.c.created_at,
                    knowledge_resolve_ranking_profiles.c.updated_at,
                )
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="ranking_profile",
            entity_id=row.id,
            entity_key=row.profile_key,
            action="upsert",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_resolve_ranking_profile_row_to_out(row)

    @app.put("/v1/knowledge/resolve/ranking-profiles/{profile_id}", response_model=KnowledgeResolveRankingProfileOut)
    async def update_knowledge_resolve_ranking_profile(
        profile_id: UUID,
        body: KnowledgeResolveRankingProfileUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeResolveRankingProfileOut:
        before_row = (
            await session.execute(
                sa.select(
                    knowledge_resolve_ranking_profiles.c.id,
                    knowledge_resolve_ranking_profiles.c.profile_key,
                    knowledge_resolve_ranking_profiles.c.description,
                    knowledge_resolve_ranking_profiles.c.enabled,
                    knowledge_resolve_ranking_profiles.c.base_score,
                    knowledge_resolve_ranking_profiles.c.lexical_weight,
                    knowledge_resolve_ranking_profiles.c.semantic_weight,
                    knowledge_resolve_ranking_profiles.c.tag_weight,
                    knowledge_resolve_ranking_profiles.c.validation_confidence_weight,
                    knowledge_resolve_ranking_profiles.c.approved_bonus,
                    knowledge_resolve_ranking_profiles.c.preferred_bonus,
                    knowledge_resolve_ranking_profiles.c.deprecated_penalty,
                    knowledge_resolve_ranking_profiles.c.created_at,
                    knowledge_resolve_ranking_profiles.c.updated_at,
                ).where(knowledge_resolve_ranking_profiles.c.id == profile_id)
            )
        ).first()
        if not before_row:
            raise HTTPException(status_code=404, detail="ranking profile not found")
        now = datetime.utcnow()
        await session.execute(
            knowledge_resolve_ranking_profiles.update()
            .where(knowledge_resolve_ranking_profiles.c.id == profile_id)
            .values(
                profile_key=str(body.profile_key or "").strip().lower(),
                description=body.description,
                enabled=body.enabled,
                base_score=body.base_score,
                lexical_weight=body.lexical_weight,
                semantic_weight=body.semantic_weight,
                tag_weight=body.tag_weight,
                validation_confidence_weight=body.validation_confidence_weight,
                approved_bonus=body.approved_bonus,
                preferred_bonus=body.preferred_bonus,
                deprecated_penalty=body.deprecated_penalty,
                updated_at=now,
            )
        )
        row = (
            await session.execute(
                sa.select(
                    knowledge_resolve_ranking_profiles.c.id,
                    knowledge_resolve_ranking_profiles.c.profile_key,
                    knowledge_resolve_ranking_profiles.c.description,
                    knowledge_resolve_ranking_profiles.c.enabled,
                    knowledge_resolve_ranking_profiles.c.base_score,
                    knowledge_resolve_ranking_profiles.c.lexical_weight,
                    knowledge_resolve_ranking_profiles.c.semantic_weight,
                    knowledge_resolve_ranking_profiles.c.tag_weight,
                    knowledge_resolve_ranking_profiles.c.validation_confidence_weight,
                    knowledge_resolve_ranking_profiles.c.approved_bonus,
                    knowledge_resolve_ranking_profiles.c.preferred_bonus,
                    knowledge_resolve_ranking_profiles.c.deprecated_penalty,
                    knowledge_resolve_ranking_profiles.c.created_at,
                    knowledge_resolve_ranking_profiles.c.updated_at,
                ).where(knowledge_resolve_ranking_profiles.c.id == profile_id)
            )
        ).one()
        await _record_policy_change_event(
            session,
            entity_type="ranking_profile",
            entity_id=row.id,
            entity_key=row.profile_key,
            action="update",
            actor=None,
            before_state=_serialize_row_state(before_row),
            after_state=_serialize_row_state(row),
            now=now,
        )
        await session.commit()
        return _knowledge_resolve_ranking_profile_row_to_out(row)

    @app.get("/v1/knowledge/validation-policy/change-events", response_model=list[KnowledgeValidationPolicyChangeEventOut])
    async def list_knowledge_validation_policy_change_events(
        entity_type: str | None = None,
        entity_key: str | None = None,
        limit: int = 100,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeValidationPolicyChangeEventOut]:
        stmt = sa.select(
            knowledge_validation_policy_change_events.c.id,
            knowledge_validation_policy_change_events.c.entity_type,
            knowledge_validation_policy_change_events.c.entity_id,
            knowledge_validation_policy_change_events.c.entity_key,
            knowledge_validation_policy_change_events.c.action,
            knowledge_validation_policy_change_events.c.actor,
            knowledge_validation_policy_change_events.c.before_state,
            knowledge_validation_policy_change_events.c.after_state,
            knowledge_validation_policy_change_events.c.created_at,
        )
        if entity_type:
            stmt = stmt.where(knowledge_validation_policy_change_events.c.entity_type == entity_type)
        if entity_key:
            stmt = stmt.where(knowledge_validation_policy_change_events.c.entity_key == entity_key)
        stmt = stmt.order_by(knowledge_validation_policy_change_events.c.created_at.desc()).limit(min(max(limit, 1), 500))
        rows = (await session.execute(stmt)).all()
        return [_knowledge_validation_policy_change_event_row_to_out(row) for row in rows]

    @app.get("/v1/knowledge/validation-policy/observability/summary", response_model=KnowledgeValidationPolicyObservabilityOut)
    async def get_knowledge_validation_policy_observability_summary(
        days: int = 7,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyObservabilityOut:
        window_days = min(max(int(days or 7), 1), 90)
        cutoff = datetime.utcnow() - timedelta(days=window_days)
        rows = (
            await session.execute(
                sa.select(
                    knowledge_resolve_audits.c.selected_count,
                    knowledge_resolve_audits.c.rejected_count,
                    knowledge_resolve_audits.c.payload,
                ).where(knowledge_resolve_audits.c.created_at >= cutoff)
            )
        ).all()
        bundle_groups: dict[str, list[dict]] = {}
        rollout_groups: dict[str, list[dict]] = {}
        ranking_groups: dict[str, list[dict]] = {}
        normalized_rows: list[dict] = []
        for row in rows:
            payload = row.payload or {}
            meta = payload.get("policy_metadata") or {}
            normalized = {
                "selected_count": int(row.selected_count or 0),
                "rejected_count": int(row.rejected_count or 0),
                "bundle_key": str(meta.get("bundle_key") or "unknown"),
                "rollout_key": str(meta.get("rollout_key") or "none"),
                "ranking_profile": str(payload.get("ranking_profile") or "unknown"),
            }
            normalized_rows.append(normalized)
            bundle_groups.setdefault(normalized["bundle_key"], []).append(normalized)
            rollout_groups.setdefault(normalized["rollout_key"], []).append(normalized)
            ranking_groups.setdefault(normalized["ranking_profile"], []).append(normalized)
        return KnowledgeValidationPolicyObservabilityOut(
            window_days=window_days,
            total_resolve_requests=len(normalized_rows),
            bundles=[_observability_item(key, items) for key, items in sorted(bundle_groups.items())],
            rollouts=[_observability_item(key, items) for key, items in sorted(rollout_groups.items())],
            ranking_profiles=[_observability_item(key, items) for key, items in sorted(ranking_groups.items())],
        )

    @app.post("/v1/knowledge/validation-policy/bundles/{bundle_id}/rollback", response_model=KnowledgeValidationPolicyRollbackOut)
    async def rollback_knowledge_validation_policy_bundle(
        bundle_id: UUID,
        body: KnowledgeValidationPolicyRollbackIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyRollbackOut:
        result = await _rollback_policy_entity(session, entity_type="bundle", entity_id=bundle_id, actor=body.actor)
        await session.commit()
        return result

    @app.post("/v1/knowledge/validation-policy/rollouts/{rollout_id}/rollback", response_model=KnowledgeValidationPolicyRollbackOut)
    async def rollback_knowledge_validation_policy_rollout(
        rollout_id: UUID,
        body: KnowledgeValidationPolicyRollbackIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyRollbackOut:
        result = await _rollback_policy_entity(session, entity_type="rollout", entity_id=rollout_id, actor=body.actor)
        await session.commit()
        return result

    @app.put("/v1/knowledge/validation-policy/{risk_level}", response_model=KnowledgeValidationPolicyOut)
    async def upsert_knowledge_validation_policy(
        risk_level: str,
        body: KnowledgeValidationPolicyUpsertIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeValidationPolicyOut:
        normalized = str(risk_level or "").strip().lower()
        if normalized not in KNOWLEDGE_RISK_ORDER:
            raise HTTPException(status_code=422, detail=f"risk_level must be one of {sorted(KNOWLEDGE_RISK_ORDER.keys())}")

        now = datetime.utcnow()
        stmt = (
            pg_insert(knowledge_validation_policies)
            .values(
                id=uuid4(),
                risk_level=normalized,
                strict_mode=body.strict_mode,
                require_validation=body.require_validation,
                require_approved=body.require_approved,
                require_not_expired=body.require_not_expired,
                min_confidence=body.min_confidence,
                max_validation_age_days=body.max_validation_age_days,
                created_at=now,
                updated_at=now,
            )
            .on_conflict_do_update(
                index_elements=[knowledge_validation_policies.c.risk_level],
                set_={
                    "strict_mode": body.strict_mode,
                    "require_validation": body.require_validation,
                    "require_approved": body.require_approved,
                    "require_not_expired": body.require_not_expired,
                    "min_confidence": body.min_confidence,
                    "max_validation_age_days": body.max_validation_age_days,
                    "updated_at": now,
                },
            )
            .returning(
                knowledge_validation_policies.c.risk_level,
                knowledge_validation_policies.c.strict_mode,
                knowledge_validation_policies.c.require_validation,
                knowledge_validation_policies.c.require_approved,
                knowledge_validation_policies.c.require_not_expired,
                knowledge_validation_policies.c.min_confidence,
                knowledge_validation_policies.c.max_validation_age_days,
                knowledge_validation_policies.c.created_at,
                knowledge_validation_policies.c.updated_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()
        return _knowledge_validation_policy_row_to_out(row)

    @app.post("/v1/knowledge/resolve", response_model=KnowledgeResolveOut)
    async def resolve_knowledge_package(
        body: KnowledgeResolveIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeResolveOut:
        requested_rank = _risk_rank(body.risk_level)
        limit = min(max(int(body.limit or 10), 1), 100)
        retrieval_mode = str(body.retrieval_mode or "lexical").strip().lower() or "lexical"
        if retrieval_mode not in {"lexical", "semantic", "hybrid"}:
            raise HTTPException(status_code=422, detail="retrieval_mode must be one of lexical, semantic, hybrid")
        ranking_profile = str(body.ranking_profile or "balanced").strip().lower() or "balanced"
        ranking_profile_row = await _load_ranking_profile(session, ranking_profile)

        stmt = sa.select(
            knowledge_units.c.id,
            knowledge_units.c.source_id,
            knowledge_units.c.unit_key,
            knowledge_units.c.title,
            knowledge_units.c.content,
            knowledge_units.c.content_sha256,
            knowledge_units.c.tags,
            knowledge_units.c.agent_scope,
            knowledge_units.c.risk_level,
            knowledge_units.c.status,
            knowledge_units.c.lifecycle_stage,
            knowledge_units.c.superseded_by_unit_id,
            knowledge_units.c.retired_at,
            knowledge_units.c.meta,
            knowledge_units.c.created_at,
            knowledge_units.c.updated_at,
            knowledge_sources.c.source_type.label("source_type"),
        ).select_from(
            knowledge_units.outerjoin(knowledge_sources, knowledge_sources.c.id == knowledge_units.c.source_id)
        ).where(knowledge_units.c.status == "active")

        if body.source_type:
            stmt = stmt.where(knowledge_sources.c.source_type == body.source_type)

        if body.tags:
            stmt = stmt.where(knowledge_units.c.tags.op("&&")(body.tags))
        if body.agent_slug:
            stmt = stmt.where(
                sa.or_(
                    sa.func.cardinality(knowledge_units.c.agent_scope) == 0,
                    knowledge_units.c.agent_scope.any(body.agent_slug),
                )
            )

        stmt = stmt.order_by(knowledge_units.c.updated_at.desc()).limit(500)
        lexical_rows = (await session.execute(stmt)).all() if retrieval_mode in {"lexical", "hybrid"} else []
        policy, policy_metadata = await _resolve_validation_policy(
            session,
            task=body.task,
            agent_slug=body.agent_slug,
            source_type=body.source_type,
            risk_level=body.risk_level,
        )

        semantic_rows: list[dict] = []
        semantic_similarity_by_unit_id: dict[str, float] = {}
        semantic_query = str(body.semantic_query or body.task or "").strip()
        if retrieval_mode in {"semantic", "hybrid"}:
            try:
                embedding_config = _knowledge_embedding_runtime_config(settings)
            except ValueError as exc:
                raise HTTPException(status_code=503, detail=str(exc)) from exc
            if not embedding_config:
                raise HTTPException(status_code=503, detail="knowledge embedding is not enabled")
            if not semantic_query:
                raise HTTPException(status_code=422, detail="semantic_query must be non-empty for semantic or hybrid retrieval")

            configured_model = str(embedding_config["model"]).strip()
            requested_model = str(body.embedding_model or "").strip() or configured_model
            if requested_model != configured_model:
                raise HTTPException(
                    status_code=422,
                    detail=f"embedding_model must match configured runtime model: {configured_model}",
                )

            try:
                vectors, model_name, actual_dimensions = await _request_knowledge_embeddings(
                    [semantic_query],
                    config=embedding_config,
                )
            except Exception as exc:
                raise HTTPException(status_code=503, detail=str(exc)) from exc

            selected_model = str(model_name or configured_model).strip() or configured_model
            if selected_model != requested_model:
                raise HTTPException(
                    status_code=422,
                    detail=f"runtime embedding model mismatch: requested={requested_model} actual={selected_model}",
                )

            requested_dimensions = body.embedding_dimensions
            if requested_dimensions is not None and int(requested_dimensions) != int(actual_dimensions):
                raise HTTPException(
                    status_code=422,
                    detail=(
                        f"embedding_dimensions must match runtime query embedding dimensions: "
                        f"requested={int(requested_dimensions)} actual={int(actual_dimensions)}"
                    ),
                )

            require_approved_validation = body.require_approved_validation
            if require_approved_validation is None:
                require_approved_validation = requested_rank >= KNOWLEDGE_RISK_ORDER["high"]

            semantic_limit = min(max(int(body.semantic_limit or max(limit * 2, 20)), limit), 100)
            semantic_rows = await _search_knowledge_units_by_embedding(
                session,
                query_embedding=vectors[0],
                embedding_model=selected_model,
                embedding_dimensions=int(actual_dimensions),
                limit=semantic_limit,
                source_id=None,
                source_type=body.source_type,
                agent_slug=body.agent_slug,
                risk_level=body.risk_level,
                require_approved_validation=bool(require_approved_validation),
                tags=body.tags,
            )
            semantic_similarity_by_unit_id = {
                str(row["id"]): _similarity_from_cosine_distance(float(row["distance"]))
                for row in semantic_rows
            }

        candidate_rows_by_unit_id: dict[str, dict] = {}
        retrieval_channels_by_unit_id: dict[str, set[str]] = {}
        for row in lexical_rows:
            mapping = _coerce_row_mapping(row)
            unit_id = str(mapping["id"])
            candidate_rows_by_unit_id[unit_id] = mapping
            retrieval_channels_by_unit_id.setdefault(unit_id, set()).add("lexical")
        for row in semantic_rows:
            mapping = _coerce_row_mapping(row)
            unit_id = str(mapping["id"])
            candidate_rows_by_unit_id[unit_id] = mapping
            retrieval_channels_by_unit_id.setdefault(unit_id, set()).add("semantic")

        rows = list(candidate_rows_by_unit_id.values())

        now = datetime.now(timezone.utc)
        selected_items: list[KnowledgeResolveItemOut] = []
        selected_payload: list[dict] = []
        rejected_payload: list[dict] = []
        min_semantic_similarity = body.min_semantic_similarity
        min_score = body.min_score

        for row in rows:
            unit_id = str(row["id"])
            reject_reason: str | None = None
            if _risk_rank(row["risk_level"]) > requested_rank:
                reject_reason = "unit_risk_exceeds_requested"

            latest_validation = (
                await session.execute(
                    sa.select(
                        knowledge_validations.c.validation_status,
                        knowledge_validations.c.expires_at,
                        knowledge_validations.c.confidence,
                        knowledge_validations.c.validated_at,
                    )
                    .where(knowledge_validations.c.unit_id == row["id"])
                    .order_by(knowledge_validations.c.validated_at.desc())
                    .limit(1)
                )
            ).first()

            validation_status = latest_validation.validation_status if latest_validation else None
            validation_expires_at = latest_validation.expires_at if latest_validation else None
            validation_confidence_raw = latest_validation.confidence if latest_validation else None
            validation_confidence = float(validation_confidence_raw) if validation_confidence_raw is not None else 0.0
            validation_validated_at = latest_validation.validated_at if latest_validation else None

            compare_expires_at = validation_expires_at
            if isinstance(compare_expires_at, datetime) and compare_expires_at.tzinfo is None:
                compare_expires_at = compare_expires_at.replace(tzinfo=timezone.utc)
            expired = bool(compare_expires_at and compare_expires_at <= now)

            compare_validated_at = validation_validated_at
            if isinstance(compare_validated_at, datetime) and compare_validated_at.tzinfo is None:
                compare_validated_at = compare_validated_at.replace(tzinfo=timezone.utc)

            strict_mode = bool(policy.get("strict_mode"))

            if reject_reason is None and policy.get("require_validation") and not latest_validation:
                reject_reason = "missing_validation"
            if reject_reason is None and strict_mode and policy.get("require_approved") and not validation_status:
                reject_reason = "missing_validation_status"
            if reject_reason is None and strict_mode and policy.get("require_not_expired") and not compare_expires_at:
                reject_reason = "missing_validation_expires_at"
            if reject_reason is None and strict_mode and policy.get("min_confidence") is not None and validation_confidence_raw is None:
                reject_reason = "missing_validation_confidence"
            if reject_reason is None and strict_mode and policy.get("max_validation_age_days") is not None and not compare_validated_at:
                reject_reason = "missing_validation_validated_at"
            if reject_reason is None and policy.get("require_approved") and validation_status != "approved":
                reject_reason = "validation_not_approved"
            if reject_reason is None and policy.get("require_not_expired"):
                if not compare_expires_at:
                    reject_reason = "validation_expiry_required"
                elif expired:
                    reject_reason = "validation_expired"
            if reject_reason is None and policy.get("min_confidence") is not None:
                if validation_confidence < float(policy.get("min_confidence") or 0.0):
                    reject_reason = "validation_confidence_too_low"
            if reject_reason is None and policy.get("max_validation_age_days") is not None:
                max_age_days = int(policy.get("max_validation_age_days") or 0)
                cutoff = now - timedelta(days=max_age_days)
                if not compare_validated_at or compare_validated_at < cutoff:
                    reject_reason = "validation_too_old"

            if reject_reason:
                rejected_payload.append(
                    {
                        "unit_id": unit_id,
                        "unit_key": row["unit_key"],
                        "source_id": str(row["source_id"]) if row["source_id"] else None,
                        "reason": reject_reason,
                    }
                )
                continue

            semantic_similarity = semantic_similarity_by_unit_id.get(unit_id)
            if semantic_similarity is not None and min_semantic_similarity is not None and semantic_similarity < float(min_semantic_similarity):
                rejected_payload.append(
                    {
                        "unit_id": unit_id,
                        "unit_key": row["unit_key"],
                        "source_id": str(row["source_id"]) if row["source_id"] else None,
                        "reason": "semantic_similarity_below_threshold",
                    }
                )
                continue

            matched_tags = len(set(body.tags).intersection(set(row["tags"] or []))) if body.tags else 0
            lexical_overlap = _compute_lexical_overlap(body.task, row["title"], row["content"], list(row["tags"] or []))
            lifecycle_adjustment = 0.0
            if str(row.get("lifecycle_stage") or "active").strip().lower() == "preferred":
                lifecycle_adjustment = float(ranking_profile_row["preferred_bonus"])
            elif str(row.get("lifecycle_stage") or "active").strip().lower() == "deprecated":
                lifecycle_adjustment = -float(ranking_profile_row["deprecated_penalty"])
            validation_bonus = float(ranking_profile_row["approved_bonus"]) if validation_status == "approved" else 0.0
            semantic_component = max(semantic_similarity or 0.0, 0.0)
            lexical_weight = float(ranking_profile_row["lexical_weight"])
            if retrieval_mode == "lexical":
                lexical_weight += 0.4
            semantic_weight = float(ranking_profile_row["semantic_weight"])
            confidence_weight = float(ranking_profile_row["validation_confidence_weight"])
            tag_weight = float(ranking_profile_row["tag_weight"])

            score_breakdown = {
                "base": float(ranking_profile_row["base_score"]),
                "lexical_overlap": lexical_overlap * lexical_weight,
                "matched_tags": float(matched_tags) * tag_weight,
                "validation_bonus": validation_bonus,
                "validation_confidence": validation_confidence * confidence_weight,
                "semantic_similarity": semantic_component * semantic_weight,
                "lifecycle_adjustment": lifecycle_adjustment,
            }
            score = sum(float(value) for value in score_breakdown.values())
            if min_score is not None and score < float(min_score):
                rejected_payload.append(
                    {
                        "unit_id": unit_id,
                        "unit_key": row["unit_key"],
                        "source_id": str(row["source_id"]) if row["source_id"] else None,
                        "reason": "score_below_threshold",
                    }
                )
                continue

            unit = KnowledgeUnitOut(
                id=row["id"],
                source_id=row["source_id"],
                unit_key=row["unit_key"],
                title=row["title"],
                content=row["content"],
                content_sha256=row["content_sha256"],
                tags=list(row["tags"] or []),
                agent_scope=list(row["agent_scope"] or []),
                risk_level=row["risk_level"],
                status=row["status"],
                lifecycle_stage=row.get("lifecycle_stage") or "active",
                superseded_by_unit_id=row.get("superseded_by_unit_id"),
                retired_at=row.get("retired_at"),
                meta=row["meta"] or {},
                created_at=row["created_at"],
                updated_at=row["updated_at"],
            )

            selected_items.append(
                KnowledgeResolveItemOut(
                    unit=unit,
                    validation_status=validation_status,
                    validation_expires_at=validation_expires_at,
                    score=score,
                    semantic_similarity=semantic_similarity,
                    retrieval_channels=sorted(retrieval_channels_by_unit_id.get(unit_id, set())),
                    score_breakdown=score_breakdown,
                )
            )

            selected_payload.append(
                {
                    "unit_id": unit_id,
                    "unit_key": row["unit_key"],
                    "source_id": str(row["source_id"]) if row["source_id"] else None,
                    "score": score,
                    "semantic_similarity": semantic_similarity,
                    "score_breakdown": score_breakdown,
                    "retrieval_channels": sorted(retrieval_channels_by_unit_id.get(unit_id, set())),
                }
            )

        selected_items.sort(key=lambda item: item.score, reverse=True)
        selected = selected_items[:limit]
        selected_ids = {str(item.unit.id) for item in selected}
        selected_payload = [item for item in selected_payload if item["unit_id"] in selected_ids]

        await session.execute(
            knowledge_resolve_audits.insert().values(
                id=uuid4(),
                task=body.task,
                agent_slug=body.agent_slug,
                requested_risk_level=(body.risk_level or "normal").strip().lower() or "normal",
                tags=body.tags,
                selected_count=len(selected_payload),
                rejected_count=len(rejected_payload),
                payload={
                    "retrieval_mode": retrieval_mode,
                    "ranking_profile": ranking_profile,
                    "ranking_profile_weights": {
                        "base_score": float(ranking_profile_row["base_score"]),
                        "lexical_weight": float(ranking_profile_row["lexical_weight"]),
                        "semantic_weight": float(ranking_profile_row["semantic_weight"]),
                        "tag_weight": float(ranking_profile_row["tag_weight"]),
                        "validation_confidence_weight": float(ranking_profile_row["validation_confidence_weight"]),
                        "approved_bonus": float(ranking_profile_row["approved_bonus"]),
                        "preferred_bonus": float(ranking_profile_row["preferred_bonus"]),
                        "deprecated_penalty": float(ranking_profile_row["deprecated_penalty"]),
                    },
                    "semantic_query": semantic_query if retrieval_mode in {"semantic", "hybrid"} else None,
                    "semantic_candidate_count": len(semantic_rows),
                    "min_semantic_similarity": min_semantic_similarity,
                    "min_score": min_score,
                    "policy": policy,
                    "policy_metadata": policy_metadata,
                    "selected": selected_payload,
                    "rejected": rejected_payload,
                },
                created_at=now,
            )
        )
        await session.commit()

        return KnowledgeResolveOut(
            task=body.task,
            agent_slug=body.agent_slug,
            risk_level=(body.risk_level or "normal").strip().lower() or "normal",
            total=len(selected),
            rejected_count=len(rejected_payload),
            items=selected,
            rejected=[KnowledgeResolveRejectedOut(**item) for item in rejected_payload] if body.include_rejected else [],
        )

    @app.post("/v1/knowledge/search", response_model=KnowledgeSearchOut)
    async def search_knowledge_units(
        body: KnowledgeSearchIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeSearchOut:
        query_text = str(body.query or "").strip()
        if not query_text:
            raise HTTPException(status_code=422, detail="query must be non-empty")

        try:
            embedding_config = _knowledge_embedding_runtime_config(settings)
        except ValueError as exc:
            raise HTTPException(status_code=503, detail=str(exc)) from exc
        if not embedding_config:
            raise HTTPException(status_code=503, detail="knowledge embedding is not enabled")

        configured_model = str(embedding_config["model"]).strip()
        requested_model = str(body.embedding_model or "").strip() or configured_model
        if requested_model != configured_model:
            raise HTTPException(
                status_code=422,
                detail=f"embedding_model must match configured runtime model: {configured_model}",
            )

        limit = min(max(int(body.limit or 10), 1), 50)

        try:
            vectors, model_name, actual_dimensions = await _request_knowledge_embeddings(
                [query_text],
                config=embedding_config,
            )
        except Exception as exc:
            raise HTTPException(status_code=503, detail=str(exc)) from exc

        selected_model = str(model_name or configured_model).strip() or configured_model
        if selected_model != requested_model:
            raise HTTPException(
                status_code=422,
                detail=f"runtime embedding model mismatch: requested={requested_model} actual={selected_model}",
            )

        requested_dimensions = body.embedding_dimensions
        if requested_dimensions is not None and int(requested_dimensions) != int(actual_dimensions):
            raise HTTPException(
                status_code=422,
                detail=(
                    f"embedding_dimensions must match runtime query embedding dimensions: "
                    f"requested={int(requested_dimensions)} actual={int(actual_dimensions)}"
                ),
            )

        rows = await _search_knowledge_units_by_embedding(
            session,
            query_embedding=vectors[0],
            embedding_model=selected_model,
            embedding_dimensions=int(actual_dimensions),
            limit=limit,
            source_id=body.source_id,
            source_type=body.source_type,
            agent_slug=body.agent_slug,
            risk_level=body.risk_level,
            require_approved_validation=bool(body.require_approved_validation),
            tags=body.tags,
        )

        items: list[KnowledgeSearchItemOut] = []
        for row in rows:
            unit = KnowledgeUnitOut(
                id=row["id"],
                source_id=row["source_id"],
                unit_key=row["unit_key"],
                title=row["title"],
                content=row["content"],
                content_sha256=row["content_sha256"],
                tags=list(row["tags"] or []),
                agent_scope=list(row["agent_scope"] or []),
                risk_level=row["risk_level"],
                status=row["status"],
                lifecycle_stage=row.get("lifecycle_stage") or "active",
                superseded_by_unit_id=row.get("superseded_by_unit_id"),
                retired_at=row.get("retired_at"),
                meta=row["meta"] or {},
                created_at=row["created_at"],
                updated_at=row["updated_at"],
            )
            distance = float(row["distance"])
            items.append(
                KnowledgeSearchItemOut(
                    unit=unit,
                    embedding_model=str(row["embedding_model"]),
                    embedding_dimensions=int(row["embedding_dimensions"]),
                    source_type=row.get("source_type"),
                    distance=distance,
                    similarity=_similarity_from_cosine_distance(distance),
                )
            )

        return KnowledgeSearchOut(
            query=query_text,
            total=len(items),
            embedding_model=selected_model,
            embedding_dimensions=int(actual_dimensions),
            items=items,
        )

    @app.post("/v1/knowledge/feedback", response_model=KnowledgeFeedbackOut)
    async def create_knowledge_feedback(
        body: KnowledgeFeedbackIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeFeedbackOut:
        feedback_type = str(body.feedback_type or "").strip().lower()
        if feedback_type not in KNOWLEDGE_FEEDBACK_TYPES:
            raise HTTPException(status_code=422, detail=f"feedback_type must be one of: {sorted(KNOWLEDGE_FEEDBACK_TYPES)}")

        unit_exists = (
            await session.execute(sa.select(knowledge_units.c.id).where(knowledge_units.c.id == body.unit_id))
        ).first()
        if not unit_exists:
            raise HTTPException(status_code=404, detail="knowledge unit not found")

        now = datetime.utcnow()
        stmt = (
            knowledge_feedback_events.insert()
            .values(
                id=uuid4(),
                unit_id=body.unit_id,
                agent=(body.agent or "").strip() or None,
                feedback_type=feedback_type,
                severity=(body.severity or "info").strip().lower() or "info",
                payload=body.payload,
                created_at=now,
            )
            .returning(
                knowledge_feedback_events.c.id,
                knowledge_feedback_events.c.unit_id,
                knowledge_feedback_events.c.feedback_type,
                knowledge_feedback_events.c.agent,
                knowledge_feedback_events.c.severity,
                knowledge_feedback_events.c.payload,
                knowledge_feedback_events.c.created_at,
            )
        )
        row = (await session.execute(stmt)).one()

        if feedback_type == "invalidation":
            await session.execute(
                knowledge_units.update()
                .where(knowledge_units.c.id == body.unit_id)
                .values(status="inactive", lifecycle_stage="inactive", retired_at=now, updated_at=now)
            )
            await _record_lifecycle_event(
                session,
                unit_id=body.unit_id,
                action="invalidate",
                actor=body.agent,
                payload={"source": "feedback", "feedback_type": feedback_type, **(body.payload or {})},
                now=now,
            )
        elif feedback_type == "promotion":
            await session.execute(
                knowledge_units.update()
                .where(knowledge_units.c.id == body.unit_id)
                .values(status="active", lifecycle_stage="preferred", retired_at=None, updated_at=now)
            )
            await _record_lifecycle_event(
                session,
                unit_id=body.unit_id,
                action="promote",
                actor=body.agent,
                payload={"source": "feedback", "feedback_type": feedback_type, **(body.payload or {})},
                now=now,
            )

        await session.commit()
        return _knowledge_feedback_row_to_out(row)

    @app.post("/v1/knowledge/units/{unit_id}/lifecycle", response_model=KnowledgeUnitOut)
    async def apply_knowledge_unit_lifecycle_action(
        unit_id: UUID,
        body: KnowledgeLifecycleActionIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeUnitOut:
        action = str(body.action or "").strip().lower()
        if action not in KNOWLEDGE_LIFECYCLE_ACTIONS:
            raise HTTPException(status_code=422, detail=f"action must be one of {sorted(KNOWLEDGE_LIFECYCLE_ACTIONS)}")

        existing = (
            await session.execute(
                sa.select(
                    knowledge_units.c.id,
                    knowledge_units.c.source_id,
                    knowledge_units.c.unit_key,
                    knowledge_units.c.title,
                    knowledge_units.c.content,
                    knowledge_units.c.content_sha256,
                    knowledge_units.c.tags,
                    knowledge_units.c.agent_scope,
                    knowledge_units.c.risk_level,
                    knowledge_units.c.status,
                    knowledge_units.c.lifecycle_stage,
                    knowledge_units.c.superseded_by_unit_id,
                    knowledge_units.c.retired_at,
                    knowledge_units.c.meta,
                    knowledge_units.c.created_at,
                    knowledge_units.c.updated_at,
                ).where(knowledge_units.c.id == unit_id)
            )
        ).first()
        if not existing:
            raise HTTPException(status_code=404, detail="knowledge unit not found")

        superseded_by_unit_id = body.superseded_by_unit_id
        if action == "supersede":
            if superseded_by_unit_id is None:
                raise HTTPException(status_code=422, detail="superseded_by_unit_id is required for supersede action")
            target_exists = (
                await session.execute(sa.select(knowledge_units.c.id).where(knowledge_units.c.id == superseded_by_unit_id))
            ).first()
            if not target_exists:
                raise HTTPException(status_code=404, detail="superseded target knowledge unit not found")

        now = datetime.utcnow()
        values: dict[str, object] = {"updated_at": now}
        if action == "promote":
            values.update({"status": "active", "lifecycle_stage": "preferred", "retired_at": None, "superseded_by_unit_id": None})
        elif action == "demote":
            values.update({"status": "active", "lifecycle_stage": "deprecated"})
        elif action == "invalidate":
            values.update({"status": "inactive", "lifecycle_stage": "inactive", "retired_at": now})
        elif action == "reactivate":
            values.update({"status": "active", "lifecycle_stage": "active", "retired_at": None, "superseded_by_unit_id": None})
        elif action == "archive":
            values.update({"status": "inactive", "lifecycle_stage": "archived", "retired_at": now})
        elif action == "supersede":
            values.update({"status": "inactive", "lifecycle_stage": "superseded", "retired_at": now, "superseded_by_unit_id": superseded_by_unit_id})

        await session.execute(
            knowledge_units.update().where(knowledge_units.c.id == unit_id).values(**values)
        )
        await _record_lifecycle_event(
            session,
            unit_id=unit_id,
            action=action,
            actor=body.actor,
            payload=body.payload | ({"superseded_by_unit_id": str(superseded_by_unit_id)} if superseded_by_unit_id else {}),
            now=now,
        )
        row = (
            await session.execute(
                sa.select(
                    knowledge_units.c.id,
                    knowledge_units.c.source_id,
                    knowledge_units.c.unit_key,
                    knowledge_units.c.title,
                    knowledge_units.c.content,
                    knowledge_units.c.content_sha256,
                    knowledge_units.c.tags,
                    knowledge_units.c.agent_scope,
                    knowledge_units.c.risk_level,
                    knowledge_units.c.status,
                    knowledge_units.c.lifecycle_stage,
                    knowledge_units.c.superseded_by_unit_id,
                    knowledge_units.c.retired_at,
                    knowledge_units.c.meta,
                    knowledge_units.c.created_at,
                    knowledge_units.c.updated_at,
                ).where(knowledge_units.c.id == unit_id)
            )
        ).one()
        await session.commit()
        return _knowledge_unit_row_to_out(row)

    @app.get("/v1/knowledge/units/{unit_id}/lifecycle-events", response_model=list[KnowledgeLifecycleEventOut])
    async def list_knowledge_unit_lifecycle_events(
        unit_id: UUID,
        limit: int = 50,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeLifecycleEventOut]:
        rows = (
            await session.execute(
                sa.select(
                    knowledge_unit_lifecycle_events.c.id,
                    knowledge_unit_lifecycle_events.c.unit_id,
                    knowledge_unit_lifecycle_events.c.action,
                    knowledge_unit_lifecycle_events.c.actor,
                    knowledge_unit_lifecycle_events.c.payload,
                    knowledge_unit_lifecycle_events.c.created_at,
                )
                .where(knowledge_unit_lifecycle_events.c.unit_id == unit_id)
                .order_by(knowledge_unit_lifecycle_events.c.created_at.desc())
                .limit(min(max(limit, 1), 200))
            )
        ).all()
        return [_knowledge_lifecycle_event_row_to_out(row) for row in rows]

    @app.get("/v1/knowledge/feedback/summary", response_model=list[KnowledgeFeedbackSummaryOut])
    async def get_knowledge_feedback_summary(
        days: int = 7,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeFeedbackSummaryOut]:
        window_days = min(max(int(days or 7), 1), 90)
        cutoff = datetime.utcnow() - timedelta(days=window_days)
        rows = (
            await session.execute(
                sa.select(
                    knowledge_feedback_events.c.feedback_type,
                    sa.func.count().label("count"),
                )
                .where(knowledge_feedback_events.c.created_at >= cutoff)
                .group_by(knowledge_feedback_events.c.feedback_type)
                .order_by(sa.func.count().desc())
            )
        ).all()
        return [KnowledgeFeedbackSummaryOut(feedback_type=row.feedback_type, count=int(row.count or 0)) for row in rows]

    @app.get("/v1/knowledge/resolve/audits", response_model=list[KnowledgeResolveAuditOut])
    async def list_knowledge_resolve_audits(
        risk_level: str | None = None,
        limit: int = 50,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeResolveAuditOut]:
        stmt = sa.select(
            knowledge_resolve_audits.c.id,
            knowledge_resolve_audits.c.task,
            knowledge_resolve_audits.c.agent_slug,
            knowledge_resolve_audits.c.requested_risk_level,
            knowledge_resolve_audits.c.tags,
            knowledge_resolve_audits.c.selected_count,
            knowledge_resolve_audits.c.rejected_count,
            knowledge_resolve_audits.c.payload,
            knowledge_resolve_audits.c.created_at,
        )
        if risk_level:
            stmt = stmt.where(knowledge_resolve_audits.c.requested_risk_level == risk_level.strip().lower())
        stmt = stmt.order_by(knowledge_resolve_audits.c.created_at.desc()).limit(min(max(limit, 1), 500))
        rows = (await session.execute(stmt)).all()
        return [_knowledge_resolve_audit_row_to_out(row) for row in rows]

    @app.get("/v1/knowledge/resolve/rejections/summary", response_model=list[KnowledgeResolveRejectSummaryOut])
    async def get_knowledge_resolve_rejection_summary(
        days: int = 7,
        risk_level: str | None = None,
        limit: int = 20,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[KnowledgeResolveRejectSummaryOut]:
        window_days = min(max(int(days or 7), 1), 90)
        max_items = min(max(int(limit or 20), 1), 200)
        cutoff = datetime.utcnow() - timedelta(days=window_days)

        stmt = sa.select(knowledge_resolve_audits.c.payload).where(knowledge_resolve_audits.c.created_at >= cutoff)
        if risk_level:
            stmt = stmt.where(knowledge_resolve_audits.c.requested_risk_level == risk_level.strip().lower())
        stmt = stmt.order_by(knowledge_resolve_audits.c.created_at.desc()).limit(2000)

        rows = (await session.execute(stmt)).all()
        counts: dict[str, int] = {}
        for row in rows:
            payload = row.payload or {}
            rejected = payload.get("rejected") or []
            if not isinstance(rejected, list):
                continue
            for item in rejected:
                if not isinstance(item, dict):
                    continue
                reason = str(item.get("reason") or "unknown").strip() or "unknown"
                counts[reason] = counts.get(reason, 0) + 1

        ordered = sorted(counts.items(), key=lambda pair: pair[1], reverse=True)
        return [
            KnowledgeResolveRejectSummaryOut(reason=reason, count=count)
            for reason, count in ordered[:max_items]
        ]

    @app.get("/v1/knowledge/resolve/metrics", response_model=KnowledgeResolveMetricsOut)
    async def get_knowledge_resolve_metrics(
        days: int = 7,
        risk_level: str | None = None,
        top_reasons: int = 5,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> KnowledgeResolveMetricsOut:
        window_days = min(max(int(days or 7), 1), 90)
        top_n = min(max(int(top_reasons or 5), 1), 20)
        cutoff = datetime.utcnow() - timedelta(days=window_days)

        stmt = sa.select(
            knowledge_resolve_audits.c.requested_risk_level,
            knowledge_resolve_audits.c.selected_count,
            knowledge_resolve_audits.c.rejected_count,
            knowledge_resolve_audits.c.payload,
        ).where(knowledge_resolve_audits.c.created_at >= cutoff)
        if risk_level:
            stmt = stmt.where(knowledge_resolve_audits.c.requested_risk_level == risk_level.strip().lower())

        rows = (await session.execute(stmt)).all()

        total_resolve_requests = len(rows)
        requests_with_hits = 0
        requests_with_rejections = 0
        total_selected_units = 0
        total_rejected_units = 0
        expired_rejection_count = 0
        reason_counts: dict[str, int] = {}
        risk_counts: dict[str, dict[str, int]] = {}

        for row in rows:
            requested_risk_level = str(row.requested_risk_level or "normal").strip().lower() or "normal"
            selected_count = int(row.selected_count or 0)
            rejected_count = int(row.rejected_count or 0)

            total_selected_units += selected_count
            total_rejected_units += rejected_count

            risk_bucket = risk_counts.setdefault(
                requested_risk_level,
                {
                    "total_resolve_requests": 0,
                    "requests_with_hits": 0,
                    "requests_with_rejections": 0,
                    "total_selected_units": 0,
                    "total_rejected_units": 0,
                },
            )
            risk_bucket["total_resolve_requests"] += 1
            risk_bucket["total_selected_units"] += selected_count
            risk_bucket["total_rejected_units"] += rejected_count

            if selected_count > 0:
                requests_with_hits += 1
                risk_bucket["requests_with_hits"] += 1
            if rejected_count > 0:
                requests_with_rejections += 1
                risk_bucket["requests_with_rejections"] += 1

            payload = row.payload or {}
            rejected = payload.get("rejected") or []
            if not isinstance(rejected, list):
                continue
            for item in rejected:
                if not isinstance(item, dict):
                    continue
                reason = str(item.get("reason") or "unknown").strip() or "unknown"
                reason_counts[reason] = reason_counts.get(reason, 0) + 1
                if reason in {"validation_expired", "validation_too_old"}:
                    expired_rejection_count += 1

        hit_denominator = max(total_resolve_requests, 1)
        rejection_denominator = max(total_selected_units + total_rejected_units, 1)
        resolve_hit_rate = float(requests_with_hits) / float(hit_denominator)
        resolve_error_rate = float(total_rejected_units) / float(rejection_denominator)
        resolve_rejection_rate = float(requests_with_rejections) / float(hit_denominator)
        unit_selection_rate = float(total_selected_units) / float(rejection_denominator)
        expired_rejection_rate = float(expired_rejection_count) / float(max(total_rejected_units, 1))
        requests_without_hits = total_resolve_requests - requests_with_hits
        requests_without_rejections = total_resolve_requests - requests_with_rejections

        ordered = sorted(reason_counts.items(), key=lambda pair: pair[1], reverse=True)
        top = [
            KnowledgeResolveRejectSummaryOut(
                reason=reason,
                count=count,
                rate=float(count) / float(max(total_rejected_units, 1)),
            )
            for reason, count in ordered[:top_n]
        ]

        risk_breakdown: list[KnowledgeResolveRiskMetricsOut] = []
        for risk_level_name in sorted(risk_counts.keys(), key=_risk_rank):
            bucket = risk_counts[risk_level_name]
            risk_total = int(bucket["total_resolve_requests"])
            risk_breakdown.append(
                KnowledgeResolveRiskMetricsOut(
                    risk_level=risk_level_name,
                    total_resolve_requests=risk_total,
                    requests_with_hits=int(bucket["requests_with_hits"]),
                    requests_with_rejections=int(bucket["requests_with_rejections"]),
                    total_selected_units=int(bucket["total_selected_units"]),
                    total_rejected_units=int(bucket["total_rejected_units"]),
                    resolve_hit_rate=float(bucket["requests_with_hits"]) / float(max(risk_total, 1)),
                    resolve_rejection_rate=float(bucket["requests_with_rejections"]) / float(max(risk_total, 1)),
                )
            )

        source_stmt = sa.select(knowledge_sources.c.meta).where(knowledge_sources.c.updated_at >= cutoff)
        source_rows = (await session.execute(source_stmt)).all()
        total_parsed_sources = 0
        ocr_fallback_count = 0
        ocr_failure_count = 0
        ocr_page_truncation_count = 0

        for source_row in source_rows:
            meta = source_row.meta or {}
            if not isinstance(meta, dict):
                continue
            parse_status = str(meta.get("parse_status") or "").strip().lower()
            if parse_status:
                total_parsed_sources += 1
            if bool(meta.get("ocr_fallback_used")):
                ocr_fallback_count += 1
            if parse_status == "failed" and bool(meta.get("ocr_used")):
                ocr_failure_count += 1
            if int(meta.get("ocr_pages_truncated") or 0) > 0:
                ocr_page_truncation_count += 1

        ocr_failure_rate = float(ocr_failure_count) / float(max(total_parsed_sources, 1))

        return KnowledgeResolveMetricsOut(
            window_days=window_days,
            total_resolve_requests=total_resolve_requests,
            requests_with_hits=requests_with_hits,
            requests_with_rejections=requests_with_rejections,
            requests_without_hits=requests_without_hits,
            requests_without_rejections=requests_without_rejections,
            total_selected_units=total_selected_units,
            total_rejected_units=total_rejected_units,
            resolve_hit_rate=resolve_hit_rate,
            resolve_error_rate=resolve_error_rate,
            resolve_rejection_rate=resolve_rejection_rate,
            unit_selection_rate=unit_selection_rate,
            expired_rejection_count=expired_rejection_count,
            expired_rejection_rate=expired_rejection_rate,
            total_parsed_sources=total_parsed_sources,
            ocr_fallback_count=ocr_fallback_count,
            ocr_failure_count=ocr_failure_count,
            ocr_failure_rate=ocr_failure_rate,
            ocr_page_truncation_count=ocr_page_truncation_count,
            top_reject_reasons=top,
            risk_breakdown=risk_breakdown,
        )

    @app.post("/v1/agents/{agent}/control", response_model=AgentControlActionOut)
    async def control_agent_container(
        agent: str,
        body: AgentControlActionIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
    ) -> AgentControlActionOut:
        slug = str(agent or "").strip()
        if not slug:
            raise HTTPException(status_code=400, detail="agent is required")

        known_agents = _known_agent_slugs(settings)
        if known_agents and slug not in known_agents:
            raise HTTPException(status_code=404, detail=f"unknown agent: {slug}")

        action = str(body.action or "").strip().lower()
        if action not in ALLOWED_AGENT_CONTROL_ACTIONS:
            raise HTTPException(status_code=400, detail=f"invalid action: {action}")

        return await _forward_agent_control(settings, agent=slug, action=action)

    @app.get("/v1/skills/mappings", response_model=list[AgentSkillMappingOut])
    async def get_skill_mappings(
        agent_slug: str | None = None,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[AgentSkillMappingOut]:
        stmt = sa.select(
            agent_skill_mappings.c.id,
            agent_skill_mappings.c.agent_slug,
            agent_skill_mappings.c.skill_slug,
            agent_skill_mappings.c.created_at,
        ).order_by(agent_skill_mappings.c.agent_slug.asc(), agent_skill_mappings.c.skill_slug.asc())
        if agent_slug:
            stmt = stmt.where(agent_skill_mappings.c.agent_slug == agent_slug)
        rows = (await session.execute(stmt)).all()
        return [AgentSkillMappingOut(**row._asdict()) for row in rows]

    @app.patch("/v1/skills/mappings", response_model=SkillsMappingPatchOut)
    async def patch_skill_mappings(
        body: SkillsMappingPatchIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> SkillsMappingPatchOut:
        agent_slugs = sorted({slug.strip() for slug in body.agent_slugs if slug and slug.strip()})
        add_skill_slugs = sorted({slug.strip() for slug in body.add_skill_slugs if slug and slug.strip()})
        remove_skill_slugs = sorted({slug.strip() for slug in body.remove_skill_slugs if slug and slug.strip()})

        if not agent_slugs:
            raise HTTPException(status_code=400, detail="agent_slugs is required")
        if not add_skill_slugs and not remove_skill_slugs:
            raise HTTPException(status_code=400, detail="nothing to update")

        failed: list[SkillsMappingFailureItem] = []
        known_agents = set(_scan_agent_slugs(settings))
        valid_agents = agent_slugs
        if known_agents:
            unknown_agents = [slug for slug in agent_slugs if slug not in known_agents]
            valid_agents = [slug for slug in agent_slugs if slug in known_agents]
            for agent in unknown_agents:
                for skill in add_skill_slugs:
                    failed.append(
                        SkillsMappingFailureItem(
                            action="add",
                            agent_slug=agent,
                            skill_slug=skill,
                            reason="unknown_agent",
                        )
                    )
                for skill in remove_skill_slugs:
                    failed.append(
                        SkillsMappingFailureItem(
                            action="remove",
                            agent_slug=agent,
                            skill_slug=skill,
                            reason="unknown_agent",
                        )
                    )

        known_global = {item.slug for item in _scan_global_skills(settings)}
        unknown_add = [slug for slug in add_skill_slugs if slug not in known_global]
        unknown_remove = [slug for slug in remove_skill_slugs if slug not in known_global]
        valid_add_skill_slugs = [slug for slug in add_skill_slugs if slug in known_global]
        valid_remove_skill_slugs = [slug for slug in remove_skill_slugs if slug in known_global]

        for skill in unknown_add:
            for agent in valid_agents:
                failed.append(
                    SkillsMappingFailureItem(
                        action="add",
                        agent_slug=agent,
                        skill_slug=skill,
                        reason="unknown_global_skill",
                    )
                )
        for skill in unknown_remove:
            for agent in valid_agents:
                failed.append(
                    SkillsMappingFailureItem(
                        action="remove",
                        agent_slug=agent,
                        skill_slug=skill,
                        reason="unknown_global_skill",
                    )
                )

        if not valid_agents:
            return SkillsMappingPatchOut(
                updated=0,
                affected_agents=[],
                restart_hint="No valid mappings were updated.",
                failed=failed,
            )

        updated = 0
        if valid_add_skill_slugs:
            target_pairs = {(agent, skill) for agent in valid_agents for skill in valid_add_skill_slugs}
            existing_stmt = sa.select(
                agent_skill_mappings.c.agent_slug,
                agent_skill_mappings.c.skill_slug,
            ).where(
                agent_skill_mappings.c.agent_slug.in_(valid_agents),
                agent_skill_mappings.c.skill_slug.in_(valid_add_skill_slugs),
            )
            existing_rows = (await session.execute(existing_stmt)).all()
            existing_pairs = {(str(row.agent_slug), str(row.skill_slug)) for row in existing_rows}

            for agent, skill in sorted(existing_pairs):
                failed.append(
                    SkillsMappingFailureItem(
                        action="add",
                        agent_slug=agent,
                        skill_slug=skill,
                        reason="already_mapped",
                    )
                )

            rows = [
                {"id": uuid4(), "agent_slug": agent, "skill_slug": skill}
                for (agent, skill) in sorted(target_pairs - existing_pairs)
            ]
            if rows:
                insert_stmt = pg_insert(agent_skill_mappings).values(rows)
                insert_stmt = insert_stmt.on_conflict_do_nothing(
                    constraint="uq_agent_skill_mappings_agent_skill"
                )
                insert_stmt = insert_stmt.returning(agent_skill_mappings.c.id)
                result = await session.execute(insert_stmt)
                updated += len(result.scalars().all())

        if valid_remove_skill_slugs:
            target_pairs = {(agent, skill) for agent in valid_agents for skill in valid_remove_skill_slugs}
            existing_stmt = sa.select(
                agent_skill_mappings.c.id,
                agent_skill_mappings.c.agent_slug,
                agent_skill_mappings.c.skill_slug,
            ).where(
                agent_skill_mappings.c.agent_slug.in_(valid_agents),
                agent_skill_mappings.c.skill_slug.in_(valid_remove_skill_slugs),
            )
            existing_rows = (await session.execute(existing_stmt)).all()
            existing_pairs = {(str(row.agent_slug), str(row.skill_slug)) for row in existing_rows}

            missing_pairs = sorted(target_pairs - existing_pairs)
            for agent, skill in missing_pairs:
                failed.append(
                    SkillsMappingFailureItem(
                        action="remove",
                        agent_slug=agent,
                        skill_slug=skill,
                        reason="not_mapped",
                    )
                )

            if existing_rows:
                existing_ids = [row.id for row in existing_rows]
                delete_stmt = agent_skill_mappings.delete().where(agent_skill_mappings.c.id.in_(existing_ids))
                result = await session.execute(delete_stmt)
                updated += int(result.rowcount or 0)

        await session.commit()

        return SkillsMappingPatchOut(
            updated=updated,
            affected_agents=valid_agents,
            restart_hint=(
                "Configuration saved. Restart affected agent containers to apply changes."
                if updated > 0
                else "No mapping changes were applied."
            ),
            failed=failed,
        )

    @app.post("/v1/tasks", response_model=TaskOut)
    async def create_task(
        body: TaskCreate,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> TaskOut:
        if body.status not in ALLOWED_TASK_STATUSES:
            raise HTTPException(
                status_code=422,
                detail=f"invalid task status: {body.status}; allowed={sorted(ALLOWED_TASK_STATUSES)}",
            )

        task_id = uuid4()
        now = datetime.utcnow()
        stmt = (
            tasks.insert()
            .values(
                id=task_id,
                title=body.title,
                status=body.status,
                assignee=body.assignee,
                tags=body.tags,
                created_at=now,
                updated_at=now,
            )
            .returning(
                tasks.c.id,
                tasks.c.title,
                tasks.c.status,
                tasks.c.assignee,
                tasks.c.tags,
                tasks.c.created_at,
                tasks.c.updated_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()
        return TaskOut(**row._asdict())

    @app.get("/v1/boards/default", response_model=BoardOut)
    async def get_board(
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> BoardOut:
        statuses = ["INBOX", "ASSIGNED", "IN PROGRESS", "REVIEW", "DONE"]
        columns: list[BoardColumn] = []
        for status in statuses:
            stmt = (
                sa.select(
                    tasks.c.id,
                    tasks.c.title,
                    tasks.c.status,
                    tasks.c.assignee,
                    tasks.c.tags,
                    tasks.c.created_at,
                    tasks.c.updated_at,
                )
                .where(tasks.c.status == status)
                .order_by(tasks.c.updated_at.desc())
                .limit(100)
            )
            rows = (await session.execute(stmt)).all()
            cards = [TaskOut(**r._asdict()) for r in rows]
            columns.append(BoardColumn(title=status, count=len(cards), cards=cards))
        return BoardOut(columns=columns)

    @app.post("/v1/tasks/{task_id}/comments", response_model=CommentOut)
    async def add_comment(
        task_id: UUID,
        body: CommentCreate,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> CommentOut:
        comment_id = uuid4()
        stmt = (
            comments.insert()
            .values(id=comment_id, task_id=task_id, author=body.author, body=body.body)
            .returning(
                comments.c.id,
                comments.c.task_id,
                comments.c.author,
                comments.c.body,
                comments.c.created_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()

        await publish_event(
            redis,
            stream_key=settings.redis_stream_key,
            event={
                "id": str(uuid4()),
                "type": "comment.created",
                "agent": body.author,
                "task_id": str(task_id),
                "payload": {"comment_id": str(row.id)},
                "created_at": datetime.utcnow().isoformat() + "Z",
            },
        )

        return CommentOut(**row._asdict())

    @app.post("/v1/events", response_model=EventOut)
    async def ingest_event(
        body: EventIn,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> EventOut:
        validation_errors: list[str] = []
        validation_details: dict = {}

        if body.type == "task.handoff":
            if not body.task_id:
                validation_errors.append("task.handoff requires task_id")
            known_agents = _known_agent_slugs(settings)
            validation_errors.extend(_validate_handoff_payload(body.payload, known_agents))
            validation_details["known_agents_count"] = len(known_agents)

        event_payload = dict(body.payload)
        if body.type == "task.status":
            if not body.task_id:
                validation_errors.append("task.status requires task_id")
            next_status = str(body.payload.get("new_status") or "").strip().upper()
            if not next_status:
                validation_errors.append("payload.new_status is required")
            elif next_status not in ALLOWED_TASK_STATUSES:
                validation_errors.append(
                    f"payload.new_status invalid: {next_status}; allowed={sorted(ALLOWED_TASK_STATUSES)}"
                )

            current_status = None
            if body.task_id:
                current_row = (
                    await session.execute(
                        sa.select(tasks.c.id, tasks.c.status).where(tasks.c.id == body.task_id)
                    )
                ).first()
                if not current_row:
                    validation_errors.append(f"task not found: {body.task_id}")
                else:
                    current_status = str(current_row.status)

            if not validation_errors and current_status is not None:
                allowed = TASK_STATUS_TRANSITIONS.get(current_status, set())
                if next_status != current_status and next_status not in allowed:
                    validation_errors.append(
                        f"invalid status transition: {current_status} -> {next_status}; allowed={sorted(allowed)}"
                    )

            if not validation_errors and current_status is not None:
                now = datetime.utcnow()
                await session.execute(
                    tasks.update()
                    .where(tasks.c.id == body.task_id)
                    .values(status=next_status, updated_at=now)
                )
                event_payload["previous_status"] = current_status
                event_payload["new_status"] = next_status
                event_payload["transition_applied"] = True
                validation_details["transition"] = {
                    "from": current_status,
                    "to": next_status,
                }

        if validation_errors:
            await publish_validation_result(
                accepted=False,
                body=body,
                errors=validation_errors,
                details=validation_details,
            )
            raise HTTPException(status_code=422, detail={"errors": validation_errors})

        event_id = uuid4()
        stmt = (
            events.insert()
            .values(id=event_id, type=body.type, agent=body.agent, task_id=body.task_id, payload=event_payload)
            .returning(
                events.c.id,
                events.c.type,
                events.c.agent,
                events.c.task_id,
                events.c.payload,
                events.c.created_at,
            )
        )
        row = (await session.execute(stmt)).one()
        await session.commit()

        await publish_event(
            redis,
            stream_key=settings.redis_stream_key,
            event={
                "id": str(row.id),
                "type": row.type,
                "agent": row.agent,
                "task_id": str(row.task_id) if row.task_id else None,
                "payload": row.payload,
                "created_at": row.created_at.isoformat(),
            },
        )

        await publish_validation_result(
            accepted=True,
            body=body,
            errors=[],
            details=validation_details,
        )

        return EventOut(**row._asdict())

    @app.get("/v1/feed", response_model=list[EventOut])
    async def get_feed(
        limit: int = 50,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[EventOut]:
        stmt = (
            sa.select(
                events.c.id,
                events.c.type,
                events.c.agent,
                events.c.task_id,
                events.c.payload,
                events.c.created_at,
            )
            .order_by(events.c.created_at.desc())
            .limit(min(limit, 200))
        )
        rows = (await session.execute(stmt)).all()
        return [EventOut(**r._asdict()) for r in rows]

    @app.get("/v1/feed-lite", response_model=list[EventLiteOut])
    async def get_feed_lite(
        limit: int = 50,
        _auth: None = Depends(lambda authorization=Header(default=None): require_auth(settings, authorization)),
        session=Depends(get_session),
    ) -> list[EventLiteOut]:
        stmt = (
            sa.select(
                events.c.id,
                events.c.type,
                events.c.agent,
                events.c.task_id,
                events.c.created_at,
                sa.func.jsonb_extract_path_text(events.c.payload, "method").label("method"),
                sa.func.jsonb_extract_path_text(events.c.payload, "path").label("path"),
                sa.cast(sa.func.nullif(sa.func.jsonb_extract_path_text(events.c.payload, "status_code"), ""), sa.Integer).label("status_code"),
                sa.func.jsonb_extract_path_text(events.c.payload, "error_type").label("error_type"),
                sa.func.jsonb_extract_path_text(events.c.payload, "test_id").label("test_id"),
                sa.cast(sa.func.nullif(sa.func.jsonb_extract_path_text(events.c.payload, "round"), ""), sa.Integer).label("round"),
            )
            .order_by(events.c.created_at.desc())
            .limit(min(limit, 500))
        )
        rows = (await session.execute(stmt)).all()
        return [EventLiteOut(**r._asdict()) for r in rows]

    @app.websocket("/ws/events")
    async def ws_events(websocket: WebSocket):
        await websocket.accept()

        # 可選 Bearer token（如果設定了 MC_AUTH_TOKEN）
        if settings.auth_token:
            auth = websocket.headers.get("authorization")
            if not auth or not auth.lower().startswith("bearer "):
                await websocket.close(code=4401)
                return
            token = auth.split(" ", 1)[1].strip()
            if token != settings.auth_token:
                await websocket.close(code=4403)
                return

        # Start from the latest stream ID to avoid replaying long history on each
        # new websocket connection, which can significantly increase perceived
        # real-time latency on the dashboard.
        last_id = "$"
        try:
            latest = await redis.xrevrange(settings.redis_stream_key, count=1)
            if latest:
                last_id = str(latest[0][0])
        except Exception:
            # Fall back to '$' (new entries only).
            last_id = "$"
        try:
            while True:
                result = await redis.xread({settings.redis_stream_key: last_id}, block=25_000, count=50)
                if not result:
                    await websocket.send_text(json.dumps({"type": "ping"}))
                    continue

                _, entries = result[0]
                for entry_id, fields in entries:
                    last_id = entry_id
                    # fields["event"] 會是一個 JSON 字串
                    event_json = fields.get("event")
                    if event_json:
                        await websocket.send_text(event_json)
        except WebSocketDisconnect:
            return

    
    @app.api_route("/chat/{agent}/{path:path}", methods=["GET", "POST", "PUT", "DELETE", "PATCH", "OPTIONS", "HEAD"])
    async def chat_proxy(agent: str, path: str, request: Request):
        started = time.perf_counter()
        upgrade = request.headers.get("upgrade", "").lower() == "websocket"
            
        target_url = f"http://openclaw-{agent}:{settings.chat_upstream_port}/{path}"
        query = request.url.query
        if query:
            target_url += f"?{query}"
            
        async with httpx.AsyncClient(timeout=3600.0) as client:
            headers = dict(request.headers.items())
            headers.pop("host", None)
            if settings.chat_force_loopback_headers:
                headers["x-real-ip"] = "127.0.0.1"
                headers["x-forwarded-for"] = "127.0.0.1"
            
            token = settings.agent_token_map.get(agent)
            if token:
                headers["authorization"] = f"Bearer {token}"
                
            req = client.build_request(
                request.method,
                target_url,
                headers=headers,
                content=request.stream(),
            )
            
            try:
                resp = await client.send(req, stream=True)
            except Exception as e:
                if not upgrade:
                    elapsed = max(0.0, time.perf_counter() - started)
                    await publish_event(
                        redis,
                        stream_key=settings.redis_stream_key,
                        event={
                            "type": "chat.gateway.access",
                            "agent": agent,
                            "payload": {
                                "path": f"/{path}",
                                "query": str(request.url.query)[:256],
                                "method": request.method,
                                "status_code": 502,
                                "request_time": f"{elapsed:.4f}",
                                "upstream_status": "error",
                                "error_type": e.__class__.__name__,
                                "is_ws_upgrade": False,
                                "source": "api_gateway_proxy",
                                "ts": datetime.utcnow().isoformat() + "Z",
                            },
                        },
                    )
                raise HTTPException(status_code=502, detail=f"Proxy error: {e}")
            
            is_html = "text/html" in resp.headers.get("content-type", "")
            if is_html:
                content = await resp.aread()
                text = content.decode("utf-8", errors="replace")
                
                if settings.chat_inject_script_enabled:
                    inject_script = _build_chat_inject_script(
                        agent,
                        token,
                        clear_device_auth_storage=settings.chat_clear_device_auth_storage,
                        inject_gateway_settings=settings.chat_inject_gateway_settings,
                        dom_avatar_rewrite=settings.chat_dom_avatar_rewrite,
                    )

                    if 'window.__OPENCLAW_CONTROL_UI_BASE_PATH__="";' in text:
                        text = text.replace('window.__OPENCLAW_CONTROL_UI_BASE_PATH__="";', inject_script)
                    else:
                        injected = f"<script>{inject_script}</script>"
                        if '<script type="module"' in text:
                            text = text.replace('<script type="module"', f"{injected}<script type=\"module\"", 1)
                        elif "</head>" in text:
                            text = text.replace("</head>", f"{injected}</head>", 1)
                        elif "<body>" in text:
                            text = text.replace("<body>", f"<body>{injected}", 1)
                        else:
                            text = f"{injected}{text}"
                    if settings.chat_dom_avatar_rewrite:
                        text = re.sub(
                            r'window\.__OPENCLAW_ASSISTANT_AVATAR__=("|\')/avatar/[^"\']*("|\')',
                            'window.__OPENCLAW_ASSISTANT_AVATAR__=""',
                            text,
                        )
                await resp.aclose()
                
                return HTMLResponse(content=text, status_code=resp.status_code)

            content = await resp.aread()
            await resp.aclose()

            if (
                settings.chat_rewrite_control_ui_config
                and request.method.upper() == "GET"
                and path.endswith("control-ui-config.json")
                and resp.status_code == 200
            ):
                content = _rewrite_control_ui_config(content, agent, rewrite_avatar=settings.chat_rewrite_avatar_payloads)

            if request.method.upper() == "GET" and path.startswith("avatar/"):
                is_meta = request.query_params.get("meta") == "1"

                if is_meta and resp.status_code == 200 and "application/json" in resp.headers.get("content-type", ""):
                    if settings.chat_rewrite_avatar_payloads:
                        content = _rewrite_avatar_meta(content, agent, query)

                if not is_meta and resp.status_code == 404:
                    svg = _avatar_fallback_svg(agent)
                    return Response(content=svg, status_code=200, media_type="image/svg+xml")

            r = Response(content=content, status_code=resp.status_code)
            for k, v in resp.headers.items():
                if k.lower() not in ("content-length", "transfer-encoding", "connection"):
                    r.headers[k] = v

            if not upgrade:
                elapsed = max(0.0, time.perf_counter() - started)
                error_type = None
                if resp.status_code >= 400:
                    error_type = f"http_{resp.status_code}"
                await publish_event(
                    redis,
                    stream_key=settings.redis_stream_key,
                    event={
                        "type": "chat.gateway.access",
                        "agent": agent,
                        "payload": {
                            "path": f"/{path}",
                            "query": str(request.url.query)[:256],
                            "method": request.method,
                            "status_code": int(resp.status_code),
                            "request_time": f"{elapsed:.4f}",
                            "upstream_status": "ok" if resp.status_code < 500 else "error",
                            "error_type": error_type,
                            "is_ws_upgrade": False,
                            "source": "api_gateway_proxy",
                            "ts": datetime.utcnow().isoformat() + "Z",
                        },
                    },
                )
            return r

    async def _chat_ws_proxy_impl(agent: str, path: str, websocket: WebSocket):
        await websocket.accept()
        
        payload = {
            "path": f"/{path}",
            "query": "",
            "method": "GET",
            "status_code": 101,
            "request_time": "0.0",
            "upstream_status": "proxy",
            "is_ws_upgrade": True,
            "source": "api_gateway_proxy",
            "ts": datetime.utcnow().isoformat() + "Z",
        }
        await publish_event(
            redis,
            stream_key=settings.redis_stream_key,
            event={
                "type": "chat.gateway.access",
                "agent": agent,
                "payload": payload,
            },
        )
            
        target_ws_url = f"ws://openclaw-{agent}:{settings.chat_upstream_port}/{path}"
        ws_query = websocket.scope.get("query_string", b"")
        if isinstance(ws_query, (bytes, bytearray)) and ws_query:
            target_ws_url = f"{target_ws_url}?{ws_query.decode('utf-8', errors='ignore')}"
        import websockets
        from websockets.exceptions import ConnectionClosed
        
        token = settings.agent_token_map.get(agent)
        headers = {}
        if token:
            headers["Authorization"] = f"Bearer {token}"
        if settings.chat_force_loopback_headers:
            headers["X-Real-IP"] = "127.0.0.1"
            headers["X-Forwarded-For"] = "127.0.0.1"
        upstream_origin = _normalize_control_ui_origin(websocket.headers.get("origin"))
            
        try:
            async with websockets.connect(
                target_ws_url,
                extra_headers=headers or None,
                origin=upstream_origin or None,
            ) as upstream_ws:
                async def forward_to_upstream():
                    try:
                        while True:
                            msg = await websocket.receive_text()
                            if settings.chat_sanitize_connect_auth:
                                msg = _sanitize_connect_auth(
                                    msg,
                                    token,
                                    strip_stale_device_fields=settings.chat_strip_stale_device_fields,
                                    force_token_in_connect=settings.chat_force_token_in_connect,
                                )
                            await upstream_ws.send(msg)
                    except WebSocketDisconnect:
                        pass
                    except Exception:
                        pass
                        
                async def forward_to_client():
                    try:
                        while True:
                            msg = await upstream_ws.recv()
                            if settings.chat_rewrite_avatar_payloads and isinstance(msg, str) and "/avatar/" in msg:
                                try:
                                    payload = json.loads(msg)
                                    rewritten = _rewrite_avatar_paths(payload, agent)
                                    msg = json.dumps(rewritten, ensure_ascii=False)
                                except Exception:
                                    pass
                            await websocket.send_text(msg)
                    except ConnectionClosed:
                        pass
                    except Exception:
                        pass
                        
                await asyncio.gather(
                    forward_to_upstream(),
                    forward_to_client()
                )
        except Exception as e:
            await websocket.close(code=1011, reason=str(e))

    @app.websocket("/chat/{agent}")
    async def chat_ws_proxy_root(agent: str, websocket: WebSocket):
        await _chat_ws_proxy_impl(agent, "", websocket)

    @app.websocket("/chat/{agent}/{path:path}")
    async def chat_ws_proxy(agent: str, path: str, websocket: WebSocket):
        await _chat_ws_proxy_impl(agent, path, websocket)

    @app.on_event("shutdown")
    async def _shutdown():
        await redis.aclose()
        await engine.dispose()

    return app


async def publish_event(redis: Redis, stream_key: str, event: dict) -> None:
    await redis.xadd(stream_key, {"event": json.dumps(event, ensure_ascii=False)})


app = create_app()
